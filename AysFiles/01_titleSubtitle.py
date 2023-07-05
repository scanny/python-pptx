from pptx import Presentation

# Create new powerpoint presentation
presentation = Presentation()
# Slide Layout seçenekleri için toplamda 11 tane seçenek bulunmaktadır. Bunu aşağıdaki kod ile gözlemleyebiliriz.

# presentation = Presentation()
# slide_layouts = presentation.slide_layouts

print(len(slide_layouts))
# Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide2.shapes.title
subtitle = slide2.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 3
slide3 = presentation.slides.add_slide(presentation.slide_layouts[2])
title = slide3.shapes.title
subtitle = slide3.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 4
slide4 = presentation.slides.add_slide(presentation.slide_layouts[3])
title = slide4.shapes.title
subtitle = slide4.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 5
slide5 = presentation.slides.add_slide(presentation.slide_layouts[4])
title = slide5.shapes.title
subtitle = slide5.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 6
slide6 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide6.shapes.title
subtitle = slide6.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 7
slide7 = presentation.slides.add_slide(presentation.slide_layouts[6])
title = slide7.shapes.title
subtitle = slide7.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 8
slide8 = presentation.slides.add_slide(presentation.slide_layouts[7])
title = slide8.shapes.title
subtitle = slide8.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 9
slide9 = presentation.slides.add_slide(presentation.slide_layouts[8])
title = slide9.shapes.title
subtitle = slide9.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 10
slide10 = presentation.slides.add_slide(presentation.slide_layouts[9])
title = slide10.shapes.title
subtitle = slide10.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

# Slide 11
slide11 = presentation.slides.add_slide(presentation.slide_layouts[10])
title = slide11.shapes.title
subtitle = slide11.placeholders[1]
title.text = "this is a title example"
subtitle.text = "Name, Department, Date"

presentation.save("01_title_subtitle.pptx")
