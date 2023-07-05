from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

# Create an empty presentation
prs = Presentation()

# Add a slide to the presentation
slide_layout = prs.slide_layouts[5]
slide1 = prs.slides.add_slide(slide_layout)

# Add a title to the slide
title = slide1.shapes.title
title.text = "this is a title example"

# Create chart data structure
chart_data = CategoryChartData()

# Enter the categories, i.e. x-axis data.
chart_data.categories = ['East', 'West', 'Midwest']

# Add data values
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

# Add a column chart to the slide
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide1.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Add a slide
slide2 = prs.slides.add_slide(slide_layout)

# Create line chart data
categories = ["Kategori 1", "Kategori 2", "Kategori 3", "Kategori 4"]
data_series1 = [10, 15, 7, 12]
data_series2 = [5, 8, 11, 6]

# Create data structure
chart_data = ChartData()
chart_data.categories = categories
chart_data.add_series("Veri Serisi 1", data_series1)
chart_data.add_series("Veri Serisi 2", data_series2)

# Create a chart on the slide
x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
chart = slide2.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# Customize line chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

# Save PowerPoint presentation
prs.save("02_chart.pptx")
