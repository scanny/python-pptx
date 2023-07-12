from pptx import Presentation
from pptx.util import Inches

# Yeni bir sunum oluştur
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Boş bir slayt ekle
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Başlık ekle
title = slide.shapes.title
title.text = "Side Title"


# Tablo özellikleri
table_rows = 10
table_cols = 3
table_width = Inches(3)
table_height = Inches(0.5)
table_top = Inches(3)
table_horizontal_spacing = Inches(0.5)

# İlk tablo konumunu belirle
table_left = Inches(1)

# 4 tabloyu oluştur
for i in range(4):
    # Slayta tabloyu ekle
    table = slide.shapes.add_table(table_rows, table_cols, table_left, table_top, table_width, table_height).table

    # Tablonun başlıklarını ayarla
    for col in range(table_cols):
        table.cell(0, col).text = f"Header {col+1}"
    
    # Tablonun içeriğini doldur
    for row in range(1, table_rows):
        for col in range(table_cols):
            table.cell(row, col).text = f"data {row + col + 1}"

    # Tablo stilini ayarla
    table.style = "Table Grid"

    # Sonraki tablo için tablo konumunu güncelle
    table_left += table_width + table_horizontal_spacing

# Yeni tablo özellikleri
new_table_rows = 1
new_table_cols = 2
new_table_width = Inches(13.5)
new_table_height = Inches(0.5)
new_table_top = Inches(2)  # İlk tablonun altına yerleştirildi
new_table_left = Inches(1)  # İlk tablonun sağ tarafında yerleştirildi

# Slayta yeni tabloyu ekle
new_table = slide.shapes.add_table(new_table_rows, new_table_cols, new_table_left, new_table_top, new_table_width, new_table_height).table

# Yeni tablonun başlıklarını ayarla
new_table.cell(0, 0).text = "Customer Cam"
new_table.cell(0, 1).text = "Referance Cam"

# Sunumu kaydet
presentation.save("03_1_Table.pptx")
