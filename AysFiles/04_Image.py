from pptx import Presentation
from pptx.util import Inches

# Yeni bir sunum oluştur
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Boş bir slayt ekle
slide_layout = presentation.slide_layouts[2]
slide1 = presentation.slides.add_slide(slide_layout)

# Başlık ve alt başlık ekle
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Slide Title"
subtitle.text = "Name, Department, Date"

# Resim eklemek için konum ve boyut ayarları
left = 0
top = 0
width = presentation.slide_width
height = 3300000

# Resmi slayta ekle
slide1.shapes.add_picture("resim.png", left, top, width, height)

# Sunumu kaydet
presentation.save("04_Image.pptx")
