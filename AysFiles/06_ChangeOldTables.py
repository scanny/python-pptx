from pptx import Presentation

presentation = Presentation('deneme.pptx')

table_count = 0
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_table:
            
            table = shape.table
            print(table)
            if table and table.cell(0, 0).text:
                print("{} tablonun 1. Hücre değeri: {}".format(table_count, table.cell(0, 0).text) )
                table.cell(0,0).text = "ays was here"
                print("{} tablonun 1. Hücre değeri: {}".format(table_count, table.cell(0, 0).text) )
                table_count = table_count + 1

print("Toplam tablo sayıisi:", table_count)
presentation.save("06_ChangeOldTables.pptx")
