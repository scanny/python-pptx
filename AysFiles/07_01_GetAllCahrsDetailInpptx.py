from pptx import Presentation
import openpyxl

# Load the PowerPoint file
presentation = Presentation('ornek_cam.pptx')

# Create a new Excel workbook
wb = openpyxl.Workbook()

# Add a new sheet to the workbook
ws = wb.active
ws.title = "Chart Details"

# Write the headers for chart details
ws.append(["Slide ID", "Chart Type", "Chart Title"])

# Loop through each slide
for slide in presentation.slides:
    # Loop through each shape on the slide
    for shape in slide.shapes:
        # Check if the shape is a chart
        if shape.has_chart:
            # Get chart details
            chart = shape.chart
            chart_type = chart.chart_type
            chart_title = chart.chart_title.text_frame.text

            # Write chart details to the Excel sheet
            ws.append([slide.slide_id, chart_type, chart_title])

# Save the Excel file
wb.save('chart_details.xlsx')
