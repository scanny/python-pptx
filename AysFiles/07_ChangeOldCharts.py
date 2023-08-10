from pptx import Presentation
from pptx.chart.data import ChartData
import openpyxl

# Presentation isimli değişkenimize üzerinde değişiklik yapacağımız powerpoint dosyasını atıyoruzz.
presentation = Presentation('Deneme.pptx')

# Burada powerpoint dosyasında nerelerde chart olduğunu tespit etmek için bir bir workbook oluşturuyoruz.
wb_chart_details = openpyxl.Workbook()
ws_chart_details = wb_chart_details.active
ws_chart_details.title = "Chart Details"
ws_chart_details.append(["Slide ID", "Chart Type", "Chart Title"])

# Bu kısım önemli! Chartların kendine ait bir ID numarası olup olmadığını kontrol ettik ve olmadığını keşfettik. 
# Bu kısımda chart'a bir title atıyoruz ve bu title o chart'ın ID numarası gibi kullanlıyor
target_chart_title = "Second"

# Bu workbook chart için manuel olarak girdiğimiz excel dosyasını içeriyor. İsmini ve sayfa ismini giriyoruz.
wb_chart_change = openpyxl.load_workbook("deneme.xlsx")
sheet = wb_chart_change['Tabelle1']

# Hangi sütunları kullanacağımızı belirliyoruz. 
# A sütunu Kategorileri belirler B sütunu 1. line ın verilerini, C sütunu 2. line'ın verilerini belirliyor.
new_categories = [cell.value for cell in sheet['A'][1:]]
new_values = [cell.value for cell in sheet['B'][1:]]
new_values2 = [cell.value for cell in sheet['C'][1:]]

for slide in presentation.slides:
    for shape in slide.shapes:
        print(shape.shape_id)
        if shape.has_chart:
            chart = shape.chart

            #  Detect Charts
            chart_type = chart.chart_type
            chart_title = chart.chart_title.text_frame.text
            ws_chart_details.append([slide.slide_id, chart_type, chart_title])
            if chart_title == target_chart_title:

                # Change Charts
                chart_data = ChartData()
                chart_data.categories = new_categories
                chart_data.add_series('Updated Data', new_values)
                chart_data.add_series('Updated Data', new_values2)
                chart.replace_data(chart_data)
                print(f"Updated data for the chart with title: {target_chart_title} on slide {slide.slide_id}")

wb_chart_details.save('deneme_chart_details.xlsx')
presentation.save('Deneme_sonuc1.pptx')
