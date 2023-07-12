# Libraries
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

# Empty presentation and default settings
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]

# Slide 1 Title, Subtitle
slide1 = presentation.slides.add_slide(presentation.slide_layouts[2])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Cam HDPX_XXX customer slide"
subtitle.text = "Name, Department, Date"

left1 = 0
top1 = 0
width1 = presentation.slide_width
height1 = 3300000

slide1.shapes.add_picture("resim.png", left1, top1, width1, height1) #Image shold be in the same folder or ypu should write path

# Slide 2 Graph
slide2 = presentation.slides.add_slide(slide_layout)

title = slide2.shapes.title
title.text = "Cam HDPX_XXX customer slide"

# Data creating
categories = ["20", "40", "60", "80", "100", "120"]
data_series1 = [0, 15, 20, 20, 10, -5, 0]
data_series2 = [0, 3, 5, 5, 0, -2, 0]

# Data structure
chart_data = CategoryChartData()
chart_data.categories = categories
chart_data.add_series("Veri Serisi 1", data_series1)
chart_data.add_series("Veri Serisi 2", data_series2)

# Linechart
x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
chart = slide2.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

# Slide 3 Table
slide3 = presentation.slides.add_slide(slide_layout)
title = slide3.shapes.title
title.text = "Cam HDPX_XXX customer slide"

table_rows = 10
table_cols = 3
table_width = Inches(3)
table_height = Inches(4)
table_top = Inches(3)
table_horizontal_spacing = Inches(0.5)

# Firs table location
table_left = Inches(1)

# 4 table
for i in range(4):
    # Add table to slide
    table = slide3.shapes.add_table(table_rows, table_cols, table_left, table_top, table_width, table_height).table

    # Tablonun başlıklarını ayarla
    for col in range(table_cols):
        table.cell(0, col).text = f"Header {col+1}"
    
    # Tablonun içeriğini doldur
    for row in range(1, table_rows):
        for col in range(table_cols):
            table.cell(row, col).text = f"data {row + col + 1}"

    # Tablo stilini ayarla
    table.style = "Table Grid"

    # Sonraki tablo için tablo konumunu güncelle
    table_left += table_width + table_horizontal_spacing

# Yeni tablo özellikleri
new_table_rows = 1
new_table_cols = 2
new_table_width = Inches(13.5)
new_table_height = Inches(0.5)
new_table_top = Inches(2)  
new_table_left = Inches(1)  

# Slayta yeni tabloyu ekle
new_table = slide3.shapes.add_table(new_table_rows, new_table_cols, new_table_left, new_table_top, new_table_width, new_table_height).table

# Yeni tablonun başlıklarını ayarla
new_table.cell(0, 0).text = "Customer Cam"
new_table.cell(0, 1).text = "Referance Cam"

#Save the presentation
presentation.save("draft.pptx")
