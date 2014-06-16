import os, sys, time, math, zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# simple class to create objects
class SimpleClass(object):
	pass

# global variables
OVAL_HEIGHT = 950000 # height of oval shape
OVAL_WIDTH = 950000 # width of oval shape
HALF_OVAL_HEIGHT = round(OVAL_HEIGHT / 2)
HALF_OVAL_WIDTH = round(OVAL_WIDTH / 2)
SLIDE_HEIGHT = 6858000 # assume standard 10" x 7.5" slide (4:3 aspect ratio)
SLIDE_WIDTH = 9144000 # assume standard 10" x 7.5" slide (4:3 aspect ratio)
MAX_TOP = SLIDE_HEIGHT - OVAL_HEIGHT # furthest top of oval shape where oval is completely on slide
MAX_LEFT = SLIDE_WIDTH - OVAL_WIDTH # furthest left of oval shape where oval is completely on slide

# check for proper number of command line arguments
if len(sys.argv) < 2:
	print "Usage: python export.py <output_path>"
	print "Example: python export.py /Users/oscar/Documents/python-pptx/export.pptx"
	exit()

# get path from command line
output_path = sys.argv[1]

# parse path and filename
index = output_path.rfind("/") + 1
output_folder = output_path[0:index]
filename = output_path[index:]

# constructs are just circle shapes, links are connectors
constructs = [[0, "A", "emotional", 1000000, 150000], [1, "B", "psychosocial", 1750000, 7000000], [2, "C", "attribute", 4000000, 750000], [3, "D", "functional", 5500000, 6000000]]
links = [[0, 2], [1, 0], [1, 2], [3, 0]]
shapes = []
connectors = []

# create a new presentation
pptx = Presentation()
current_link_id = 1000

# set up new slide
BLANK_SLIDE_LAYOUT = pptx.slide_layouts[6]
slide = pptx.slides.add_slide(BLANK_SLIDE_LAYOUT)

# create each connector object and add to list
for l in links:

	c = SimpleClass()
	c.start = l[0]
	c.end = l[1]
	c.id = "-1"
	connectors.append(c)

# create each shape object and add to list
for c in constructs:

	s = SimpleClass()
	s.id = c[0]
	s.word = c[1]
	s.type = c[2]
	s.top = c[3]
	s.left = c[4]
	shapes.append(s)

# add shapes
for s in shapes:

	# create oval shape and set id
	shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, s.left, s.top, OVAL_WIDTH, OVAL_HEIGHT)

	# not sure if this ever worked
	s.xml_id = shape.id
	current_link_id = shape.id + 1

	# clear current text and add new placeholder
	shape.textframe.clear()
	p = shape.textframe.paragraphs[0]
	run = p.add_run()
	run.text = s.word

	# set font, size, and color
	font = run.font
	font.name = "Calibri"
	font.size = Pt(10.5)
	font.color.rgb = RGBColor(50, 50, 50) # grey

	# set shape fill
	fill = shape.fill
	fill.solid()
	if s.type == "emotional": # blue
		fill.fore_color.rgb = RGBColor(204, 224, 255)
	elif s.type == "psychosocial": # green
		fill.fore_color.rgb = RGBColor(217, 228, 191)
	elif s.type == "functional": # peach
		fill.fore_color.rgb = RGBColor(255, 231, 219)
	elif s.type == "attribute": # purple
		fill.fore_color.rgb = RGBColor(226, 215, 243)
	elif s.type == "other": # white
		fill.fore_color.rgb = RGBColor(255, 255, 255)
	else: # white
		fill.fore_color.rgb = RGBColor(255, 255, 255)

# save file
os.chdir(output_folder)
pptx.save(filename)

# unzip .pptx to temporary space and remove file
timestamp = str(time.time()).replace(".", "")
zip_file = zipfile.ZipFile(filename, "r")
zip_file.extractall(os.path.join("/tmp", timestamp))
zip_file.close
os.remove(filename)

# register necessary namespaces
ET.register_namespace("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
ET.register_namespace("p", "http://schemas.openxmlformats.org/presentationml/2006/main")
ET.register_namespace("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")

# parse xml document and find shape tree
tree = ET.parse(os.path.join("/tmp", timestamp, "ppt", "slides", "slide1.xml"))
root = tree.getroot()

# get necessary elements
iterator = root.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}spTree")
shapeTree = iterator.next()
iterator = root.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr")
groupShapeProperties = iterator.next()

# add properties (these seem to not change)
twodTransform = ET.SubElement(groupShapeProperties, "a:xfrm")
offset = ET.SubElement(twodTransform, "a:off")
offset.set("x", "0")
offset.set("y", "0")
extents = ET.SubElement(twodTransform, "a:ext")
extents.set("cx", "0")
extents.set("cy", "0")
childOffset = ET.SubElement(twodTransform, "a:chOff")
childOffset.set("x", "0")
childOffset.set("y", "0")
childExtents = ET.SubElement(twodTransform, "a:chExt")
childExtents.set("cx", "0")
childExtents.set("cy", "0")

# for each connector
for connector in connectors:

	# variables
	start_shape = None
	end_shape = None
	flip_v = False
	flip_h = False

	# assign and increment link id
	connector.id = current_link_id
	current_link_id += 1

	# get start and end constructs
	for s in shapes:
		if (s.id == connector.start):
			start_shape = s
		elif (s.id == connector.end):
			end_shape = s
		if start_shape != None and end_shape != None:
			break

	# if start and end constructs can't be found, do not run this iteration
	if start_shape == None or end_shape == None:
		continue

	# calculate x and y differences
	x_diff = abs(start_shape.left - end_shape.left)
	y_diff = abs(start_shape.top - end_shape.top)

	# start node is completely above end node
	if (start_shape.top + OVAL_HEIGHT) < end_shape.top:

		# start node is completely above and completely left of end node
		if (start_shape.left + OVAL_WIDTH) < end_shape.left:

			# further apart left-right than up-down
			if x_diff > y_diff:

				start_idx = "6"
				end_idx = "2"
				cxn_x = start_shape.left + OVAL_WIDTH
				cxn_y = start_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = end_shape.left - start_shape.left - OVAL_WIDTH
				cxn_cy = end_shape.top - start_shape.top

			# further apart up-down than left-right
			else:

				start_idx = "4"
				end_idx = "0"
				cxn_x = start_shape.left + HALF_OVAL_HEIGHT
				cxn_y = start_shape.top + OVAL_HEIGHT
				cxn_cx = end_shape.left - start_shape.left
				cxn_cy = end_shape.top - start_shape.top - OVAL_HEIGHT

		# start node is completely above and completely right of end node
		elif (end_shape.left + OVAL_WIDTH) < start_shape.left:

			flip_h = True
			
			# further apart left-right than up-down
			if x_diff > y_diff:

				start_idx = "2"
				end_idx = "6"
				cxn_x = end_shape.left + OVAL_WIDTH
				cxn_y = start_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = start_shape.left - end_shape.left - OVAL_WIDTH
				cxn_cy = end_shape.top - start_shape.top

			# further apart up-down than left-right
			else:

				start_idx = "4"
				end_idx = "0"
				cxn_x = end_shape.left + HALF_OVAL_HEIGHT
				cxn_y = start_shape.top + OVAL_HEIGHT
				cxn_cx = start_shape.left - end_shape.left
				cxn_cy = end_shape.top - start_shape.top - OVAL_HEIGHT

		# start node is completely above and partially left of end node
		elif start_shape.left < end_shape.left:

			start_idx = "4"
			end_idx = "0"
			cxn_x = start_shape.left + HALF_OVAL_HEIGHT
			cxn_y = start_shape.top + OVAL_HEIGHT
			cxn_cx = end_shape.left - start_shape.left
			cxn_cy = end_shape.top - start_shape.top - OVAL_HEIGHT

		# start node is completely above and partially right of end node
		else:

			flip_h = True
			start_idx = "4"
			end_idx = "0"
			cxn_x = end_shape.left + HALF_OVAL_HEIGHT
			cxn_y = start_shape.top + OVAL_HEIGHT
			cxn_cx = start_shape.left - end_shape.left
			cxn_cy = end_shape.top - start_shape.top - OVAL_HEIGHT

	# start node is completely below end node
	elif (end_shape.top + OVAL_HEIGHT) < start_shape.top:

		flip_v = True

		# start node is completely below and completely left of end node
		if (start_shape.left + OVAL_WIDTH) < end_shape.left:
			
			# further apart left-right than up-down
			if x_diff > y_diff:

				start_idx = "6"
				end_idx = "2"
				cxn_x = start_shape.left + OVAL_WIDTH
				cxn_y = end_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = end_shape.left - start_shape.left - OVAL_WIDTH
				cxn_cy = start_shape.top - end_shape.top

			# further apart up-down than left-right
			else:

				start_idx = "0"
				end_idx = "4"
				cxn_x = start_shape.left + HALF_OVAL_HEIGHT
				cxn_y = end_shape.top + OVAL_HEIGHT
				cxn_cx = end_shape.left - start_shape.left
				cxn_cy = start_shape.top - end_shape.top - OVAL_HEIGHT

		# start node is completely below and completely right of end node
		elif (end_shape.left + OVAL_WIDTH) < start_shape.left:

			flip_h = True
			
			# further apart left-right than up-down
			if x_diff > y_diff:

				start_idx = "2"
				end_idx = "6"
				cxn_x = end_shape.left + OVAL_WIDTH
				cxn_y = end_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = start_shape.left - end_shape.left - OVAL_WIDTH
				cxn_cy = start_shape.top - end_shape.top

			# further apart up-down than left-right
			else:

				start_idx = "0"
				end_idx = "4"
				cxn_x = end_shape.left + HALF_OVAL_HEIGHT
				cxn_y = end_shape.top + OVAL_HEIGHT
				cxn_cx = start_shape.left - end_shape.left
				cxn_cy = start_shape.top - end_shape.top - OVAL_HEIGHT

		# start node is completely below and partially left of end node
		elif start_shape.left < end_shape.left:

			start_idx = "0"
			end_idx = "4"
			cxn_x = start_shape.left + HALF_OVAL_HEIGHT
			cxn_y = end_shape.top + OVAL_HEIGHT
			cxn_cx = end_shape.left - start_shape.left
			cxn_cy = start_shape.top - end_shape.top - OVAL_HEIGHT

		# start node is completely below and partially right of end node
		else:

			flip_h = True
			start_idx = "0"
			end_idx = "4"
			cxn_x = end_shape.left + HALF_OVAL_HEIGHT
			cxn_y = end_shape.top + OVAL_HEIGHT
			cxn_cx = start_shape.left - end_shape.left
			cxn_cy = start_shape.top - end_shape.top - OVAL_HEIGHT

	# start node is partially above end node
	elif start_shape.top < end_shape.top:

		# start node is partially above and completely left of end node
		if (start_shape.left + OVAL_WIDTH) < end_shape.left:

			start_idx = "6"
			end_idx = "2"
			cxn_x = start_shape.left + OVAL_WIDTH
			cxn_y = start_shape.top + HALF_OVAL_HEIGHT
			cxn_cx = end_shape.left - start_shape.left - OVAL_WIDTH
			cxn_cy = end_shape.top - start_shape.top

		# start node is partially above and completely right of end node
		elif (end_shape.left + OVAL_WIDTH) < start_shape.left:

			flip_h = True
			start_idx = "2"
			end_idx = "6"
			cxn_x = end_shape.left + OVAL_WIDTH
			cxn_y = start_shape.top + HALF_OVAL_HEIGHT
			cxn_cx = start_shape.left - end_shape.left - OVAL_WIDTH
			cxn_cy = end_shape.top - start_shape.top

		# start node is partially above and partially left of end node
		elif start_shape.left < end_shape.left:

			# further apart left-right than up-down
			if x_diff > y_diff:
				
				flip_h = True
				start_idx = "6"
				end_idx = "2"
				cxn_x = end_shape.left
				cxn_y = start_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = (start_shape.left + OVAL_WIDTH) - end_shape.left
				cxn_cy = end_shape.top - start_shape.top

			# further apart up-down than left-right
			else:

				flip_v = True
				start_idx = "4"
				end_idx = "0"
				cxn_x = start_shape.left + HALF_OVAL_HEIGHT
				cxn_y = end_shape.top
				cxn_cx = end_shape.left - start_shape.left
				cxn_cy = (start_shape.top + OVAL_HEIGHT) - end_shape.top

		# start node is partially above and partially right of end node
		else:

			# further apart left-right than up-down
			if x_diff > y_diff:
				
				start_idx = "2"
				end_idx = "6"
				cxn_x = start_shape.left
				cxn_y = end_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = (end_shape.left + OVAL_WIDTH) - start_shape.left
				cxn_cy = end_shape.top - start_shape.top

			# further apart up-down than left-right
			else:

				flip_v = True
				flip_h = True
				start_idx = "4"
				end_idx = "0"
				cxn_x = end_shape.left + HALF_OVAL_WIDTH
				cxn_y = end_shape.top
				cxn_cx = start_shape.left - end_shape.left
				cxn_cy = (start_shape.top + OVAL_HEIGHT) - end_shape.top

	# start node is partially below end node
	else:

		# start node is partially below and completely left of end node
		if (start_shape.left + OVAL_WIDTH) < end_shape.left:

			flip_v = True
			start_idx = "6"
			end_idx = "2"
			cxn_x = start_shape.left + OVAL_WIDTH
			cxn_y = end_shape.top + HALF_OVAL_HEIGHT
			cxn_cx = end_shape.left - (start_shape.left + OVAL_WIDTH)
			cxn_cy = start_shape.top - end_shape.top

		# start node is partially below and completely right of end node
		elif (end_shape.left + OVAL_WIDTH) < start_shape.left:

			flip_v = True
			flip_h = True
			start_idx = "2"
			end_idx = "6"
			cxn_x = end_shape.left + OVAL_WIDTH
			cxn_y = end_shape.top + HALF_OVAL_HEIGHT
			cxn_cx = start_shape.left - (end_shape.left + OVAL_WIDTH)
			cxn_cy = start_shape.top - end_shape.top

		# start node is partially below and partially left of end node
		elif start_shape.left < end_shape.left:

			# further apart left-right than up-down
			if x_diff > y_diff:

				flip_v = True
				flip_h = True
				start_idx = "6"
				end_idx = "2"
				cxn_x = end_shape.left
				cxn_y = end_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = (start_shape.left + OVAL_WIDTH) - end_shape.left
				cxn_cy = start_shape.top - end_shape.top

			# further apart up-down than left-right
			else:

				start_idx = "0"
				end_idx = "4"
				cxn_x = start_shape.left + HALF_OVAL_HEIGHT
				cxn_y = start_shape.top
				cxn_cx = end_shape.left - start_shape.left
				cxn_cy = (end_shape.top + OVAL_HEIGHT) - start_shape.top

		# start node is partially below and partially right of end node
		else:

			# further apart left-right than up-down
			if x_diff > y_diff:

				flip_v = True
				start_idx = "2"
				end_idx = "6"
				cxn_x = start_shape.left
				cxn_y = end_shape.top + HALF_OVAL_HEIGHT
				cxn_cx = (end_shape.left + OVAL_WIDTH) - start_shape.left
				cxn_cy = start_shape.top - end_shape.top

			# further apart up-down than left-right
			else:

				flip_h = True
				start_idx = "0"
				end_idx = "4"
				cxn_x = end_shape.left + HALF_OVAL_WIDTH
				cxn_y = start_shape.top
				cxn_cx = start_shape.left - end_shape.left
				cxn_cy = (end_shape.top + OVAL_HEIGHT) - start_shape.top

	connectionShape = ET.SubElement(shapeTree, "p:cxnSp")
	nonVisualConnectorShapeDrawingProperties = ET.SubElement(connectionShape, "p:nvCxnSpPr")
	cNonVisualProperties = ET.SubElement(nonVisualConnectorShapeDrawingProperties, "p:cNvPr")
	cNonVisualProperties.set("id", str(connector.id))
	cNonVisualProperties.set("name", "Connector " + str(connector.id))
	cNonVisualConnectorShapeDrawingProperties = ET.SubElement(nonVisualConnectorShapeDrawingProperties, "p:cNvCxnSpPr")
	ET.SubElement(nonVisualConnectorShapeDrawingProperties, "p:nvPr")

	connectionStart = ET.SubElement(cNonVisualConnectorShapeDrawingProperties, "a:stCxn")
	connectionStart.set("id", str(start_shape.xml_id)) # shape index from which connector starts (param)
	connectionStart.set("idx", start_idx) # connector spawn point index
	connectionEnd = ET.SubElement(cNonVisualConnectorShapeDrawingProperties, "a:endCxn")
	connectionEnd.set("id", str(end_shape.xml_id)) # shape index at which connector ends (param)
	connectionEnd.set("idx", end_idx) # connector termination point index

	shapeProperties = ET.SubElement(connectionShape, "p:spPr")
	twodTransform = ET.SubElement(shapeProperties, "a:xfrm")
	
	# check flipV and flipH
	if flip_v:
		twodTransform.set("flipV", "1")
	if flip_h:
		twodTransform.set("flipH", "1")

	# location of bounding box (params)
	offset = ET.SubElement(twodTransform, "a:off")
	offset.set("x", str(cxn_x).replace(".0", ""))
	offset.set("y", str(cxn_y).replace(".0", ""))

	# height and width of bounding box (params)
	extents = ET.SubElement(twodTransform, "a:ext")
	extents.set("cx", str(cxn_cx).replace(".0", ""))
	extents.set("cy", str(cxn_cy).replace(".0", ""))

	presetGeometry = ET.SubElement(shapeProperties, "a:prstGeom")
	presetGeometry.set("prst", "line")
	ET.SubElement(presetGeometry, "a:avLst")

	style = ET.SubElement(connectionShape, "p:style")

	lineReference = ET.SubElement(style, "a:lnRef")
	lineReference.set("idx", "1")
	schemeColor = ET.SubElement(lineReference, "a:schemeClr")
	schemeColor.set("val", "dk1")

	fillReference = ET.SubElement(style, "a:fillRef")
	fillReference.set("idx", "0")
	schemeColor = ET.SubElement(fillReference, "a:schemeClr")
	schemeColor.set("val", "dk1")

	effectReference = ET.SubElement(style, "a:effectRef")
	effectReference.set("idx", "0")
	schemeColor = ET.SubElement(effectReference, "a:schemeClr")
	schemeColor.set("val", "dk1")

	fontReference = ET.SubElement(style, "a:fontRef")
	fontReference.set("idx", "minor")
	schemeColor = ET.SubElement(fontReference, "a:schemeClr")
	schemeColor.set("val", "tx1")

# write out the final tree
f = open(os.path.join("/tmp", timestamp, "ppt", "slides", "slide1.xml"), "w")
tree.write(f, encoding="UTF-8", xml_declaration=True)
f.close()

# zip the file
zip_file = zipfile.ZipFile(filename, "w", zipfile.ZIP_DEFLATED)
os.chdir(os.path.join("/tmp", timestamp))
for root, dirs, files in os.walk("."):
	for file in files:
		zip_file.write(os.path.join(root, file))
zip_file.close
