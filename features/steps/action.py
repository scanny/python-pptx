# encoding: utf-8

"""
Gherkin step implementations for click action-related features.
"""

from __future__ import absolute_import, print_function

from behave import given, then, when

from pptx import Presentation
from pptx.action import Hyperlink
from pptx.enum.action import PP_ACTION

from helpers import test_pptx


# given ===================================================

@given('an ActionSetting object having action {action} as click_action')
def given_an_ActionSetting_object_as_click_action(context, action):
    shape_idx = {'NONE': 0, 'NAMED_SLIDE': 6}[action]
    slides = Presentation(test_pptx('act-props')).slides
    context.slides = slides
    context.click_action = slides[2].shapes[shape_idx].click_action


@given('another slide in the deck as slide')
def given_another_slide_in_the_deck_as_slide(context):
    context.slide = context.slides[1]


@given('a shape having click action {action}')
def given_a_shape_having_click_action_action(context, action):
    shape_idx = (
        'none',
        'first slide',
        'last slide',
        'previous slide',
        'next slide',
        'last slide viewed',
        'named slide',
        'end show',
        'hyperlink',
        'other presentation',
        'open file',
        'custom slide show',
        'OLE action',
        'run macro',
        'run program',
    ).index(action)
    slides = Presentation(test_pptx('act-props')).slides
    context.slides = slides
    context.click_action = slides[2].shapes[shape_idx].click_action


# when ====================================================

@when('I assign {value} to click_action.hyperlink.address')
def when_I_assign_value_to_click_action_hyperlink_address(context, value):
    value = None if value == 'None' else value
    context.click_action.hyperlink.address = value


@when('I assign {value} to click_action.target_slide')
def when_I_assign_value_to_click_action_target_slide(context, value):
    rhs = {'None': None, 'slide': context.slide}[value]
    context.click_action.target_slide = rhs


# then ====================================================

@then('click_action.action is {member_name}')
def then_click_action_action_is_value(context, member_name):
    click_action = context.click_action
    expected_value = getattr(PP_ACTION, member_name)
    assert click_action.action == expected_value


@then('click_action.hyperlink is a Hyperlink object')
def then_click_action_hyperlink_is_a_Hyperlink_object(context):
    hyperlink = context.click_action.hyperlink
    assert isinstance(hyperlink, Hyperlink)


@then('click_action.hyperlink.address is {value}')
def then_click_action_hyperlink_address_is_value(context, value):
    expected_value = None if value == 'None' else value
    hyperlink = context.click_action.hyperlink
    print('expected value %s != %s' % (expected_value, hyperlink.address))
    assert hyperlink.address == expected_value


@then('click_action.target_slide is {value}')
def then_click_action_target_slide_is_value(context, value):
    if value.startswith('slides['):
        idx = value[7]
        expected_value = context.slides[int(idx)]
    elif value == 'None':
        expected_value = None
    else:
        expected_value = context.slide

    click_action = context.click_action
    assert click_action.target_slide == expected_value
