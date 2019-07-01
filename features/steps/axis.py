# encoding: utf-8

"""Gherkin step implementations for chart axis features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation
from pptx.enum.chart import XL_AXIS_CROSSES, XL_CATEGORY_TYPE

from helpers import test_pptx


# given ===================================================


@given("a {axis_type} axis")
def given_a_axis_type_axis(context, axis_type):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[0].shapes[0].chart
    context.axis = {"category": chart.category_axis, "value": chart.value_axis}[
        axis_type
    ]


@given("a major gridlines")
def given_a_major_gridlines(context):
    prs = Presentation(test_pptx("cht-gridlines-props"))
    axis = prs.slides[0].shapes[0].chart.value_axis
    context.gridlines = axis.major_gridlines


@given("a value axis having category axis crossing of {crossing}")
def given_a_value_axis_having_cat_ax_crossing_of(context, crossing):
    slide_idx = {"automatic": 0, "maximum": 2, "minimum": 3, "2.75": 4, "-1.5": 5}[
        crossing
    ]
    prs = Presentation(test_pptx("cht-axis-props"))
    context.value_axis = prs.slides[slide_idx].shapes[0].chart.value_axis


@given("an axis")
def given_an_axis(context):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.value_axis


@given("an axis having {a_or_no} title")
def given_an_axis_having_a_or_no_title(context, a_or_no):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[7].shapes[0].chart
    context.axis = {"a": chart.value_axis, "no": chart.category_axis}[a_or_no]


@given("an axis having {major_or_minor} gridlines")
def given_an_axis_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.value_axis


@given("an axis having {major_or_minor} unit of {value}")
def given_an_axis_having_major_or_minor_unit_of_value(context, major_or_minor, value):
    slide_idx = 0 if value == "Auto" else 1
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.axis = chart.value_axis


@given("an axis of type {cls_name}")
def given_an_axis_of_type_cls_name(context, cls_name):
    slide_idx = {"CategoryAxis": 0, "DateAxis": 6}[cls_name]
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.axis = chart.category_axis


@given("an axis not having {major_or_minor} gridlines")
def given_an_axis_not_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.category_axis


@given("an axis title")
def given_an_axis_title(context):
    prs = Presentation(test_pptx("cht-axis-props"))
    context.axis_title = prs.slides[7].shapes[0].chart.value_axis.axis_title


@given("an axis title having {a_or_no} text frame")
def given_an_axis_title_having_a_or_no_text_frame(context, a_or_no):
    prs = Presentation(test_pptx("cht-axis-props"))
    chart = prs.slides[7].shapes[0].chart
    axis = {"a": chart.value_axis, "no": chart.category_axis}[a_or_no]
    context.axis_title = axis.axis_title


@given("tick labels having an offset of {setting}")
def given_tick_labels_having_an_offset_of_setting(context, setting):
    slide_idx = {"no explicit setting": 0, "420": 1}[setting]
    prs = Presentation(test_pptx("cht-ticklabels-props"))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.tick_labels = chart.category_axis.tick_labels


# when ====================================================


@when("I assign {value} to axis.has_title")
def when_I_assign_value_to_axis_has_title(context, value):
    context.axis.has_title = {"True": True, "False": False}[value]


@when("I assign {value} to axis.has_{major_or_minor}_gridlines")
def when_I_assign_value_to_axis_has_major_or_minor_gridlines(
    context, value, major_or_minor
):
    axis = context.axis
    propname = "has_%s_gridlines" % major_or_minor
    new_value = {"True": True, "False": False}[value]
    setattr(axis, propname, new_value)


@when("I assign {value} to axis.{major_or_minor}_unit")
def when_I_assign_value_to_axis_major_or_minor_unit(context, value, major_or_minor):
    axis = context.axis
    propname = "%s_unit" % major_or_minor
    new_value = {"8.4": 8.4, "5": 5, "None": None}[value]
    setattr(axis, propname, new_value)


@when("I assign {value} to axis_title.has_text_frame")
def when_I_assign_value_to_axis_title_has_text_frame(context, value):
    context.axis_title.has_text_frame = {"True": True, "False": False}[value]


@when("I assign {value} to tick_labels.offset")
def when_I_assign_value_to_tick_labels_offset(context, value):
    new_value = int(value)
    context.tick_labels.offset = new_value


@when("I assign {member} to value_axis.crosses")
def when_I_assign_member_to_value_axis_crosses(context, member):
    value_axis = context.value_axis
    value_axis.crosses = getattr(XL_AXIS_CROSSES, member)


@when("I assign {value} to value_axis.crosses_at")
def when_I_assign_value_to_value_axis_crosses_at(context, value):
    new_value = None if value == "None" else float(value)
    context.value_axis.crosses_at = new_value


# then ====================================================


@then("axis.axis_title is an AxisTitle object")
def then_axis_axis_title_is_an_AxisTitle_object(context):
    class_name = type(context.axis.axis_title).__name__
    assert class_name == "AxisTitle", "got %s" % class_name


@then("axis.category_type is XL_CATEGORY_TYPE.{member}")
def then_axis_category_type_is_XL_CATEGORY_TYPE_member(context, member):
    expected_value = getattr(XL_CATEGORY_TYPE, member)
    category_type = context.axis.category_type
    assert category_type is expected_value, "got %s" % category_type


@then("axis.format is a ChartFormat object")
def then_axis_format_is_a_ChartFormat_object(context):
    axis = context.axis
    assert type(axis.format).__name__ == "ChartFormat"


@then("axis.format.fill is a FillFormat object")
def then_axis_format_fill_is_a_FillFormat_object(context):
    axis = context.axis
    assert type(axis.format.fill).__name__ == "FillFormat"


@then("axis.format.line is a LineFormat object")
def then_axis_format_line_is_a_LineFormat_object(context):
    axis = context.axis
    assert type(axis.format.line).__name__ == "LineFormat"


@then("axis.has_title is {value}")
def then_axis_has_title_is_value(context, value):
    axis = context.axis
    actual_value = axis.has_title
    expected_value = {"True": True, "False": False}[value]
    assert actual_value is expected_value, "got %s" % actual_value


@then("axis.has_{major_or_minor}_gridlines is {value}")
def then_axis_has_major_or_minor_gridlines_is_expected_value(
    context, major_or_minor, value
):
    axis = context.axis
    actual_value = {
        "major": axis.has_major_gridlines,
        "minor": axis.has_minor_gridlines,
    }[major_or_minor]
    expected_value = {"True": True, "False": False}[value]
    assert actual_value is expected_value, "got %s" % actual_value


@then("axis.major_gridlines is a MajorGridlines object")
def then_axis_major_gridlines_is_a_MajorGridlines_object(context):
    axis = context.axis
    assert type(axis.major_gridlines).__name__ == "MajorGridlines"


@then("axis.{major_or_minor}_unit is {value}")
def then_axis_major_or_minor_unit_is_value(context, major_or_minor, value):
    axis = context.axis
    propname = "%s_unit" % major_or_minor
    actual_value = getattr(axis, propname)
    expected_value = {"20.0": 20.0, "8.4": 8.4, "5.0": 5.0, "4.2": 4.2, "None": None}[
        value
    ]
    assert actual_value == expected_value, "got %s" % actual_value


@then("axis_title.format is a ChartFormat object")
def then_axis_title_format_is_a_ChartFormat_object(context):
    class_name = type(context.axis_title.format).__name__
    assert class_name == "ChartFormat", "got %s" % class_name


@then("axis_title.format.fill is a FillFormat object")
def then_axis_title_format_fill_is_a_FillFormat_object(context):
    class_name = type(context.axis_title.format.fill).__name__
    assert class_name == "FillFormat", "got %s" % class_name


@then("axis_title.format.line is a LineFormat object")
def then_axis_title_format_line_is_a_LineFormat_object(context):
    class_name = type(context.axis_title.format.line).__name__
    assert class_name == "LineFormat", "got %s" % class_name


@then("axis_title.has_text_frame is {value}")
def then_axis_title_has_text_frame_is_value(context, value):
    actual_value = context.axis_title.has_text_frame
    expected_value = {"True": True, "False": False}[value]
    assert actual_value is expected_value, "got %s" % actual_value


@then("axis_title.text_frame is a TextFrame object")
def then_axis_title_text_frame_is_a_TextFrame_object(context):
    class_name = type(context.axis_title.text_frame).__name__
    assert class_name == "TextFrame", "got %s" % class_name


@then("gridlines.format is a ChartFormat object")
def then_gridlines_format_is_a_ChartFormat_object(context):
    gridlines = context.gridlines
    assert type(gridlines.format).__name__ == "ChartFormat"


@then("gridlines.format.fill is a FillFormat object")
def then_gridlines_format_fill_is_a_FillFormat_object(context):
    gridlines = context.gridlines
    assert type(gridlines.format.fill).__name__ == "FillFormat"


@then("gridlines.format.line is a LineFormat object")
def then_gridlines_format_line_is_a_LineFormat_object(context):
    gridlines = context.gridlines
    assert type(gridlines.format.line).__name__ == "LineFormat"


@then("tick_labels.offset is {value}")
def then_tick_labels_offset_is_expected_value(context, value):
    expected_value = int(value)
    tick_labels = context.tick_labels
    assert tick_labels.offset == expected_value, "got %s" % tick_labels.offset


@then("value_axis.crosses is {member}")
def then_value_axis_crosses_is_value(context, member):
    value_axis = context.value_axis
    expected_value = getattr(XL_AXIS_CROSSES, member)
    assert value_axis.crosses == expected_value, "got %s" % value_axis.crosses


@then("value_axis.crosses_at is {value}")
def then_value_axis_crosses_at_is_value(context, value):
    value_axis = context.value_axis
    expected_value = None if value == "None" else float(value)
    assert value_axis.crosses_at == expected_value, "got %s" % value_axis.crosses_at
