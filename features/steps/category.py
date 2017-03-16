# encoding: utf-8

"""Gherkin step implementations for chart category features."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

from behave import given, then

from pptx import Presentation

from helpers import test_pptx


# given ===================================================

@given('a Categories object containing 3 categories')
def given_a_Categories_object_containing_3_categories(context):
    prs = Presentation(test_pptx('cht-category-access'))
    context.categories = prs.slides[0].shapes[0].chart.plots[0].categories


@given('a Categories object having {count} category levels')
def given_a_Categories_object_having_count_category_levels(context, count):
    slide_idx = [2, 0, 3, 1][int(count)]
    slide = Presentation(test_pptx('cht-category-access')).slides[slide_idx]
    context.categories = slide.shapes[0].chart.plots[0].categories


@given('a Categories object having {leafs} categories and {levels} levels')
def given_a_Categories_obj_having_leafs_and_levels(context, leafs, levels):
    slide_idx = {
        (3, 1): 0,
        (8, 3): 1,
        (0, 0): 2,
        (4, 2): 3,
    }[(int(leafs), int(levels))]
    slide = Presentation(test_pptx('cht-category-access')).slides[slide_idx]
    context.categories = slide.shapes[0].chart.plots[0].categories


@given('a Category object having idx value {idx}')
def given_a_Category_object_having_idx_value_idx(context, idx):
    cat_offset = int(idx)
    slide = Presentation(test_pptx('cht-category-access')).slides[0]
    context.category = slide.shapes[0].chart.plots[0].categories[cat_offset]


@given('a Category object having {label}')
def given_a_Category_object_having_label(context, label):
    cat_offset = {'label \'Foo\'': 0, 'no label': 1}[label]
    slide = Presentation(test_pptx('cht-category-access')).slides[0]
    context.category = slide.shapes[0].chart.plots[0].categories[cat_offset]


@given('a CategoryLevel object containing 4 categories')
def given_a_CategoryLevel_object_containing_4_categories(context):
    slide = Presentation(test_pptx('cht-category-access')).slides[1]
    chart = slide.shapes[0].chart
    context.category_level = chart.plots[0].categories.levels[1]


# then ====================================================

@then('categories[2] is a Category object')
def then_categories_2_is_a_Category_object(context):
    type_name = type(context.categories[2]).__name__
    assert type_name == 'Category', 'got %s' % type_name


@then('categories.depth is {value}')
def then_categories_depth_is_value(context, value):
    expected_value = int(value)
    depth = context.categories.depth
    assert depth == expected_value, 'got %s' % expected_value


@then('categories.flattened_labels is a tuple of {leafs} tuples')
def then_categories_flattened_labels_is_tuple_of_tuples(context, leafs):
    flattened_labels = context.categories.flattened_labels
    type_name = type(flattened_labels).__name__
    length = len(flattened_labels)
    assert type_name == 'tuple', 'got %s' % type_name
    assert length == int(leafs), 'got %s' % length
    for labels in flattened_labels:
        type_name = type(labels).__name__
        assert type_name == 'tuple', 'got %s' % type_name


@then('categories.levels contains {count} CategoryLevel objects')
def then_categories_levels_contains_count_CategoryLevel_objs(context, count):
    expected_idx = int(count) - 1
    idx = -1
    for idx, category_level in enumerate(context.categories.levels):
        type_name = type(category_level).__name__
        assert type_name == 'CategoryLevel', 'got %s' % type_name
    assert idx == expected_idx, 'got %s' % idx


@then('category.idx is {value}')
def then_category_idx_is_value(context, value):
    expected_value = None if value == 'None' else int(value)
    idx = context.category.idx
    assert idx == expected_value, 'got %s' % idx


@then('category.label is {value}')
def then_category_label_is_value(context, value):
    expected_value = {'\'Foo\'': 'Foo', '\'\'': ''}[value]
    label = context.category.label
    assert label == expected_value, 'got %s' % label


@then('category_level[2] is a Category object')
def then_category_level_2_is_a_Category_object(context):
    type_name = type(context.category_level[2]).__name__
    assert type_name == 'Category', 'got %s' % type_name


@then('each label tuple contains {levels} labels')
def then_each_label_tuple_contains_levels_labels(context, levels):
    flattened_labels = context.categories.flattened_labels
    for labels in flattened_labels:
        length = len(labels)
        assert length == int(levels), 'got %s' % levels
        for label in labels:
            type_name = type(label).__name__
            assert type_name == 'str', 'got %s' % type_name


@then('iterating categories produces 3 Category objects')
def then_iterating_categories_produces_3_category_objects(context):
    categories = context.categories
    idx = -1
    for idx, category in enumerate(categories):
        assert type(category).__name__ == 'Category'
    assert idx == 2, 'got %s' % idx


@then('iterating category_level produces 4 Category objects')
def then_iterating_category_level_produces_4_Category_objects(context):
    idx = -1
    for idx, category in enumerate(context.category_level):
        type_name = type(category).__name__
        assert type_name == 'Category', 'got %s' % type_name
    assert idx == 3, 'got %s' % idx


@then('len(categories) is {count}')
def then_len_categories_is_count(context, count):
    expected_count = int(count)
    assert len(context.categories) == expected_count


@then('len(category_level) is {count}')
def then_len_category_level_is_count(context, count):
    expected_count = int(count)
    actual_count = len(context.category_level)
    assert actual_count == expected_count, 'got %s' % actual_count


@then('list(categories) == [\'Foo\', \'\', \'Baz\']')
def then_list_categories_is_Foo_empty_Baz(context):
    cats_list = list(context.categories)
    assert cats_list == ['Foo', '', 'Baz'], 'got %s' % cats_list
