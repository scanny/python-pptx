# encoding: utf-8

"""
Gherkin step implementations for chart features.
"""

from __future__ import absolute_import, print_function

import hashlib

from itertools import islice

from behave import given, then, when

from pptx import Presentation
from pptx.chart.chart import Legend
from pptx.chart.data import (
    BubbleChartData, CategoryChartData, ChartData, XyChartData
)
from pptx.enum.chart import XL_CHART_TYPE
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.util import Inches

from helpers import count, test_pptx


# given ===================================================

@given('a Chart object as chart')
def given_a_Chart_object_as_chart(context):
    slide = Presentation(test_pptx('shp-common-props')).slides[0]
    context.chart = slide.shapes[6].chart


@given('a chart having {a_or_no} title')
def given_a_chart_having_a_or_no_title(context, a_or_no):
    shape_idx = {'no': 0, 'a': 1}[a_or_no]
    prs = Presentation(test_pptx('cht-chart-props'))
    context.chart = prs.slides[0].shapes[shape_idx].chart


@given('a chart {having_or_not} a legend')
def given_a_chart_having_or_not_a_legend(context, having_or_not):
    slide_idx = {
        'having':     0,
        'not having': 1,
    }[having_or_not]
    prs = Presentation(test_pptx('cht-legend'))
    context.chart = prs.slides[slide_idx].shapes[0].chart


@given('a chart of size and type {spec}')
def given_a_chart_of_size_and_type_spec(context, spec):
    slide_idx = {
        '2x2 Clustered Bar':    0,
        '2x2 100% Stacked Bar': 1,
        '2x2 Clustered Column': 2,
        '4x3 Line':             3,
        '3x1 Pie':              4,
        '3x2 XY':               5,
        '3x2 Bubble':           6,
    }[spec]
    prs = Presentation(test_pptx('cht-replace-data'))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.chart = chart
    context.xlsx_sha1 = hashlib.sha1(
        chart._workbook.xlsx_part.blob
    ).hexdigest()


@given('a chart of type {chart_type}')
def given_a_chart_of_type_chart_type(context, chart_type):
    slide_idx, shape_idx = {
        'Area':                        (0, 0),
        'Stacked Area':                (0, 1),
        '100% Stacked Area':           (0, 2),
        '3-D Area':                    (0, 3),
        '3-D Stacked Area':            (0, 4),
        '3-D 100% Stacked Area':       (0, 5),
        'Clustered Bar':               (1, 0),
        'Stacked Bar':                 (1, 1),
        '100% Stacked Bar':            (1, 2),
        'Clustered Column':            (1, 3),
        'Stacked Column':              (1, 4),
        '100% Stacked Column':         (1, 5),
        'Line':                        (2, 0),
        'Stacked Line':                (2, 1),
        '100% Stacked Line':           (2, 2),
        'Marked Line':                 (2, 3),
        'Stacked Marked Line':         (2, 4),
        '100% Stacked Marked Line':    (2, 5),
        'Pie':                         (3, 0),
        'Exploded Pie':                (3, 1),
        'XY (Scatter)':                (4, 0),
        'XY Lines':                    (4, 1),
        'XY Lines No Markers':         (4, 2),
        'XY Smooth Lines':             (4, 3),
        'XY Smooth No Markers':        (4, 4),
        'Bubble':                      (5, 0),
        '3D-Bubble':                   (5, 1),
        'Radar':                       (6, 0),
        'Marked Radar':                (6, 1),
        'Filled Radar':                (6, 2),
        'Line (with date categories)': (7, 0),
    }[chart_type]
    prs = Presentation(test_pptx('cht-chart-type'))
    context.chart = prs.slides[slide_idx].shapes[shape_idx].chart


@given('a chart title')
def given_a_chart_title(context):
    prs = Presentation(test_pptx('cht-chart-props'))
    context.chart_title = prs.slides[0].shapes[1].chart.chart_title


@given('a chart title having {a_or_no} text frame')
def given_a_chart_title_having_a_or_no_text_frame(context, a_or_no):
    prs = Presentation(test_pptx('cht-chart-props'))
    shape_idx = {'no': 0, 'a': 1}[a_or_no]
    context.chart_title = prs.slides[1].shapes[shape_idx].chart.chart_title


# when ====================================================

@when('I add a Clustered bar chart with multi-level categories')
def when_I_add_a_clustered_bar_chart_with_multi_level_categories(context):
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    chart_data = CategoryChartData()

    WEST = chart_data.add_category('WEST')
    WEST.add_sub_category('SF')
    WEST.add_sub_category('LA')
    EAST = chart_data.add_category('EAST')
    EAST.add_sub_category('NY')
    EAST.add_sub_category('NJ')

    chart_data.add_series('Series 1', (1, 2, None, 4))
    chart_data.add_series('Series 2', (5, None, 7, 8))

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I add a {kind} chart with {cats} categories and {sers} series')
def when_I_add_a_chart_with_categories_and_series(context, kind, cats, sers):
    chart_type = {
        'Area':                      XL_CHART_TYPE.AREA,
        'Stacked Area':              XL_CHART_TYPE.AREA_STACKED,
        '100% Stacked Area':         XL_CHART_TYPE.AREA_STACKED_100,
        'Clustered Bar':             XL_CHART_TYPE.BAR_CLUSTERED,
        'Stacked Bar':               XL_CHART_TYPE.BAR_STACKED,
        '100% Stacked Bar':          XL_CHART_TYPE.BAR_STACKED_100,
        'Clustered Column':          XL_CHART_TYPE.COLUMN_CLUSTERED,
        'Stacked Column':            XL_CHART_TYPE.COLUMN_STACKED,
        '100% Stacked Column':       XL_CHART_TYPE.COLUMN_STACKED_100,
        'Doughnut':                  XL_CHART_TYPE.DOUGHNUT,
        'Exploded Doughnut':         XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        'Line':                      XL_CHART_TYPE.LINE,
        'Line with Markers':         XL_CHART_TYPE.LINE_MARKERS,
        'Line Markers Stacked':      XL_CHART_TYPE.LINE_MARKERS_STACKED,
        '100% Line Markers Stacked': XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
        'Line Stacked':              XL_CHART_TYPE.LINE_STACKED,
        '100% Line Stacked':         XL_CHART_TYPE.LINE_STACKED_100,
        'Pie':                       XL_CHART_TYPE.PIE,
        'Exploded Pie':              XL_CHART_TYPE.PIE_EXPLODED,
        'Radar':                     XL_CHART_TYPE.RADAR,
        'Filled Radar':              XL_CHART_TYPE.RADAR_FILLED,
        'Radar with markers':        XL_CHART_TYPE.RADAR_MARKERS,
    }[kind]
    category_count, series_count = int(cats), int(sers)
    category_source = ('Foo', 'Bar', 'Baz', 'Boo', 'Far', 'Faz')
    series_value_source = count(1.1, 1.1)

    chart_data = CategoryChartData()
    chart_data.categories = category_source[:category_count]
    for idx in range(series_count):
        series_title = 'Series %d' % (idx+1)
        series_values = tuple(islice(series_value_source, category_count))
        chart_data.add_series(series_title, series_values)

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I add a {bubble_type} chart having 2 series of 3 points each')
def when_I_add_a_bubble_chart_having_2_series_of_3_pts(context, bubble_type):
    chart_type = getattr(XL_CHART_TYPE, bubble_type)
    data = (
        ('Series 1', ((-0.1, 0.5, 1.0), (16.2, 0.0, 2.0), (8.0, -0.2, 3.0))),
        ('Series 2', ((12.4, 0.8, 4.0), (-7.5, 0.5, 5.0), (5.1, -0.5, 6.0))),
    )

    chart_data = BubbleChartData()

    for series_data in data:
        series_label, points = series_data
        series = chart_data.add_series(series_label)
        for point in points:
            x, y, size = point
            series.add_data_point(x, y, size)

    context.chart = context.slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    ).chart


@when('I assign {value} to chart.has_legend')
def when_I_assign_value_to_chart_has_legend(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.chart.has_legend = new_value


@when('I assign {value} to chart.has_title')
def when_I_assign_value_to_chart_has_title(context, value):
    context.chart.has_title = {'True': True, 'False': False}[value]


@when('I assign {value} to chart_title.has_text_frame')
def when_I_assign_value_to_chart_title_has_text_frame(context, value):
    context.chart_title.has_text_frame = {
        'True':  True,
        'False': False
    }[value]


@when('I replace its data with {cats} categories and {sers} series')
def when_I_replace_its_data_with_categories_and_series(context, cats, sers):
    category_count, series_count = int(cats), int(sers)
    category_source = ('Foo', 'Bar', 'Baz', 'Boo', 'Far', 'Faz')
    series_value_source = count(1.1, 1.1)

    chart_data = ChartData()
    chart_data.categories = category_source[:category_count]
    for idx in range(series_count):
        series_title = 'New Series %d' % (idx+1)
        series_values = tuple(islice(series_value_source, category_count))
        chart_data.add_series(series_title, series_values)

    context.chart.replace_data(chart_data)


@when('I replace its data with 3 series of 3 bubble points each')
def when_I_replace_its_data_with_3_series_of_three_bubble_pts_each(context):
    chart_data = BubbleChartData()
    for idx in range(3):
        series_title = 'New Series %d' % (idx+1)
        series = chart_data.add_series(series_title)
        for jdx in range(3):
            x, y, size = idx * 3 + jdx, idx * 2 + jdx, idx + jdx
            series.add_data_point(x, y, size)

    context.chart.replace_data(chart_data)


@when('I replace its data with 3 series of 3 points each')
def when_I_replace_its_data_with_3_series_of_three_points_each(context):
    chart_data = XyChartData()
    x = y = 0
    for idx in range(3):
        series_title = 'New Series %d' % (idx+1)
        series = chart_data.add_series(series_title)
        for jdx in range(3):
            x, y = idx * 3 + jdx, idx * 2 + jdx
            series.add_data_point(x, y)

    context.chart.replace_data(chart_data)


# then ====================================================

@then('chart.category_axis is a {cls_name} object')
def then_chart_category_axis_is_a_cls_name_object(context, cls_name):
    category_axis = context.chart.category_axis
    type_name = type(category_axis).__name__
    assert type_name == cls_name, 'got %s' % type_name


@then('chart.chart_title is a ChartTitle object')
def then_chart_chart_title_is_a_ChartTitle_object(context):
    class_name = type(context.chart.chart_title).__name__
    assert class_name == 'ChartTitle', 'got %s' % class_name


@then('chart.chart_type is {enum_member}')
def then_chart_chart_type_is_value(context, enum_member):
    expected_value = getattr(XL_CHART_TYPE, enum_member)
    chart = context.chart
    assert chart.chart_type is expected_value, 'got %s' % chart.chart_type


@then('chart.font is a Font object')
def then_chart_font_is_a_Font_object(context):
    actual = type(context.chart.font).__name__
    expected = 'Font'
    assert actual == expected, 'chart.font is a %s object' % actual


@then('chart.has_legend is {value}')
def then_chart_has_legend_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    chart = context.chart
    assert chart.has_legend is expected_value


@then('chart.has_title is {value}')
def then_chart_has_title_is_value(context, value):
    chart = context.chart
    actual_value = chart.has_title
    expected_value = {'True': True, 'False': False}[value]
    assert actual_value is expected_value, 'got %s' % actual_value


@then('chart.legend is a legend object')
def then_chart_legend_is_a_legend_object(context):
    chart = context.chart
    assert isinstance(chart.legend, Legend)


@then('chart.series is a SeriesCollection object')
def then_chart_series_is_a_SeriesCollection_object(context):
    type_name = type(context.chart.series).__name__
    assert type_name == 'SeriesCollection', 'got %s' % type_name


@then('chart.value_axis is a ValueAxis object')
def then_chart_value_axis_is_a_ValueAxis_object(context):
    value_axis = context.chart.value_axis
    assert type(value_axis).__name__ == 'ValueAxis'


@then('chart_title.format is a ChartFormat object')
def then_chart_title_format_is_a_ChartFormat_object(context):
    class_name = type(context.chart_title.format).__name__
    assert class_name == 'ChartFormat', 'got %s' % class_name


@then('chart_title.format.fill is a FillFormat object')
def then_chart_title_format_fill_is_a_FillFormat_object(context):
    class_name = type(context.chart_title.format.fill).__name__
    assert class_name == 'FillFormat', 'got %s' % class_name


@then('chart_title.format.line is a LineFormat object')
def then_chart_title_format_line_is_a_LineFormat_object(context):
    class_name = type(context.chart_title.format.line).__name__
    assert class_name == 'LineFormat', 'got %s' % class_name


@then('chart_title.has_text_frame is {value}')
def then_chart_title_has_text_frame_is_value(context, value):
    actual_value = context.chart_title.has_text_frame
    expected_value = {'True': True, 'False': False}[value]
    assert actual_value is expected_value, 'got %s' % actual_value


@then('chart_title.text_frame is a TextFrame object')
def then_chart_title_text_frame_is_a_TextFrame_object(context):
    class_name = type(context.chart_title.text_frame).__name__
    assert class_name == 'TextFrame', 'got %s' % class_name


@then('each series has a new name')
def then_each_series_has_a_new_name(context):
    for series in context.chart.plots[0].series:
        assert series.name.startswith('New ')


@then('each series has {count} values')
def then_each_series_has_count_values(context, count):
    expected_count = int(count)
    for series in context.chart.plots[0].series:
        actual_value_count = len(series.values)
        assert actual_value_count == expected_count


@then('len(chart.series) is {count}')
def then_len_chart_series_is_count(context, count):
    expected_count = int(count)
    assert len(context.chart.series) == expected_count


@then('the chart has an Excel data worksheet')
def then_the_chart_has_an_Excel_data_worksheet(context):
    xlsx_part = context.chart._workbook.xlsx_part
    assert isinstance(xlsx_part, EmbeddedXlsxPart)


@then('the chart has new chart data')
def then_the_chart_has_new_chart_data(context):
    orig_xlsx_sha1 = context.xlsx_sha1
    new_xlsx_sha1 = hashlib.sha1(
        context.chart._workbook.xlsx_part.blob
    ).hexdigest()
    assert new_xlsx_sha1 != orig_xlsx_sha1
