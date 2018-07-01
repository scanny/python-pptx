# encoding: utf-8

"""Gherkin step implementations for ColorFormat-related features."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

from behave import given, when, then

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

from helpers import test_pptx


# given ====================================================

@given('a ColorFormat object as color')
def given_a_ColorFormat_object_as_color(context):
    shape = Presentation(test_pptx('dml-fill')).slides[0].shapes[2]
    color = shape.fill.fore_color
    context.color = color


# when =====================================================

@when("I assign MSO_THEME_COLOR.ACCENT_6 to color.theme_color")
def when_assign_MSO_THEME_COLOR_ACCENT_6_to_color_theme_color(context):
    context.color.theme_color = MSO_THEME_COLOR.ACCENT_6


@when("I assign RGBColor(12, 34, 56) to color.rgb")
def when_I_assign_RGBColor_to_color_rgb(context):
    context.color.rgb = RGBColor(12, 34, 56)


@when("I assign 0.42 to color.brightness")
def when_I_assign_0_42_to_color_brightness(context):
    context.color.brightness = 0.42


# then =====================================================

@then('color.brightness is 0.42')
def then_color_brightness_is_value(context):
    brightness = context.color.brightness
    expected_value = 0.42
    assert brightness == expected_value, (
        'expected %s, got %s' % (expected_value, brightness)
    )


@then('color.rgb is RGBColor(12, 34, 56)')
def then_color_rgb_is_RGBColor_12_34_56(context):
    rgb = context.color.rgb
    expected_value = RGBColor(12, 34, 56)
    assert rgb == expected_value, (
        'expected %s, got %s' % (repr(expected_value), repr(rgb))
    )


@then('color.theme_color is MSO_THEME_COLOR.ACCENT_6')
def then_color_theme_color_is_MSO_THEME_COLOR_ACCENT_6(context):
    theme_color = context.color.theme_color
    expected_value = MSO_THEME_COLOR.ACCENT_6
    assert theme_color == expected_value, (
        'expected %s, got %s' % (expected_value, theme_color)
    )
