# encoding: utf-8

"""Gherkin step implementations for chart data label features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation
from pptx.enum.chart import XL_DATA_LABEL_POSITION

from helpers import test_pptx


# given ===================================================


@given("a DataLabels object {showing_or_not} category-name as data_labels")
def given_a_DataLabels_object_showing_or_not_cat_name(context, showing_or_not):
    series_idx = {"not showing": 0, "showing": 1}[showing_or_not]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[2].shapes[0].chart
    context.data_labels = chart.plots[0].series[series_idx].data_labels


@given("a DataLabels object {showing_or_not} legend-key as data_labels")
def given_a_DataLabels_object_showing_or_not_leg_key(context, showing_or_not):
    series_idx = {"not showing": 0, "showing": 1}[showing_or_not]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[2].shapes[0].chart
    context.data_labels = chart.plots[0].series[series_idx].data_labels


@given("a DataLabels object {showing_or_not} percentage as data_labels")
def given_a_DataLabels_object_showing_or_not_percent(context, showing_or_not):
    slide_idx = {"not showing": 4, "showing": 3}[showing_or_not]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.data_labels = chart.plots[0].series[0].data_labels


@given("a DataLabels object {showing_or_not} series-name as data_labels")
def given_a_DataLabels_object_showing_or_not_ser_name(context, showing_or_not):
    series_idx = {"not showing": 0, "showing": 1}[showing_or_not]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[2].shapes[0].chart
    context.data_labels = chart.plots[0].series[series_idx].data_labels


@given("a DataLabels object {showing_or_not} value as data_labels")
def given_a_DataLabels_object_showing_or_not_value(context, showing_or_not):
    series_idx = {"not showing": 0, "showing": 1}[showing_or_not]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[2].shapes[0].chart
    context.data_labels = chart.plots[0].series[series_idx].data_labels


@given("a DataLabels object with {pos} position as data_labels")
def given_a_DataLabels_object_with_pos_position(context, pos):
    slide_idx = {"inherited": 0, "inside-base": 1}[pos]
    prs = Presentation(test_pptx("cht-datalabels"))
    chart = prs.slides[slide_idx].shapes[0].chart
    context.data_labels = chart.plots[0].data_labels


@given("a data label {having_or_not} custom font as data_label")
def given_a_data_label_having_or_not_custom_font(context, having_or_not):
    point_idx = {"having a": 0, "having no": 1}[having_or_not]
    prs = Presentation(test_pptx("cht-point-props"))
    points = prs.slides[2].shapes[0].chart.plots[0].series[0].points
    context.data_label = points[point_idx].data_label


@given("a data label {having_or_not} custom text as data_label")
def given_a_data_label_having_or_not_custom_text(context, having_or_not):
    point_idx = {"having": 0, "having no": 1}[having_or_not]
    prs = Presentation(test_pptx("cht-point-props"))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.data_label = plot.series[0].points[point_idx].data_label


@given("a data label with {pos} position as data_label")
def given_a_data_label_with_pos_position_as_data_label(context, pos):
    point_idx = {"inherited": 0, "centered": 1, "below": 2}[pos]
    prs = Presentation(test_pptx("cht-point-props"))
    plot = prs.slides[1].shapes[0].chart.plots[0]
    context.data_label = plot.series[0].points[point_idx].data_label


# when ====================================================


@when("I assign {value} to data_label.has_text_frame")
def when_I_assign_value_to_data_label_has_text_frame(context, value):
    new_value = {"True": True, "False": False}[value]
    context.data_label.has_text_frame = new_value


@when("I assign {value} to data_label.position")
def when_I_assign_value_to_data_label_position(context, value):
    new_value = None if value == "None" else getattr(XL_DATA_LABEL_POSITION, value)
    context.data_label.position = new_value


@when("I assign {value} to data_labels.position")
def when_I_assign_value_to_data_labels_position(context, value):
    new_value = None if value == "None" else getattr(XL_DATA_LABEL_POSITION, value)
    context.data_labels.position = new_value


@when("I assign {value} to data_labels.show_category_name")
def when_I_assign_value_to_data_labels_show_category_name(context, value):
    context.data_labels.show_category_name = eval(value)


@when("I assign {value} to data_labels.show_legend_key")
def when_I_assign_value_to_data_labels_show_legend_key(context, value):
    context.data_labels.show_legend_key = eval(value)


@when("I assign {value} to data_labels.show_percentage")
def when_I_assign_value_to_data_labels_show_percentage(context, value):
    context.data_labels.show_percentage = eval(value)


@when("I assign {value} to data_labels.show_series_name")
def when_I_assign_value_to_data_labels_show_series_name(context, value):
    context.data_labels.show_series_name = eval(value)


@when("I assign {value} to data_labels.show_value")
def when_I_assign_value_to_data_labels_show_value(context, value):
    context.data_labels.show_value = eval(value)


# then ====================================================


@then("data_label.font is a Font object")
def then_data_label_font_is_a_Font_object(context):
    font = context.data_label.font
    assert type(font).__name__ == "Font"


@then("data_label.has_text_frame is {value}")
def then_data_label_has_text_frame_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    data_label = context.data_label
    assert data_label.has_text_frame is expected_value


@then("data_label.position is {value}")
def then_data_label_position_is_value(context, value):
    expected_value = None if value == "None" else getattr(XL_DATA_LABEL_POSITION, value)
    data_label = context.data_label
    assert data_label.position is expected_value, "got %s" % data_label.position


@then("data_label.text_frame is a TextFrame object")
def then_data_label_text_frame_is_a_TextFrame_object(context):
    text_frame = context.data_label.text_frame
    assert type(text_frame).__name__ == "TextFrame"


@then("data_labels.position is {value}")
def then_data_labels_position_is_value(context, value):
    expected_value = None if value == "None" else getattr(XL_DATA_LABEL_POSITION, value)
    data_labels = context.data_labels
    assert data_labels.position is expected_value, "got %s" % data_labels.position


@then("data_labels.show_category_name is {value}")
def then_data_labels_show_category_name_is_value(context, value):
    actual, expected = context.data_labels.show_category_name, eval(value)
    assert actual is expected, "data_labels.show_category_name is %s" % actual


@then("data_labels.show_legend_key is {value}")
def then_data_labels_show_legend_key_is_value(context, value):
    actual, expected = context.data_labels.show_legend_key, eval(value)
    assert actual is expected, "data_labels.show_legend_key is %s" % actual


@then("data_labels.show_percentage is {value}")
def then_data_labels_show_percentage_is_value(context, value):
    actual, expected = context.data_labels.show_percentage, eval(value)
    assert actual is expected, "data_labels.show_percentage is %s" % actual


@then("data_labels.show_series_name is {value}")
def then_data_labels_show_series_name_is_value(context, value):
    actual, expected = context.data_labels.show_series_name, eval(value)
    assert actual is expected, "data_labels.show_series_name is %s" % actual


@then("data_labels.show_value is {value}")
def then_data_labels_show_value_is_value(context, value):
    actual, expected = context.data_labels.show_value, eval(value)
    assert actual is expected, "data_labels.show_value is %s" % actual
