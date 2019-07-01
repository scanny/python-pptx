# encoding: utf-8

"""Gherkin step implementations for ShadowFormat-related features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation

from helpers import test_pptx


# given ====================================================


@given("a ShadowFormat object that {inherits} as shadow")
def given_a_ShadowFormat_object_that_inherits_or_not(context, inherits):
    shape_idx = {"inherits": 0, "does not inherit": 1}[inherits]
    shape = Presentation(test_pptx("dml-effect")).slides[0].shapes[shape_idx]
    context.shadow = shape.shadow


# when =====================================================


@when("I assign {value} to shadow.inherit")
def when_I_assign_value_to_shadow_inherit(context, value):
    context.shadow.inherit = eval(value)


# then =====================================================


@then("shadow.inherit is {bool_str}")
def then_shadow_inherit_is_bool_val(context, bool_str):
    expected_value = eval(bool_str)
    actual_value = context.shadow.inherit
    assert actual_value is expected_value, "shadow.inherit is %s" % actual_value
