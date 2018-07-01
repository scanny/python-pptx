# encoding: utf-8

"""Gherkin step implementations for FillFormat-related features."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

from behave import given, then, when

from pptx import Presentation
from pptx.enum.dml import MSO_FILL, MSO_PATTERN  # noqa

from helpers import test_pptx


# given ====================================================

@given('a FillFormat object as fill')
def given_a_FillFormat_object_as_fill(context):
    fill = Presentation(test_pptx('dml-fill')).slides[0].shapes[0].fill
    context.fill = fill


@given('a FillFormat object as fill having {pattern} fill')
def given_a_FillFormat_object_as_fill_having_pattern(context, pattern):
    shape_idx = {
        'no pattern':        0,
        'MSO_PATTERN.DIVOT': 1,
        'MSO_PATTERN.WAVE':  2,
    }[pattern]
    slide = Presentation(test_pptx('dml-fill')).slides[1]
    fill = slide.shapes[shape_idx].fill
    context.fill = fill


@given('{type} FillFormat object as fill')
def given_type_FillFormat_object_as_fill(context, type):
    shape_idx = {
        'an inheriting': 0, 'a no-fill': 1, 'a solid': 2, 'a picture': 3,
        'a gradient': 4, 'a patterned': 5,
    }[type]
    shape = Presentation(test_pptx('dml-fill')).slides[0].shapes[shape_idx]
    context.fill = shape.fill


@given('a _GradientStop object as stop')
def given_a_GradientStop_object_as_stop(context):
    shape = Presentation(test_pptx('dml-fill')).slides[0].shapes[4]
    context.stop = shape.fill.gradient_stops[0]


# when =====================================================

@when("I assign {value} to fill.gradient_angle")
def when_I_assign_value_to_fill_gradient_angle(context, value):
    context.fill.gradient_angle = eval(value)


@when("I assign {value} to fill.pattern")
def when_I_assign_value_to_fill_pattern(context, value):
    pattern = {
        'None':              None,
        'MSO_PATTERN.CROSS': MSO_PATTERN.CROSS,
        'MSO_PATTERN.DIVOT': MSO_PATTERN.DIVOT,
        'MSO_PATTERN.WAVE':  MSO_PATTERN.WAVE,
    }[value]
    context.fill.pattern = pattern


@when("I assign {value} to stop.position")
def when_I_assign_value_to_stop_position(context, value):
    context.stop.position = eval(value)


@when("I call fill.background()")
def when_I_call_fill_background(context):
    context.fill.background()


@when("I call fill.gradient()")
def when_I_call_fill_gradient(context):
    context.fill.gradient()


@when("I call fill.patterned()")
def when_I_call_fill_patterned(context):
    context.fill.patterned()


@when("I call fill.solid()")
def when_I_call_fill_solid(context):
    context.fill.solid()


# then =====================================================

@then('fill.back_color is a ColorFormat object')
def then_fill_back_color_is_a_ColorFormat_object(context):
    actual_value = context.fill.back_color.__class__.__name__
    expected_value = 'ColorFormat'
    assert actual_value == expected_value, (
        'fill.back_color is a %s object' % actual_value
    )


@then('fill.fore_color is a ColorFormat object')
def then_fill_fore_color_is_a_ColorFormat_object(context):
    actual_value = context.fill.fore_color.__class__.__name__
    expected_value = 'ColorFormat'
    assert actual_value == expected_value, (
        'fill.fore_color is a %s object' % actual_value
    )


@then('fill.gradient_angle == {value}')
def then_fill_gradient_angle_eq_value(context, value):
    expected_value = round(eval(value), 2)
    actual_value = round(context.fill.gradient_angle, 2)
    assert actual_value == expected_value, (
        'fill.gradient_angle == %s' % actual_value
    )


@then('fill.gradient_stops is a _GradientStops object')
def then_fill_gradient_stops_is_a_GradientStops_object(context):
    expected_value = '_GradientStops'
    actual_value = context.fill.gradient_stops.__class__.__name__
    assert actual_value == expected_value, (
        'fill.gradient_stops is a %s object' % actual_value
    )


@then('fill.pattern is {value}')
def then_fill_pattern_is_value(context, value):
    fill_pattern = context.fill.pattern
    expected_value = {
        'None':              None,
        'MSO_PATTERN.CROSS': MSO_PATTERN.CROSS,
        'MSO_PATTERN.DIVOT': MSO_PATTERN.DIVOT,
        'MSO_PATTERN.WAVE':  MSO_PATTERN.WAVE,
    }[value]
    assert fill_pattern == expected_value, (
        'expected fill pattern %s, got %s' % (expected_value, fill_pattern)
    )


@then('fill.type == {value}')
def then_fill_type_eq_value(context, value):
    expected_value = eval(value)
    actual_value = context.fill.type
    assert actual_value == expected_value, ('fill.type is %s' % actual_value)


@then('stop.color is a ColorFormat object')
def then_stop_color_is_a_ColorFormat_object(context):
    expected_value = 'ColorFormat'
    actual_value = context.stop.color.__class__.__name__
    assert actual_value == expected_value, (
        'stop.color is a %s object' % actual_value
    )


@then('stop.position == {value}')
def then_stop_position_eq_value(context, value):
    expected_value = eval(value)
    actual_value = context.stop.position
    assert actual_value == expected_value, (
        'stop.position == %s' % actual_value
    )
