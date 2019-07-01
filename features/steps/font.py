# encoding: utf-8

"""
Step implementations for run property (font)-related features
"""

from __future__ import absolute_import

from behave import given, then, when

from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_UNDERLINE

from helpers import test_pptx


# given ===================================================


@given("a font")
def given_a_font(context):
    prs = Presentation(test_pptx("txt-font-props"))
    slide = prs.slides[1]
    textbox = slide.shapes[0]
    run = textbox.text_frame.paragraphs[0].runs[0]
    context.font = run.font


@given("a font having language id {value}")
def given_a_font_having_language_id_value(context, value):
    shape_idx = {
        "of no explicit setting": 0,
        "MSO_LANGUAGE_ID.FRENCH": 1,
        "MSO_LANGUAGE_ID.POLISH": 2,
    }[value]
    prs = Presentation(test_pptx("txt-font-props"))
    textbox = prs.slides[4].shapes[shape_idx]
    run = textbox.text_frame.paragraphs[0].runs[0]
    context.font = run.font


@given("a font having size of {value}")
def given_a_font_having_size_of_value(context, value):
    shape_idx = {"no explicit value": 0, "42pt": 1}[value]
    prs = Presentation(test_pptx("txt-font-props"))
    slide = prs.slides[1]
    textbox = slide.shapes[shape_idx]
    run = textbox.text_frame.paragraphs[0].runs[0]
    context.font = run.font


@given("a font with bold set {state}")
def given_a_font_with_bold_set_state(context, state):
    shape_idx = ["on", "off", "to inherit"].index(state)
    prs = Presentation(test_pptx("txt-font-props"))
    paragraph = prs.slides[2].shapes[shape_idx].text_frame.paragraphs[0]
    context.font = paragraph.runs[0].font


@given("a font with italic set {state}")
def given_run_with_italic_set_to_state(context, state):
    run_idx = ["on", "off", "to inherit"].index(state)
    prs = Presentation(test_pptx("txt-font-props"))
    runs = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs
    context.font = runs[run_idx].font


@given("a font with underline set {state}")
def given_run_with_underline_set_to_state(context, state):
    run_idx = ["on", "off", "to inherit", "to DOUBLE_LINE", "to WAVY_LINE"].index(state)
    prs = Presentation(test_pptx("txt-font-props"))
    runs = prs.slides[3].shapes[0].text_frame.paragraphs[0].runs
    print(runs[run_idx]._r.xml)
    context.font = runs[run_idx].font


# when ===================================================


@when("I assign {value} to font.bold")
def when_I_assign_value_to_font_bold(context, value):
    new_value = {"True": True, "False": False, "None": None}[value]
    context.font.bold = new_value


@when("I assign {value} to font.italic")
def when_I_assign_value_to_font_italic(context, value):
    new_value = {"True": True, "False": False, "None": None}[value]
    context.font.italic = new_value


@when("I assign {value} to font.language_id")
def when_I_assign_value_to_font_language_id(context, value):
    new_value = None if value == "None" else getattr(MSO_LANGUAGE_ID, value[16:])
    context.font.language_id = new_value


@when("I assign {value} to font.underline")
def when_I_assign_value_to_font_underline(context, value):
    new_value = {
        "True": True,
        "False": False,
        "None": None,
        "DOUBLE_LINE": MSO_UNDERLINE.DOUBLE_LINE,
        "NONE": MSO_UNDERLINE.NONE,
        "SINGLE_LINE": MSO_UNDERLINE.SINGLE_LINE,
    }[value]
    context.font.underline = new_value


# then ===================================================


@then("font.bold is {value}")
def then_font_bold_is_value(context, value):
    expected_value = {"True": True, "False": False, "None": None}[value]
    font = context.font
    assert font.bold is expected_value


@then("font.italic is {value}")
def then_font_italic_is_value(context, value):
    expected_value = {"True": True, "False": False, "None": None}[value]
    font = context.font
    assert font.italic is expected_value


@then("font.language_id is MSO_LANGUAGE_ID.{member}")
def then_font_language_id_is_MSO_LANGUAGE_ID_(context, member):
    expected_value = getattr(MSO_LANGUAGE_ID, member)
    font = context.font
    assert font.language_id is expected_value


@then("font.size is {value_str}")
def then_font_size_is_value(context, value_str):
    expected_value = {"42.0 points": 42.0, "None": None}[value_str]
    font = context.font
    value = font.size if font.size is None else font.size.pt
    assert value == expected_value, "expected %s, got %s" % (expected_value, value)


@then("font.underline is {value}")
def then_font_underline_is_value(context, value):
    expected_value = {
        "True": True,
        "False": False,
        "None": None,
        "DOUBLE_LINE": MSO_UNDERLINE.DOUBLE_LINE,
        "SINGLE_LINE": MSO_UNDERLINE.SINGLE_LINE,
        "WAVY_LINE": MSO_UNDERLINE.WAVY_LINE,
    }[value]
    font = context.font
    print(font._rPr.xml)
    assert font.underline is expected_value, "got %s" % font.underline
