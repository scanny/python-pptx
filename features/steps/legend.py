# encoding: utf-8

"""Gherkin step implementations for chart legend features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.text.text import Font

from helpers import test_pptx


# given ===================================================


@given("a legend")
def given_a_legend(context):
    prs = Presentation(test_pptx("cht-legend-props"))
    context.legend = prs.slides[0].shapes[0].chart.legend


@given("a legend having horizontal offset of {value}")
def given_a_legend_having_horizontal_offset_of_value(context, value):
    slide_idx = {"none": 0, "-0.5": 1, "0.42": 2}[value]
    prs = Presentation(test_pptx("cht-legend-props"))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


@given("a legend positioned {location} the chart")
def given_a_legend_positioned_location_the_chart(context, location):
    slide_idx = {"at an unspecified location of": 0, "below": 1, "to the right of": 2}[
        location
    ]
    prs = Presentation(test_pptx("cht-legend-props"))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


@given("a legend with overlay setting of {setting}")
def given_a_legend_with_overlay_setting_of_setting(context, setting):
    slide_idx = {"no explicit setting": 0, "True": 1, "False": 2}[setting]
    prs = Presentation(test_pptx("cht-legend-props"))
    context.legend = prs.slides[slide_idx].shapes[0].chart.legend


# when ====================================================


@when("I assign {value} to legend.horz_offset")
def when_I_assign_value_to_legend_horz_offset(context, value):
    new_value = float(value)
    context.legend.horz_offset = new_value


@when("I assign {value} to legend.include_in_layout")
def when_I_assign_value_to_legend_include_in_layout(context, value):
    new_value = {"True": True, "False": False}[value]
    context.legend.include_in_layout = new_value


@when("I assign {value} to legend.position")
def when_I_assign_value_to_legend_position(context, value):
    enum_value = getattr(XL_LEGEND_POSITION, value)
    context.legend.position = enum_value


# then ====================================================


@then("legend.font is a Font object")
def then_legend_font_is_a_Font_object(context):
    legend = context.legend
    assert isinstance(legend.font, Font)


@then("legend.horz_offset is {value}")
def then_legend_horz_offset_is_value(context, value):
    expected_value = float(value)
    legend = context.legend
    assert legend.horz_offset == expected_value


@then("legend.include_in_layout is {value}")
def then_legend_include_in_layout_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    legend = context.legend
    assert legend.include_in_layout is expected_value


@then("legend.position is {value}")
def then_legend_position_is_value(context, value):
    expected_position = getattr(XL_LEGEND_POSITION, value)
    legend = context.legend
    assert legend.position is expected_position, "got %s" % legend.position
