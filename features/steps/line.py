# encoding: utf-8

"""Step implementations for LineFormat-related features."""

from __future__ import absolute_import, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation
from pptx.enum.dml import MSO_LINE
from pptx.util import Length, Pt

from helpers import test_pptx


# given ===================================================


@given("a LineFormat object as line")
def given_a_LineFormat_object_as_line(context):
    line = Presentation(test_pptx("dml-line")).slides[0].shapes[0].line
    context.line = line


@given("a LineFormat object as line having {current} dash style")
def given_a_LineFormat_object_as_line_having_dash_style(context, current):
    shape_idx = {"no explicit": 0, "solid": 1, "dashed": 2, "dash-dot": 3}[current]
    shape = Presentation(test_pptx("dml-line")).slides[3].shapes[shape_idx]
    context.line = shape.line


@given("a LineFormat object as line having {line_width} width")
def given_a_LineFormat_object_as_line_having_width(context, line_width):
    shape_idx = {"no explicit": 0, "1 pt": 1}[line_width]
    prs = Presentation(test_pptx("dml-line"))
    shape = prs.slides[2].shapes[shape_idx]
    context.line = shape.line


# when ====================================================


@when("I assign {value_key} to line.dash_style")
def when_I_assign_value_to_line_dash_style(context, value_key):
    value = {
        "None": None,
        "MSO_LINE.DASH": MSO_LINE.DASH,
        "MSO_LINE.DASH_DOT": MSO_LINE.DASH_DOT,
        "MSO_LINE.SOLID": MSO_LINE.SOLID,
    }[value_key]
    context.line.dash_style = value


@when("I assign {line_width} to line.width")
def when_I_assign_value_to_line_width(context, line_width):
    value = {"None": None, "1 pt": Pt(1), "2.34 pt": Pt(2.34)}[line_width]
    context.line.width = value


# then ====================================================


@then("line.color is a ColorFormat object")
def then_line_color_is_a_ColorFormat_object(context):
    class_name = context.line.color.__class__.__name__
    expected_value = "ColorFormat"
    assert class_name == expected_value, "expected '%s', got '%s'" % (
        expected_value,
        class_name,
    )


@then("line.dash_style is {dash_style}")
def then_line_dash_style_is_value(context, dash_style):
    expected_value = {
        "None": None,
        "MSO_LINE.DASH": MSO_LINE.DASH,
        "MSO_LINE.DASH_DOT": MSO_LINE.DASH_DOT,
        "MSO_LINE.SOLID": MSO_LINE.SOLID,
    }[dash_style]
    actual_value = context.line.dash_style
    assert actual_value == expected_value, "expected %s, got %s" % (
        expected_value,
        actual_value,
    )


@then("line.fill is a FillFormat object")
def then_line_fill_is_a_FillFormat_object(context):
    class_name = context.line.fill.__class__.__name__
    expected_value = "FillFormat"
    assert class_name == expected_value, "expected '%s', got '%s'" % (
        expected_value,
        class_name,
    )


@then("line.width is {line_width}")
def then_line_width_is_value(context, line_width):
    expected_value = {"0": 0, "1 pt": Pt(1), "2.34 pt": Pt(2.34)}[line_width]
    line_width = context.line.width
    assert line_width == expected_value
    assert isinstance(line_width, Length)
