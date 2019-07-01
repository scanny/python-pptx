# encoding: utf-8

"""
Gherkin step implementations for picture-related features.
"""

from __future__ import absolute_import

from behave import given, when, then

from pptx import Presentation
from pptx.compat import BytesIO
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.package import Package
from pptx.util import Inches

from helpers import saved_pptx_path, test_image, test_pptx


# given ===================================================


@given("a Picture object masked by a {shape} as picture")
def given_a_picture_object_masked_by_shape_as_picture(context, shape):
    shape_idx = {"rectangle": 0, "circle": 1}[shape]
    prs = Presentation(test_pptx("shp-picture"))
    context.picture = prs.slides[1].shapes[shape_idx]


@given("a picture of known position and size")
def given_a_picture_of_known_position_and_size(context):
    prs = Presentation(test_pptx("shp-pos-and-size"))
    context.picture = prs.slides[1].shapes[0]


# when ====================================================


@when("I add the image {filename} using shapes.add_picture()")
def when_I_add_the_image_filename_using_shapes_add_picture(context, filename):
    shapes = context.slide.shapes
    shapes.add_picture(test_image(filename), Inches(1.25), Inches(1.25))


@when("I add the stream image {filename} using shapes.add_picture()")
def when_I_add_the_stream_image_filename_using_add_picture(context, filename):
    shapes = context.slide.shapes
    with open(test_image(filename), "rb") as f:
        stream = BytesIO(f.read())
    shapes.add_picture(stream, Inches(1.25), Inches(1.25))


@when("I assign MSO_AUTO_SHAPE_TYPE.{member} to picture.auto_shape_type")
def when_I_assign_member_to_picture_auto_shape_type(context, member):
    context.picture.auto_shape_type = getattr(MSO_AUTO_SHAPE_TYPE, member)


# then ====================================================


@then("a {ext} image part appears in the pptx file")
def step_then_a_ext_image_part_appears_in_the_pptx_file(context, ext):
    pkg = Package().open(saved_pptx_path)
    partnames = [part.partname for part in pkg.parts]
    image_partname = "/ppt/media/image1.%s" % ext
    assert image_partname in partnames, "got %s" % [
        p for p in partnames if "image" in p
    ]


@then("picture.auto_shape_type == MSO_AUTO_SHAPE_TYPE.{member}")
def then_picture_auto_shape_type_eq_shape_type_member(context, member):
    expected = getattr(MSO_AUTO_SHAPE_TYPE, member)
    actual = context.picture.auto_shape_type
    assert actual == expected, "shape.auto_shape_type == %s" % actual


@then("the picture appears in the slide")
def then_the_picture_appears_in_the_slide(context):
    prs = Presentation(saved_pptx_path)
    slide = prs.slides[0]
    shapes = slide.shapes
    cls_names = [sp.__class__.__name__ for sp in shapes]
    assert "Picture" in cls_names
