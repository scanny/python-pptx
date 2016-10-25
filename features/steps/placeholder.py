# encoding: utf-8

"""
Gherkin step implementations for placeholder-related features.
"""

from __future__ import absolute_import

import hashlib

from behave import given, when, then

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import _PlaceholderFormat

from helpers import saved_pptx_path, test_file, test_pptx, test_text


# given ===================================================

@given('a bullet body placeholder')
def given_a_bullet_body_placeholder(context):
    prs = Presentation(test_pptx('ph-unpopulated-placeholders'))
    context.prs = prs
    context.sld = prs.slides[2]
    context.body = prs.slides[2].shapes.placeholders[10]


@given('a known {placeholder_type} placeholder shape')
def given_a_known_placeholder_shape(context, placeholder_type):
    context.execute_steps(
        'given an unpopulated %s placeholder shape' % placeholder_type
    )


@given('a layout placeholder having directly set position and size')
def given_layout_placeholder_with_directly_set_pos_and_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slide_layouts[0].placeholders[1]


@given('a layout placeholder having no direct position or size settings')
def given_layout_placeholder_with_no_direct_pos_or_size_settings(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slide_layouts[0].placeholders[0]


@given('a master placeholder')
def given_a_master_placeholder(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.master_placeholder = prs.slide_master.placeholders[1]


@given('a notes slide placeholder having directly set position and size')
def given_notes_slide_placeholder_having_directly_set_pos_and_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[1].notes_slide.placeholders[1]


@given('a notes slide placeholder having no direct position or size settings')
def given_notes_slide_placeholder_having_no_direct_pos_or_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[0].notes_slide.placeholders[1]


@given('a slide placeholder having directly set position and size')
def given_slide_placeholder_with_directly_set_pos_and_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[0].placeholders[10]


@given('a slide placeholder having no direct position or size settings')
def given_slide_placeholder_with_no_direct_pos_or_size_settings(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[0].placeholders[0]


@given('a slide with an unpopulated {type_} placeholder')
def given_a_slide_with_an_unpopulated_type_placeholder(context, type_):
    slide_idx = [
        'title', 'content', 'text', 'chart', 'table', 'smart art', 'media',
        'clip art', 'picture',
    ].index(type_)
    prs = Presentation(test_pptx('ph-unpopulated-placeholders'))
    context.shape = prs.slides[slide_idx].shapes[0]


@given('a slide with a {type_} placeholder populated with {content}')
def given_a_slide_with_a_type_ph_with_content(context, type_, content):
    slide_idx = [
        'picture', 'clip art', 'table', 'chart', 'title', 'content', 'text',
        'smart art', 'media',
    ].index(type_)
    prs = Presentation(test_pptx('ph-populated-placeholders'))
    context.shape = prs.slides[slide_idx].shapes[0]


@given('an unpopulated {placeholder_type} placeholder shape')
def given_an_unpopulated_placeholder_shape(context, placeholder_type):
    slide_idx = [
        'title', 'content', 'text', 'chart', 'table', 'smart art', 'media',
        'clip art', 'picture',
    ].index(placeholder_type)
    prs = Presentation(test_pptx('ph-unpopulated-placeholders'))
    context.shape = prs.slides[slide_idx].shapes[0]


# when ====================================================

@when('I call placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data)')
def when_I_call_placeholder_insert_chart(context):
    chart_data = CategoryChartData()
    chart_data.categories = ['Yes', 'No']
    chart_data.add_series('Series 1', (42, 24))
    placeholder = context.shape
    context.placeholder = placeholder.insert_chart(
        XL_CHART_TYPE.PIE, chart_data
    )


@when('I call placeholder.insert_picture(\'{filename}\')')
def when_I_call_placeholder_insert_picture(context, filename):
    placeholder = context.shape
    path = test_file(filename)
    with open(path, 'rb') as f:
        context.image_sha1 = hashlib.sha1(f.read()).hexdigest()
    context.placeholder = placeholder.insert_picture(path)


@when('I call placeholder.insert_table(rows=2, cols=3)')
def when_I_call_placeholder_insert_table(context):
    placeholder = context.shape
    context.placeholder = placeholder.insert_table(2, 3)


@when('I indent the first paragraph')
def when_I_indent_the_first_paragraph(context):
    context.body.text_frame.paragraphs[0].level = 1


@when("I set the title text of the slide")
def step_when_set_slide_title_text(context):
    context.slide.shapes.title.text = test_text


# then ====================================================

@then('I can get the placeholder dimensions')
def then_I_can_get_the_placeholder_dimensions(context):
    placeholder = context.master_placeholder
    assert placeholder.width == 6923112, 'got %d' % placeholder.width
    assert placeholder.height == 3484984, 'got %d' % placeholder.height


@then('I can get the placeholder position')
def then_I_can_get_the_placeholder_position(context):
    placeholder = context.master_placeholder
    assert placeholder.left == 1110444, 'got %d' % placeholder.left
    assert placeholder.top == 1686508, 'got %d' % placeholder.top


@then('I get the direct settings when I query position and size')
def then_I_get_direct_settings_when_query_pos_and_size(context):
    placeholder = context.placeholder
    assert placeholder.left == 468312, 'got %s' % placeholder.left
    assert placeholder.top == 1700212, 'got %s' % placeholder.top
    assert placeholder.width == 8208143, 'got %s' % placeholder.width
    assert placeholder.height == 4537099, 'got %s' % placeholder.height


@then('I get inherited settings when I query position and size')
def then_I_get_inherited_settings_when_I_query_position_and_size(context):
    placeholder = context.placeholder
    assert placeholder.left == 457200, 'got %s' % placeholder.left
    assert placeholder.top == 274638, 'got %s' % placeholder.top
    assert placeholder.width == 8229600, 'got %s' % placeholder.width
    assert placeholder.height == 1143000, 'got %s' % placeholder.height


@then('placeholder_format.idx is {value}')
def then_placeholder_format_idx_is_value(context, value):
    expected_value = int(value)
    placeholder_format = context.shape.placeholder_format
    assert placeholder_format.idx == expected_value


@then('placeholder_format.type is {value}')
def then_placeholder_format_type_is_value(context, value):
    expected_value = getattr(PP_PLACEHOLDER, value.split('.')[1])
    placeholder_format = context.shape.placeholder_format
    assert placeholder_format.type == expected_value


@then('shape.placeholder_format is its _PlaceholderFormat object')
def then_shape_placeholder_format_is_its_PlaceholderFormat_object(context):
    shape = context.shape
    placeholder_format = shape.placeholder_format
    assert isinstance(placeholder_format, _PlaceholderFormat)
    assert placeholder_format.element is shape.element.ph


@then('shape.shape_type is MSO_SHAPE_TYPE.PLACEHOLDER')
def then_shape_shape_type_is_PLACEHOLDER(context):
    shape = context.shape
    assert shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER


@then('slide.shapes[0] is a {cls} proxy object for that placeholder')
def then_slide_shapes_0_is_a_cls_proxy_for_that_placeholder(context, cls):
    placeholder = context.shape
    clsname = placeholder.__class__.__name__
    assert clsname == cls, 'got %s' % clsname


@then('the chart is a pie chart')
def then_the_chart_is_a_pie_chart(context):
    chart = context.chart
    assert chart.chart_type == XL_CHART_TYPE.PIE


@then('the return value is a Placeholder{type} object')
def then_the_return_value_is_a_PlaceholderType_object(context, type):
    expected_type_name = 'Placeholder%s' % type
    placeholder_type_name = context.placeholder.__class__.__name__
    assert placeholder_type_name == expected_type_name


@then('the paragraph is indented')
def then_the_paragraph_is_indented(context):
    prs = Presentation(saved_pptx_path)
    p = prs.slides[2].shapes.placeholders[10].text_frame.paragraphs[0]
    assert p.level == 1


@then('the placeholder contains the chart')
def then_the_placeholder_contains_the_chart(context):
    placeholder_graphic_frame = context.placeholder
    assert placeholder_graphic_frame.has_chart
    context.chart = placeholder_graphic_frame.chart


@then('the placeholder contains the image')
def then_the_placeholder_contains_the_image(context):
    placeholder_picture = context.placeholder
    assert placeholder_picture.image.sha1 == context.image_sha1


@then('the placeholder contains the table')
def then_the_placeholder_contains_the_table(context):
    placeholder_graphic_frame = context.placeholder
    assert placeholder_graphic_frame.has_table
    context.table_ = placeholder_graphic_frame.table


@then('the placeholder\'s position and size are inherited from its layout')
def then_the_placeholders_position_and_size_are_inherited(context):
    placeholder = context.shape
    expected_values = (
        ('left', 2743200),
        ('top', 2057400),
        ('width', 3657600),
        ('height', 2743200),
    )
    for prop_name, expected_value in expected_values:
        value = getattr(placeholder, prop_name)
        assert value == expected_value, 'got %s' % value


@then('the {sides} crop is {value}')
def then_the_sides_crop_is_value(context, sides, value):
    side_prop_names = {
        'top and bottom': ('crop_top',  'crop_bottom'),
        'left and right': ('crop_left', 'crop_right'),
    }[sides]
    expected_value = float(value)
    placeholder_picture = context.placeholder
    for prop_name in side_prop_names:
        value = getattr(placeholder_picture, prop_name)
        difference = abs(expected_value - value)
        assert difference < 0.000002, 'got %s for %s' % (value, prop_name)


@then('the table has 2 rows and 3 columns')
def then_the_table_has_2_rows_and_3_columns(context):
    table = context.table_
    assert len(table.rows) == 2
    assert len(table.columns) == 3


@then('the text appears in the title placeholder')
def step_then_text_appears_in_title_placeholder(context):
    prs = Presentation(saved_pptx_path)
    title_shape = prs.slides[0].shapes.title
    title_text = title_shape.text_frame.paragraphs[0].runs[0].text
    assert title_text == test_text
