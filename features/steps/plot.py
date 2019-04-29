# encoding: utf-8

"""Gherkin step implementations for chart plot features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation

from helpers import test_pptx


# given ===================================================


@given("a bar plot {having_or_not} data labels")
def given_a_bar_plot_having_or_not_data_labels(context, having_or_not):
    slide_idx = {"having": 0, "not having": 1}[having_or_not]
    prs = Presentation(test_pptx("cht-plot-props"))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given("a bar plot having gap width of {width}")
def given_a_bar_plot_having_gap_width_of_width(context, width):
    slide_idx = {"no explicit value": 0, "300": 1}[width]
    prs = Presentation(test_pptx("cht-plot-props"))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given("a bar plot having overlap of {overlap}")
def given_a_bar_plot_having_overlap_of_overlap(context, overlap):
    slide_idx = {"no explicit value": 0, "42": 1, "-42": 2}[overlap]
    prs = Presentation(test_pptx("cht-plot-props"))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given("a bar plot having vary color by category set to {setting}")
def given_a_bar_plot_having_vary_color_by_category_setting(context, setting):
    slide_idx = {"no explicit setting": 0, "True": 1, "False": 2}[setting]
    prs = Presentation(test_pptx("cht-plot-props"))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given("a bubble plot having bubble scale of {percent}")
def given_a_bubble_plot_having_bubble_scale_of_percent(context, percent):
    slide_idx = {"no explicit value": 3, "70%": 4}[percent]
    prs = Presentation(test_pptx("cht-plot-props"))
    context.bubble_plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given("a category plot")
def given_a_category_plot(context):
    prs = Presentation(test_pptx("cht-plot-props"))
    context.plot = prs.slides[2].shapes[0].chart.plots[0]


# when ====================================================


@when("I assign {value} to bubble_plot.bubble_scale")
def when_I_assign_value_to_bubble_plot_bubble_scale(context, value):
    new_value = None if value == "None" else int(value)
    context.bubble_plot.bubble_scale = new_value


@when("I assign {value} to plot.gap_width")
def when_I_assign_value_to_plot_gap_width(context, value):
    new_value = int(value)
    context.plot.gap_width = new_value


@when("I assign {value} to plot.has_data_labels")
def when_I_assign_value_to_plot_has_data_labels(context, value):
    new_value = {"True": True, "False": False}[value]
    context.plot.has_data_labels = new_value


@when("I assign {value} to plot.overlap")
def when_I_assign_value_to_plot_overlap(context, value):
    new_value = int(value)
    context.plot.overlap = new_value


@when("I assign {value} to plot.vary_by_categories")
def when_I_assign_value_to_plot_vary_by_categories(context, value):
    new_value = {"True": True, "False": False}[value]
    context.plot.vary_by_categories = new_value


# then ====================================================


@then("bubble_plot.bubble_scale is {value}")
def then_bubble_plot_bubble_scale_is_value(context, value):
    expected_value = int(value)
    bubble_plot = context.bubble_plot
    assert bubble_plot.bubble_scale == expected_value, (
        "got %s" % bubble_plot.bubble_scale
    )


@then("len(plot.categories) is {count}")
def then_len_plot_categories_is_count(context, count):
    plot = context.chart.plots[0]
    expected_count = int(count)
    assert len(plot.categories) == expected_count


@then("plot.categories is a Categories object")
def then_plot_categories_is_a_Categories_object(context):
    plot = context.plot
    type_name = type(plot.categories).__name__
    assert type_name == "Categories", "got %s" % type_name


@then("plot.gap_width is {value}")
def then_plot_gap_width_is_value(context, value):
    expected_value = int(value)
    plot = context.plot
    assert plot.gap_width == expected_value, "got %s" % plot.gap_width


@then("plot.has_data_labels is {value}")
def then_plot_has_data_labels_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    assert context.plot.has_data_labels is expected_value


@then("plot.overlap is {value}")
def then_plot_overlap_is_expected_value(context, value):
    expected_value = int(value)
    plot = context.plot
    assert plot.overlap == expected_value, "got %s" % plot.overlap


@then("plot.vary_by_categories is {value}")
def then_plot_vary_by_categories_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    plot = context.plot
    assert plot.vary_by_categories is expected_value
