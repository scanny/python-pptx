"""Gherkin step implementations for presentation-level features."""

from __future__ import annotations

import io
import os
import zipfile
from typing import TYPE_CHECKING, cast

from behave import given, then, when
from behave.runner import Context
from helpers import saved_pptx_path, test_file, test_pptx

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches

if TYPE_CHECKING:
    from pptx import presentation
    from pptx.shapes.picture import Picture

# given ===================================================


@given("a clean working directory")
def given_clean_working_dir(context: Context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)


@given("a presentation")
def given_a_presentation(context: Context):
    context.presentation = Presentation(test_pptx("prs-properties"))


@given("a presentation having a notes master")
def given_a_presentation_having_a_notes_master(context: Context):
    context.prs = Presentation(test_pptx("prs-notes"))


@given("a presentation having no notes master")
def given_a_presentation_having_no_notes_master(context: Context):
    context.prs = Presentation(test_pptx("prs-properties"))


@given("a presentation with an image/jpg MIME-type")
def given_prs_with_image_jpg_MIME_type(context):
    context.prs = Presentation(test_pptx("test-image-jpg-mime"))


@given("a presentation with external relationships")
def given_prs_with_ext_rels(context: Context):
    context.prs = Presentation(test_pptx("ext-rels"))


@given("an initialized pptx environment")
def given_initialized_pptx_env(context: Context):
    pass


# when ====================================================


@when("I change the slide width and height")
def when_change_slide_width_and_height(context: Context):
    presentation = context.presentation
    presentation.slide_width = Inches(4)
    presentation.slide_height = Inches(3)


@when("I construct a Presentation instance with no path argument")
def when_construct_default_prs(context: Context):
    context.prs = Presentation()


@when("I open a basic PowerPoint presentation")
def when_open_basic_pptx(context: Context):
    context.prs = Presentation(test_pptx("test"))


@when("I open a presentation extracted into a directory")
def when_I_open_a_presentation_extracted_into_a_directory(context: Context):
    context.prs = Presentation(test_file("extracted-pptx"))


@when("I open a presentation contained in a stream")
def when_open_presentation_stream(context: Context):
    with open(test_pptx("test"), "rb") as f:
        stream = io.BytesIO(f.read())
    context.prs = Presentation(stream)
    stream.close()


@when("I save and reload the presentation")
def when_save_and_reload_prs(context: Context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)
    context.prs = Presentation(saved_pptx_path)


@when("I save that stream to a file")
def when_save_stream_to_a_file(context: Context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.stream.seek(0)
    with open(saved_pptx_path, "wb") as f:
        f.write(context.stream.read())


@when("I save the presentation")
def when_save_presentation(context: Context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)


@when("I save the presentation to a stream")
def when_save_presentation_to_stream(context: Context):
    context.stream = io.BytesIO()
    context.prs.save(context.stream)


# then ====================================================


@then("I receive a presentation based on the default template")
def then_receive_prs_based_on_def_tmpl(context: Context):
    prs = context.prs
    assert prs is not None
    slide_masters = prs.slide_masters
    assert slide_masters is not None
    assert len(slide_masters) == 1
    slide_layouts = slide_masters[0].slide_layouts
    assert slide_layouts is not None
    assert len(slide_layouts) == 11


@then("its slide height matches its known value")
def then_slide_height_matches_known_value(context: Context):
    presentation = context.presentation
    assert presentation.slide_height == 6858000


@then("its slide width matches its known value")
def then_slide_width_matches_known_value(context: Context):
    presentation = context.presentation
    assert presentation.slide_width == 9144000


@then("I see the pptx file in the working directory")
def then_see_pptx_file_in_working_dir(context: Context):
    assert os.path.isfile(saved_pptx_path)
    minimum = 30000
    actual = os.path.getsize(saved_pptx_path)
    assert actual > minimum


@then("len(notes_master.shapes) is {shape_count}")
def then_len_notes_master_shapes_is_shape_count(context: Context, shape_count: str):
    notes_master = context.prs.notes_master
    expected = int(shape_count)
    actual = len(notes_master.shapes)
    assert actual == expected, "got %s" % actual


@then("prs.notes_master is a NotesMaster object")
def then_prs_notes_master_is_a_NotesMaster_object(context: Context):
    prs = context.prs
    assert type(prs.notes_master).__name__ == "NotesMaster"


@then("prs.slides is a Slides object")
def then_prs_slides_is_a_Slides_object(context: Context):
    prs = context.presentation
    assert type(prs.slides).__name__ == "Slides"


@then("prs.slide_masters is a SlideMasters object")
def then_prs_slide_masters_is_a_SlideMasters_object(context: Context):
    prs = context.presentation
    assert type(prs.slide_masters).__name__ == "SlideMasters"


@then("the external relationships are still there")
def then_ext_rels_are_preserved(context: Context):
    prs = context.prs
    sld = prs.slides[0]
    rel = sld.part._rels["rId2"]
    assert rel.is_external
    assert rel.reltype == RT.HYPERLINK
    assert rel.target_ref == "https://github.com/scanny/python-pptx"


@then("the package has the expected number of .rels parts")
def then_the_package_has_the_expected_number_of_rels_parts(context: Context):
    with zipfile.ZipFile(saved_pptx_path, "r") as z:
        member_count = len(z.namelist())
    assert member_count == 18, "expected 18, got %d" % member_count


@then("I can access the JPEG image")
def then_I_can_access_the_JPEG_image(context):
    prs = cast("presentation.Presentation", context.prs)
    slide = prs.slides[0]
    picture = cast("Picture", slide.shapes[0])
    try:
        picture.image
    except AttributeError:
        raise AssertionError("JPEG image not recognized")


@then("the slide height matches the new value")
def then_slide_height_matches_new_value(context: Context):
    presentation = context.presentation
    assert presentation.slide_height == Inches(3)


@then("the slide width matches the new value")
def then_slide_width_matches_new_value(context: Context):
    presentation = context.presentation
    assert presentation.slide_width == Inches(4)
