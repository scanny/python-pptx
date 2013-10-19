# encoding: utf-8

import os

from behave import given, when, then
from hamcrest import (
    assert_that, equal_to, has_item, is_, is_not, greater_than
)
from StringIO import StringIO

from pptx import packaging
from pptx import Presentation
from pptx.constants import MSO_AUTO_SHAPE_TYPE as MAST, MSO
from pptx.util import Inches


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
scratch_dir = absjoin(thisdir, '../_scratch')
test_file_dir = absjoin(thisdir, '../../tests/test_files')
basic_pptx_path = absjoin(test_file_dir, 'test.pptx')
saved_pptx_path = absjoin(scratch_dir, 'test_out.pptx')
test_image_path = absjoin(test_file_dir, 'python-powered.png')

test_text = "python-pptx was here!"


# given ===================================================

@given('a clean working directory')
def step_given_clean_working_dir(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)


@given('an initialized pptx environment')
def step_given_initialized_pptx_env(context):
    pass


@given('I have a reference to a blank slide')
def step_given_ref_to_blank_slide(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[6]
    context.sld = context.prs.slides.add_slide(slidelayout)


@given('I have a reference to a chevron shape')
def step_given_ref_to_chevron_shape(context):
    context.prs = Presentation()
    blank_slidelayout = context.prs.slidelayouts[6]
    shapes = context.prs.slides.add_slide(blank_slidelayout).shapes
    x = y = cx = cy = 914400
    context.chevron_shape = shapes.add_shape(MAST.CHEVRON, x, y, cx, cy)


@given('I have a reference to a slide')
def step_given_ref_to_slide(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[0]
    context.sld = context.prs.slides.add_slide(slidelayout)


@given('I have an empty presentation open')
def step_given_empty_prs(context):
    context.prs = Presentation()


# when ====================================================

@when('I add a new slide')
def step_when_add_slide(context):
    slidelayout = context.prs.slidemasters[0].slidelayouts[0]
    context.prs.slides.add_slide(slidelayout)


@when("I add a picture stream to the slide's shape collection")
def step_when_add_picture_stream(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.25), Inches(1.25))
    with open(test_image_path) as f:
        stream = StringIO(f.read())
    shapes.add_picture(stream, x, y)


@when("I add a picture to the slide's shape collection")
def step_when_add_picture(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.25), Inches(1.25))
    shapes.add_picture(test_image_path, x, y)


@when("I add a text box to the slide's shape collection")
def step_when_add_text_box(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    sp = shapes.add_textbox(x, y, cx, cy)
    sp.text = test_text


@when("I add an auto shape to the slide's shape collection")
def step_when_add_auto_shape(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(4.00))
    sp = shapes.add_shape(MAST.ROUNDED_RECTANGLE, x, y, cx, cy)
    sp.text = test_text


@when('I construct a Presentation instance with no path argument')
def step_when_construct_default_prs(context):
    context.prs = Presentation()


@when('I open a basic PowerPoint presentation')
def step_when_open_basic_pptx(context):
    context.prs = Presentation(basic_pptx_path)


@when('I open a presentation contained in a stream')
def step_when_open_presentation_stream(context):
    with open(basic_pptx_path) as f:
        stream = StringIO(f.read())
    context.prs = Presentation(stream)
    stream.close()


@when('I save that stream to a file')
def step_when_save_stream_to_a_file(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.stream.seek(0)
    with open(saved_pptx_path, 'wb') as f:
        f.write(context.stream.read())


@when('I save the presentation')
def step_when_save_presentation(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)


@when('I save the presentation to a stream')
def step_when_save_presentation_to_stream(context):
    context.stream = StringIO()
    context.prs.save(context.stream)


@when("I set the first adjustment value to 0.15")
def step_when_set_first_adjustment_value(context):
    context.chevron_shape.adjustments[0] = 0.15


# then ====================================================


@then('I receive a presentation based on the default template')
def step_then_receive_prs_based_on_def_tmpl(context):
    prs = context.prs
    assert_that(prs, is_not(None))
    slidemasters = prs.slidemasters
    assert_that(slidemasters, is_not(None))
    assert_that(len(slidemasters), is_(1))
    slidelayouts = slidemasters[0].slidelayouts
    assert_that(slidelayouts, is_not(None))
    assert_that(len(slidelayouts), is_(11))


@then('I see the pptx file in the working directory')
def step_then_see_pptx_file_in_working_dir(context):
    assert_that(os.path.isfile(saved_pptx_path))
    minimum = 30000
    actual = os.path.getsize(saved_pptx_path)
    assert_that(actual, is_(greater_than(minimum)))


@then('the auto shape appears in the slide')
def step_then_auto_shape_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sp = prs.slides[0].shapes[0]
    sp_text = sp.textframe.paragraphs[0].runs[0].text
    assert_that(sp.shape_type, is_(equal_to(MSO.AUTO_SHAPE)))
    assert_that(sp.auto_shape_type, is_(equal_to(MAST.ROUNDED_RECTANGLE)))
    assert_that(sp_text, is_(equal_to(test_text)))


@then('the chevron shape appears with a less acute arrow head')
def step_then_chevron_shape_appears_with_less_acute_arrow_head(context):
    chevron = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(chevron.adjustments[0], is_(equal_to(0.15)))


@then('the image is saved in the pptx file')
def step_then_img_saved_in_pptx_file(context):
    pkgng_pkg = packaging.Package().open(saved_pptx_path)
    partnames = [part.partname for part in pkgng_pkg.parts
                 if part.partname.startswith('/ppt/media/')]
    assert_that(partnames, has_item('/ppt/media/image1.png'))


@then('the picture appears in the slide')
def step_then_picture_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    shapes = sld.shapes
    classnames = [sp.__class__.__name__ for sp in shapes]
    assert_that(classnames, has_item('_Picture'))


@then('the pptx file contains a single slide')
def step_then_pptx_file_contains_single_slide(context):
    prs = Presentation(saved_pptx_path)
    assert_that(len(prs.slides), is_(equal_to(1)))


@then('the text box appears in the slide')
def step_then_text_box_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    textbox = prs.slides[0].shapes[0]
    textbox_text = textbox.textframe.paragraphs[0].runs[0].text
    assert_that(textbox_text, is_(equal_to(test_text)))
