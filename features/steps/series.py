# encoding: utf-8

"""Gherkin step implementations for chart plot features."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

from ast import literal_eval

from behave import given, then, when

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR

from helpers import test_pptx


# given ===================================================

@given('a bar series having fill of {fill}')
def given_a_bar_series_having_fill_of_fill(context, fill):
    series_idx = {
        'Automatic': 0,
        'No Fill':   1,
        'Orange':    2,
        'Accent 1':  3,
    }[fill]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series having invert_if_negative of {setting}')
def given_a_bar_series_having_invert_if_negative_setting(context, setting):
    series_idx = {
        'no explicit setting': 0,
        'True':                1,
        'False':               2,
    }[setting]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series with values {values}')
def given_a_bar_series_with_values(context, values):
    prs = Presentation(test_pptx('cht-series-props'))
    series_idx = {
        '1.2, 2.3, 3.4':  0,
        '4.5, None, 6.7': 1,
    }[values]
    context.series = prs.slides[1].shapes[0].chart.plots[0].series[series_idx]


@given('a bar series having {width} line')
def given_a_bar_series_having_width_line(context, width):
    series_idx = {
        'no':      0,
        '1 point': 1,
    }[width]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a marker')
def given_a_marker(context):
    prs = Presentation(test_pptx('cht-marker-props'))
    series = prs.slides[0].shapes[0].chart.series[0]
    context.marker = series.marker


@given('a marker having size of {case}')
def given_a_marker_having_size_of_case(context, case):
    series_idx = {
        'no explicit value': 0,
        '24 points':         1,
        '36 points':         2,
    }[case]
    prs = Presentation(test_pptx('cht-marker-props'))
    series = prs.slides[0].shapes[0].chart.series[series_idx]
    context.marker = series.marker


@given('a marker having style of {case}')
def given_a_marker_having_style_of_case(context, case):
    series_idx = {
        'no explicit value': 0,
        'circle':            1,
        'triangle':          2,
    }[case]
    prs = Presentation(test_pptx('cht-marker-props'))
    series = prs.slides[0].shapes[0].chart.series[series_idx]
    context.marker = series.marker


@given('a point')
def given_a_point(context):
    prs = Presentation(test_pptx('cht-point-props'))
    chart = prs.slides[0].shapes[0].chart
    context.point = chart.plots[0].series[0].points[0]


@given('a {points_type} object containing 3 points')
def given_a_points_type_object_containing_3_points(context, points_type):
    slide_idx = {
        'XyPoints':       0,
        'BubblePoints':   1,
        'CategoryPoints': 2,
    }[points_type]
    prs = Presentation(test_pptx('cht-point-access'))
    series = prs.slides[slide_idx].shapes[0].chart.plots[0].series[0]
    context.points = series.points


@given('a series')
def given_a_series(context):
    prs = Presentation(test_pptx('cht-series-props'))
    context.series = prs.slides[0].shapes[0].chart.plots[0].series[0]


@given('a series collection for a plot having {count} series')
def given_a_series_collection_for_a_plot_having_count_series(context, count):
    prs = Presentation(test_pptx('cht-series-access'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series_collection = plot.series
    context.series_count = int(count)


@given('a series collection for a {type_} chart having {count} series')
def given_a_series_collection_for_chart_having_series(context, type_, count):
    slide_idx = {
        'single-plot': 0,
        'multi-plot':  1,
    }[type_]
    prs = Presentation(test_pptx('cht-series-access'))
    context.series_collection = prs.slides[slide_idx].shapes[0].chart.series
    context.series_count = int(count)


@given('a series of type {series_type}')
def given_a_series_of_type_series_type(context, series_type):
    slide_idx = {
        'Category': 1,
        'XY':       2,
        'Bubble':   3,
        'Line':     4,
        'Radar':    5,
    }[series_type]
    prs = Presentation(test_pptx('cht-series-props'))
    context.series = prs.slides[slide_idx].shapes[0].chart.plots[0].series[0]


# when ====================================================

@when('I add a series with number format {strval}')
def when_I_add_a_series_with_number_format(context, strval):
    chart_data = context.chart_data
    params = {'name': 'Series Foo'}
    if strval != 'None':
        params['number_format'] = int(strval)
    context.series_data = chart_data.add_series(**params)


@when('I assign {value} to marker.size')
def when_I_assign_value_to_marker_size(context, value):
    new_value = None if value == 'None' else int(value)
    context.marker.size = new_value


@when('I assign {value} to marker.style')
def when_I_assign_value_to_marker_style(context, value):
    new_value = None if value == 'None' else getattr(XL_MARKER_STYLE, value)
    context.marker.style = new_value


@when('I assign {value} to series.invert_if_negative')
def when_I_assign_value_to_series_invert_if_negative(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.series.invert_if_negative = new_value


# then ====================================================

@then('data_point.number_format is {value_str}')
def then_data_point_number_format_is(context, value_str):
    data_point = context.data_point
    number_format = value_str if value_str == 'General' else int(value_str)
    assert data_point.number_format == number_format


@then('iterating points produces 3 Point objects')
def then_iterating_points_produces_3_point_objects(context):
    points = context.points
    idx = -1
    for idx, point in enumerate(points):
        assert type(point).__name__ == 'Point'
    assert idx == 2, 'got %s' % idx


@then('iterating series_collection produces {count} Series objects')
def then_iterating_series_collection_produces_count_series(context, count):
    expected_idx = int(count) - 1
    idx = -1
    for idx, series in enumerate(context.series_collection):
        type_name = type(series).__name__
        assert type_name.endswith('Series'), 'got %s' % type_name
    assert idx == expected_idx, 'got %s' % idx


@then('len(points) is 3')
def then_len_points_is_3(context):
    points = context.points
    assert len(points) == 3


@then('len(series_collection) is {count}')
def then_len_series_collection_is_count(context, count):
    expected_len = int(count)
    actual_len = len(context.series_collection)
    assert actual_len == expected_len, 'got %s' % actual_len


@then('len(series.values) is {count} for each series')
def then_len_series_values_is_count_for_each_series(context, count):
    expected_count = int(count)
    for series in context.chart.plots[0].series:
        assert len(series.values) == expected_count


@then('marker.format is a ChartFormat object')
def then_marker_format_is_a_ChartFormat_object(context):
    marker = context.marker
    assert type(marker.format).__name__ == 'ChartFormat'


@then('marker.format.fill is a FillFormat object')
def then_marker_format_fill_is_a_FillFormat_object(context):
    marker = context.marker
    assert type(marker.format.fill).__name__ == 'FillFormat'


@then('marker.format.line is a LineFormat object')
def then_marker_format_line_is_a_LineFormat_object(context):
    marker = context.marker
    assert type(marker.format.line).__name__ == 'LineFormat'


@then('marker.size is {case}')
def then_marker_size_is_case(context, case):
    expected_value = None if case == 'None' else int(case)
    marker = context.marker
    assert marker.size == expected_value, 'got %s' % marker.size


@then('marker.style is {case}')
def then_marker_style_is_case(context, case):
    expected_value = None if case == 'None' else getattr(XL_MARKER_STYLE, case)
    marker = context.marker
    assert marker.style == expected_value, 'got %s' % marker.style


@then('point.data_label is a DataLabel object')
def then_point_data_label_is_a_DataLabel_object(context):
    point = context.point
    assert type(point.data_label).__name__ == 'DataLabel'


@then('point.format is a ChartFormat object')
def then_point_format_is_a_ChartFormat_object(context):
    point = context.point
    assert type(point.format).__name__ == 'ChartFormat'


@then('point.format.fill is a FillFormat object')
def then_point_format_fill_is_a_FillFormat_object(context):
    point = context.point
    assert type(point.format.fill).__name__ == 'FillFormat'


@then('point.format.line is a LineFormat object')
def then_point_format_line_is_a_LineFormat_object(context):
    point = context.point
    assert type(point.format.line).__name__ == 'LineFormat'


@then('point.marker is a Marker object')
def then_point_marker_is_a_Marker_object(context):
    point = context.point
    assert type(point.marker).__name__ == 'Marker'


@then('points[2] is a Point object')
def then_points_2_is_a_Point_object(context):
    points = context.points
    assert type(points[2]).__name__ == 'Point'


@then('series_collection[2] is a Series object')
def then_series_collection_2_is_a_Series_object(context):
    type_name = type(context.series_collection[2]).__name__
    assert type_name.endswith('Series'), 'got %s' % type_name


@then('series.format.fill.fore_color.rgb is FF6600')
def then_series_format_fill_fore_color_rgb_is_FF6600(context):
    rgb_color = context.series.format.fill.fore_color.rgb
    assert rgb_color == RGBColor(0xFF, 0x66, 0x00), 'got %s' % rgb_color


@then('series.format.fill.fore_color.theme_color is Accent 1')
def then_series_format_fill_fore_color_theme_color_is_Accent_1(context):
    theme_color = context.series.format.fill.fore_color.theme_color
    assert theme_color == MSO_THEME_COLOR.ACCENT_1, 'got %s' % theme_color


@then('series.format.fill.type is {fill_type}')
def then_series_format_fill_type_is_type(context, fill_type):
    expected_fill_type = {
        'None':                     None,
        'MSO_FILL_TYPE.BACKGROUND': MSO_FILL_TYPE.BACKGROUND,
        'MSO_FILL_TYPE.SOLID':      MSO_FILL_TYPE.SOLID,
    }[fill_type]
    fill_type = context.series.format.fill.type
    assert fill_type == expected_fill_type, 'got %s' % fill_type


@then('series.invert_if_negative is {value}')
def then_series_invert_if_negative_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    series = context.series
    assert series.invert_if_negative is expected_value


@then('series.format.line.width is {width}')
def then_series_format_line_width_is_width(context, width):
    expected_width = int(width)
    line_width = context.series.format.line.width
    assert line_width == expected_width, 'got %s' % line_width


@then('series.format is a ChartFormat object')
def then_series_format_is_a_ChartFormat_object(context):
    series = context.series
    assert type(series.format).__name__ == 'ChartFormat'


@then('series.marker is a Marker object')
def then_series_marker_is_a_Marker_object(context):
    series = context.series
    assert type(series.marker).__name__ == 'Marker'


@then('series.points is a {type_name} object')
def then_series_points_is_a_type_name_object(context, type_name):
    series = context.series
    assert type(series.points).__name__ == type_name


@then('series.values is {values}')
def then_series_values_is_values(context, values):
    series = context.series
    expected_values = literal_eval(values)
    assert series.values == expected_values, 'got %s' % (series.values,)
