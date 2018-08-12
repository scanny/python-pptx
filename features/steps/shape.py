# encoding: utf-8

"""Gherkin step implementations for shape-related features."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import hashlib

from behave import given, when, then

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_MEDIA_TYPE
from pptx.action import ActionSetting
from pptx.util import Emu

from helpers import cls_qname, test_file, test_pptx


# given ===================================================

@given('an autoshape')
def given_an_autoshape(context):
    prs = Presentation(test_pptx('shp-autoshape-adjustments'))
    context.shape = prs.slides[0].shapes[0]


@given('an autoshape having text')
def given_an_autoshape_having_text(context):
    prs = Presentation(test_pptx('shp-autoshape-props'))
    context.shape = prs.slides[0].shapes[0]


@given("(builder._start_x, builder._start_y) is ({x_str}, {y_str})")
def given_builder_start_x_builder_start_y_is_x_y(context, x_str, y_str):
    builder = context.builder
    builder._start_x, builder._start_y = int(x_str), int(y_str)


@given("(builder._x_scale, builder._y_scale) is ({p_str}, {q_str})")
def given_builder_x_scale_builder_y_scale_is_p_q(context, p_str, q_str):
    builder = context.builder
    builder._x_scale, builder._y_scale = float(p_str), float(q_str)


@given('a chevron shape')
def given_a_chevron_shape(context):
    prs = Presentation(test_pptx('shp-autoshape-adjustments'))
    context.shape = prs.slides[0].shapes[0]


@given('a Connector object as shape')
def given_a_Connector_object_as_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[4]


@given('a connector and a 1 inch square picture at 0, 0')
def given_a_connector_and_a_1_inch_square_picture_at_0_0(context):
    prs = Presentation(test_pptx('shp-connector-props'))
    shapes = prs.slides[1].shapes
    context.picture = shapes[0]
    context.connector = shapes[1]


@given('a connector having its begin point at ({x}, {y})')
def given_a_connector_having_its_begin_point_at_x_y(context, x, y):
    prs = Presentation(test_pptx('shp-connector-props'))
    sld = prs.slides[0]
    context.connector = sld.shapes[0]


@given('a connector having its end point at ({x}, {y})')
def given_a_connector_having_its_end_point_at_x_y(context, x, y):
    prs = Presentation(test_pptx('shp-connector-props'))
    sld = prs.slides[0]
    context.connector = sld.shapes[0]


@given('an empty GroupShape object as shape')
def given_an_empty_GroupShape_object_as_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes.add_group_shape()


@given('a FreeformBuilder object as builder')
def given_a_FreeformBuilder_object_as_builder(context):
    shapes = Presentation(test_pptx('shp-freeform')).slides[0].shapes
    builder = shapes.build_freeform()
    context.builder = builder


@given('a GraphicFrame object containing a chart as shape')
def given_a_GraphicFrame_object_containing_a_chart_as_shape(context):
    prs = Presentation(test_pptx('shp-access-chart'))
    sld = prs.slides[0]
    context.shape = sld.shapes[0]


@given('a GraphicFrame object containing a table as shape')
def given_a_GraphicFrame_object_containing_a_table_as_shape(context):
    prs = Presentation(test_pptx('shp-access-chart'))
    sld = prs.slides[1]
    context.shape = sld.shapes[0]


@given('a GraphicFrame object as shape')
def given_a_GraphicFrame_object_as_shape(context):
    # shouldn't matter, but this one contains a table
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[2]


@given('a GroupShape object as group_shape')
def given_a_GroupShape_object_as_group_shape(context):
    prs = Presentation(test_pptx('shp-groupshape'))
    sld = prs.slides[0]
    context.group_shape = sld.shapes[0]


@given('a GroupShape object as shape')
def given_a_GroupShape_object_as_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[3]


@given('a movie shape')
def given_a_movie_shape(context):
    prs = Presentation(test_pptx('shp-movie-props'))
    context.movie = prs.slides[0].shapes[0]


@given('a Picture object as picture')
def given_a_Picture_object_as_picture(context):
    slide = Presentation(test_pptx('shp-picture')).slides[0]
    context.picture = slide.shapes[0]


@given('a Picture object as shape')
def given_a_Picture_object_as_shape(context):
    slide = Presentation(test_pptx('shp-common-props')).slides[0]
    context.shape = slide.shapes[1]


@given('a Picture object with {crop_or_no} as picture')
def given_a_Picture_object_with_crop_or_no_as_picture(context, crop_or_no):
    shape_idx = {
        'no cropping': 0,
        'cropping': 1,
    }[crop_or_no]
    slide = Presentation(test_pptx('shp-picture')).slides[0]
    context.picture = slide.shapes[shape_idx]


@given('a rotated {shape_type} object as shape')
def given_a_rotated_shape_type_object_as_shape(context, shape_type):
    shape_idx = {
        'Shape':        0,
        'Picture':      1,
        'GraphicFrame': 2,
        'GroupShape':   3,
        'Connector':    4,
    }[shape_type]
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[1]
    context.shape = sld.shapes[shape_idx]


@given('a Shape object as shape')
def given_a_Shape_object_as_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[0]


@given('a {shape_type} object on a slide as shape')
def given_a_shape_on_a_slide(context, shape_type):
    shape_idx = {
        'Shape':        0,
        'Picture':      1,
        'GraphicFrame': 2,
        'GroupShape':   3,
        'Connector':    4,
    }[shape_type]
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[shape_idx]
    context.slide = sld


@given('a textbox')
def given_a_textbox(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[5]


@given('a shape of known position and size')
def given_a_shape_of_known_position_and_size(context):
    prs = Presentation(test_pptx('shp-pos-and-size'))
    context.shape = prs.slides[0].shapes[0]


# when ====================================================

@when('I add a {cx} x {cy} shape at ({x}, {y})')
def when_I_add_a_cx_cy_shape_at_x_y(context, cx, cy, x, y):
    context.shape.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(cx), int(cy)
    )


@when("I assign 0.15 to shape.adjustments[0]")
def when_I_assign_to_shape_adjustments(context):
    context.shape.adjustments[0] = 0.15


@when("I assign a string to shape.text")
def when_I_assign_a_string_to_shape_text(context):
    context.shape.text = u'F\xf8o\nBar'


@when("I assign builder.convert_to_shape() to shape")
def when_I_assign_builder_convert_to_shape_to_shape(context):
    builder = context.builder
    context.shape = builder.convert_to_shape()


@when("I assign builder.convert_to_shape({x_str}, {y_str}) to shape")
def when_I_assign_builder_convert_to_shape_origin_x_y(context, x_str, y_str):
    builder = context.builder
    origin_x, origin_y = int(x_str), int(y_str)
    context.shape = builder.convert_to_shape(origin_x, origin_y)


@when('I assign {value} to connector.begin_x')
def when_I_assign_value_to_connector_begin_x(context, value):
    context.connector.begin_x = int(value)


@when('I assign {value} to connector.begin_y')
def when_I_assign_value_to_connector_begin_y(context, value):
    context.connector.begin_y = int(value)


@when('I assign {value} to connector.end_x')
def when_I_assign_value_to_connector_end_x(context, value):
    context.connector.end_x = int(value)


@when('I assign {value} to connector.end_y')
def when_I_assign_value_to_connector_end_y(context, value):
    context.connector.end_y = int(value)


@when('I assign {value} to picture.crop_{side}')
def when_I_assign_value_to_picture_crop_side(context, value, side):
    new_value = (
        None if value == 'None' else
        float(value) if '.' in value else
        int(value)
    )
    setattr(context.picture, 'crop_%s' % side, new_value)


@when('I assign {value} to shape.height')
def when_I_assign_value_to_shape_height(context, value):
    context.shape.height = int(value)


@when('I assign {value} to shape.left')
def when_I_assign_value_to_shape_left(context, value):
    context.shape.left = int(value)


@when('I assign \'{value}\' to shape.name')
def when_I_assign_value_to_shape_name(context, value):
    context.shape.name = value


@when('I assign {value} to shape.rotation')
def when_I_assign_value_to_shape_rotation(context, value):
    context.shape.rotation = float(value)


@when('I assign {value} to shape.top')
def when_I_assign_value_to_shape_top(context, value):
    context.shape.top = int(value)


@when('I assign {value} to shape.width')
def when_I_assign_value_to_shape_width(context, value):
    context.shape.width = int(value)


@when('I call builder.add_line_segments([(100, 25), (25, 100)])')
def when_I_call_builder_add_line_segments_100_25_25_100(context):
    builder = context.builder
    builder.add_line_segments([(100, 25), (25, 100)])


@when('I call connector.begin_connect(picture, 3)')
def when_I_call_connector_begin_connect_picture_3(context):
    connector, picture = context.connector, context.picture
    connector.begin_connect(picture, 3)


@when('I call connector.end_connect(picture, 3)')
def when_I_call_connector_end_connect_picture_3(context):
    connector, picture = context.connector, context.picture
    connector.end_connect(picture, 3)


# then ====================================================

@then("accessing shape.click_action raises TypeError")
def then_accessing_shape_click_action_raises_TypeError(context):
    try:
        context.shape.click_action
    except TypeError:
        return
    except Exception as e:
        raise AssertionError(
            'Accessing GroupShape.click_action raised %s' % type(e).__name__
        )
    raise AssertionError('Accessing GroupShape.click_action did not raise')


@then("builder is a FreeformBuilder object")
def then_builder_is_a_FreeformBuilder_object(context):
    builder = context.builder
    class_name = builder.__class__.__name__
    expected_value = 'FreeformBuilder'
    assert class_name == expected_value, (
        'Expected class name \'%s\', got \'%s\'' %
        (expected_value, class_name)
    )


@then("(builder._start_x, builder._start_y) is ({x_str}, {y_str})")
def then_builder_start_x_builder_start_y_is_x_y(context, x_str, y_str):
    builder = context.builder
    actual_value = builder._start_x, builder._start_y
    expected_value = int(x_str), int(y_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then("(builder._x_scale, builder._y_scale) is ({p_str}, {q_str})")
def then_builder_x_scale_builder_y_scale_is_x_y(context, p_str, q_str):
    builder = context.builder
    actual_value = builder._x_scale, builder._y_scale
    expected_value = float(p_str), float(q_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then('connector is a Connector object')
def then_connector_is_a_Connector_object(context):
    assert type(context.connector).__name__ == 'Connector'


@then('connector.begin_x == {value}')
def then_connector_begin_x_equals_value(context, value):
    assert context.connector.begin_x == int(value)


@then('connector.begin_x is an Emu object with value {x}')
def then_connector_begin_x_is_an_Emu_object_with_value_x(context, x):
    begin_x = context.connector.begin_x
    assert isinstance(begin_x, Emu)
    assert begin_x == int(x)


@then('connector.begin_y == {value}')
def then_connector_begin_y_equals_value(context, value):
    assert context.connector.begin_y == int(value)


@then('connector.begin_y is an Emu object with value {y}')
def then_connector_begin_y_is_an_Emu_object_with_value_y(context, y):
    begin_y = context.connector.begin_y
    assert isinstance(begin_y, Emu)
    assert begin_y == int(y)


@then('connector.end_x == {value}')
def then_connector_end_x_equals_value(context, value):
    assert context.connector.end_x == int(value)


@then('connector.end_x is an Emu object with value {x}')
def then_connector_end_x_is_an_Emu_object_with_value_x(context, x):
    end_x = context.connector.end_x
    assert isinstance(end_x, Emu)
    assert end_x == int(x)


@then('connector.end_y == {value}')
def then_connector_end_y_equals_value(context, value):
    assert context.connector.end_y == int(value)


@then('connector.end_y is an Emu object with value {y}')
def then_connector_end_y_is_an_Emu_object_with_value_y(context, y):
    end_y = context.connector.end_y
    assert isinstance(end_y, Emu)
    assert end_y == int(y)


@then('group_shape.shapes is a GroupShapes object')
def then_group_shape_shapes_is_a_GroupShapes_object(context):
    class_name = context.group_shape.shapes.__class__.__name__
    assert class_name == 'GroupShapes', 'got %s' % class_name


@then('movie is a Movie object')
def then_movie_is_a_Movie_object(context):
    class_name = context.movie.__class__.__name__
    assert class_name == 'Movie', 'got %s' % class_name


@then("movie.left, movie.top == x, y")
def then_movie_left_movie_top_eq_x_y(context):
    movie = context.movie
    position = movie.left, movie.top
    assert position == (Emu(2590800), Emu(571500)), 'got %s' % position


@then('movie.media_format is a _MediaFormat object')
def then_movie_media_format_is_a_MediaFormat_object(context):
    class_name = context.movie.media_format.__class__.__name__
    assert class_name == '_MediaFormat', 'got %s' % class_name


@then('movie.media_type is PP_MEDIA_TYPE.MOVIE')
def then_movie_media_type_is_PP_MEDIA_TYPE_MOVIE(context):
    media_type = context.movie.media_type
    assert media_type == PP_MEDIA_TYPE.MOVIE, 'got %s' % media_type


@then("movie.poster_frame is the same image as poster_frame")
def then_movie_poster_frame_is_the_same_image_as_poster_frame(context):
    actual_sha1 = context.movie.poster_frame.sha1
    with open(test_file('just-two-mice.png'), 'rb') as f:
        expected_sha1 = hashlib.sha1(f.read()).hexdigest()
    assert actual_sha1 == expected_sha1, 'not the same image'


@then('movie.shape_type is MSO_SHAPE_TYPE.MEDIA')
def then_movie_shape_type_is_MSO_SHAPE_TYPE_MEDIA(context):
    shape_type = context.movie.shape_type
    assert shape_type == MSO_SHAPE_TYPE.MEDIA, 'got %s' % shape_type


@then("movie.width, movie.height == cx, cy")
def then_movie_width_movie_height_eq_cx_cy(context):
    movie = context.movie
    size = movie.width, movie.height
    assert size == (Emu(3962400), Emu(5715000)), 'got %s' % size


@then('picture.crop_{side} == {value}')
def then_picture_crop_side_eq_value(context, side, value):
    expected_value = round(float(value), 5)
    actual_value = round(getattr(context.picture, 'crop_%s' % side), 5)
    assert actual_value == expected_value, (
        'picture.crop_%s == %s' % (side, actual_value)
    )


@then('picture.image is an Image object')
def then_picture_image_is_an_Image_object(context):
    class_name = context.picture.image.__class__.__name__
    assert class_name == 'Image', (
        'picture.image is a %s object' % class_name
    )


@then('shape.adjustments[0] is 0.15')
def then_shape_adjustments_is_value(context):
    shape = context.shape
    assert shape.adjustments[0] == 0.15


@then('shape.chart is a Chart object')
def then_shape_chart_is_a_Chart_object(context):
    chart = context.shape.chart
    class_name = chart.__class__.__name__
    assert class_name == 'Chart', 'got %s' % class_name


@then("shape.click_action is an ActionSetting object")
def then_shape_click_action_is_an_ActionSetting_object(context):
    assert isinstance(context.shape.click_action, ActionSetting)


@then('shape.has_chart is {value}')
def then_shape_has_chart_is_value(context, value):
    expected_value = {'True': True, 'False': False}[value]
    actual_value = context.shape.has_chart
    assert actual_value is expected_value, (
        'shape.has_chart is %s' % actual_value
    )


@then('shape.has_table is {value}')
def then_shape_has_table_is_value(context, value):
    expected_value = {'True': True, 'False': False}[value]
    actual_value = context.shape.has_table
    assert actual_value is expected_value, (
        'shape.has_table is %s' % actual_value
    )


@then('shape.has_text_frame is {value_str}')
def then_shape_has_text_frame_is(context, value_str):
    expected_value = {'True': True, 'False': False}[value_str]
    has_text_frame = context.shape.has_text_frame
    assert has_text_frame is expected_value, 'got %s' % has_text_frame


@then('shape.height == {value}')
def then_shape_height_eq_value(context, value):
    expected_height = int(value)
    actual_height = context.shape.height
    assert actual_height == expected_height, (
        'shape.height == %s' % actual_height
    )


@then('shape.left == {value}')
def then_shape_left_eq_value(context, value):
    expected_left = int(value)
    actual_left = context.shape.left
    assert actual_left == expected_left, 'shape.left == %s' % actual_left


@then('shape.line is a LineFormat object')
def then_shape_line_is_a_LineFormat_object(context):
    shape = context.shape
    line_format = shape.line
    line_format_cls_name = cls_qname(line_format)
    expected_cls_name = 'pptx.dml.line.LineFormat'
    assert line_format_cls_name == expected_cls_name, (
        "expected '%s', got '%s'" % (expected_cls_name, line_format_cls_name)
    )


@then("shape.name == '{expected_value}'")
def then_shape_name_eq_value(context, expected_value):
    shape = context.shape
    msg = "expected shape name '%s', got '%s'" % (shape.name, expected_value)
    assert shape.name == expected_value, msg


@then('shape.part is a SlidePart object')
def then_shape_part_is_a_SlidePart_object(context):
    cls_name = type(context.shape.part).__name__
    expected_cls_name = 'SlidePart'
    assert cls_name == expected_cls_name, (
        "expected '%s', got '%s'" % (expected_cls_name, cls_name)
    )


@then('shape.part is slide.part')
def then_shape_part_is_slide_part(context):
    assert context.shape.part is context.slide.part


@then("shape.rotation == {value}")
def then_shape_rotation_eq_value(context, value):
    shape = context.shape
    expected_value = float(value)
    assert shape.rotation == expected_value, 'got %s' % expected_value


@then('shape.shadow is a ShadowFormat object')
def then_shape_shadow_is_a_ShadowFormat_object(context):
    cls_name = type(context.shape.shadow).__name__
    assert cls_name == 'ShadowFormat', (
        "shape.shadow is a '%s' object" % cls_name
    )


@then("shape.shadow raises NotImplementedError")
def then_shape_shadow_raises_NotImplementedError(context):
    try:
        context.shape.shadow
    except NotImplementedError:
        return
    except Exception as e:
        raise AssertionError(
            'shape.shadow raises %s' % type(e).__name__
        )
    raise AssertionError('shape.shadow did not raise')


@then('shape.shape_id == {value_str}')
def then_shape_shape_id_equals(context, value_str):
    expected_value = int(value_str)
    shape_id = context.shape.shape_id
    assert shape_id == expected_value, 'got %s' % shape_id


@then('shape.shape_type == MSO_SHAPE_TYPE.{member_name}')
def then_shape_shape_type_is_MSO_SHAPE_TYPE_member(context, member_name):
    expected_shape_type = getattr(MSO_SHAPE_TYPE, member_name)
    actual_shape_type = context.shape.shape_type
    assert actual_shape_type == expected_shape_type, (
        'shape.shape_type == %s' % actual_shape_type
    )


@then('shape.text is the string I assigned')
def then_shape_text_is_the_string_I_assigned(context):
    shape = context.shape
    assert shape.text == u'F\xf8o\nBar'


@then('shape.text is the text in the shape')
def then_shape_text_is_the_text_in_the_shape(context):
    shape = context.shape
    assert shape.text == u'Fee Fi\nF\xf8\xf8 Fum\nI am a shape\nwith textium'


@then('shape.top == {value}')
def then_shape_top_eq_value(context, value):
    expected_top = int(value)
    actual_top = context.shape.top
    assert actual_top == expected_top, 'shape.top == %s' % actual_top


@then('shape.width == {value}')
def then_shape_width_eq_value(context, value):
    expected_width = int(value)
    actual_width = context.shape.width
    assert actual_width == expected_width, 'shape.width == %s' % actual_width
