# encoding: utf-8

"""Gherkin step implementations for shape collections."""

import io

from behave import given, then, when

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, PP_PLACEHOLDER, PROG_ID
from pptx.shapes.base import BaseShape
from pptx.util import Emu, Inches

from helpers import saved_pptx_path, test_file, test_image, test_pptx


# given ===================================================


@given("a _BaseShapes object as shapes")
def given_a_BaseShapes_object_as_shapes(context):
    prs = Presentation()
    context.shapes = prs.slides.add_slide(prs.slide_layouts[6]).shapes


@given("a GroupShapes object as shapes")
@given("a GroupShapes object of length 3 as shapes")
def given_a_GroupShapes_object_of_length_3_as_shapes(context):
    prs = Presentation(test_pptx("shp-groupshape"))
    group_shape = prs.slides[0].shapes[0]
    context.shapes = group_shape.shapes


@given("a LayoutPlaceholders object of length 2 as shapes")
def given_a_LayoutPlaceholders_object_of_length_2_as_shapes(context):
    prs = Presentation(test_pptx("lyt-shapes"))
    context.shapes = prs.slide_layouts[0].placeholders


@given("a LayoutShapes object of length 3 as shapes")
def given_a_LayoutShapes_object_of_length_3_as_shapes(context):
    prs = Presentation(test_pptx("lyt-shapes"))
    context.shapes = prs.slide_layouts[0].shapes


@given("a MasterPlaceholders object of length 2 as shapes")
def given_a_MasterPlaceholders_object_of_length_2_as_shapes(context):
    prs = Presentation(test_pptx("mst-placeholders"))
    context.shapes = prs.slide_masters[0].placeholders


@given("a MasterShapes object of length 2 as shapes")
def given_a_MasterShapes_object_of_length_2_as_shapes(context):
    prs = Presentation(test_pptx("mst-shapes"))
    context.shapes = prs.slide_masters[0].shapes


@given("a {PROG_ID_member} file as ole_object_file")
def given_a_PROG_ID_member_file_as_ole_object_file(context, PROG_ID_member):
    filename = {
        "DOCX": "shp-embedded-docx.docx",
        "PPTX": "shp-embedded-pptx.pptx",
        "XLSX": "shp-embedded-xlsx.xlsx",
    }[PROG_ID_member]
    with open(test_file(filename), "rb") as f:
        context.ole_object_file = io.BytesIO(f.read())
    context.PROG_ID_member = PROG_ID_member


@given("a SlidePlaceholders object of length 2 as shapes")
def given_a_SlidePlaceholders_object_of_length_2_as_shapes(context):
    prs = Presentation(test_pptx("shp-shapes"))
    context.shapes = prs.slides[0].placeholders


@given("a SlideShapes object as shapes")
def given_a_SlideShapes_object_as_shapes(context):
    prs = Presentation(test_pptx("shp-shapes"))
    context.shapes = prs.slides[0].shapes


@given("a SlideShapes object containing {a_or_no} movies")
def given_a_SlideShapes_object_containing_a_or_no_movies(context, a_or_no):
    pptx = {"one or more": "shp-movie-props", "no": "shp-shapes"}[a_or_no]
    prs = Presentation(test_pptx(pptx))
    context.prs = prs
    context.shapes = prs.slides[0].shapes


@given("a SlideShapes object of length 6 shapes as shapes")
def given_a_SlideShapes_object_of_length_6_as_shapes(context):
    prs = Presentation(test_pptx("shp-shapes"))
    context.shapes = prs.slides[0].shapes


@given("a SlideShapes object having a {type} shape at offset {idx}")
def given_a_SlideShapes_obj_having_type_shape_at_off_idx(context, type, idx):
    prs = Presentation(test_pptx("shp-shapes"))
    context.shapes = prs.slides[1].shapes


# when ====================================================


@when("I add 100 shapes")
def when_I_add_100_shapes(context):
    X_ORIG = Y_ORIG = Inches(0.0625)
    X_INCR = Y_INCR = Inches(0.5)
    CX = CY = Inches(0.375)

    def iter_corner():
        y = Y_ORIG
        while True:
            for i in range(20):
                x = X_ORIG + (X_INCR * i)
                yield x, y
            y += Y_INCR

    shapes = context.shapes
    corners = iter_corner()
    for i in range(100):
        x, y = next(corners)
        shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, CX, CY)


@when("I add a table to the slide's shape collection")
def when_I_call_shapes_add_table(context):
    shapes = context.slide.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    shapes.add_table(2, 2, x, y, cx, cy)


@when("I assign shape.ole_format to ole_format")
def when_I_assign_shape_ole_format_to_ole_format(context):
    context.ole_format = context.shape.ole_format


@when("I assign shapes.add_chart() to shape")
def when_I_assign_shapes_add_chart_to_shape(context):
    chart_data = CategoryChartData()
    chart_data.categories = ("Foo", "Bar")
    chart_data.add_series("East", (1.0, 2.0))
    chart_data.add_series("West", (3.0, 4.0))

    context.shape = context.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5),
        chart_data,
    )


@when("I assign shapes.add_connector() to shape")
def when_I_assign_shapes_add_connector_to_shape(context):
    context.shape = context.shapes.add_connector(MSO_CONNECTOR.CURVE, 4, 3, 2, 1)


@when("I assign shapes.add_group_shape() to shape")
def when_I_assign_shapes_add_group_shape_to_shape(context):
    context.shape = context.shapes.add_group_shape()


@when("I assign shapes.add_ole_object(ole_object_file) to shape")
def when_I_assign_shapes_add_ole_object_to_shape(context):
    context.shape = context.shapes.add_ole_object(
        context.ole_object_file, getattr(PROG_ID, context.PROG_ID_member), 4, 3, 2, 1
    )


@when("I assign shapes.add_picture() to shape")
def when_I_assign_shapes_add_picture_to_shape(context):
    context.shape = context.shapes.add_picture(
        test_image("sonic.gif"), Inches(1), Inches(2)
    )


@when("I assign shapes.add_shape() to shape")
def when_I_assign_shapes_add_shape_to_shape(context):
    context.shape = context.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3), Inches(1), Inches(0.5)
    )


@when("I assign shapes.add_textbox() to shape")
def when_I_assign_shapes_add_textbox_to_shape(context):
    context.shape = context.shapes.add_textbox(
        Inches(1), Inches(2), Inches(3), Inches(0.5)
    )


@when("I assign shapes.build_freeform() to builder")
def when_I_assign_shapes_build_freeform_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform()
    context.builder = builder


@when("I assign shapes.build_freeform(scale=100.0) to builder")
def when_I_assign_shapes_build_freeform_scale_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(scale=100.0)
    context.builder = builder


@when("I assign shapes.build_freeform(scale=(200.0, 100.0)) to builder")
def when_I_assign_shapes_build_freeform_scale_rectnglr_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(scale=(200.0, 100.0))
    context.builder = builder


@when("I assign shapes.build_freeform(start_x=25, start_y=125) to builder")
def when_I_assign_shapes_build_freeform_start_x_start_y_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(25, 125)
    context.builder = builder


@when("I assign True to shapes.turbo_add_enabled")
def when_I_assign_True_to_shapes_turbo_add_enabled(context):
    context.shapes.turbo_add_enabled = True


@when("I call shapes.add_chart({type_}, chart_data)")
def when_I_call_shapes_add_chart(context, type_):
    chart_type = getattr(XL_CHART_TYPE, type_)
    context.chart = context.shapes.add_chart(
        chart_type, 0, 0, 0, 0, context.chart_data
    ).chart


@when("I call shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 1, 2, 3, 4)")
def when_I_call_shapes_add_connector(context):
    context.connector = context.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 1, 2, 3, 4)


@when("I call shapes.add_movie(file, x, y, cx, cy, poster_frame)")
def when_I_call_shapes_add_movie(context):
    shapes = context.shapes
    x, y, cx, cy = Emu(2590800), Emu(571500), Emu(3962400), Emu(5715000)
    context.movie = shapes.add_movie(
        test_file("just-two-mice.mp4"), x, y, cx, cy, test_file("just-two-mice.png")
    )


# then ====================================================


@then("iterating shapes produces {count} objects of type {class_name}")
def then_iterating_shapes_produces_count_objects_of_type_class_name(
    context, count, class_name
):
    shapes = context.shapes
    expected_count, expected_class_name = int(count), class_name
    idx = -1
    for idx, shape in enumerate(shapes):
        actual_class_name = shape.__class__.__name__
        assert actual_class_name == expected_class_name, (
            "shape.__class__.__name__ == %s" % actual_class_name
        )
    actual_count = idx + 1
    assert actual_count == expected_count, "got %d items" % actual_count


@then("iterating shapes produces {count} objects that subclass BaseShape")
def then_iterating_shapes_produces_count_objects_that_subclass_BaseShape(
    context, count
):
    shapes = context.shapes
    expected_count = int(count)
    idx = -1
    for idx, shape in enumerate(shapes):
        class_name = shape.__class__.__name__
        assert isinstance(shape, BaseShape), (
            "%s does not subclass BaseShape" % class_name
        )
    actual_count = idx + 1
    assert actual_count == expected_count, "got %d items" % actual_count


@then("len(shapes) == {value}")
def then_len_shapes_eq_value(context, value):
    expected_len = int(value)
    actual_len = len(context.shapes)
    assert actual_len == expected_len, "len(shapes) == %s" % actual_len


@then("shape is a {clsname} object")
def then_shape_is_a_type_object(context, clsname):
    actual_class_name = context.shape.__class__.__name__
    expected_class_name = clsname
    assert actual_class_name == expected_class_name, (
        "shape is a %s object" % actual_class_name
    )


@then("shapes[-1] == shape")
def then_shapes_minus_1_eq_shape(context):
    shapes, shape = context.shapes, context.shape
    assert shapes[-1] == shape


@then("shapes[{idx}] is a {type_} object")
def then_shapes_idx_is_a_type_object(context, idx, type_):
    shapes = context.shapes
    type_name = type(shapes[int(idx)]).__name__
    assert type_name == type_, "got %s" % type_name


@then("shapes.get(idx=10) is the body placeholder")
def then_shapes_get_10_is_the_body_placeholder(context):
    shapes = context.shapes
    title_placeholder = shapes.get(idx=0)
    body_placeholder = shapes.get(idx=10)
    assert title_placeholder._element is shapes[0]._element
    assert body_placeholder._element is shapes[1]._element


@then("shapes.get(PP_PLACEHOLDER.BODY) is the body placeholder")
def then_shapes_get_PP_PLACEHOLDER_BODY_is_the_body_ph(context):
    shapes = context.shapes
    title_placeholder = shapes.get(PP_PLACEHOLDER.TITLE)
    body_placeholder = shapes.get(PP_PLACEHOLDER.BODY)
    assert title_placeholder._element is shapes[0]._element
    assert body_placeholder._element is shapes[1]._element


@then("shapes.index(shape) for each shape matches its sequence position")
def then_shapes_index_for_each_shape_matches_sequence_position(context):
    shapes = context.shapes
    for idx, shape in enumerate(shapes):
        assert idx == shapes.index(shape), "index doesn't match for idx == %s" % idx


@then("shapes.title is the title placeholder")
def then_shapes_title_is_the_title_placeholder(context):
    shapes = context.shapes
    title_placeholder = shapes.title
    assert title_placeholder.element is shapes[0].element
    assert title_placeholder.shape_id == 4


@then("shapes.turbo_add_enabled is False")
def then_shapes_turbo_add_enabled_is_False(context):
    shapes = context.shapes
    assert shapes.turbo_add_enabled is False


@then("the table appears in the slide")
def then_the_table_appears_in_the_slide(context):
    prs = Presentation(saved_pptx_path)
    expected_table_graphic_frame = prs.slides[0].shapes[0]
    assert expected_table_graphic_frame.has_table
