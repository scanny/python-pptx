# encoding: utf-8

"""
Step implementations for shape tree-related features
"""

from __future__ import absolute_import

from behave import given, then

from pptx import Presentation
from pptx.shapes.shape import BaseShape
from pptx.shapes.shapetree import BaseShapeTree

from .helpers import test_pptx


# given ===================================================

@given('a slide having two shapes')
def given_slide_having_two_shapes(context):
    presentation = Presentation(test_pptx('sld-access-shapes'))
    context.slide = presentation.slides[0]


@given('a shape collection containing two shapes')
def given_shape_collection_containing_two_shapes(context):
    presentation = Presentation(test_pptx('sld-access-shapes'))
    slide = presentation.slides[0]
    context.shapes = slide.shapes_new


# then ====================================================

@then('I can access a shape by index')
def then_can_access_shape_by_index(context):
    shapes = context.shapes
    for idx in range(2):
        shape = shapes[idx]
        assert isinstance(shape, BaseShape)


@then('I can access the shape collection of the slide')
def then_can_access_shapes_of_slide(context):
    slide = context.slide
    shapes = slide.shapes_new
    msg = 'Slide.shapes not instance of ShapeTree'
    assert isinstance(shapes, BaseShapeTree), msg


@then('I can iterate over the shapes')
def then_can_iterate_over_the_shapes(context):
    shapes = context.shapes
    actual_count = 0
    for shape in shapes:
        actual_count += 1
        assert isinstance(shape, BaseShape)
    assert actual_count == 2


@then('the length of the shape collection is 2')
def then_len_of_shape_collection_is_2(context):
    slide = context.slide
    shapes = slide.shapes_new
    assert len(shapes) == 2, (
        'expected len(shapes) of 2, got %s' % len(shapes)
    )
