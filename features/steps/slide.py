# encoding: utf-8

"""Gherkin step implementations for slide-related features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then

from pptx import Presentation

from helpers import test_pptx


# given ===================================================


@given("a blank slide")
def given_a_blank_slide(context):
    context.prs = Presentation(test_pptx("sld-blank"))
    context.slide = context.prs.slides[0]


@given("a notes slide")
def given_a_notes_slide(context):
    prs = Presentation(test_pptx("sld-notes"))
    context.notes_slide = prs.slides[0].notes_slide


@given("a slide")
def given_a_slide(context):
    presentation = Presentation(test_pptx("shp-shapes"))
    context.slide = presentation.slides[0]


@given("a slide having a notes slide")
def given_a_slide_having_a_notes_slide(context):
    context.slide = Presentation(test_pptx("sld-notes")).slides[0]


@given("a slide having no notes slide")
def given_a_slide_having_no_notes_slide(context):
    context.slide = Presentation(test_pptx("sld-notes")).slides[1]


@given("a slide having a title")
def given_a_slide_having_a_title(context):
    prs = Presentation(test_pptx("shp-shapes"))
    context.prs, context.slide = prs, prs.slides[0]


@given("a slide having name {name}")
def given_a_slide_having_name_name(context, name):
    slide_idx = 0 if name == "Overview" else 1
    presentation = Presentation(test_pptx("sld-slide"))
    context.slide = presentation.slides[slide_idx]


@given("a slide having slide id 256")
def given_a_slide_having_slide_id_256(context):
    presentation = Presentation(test_pptx("shp-shapes"))
    context.slide = presentation.slides[0]


@given("a Slide object based on slide_layout as slide")
def given_a_Slide_object_based_on_slide_layout_as_slide(context):
    context.slide = context.prs.slides[0]


@given("a Slide object having {def_or_ovr} background as slide")
def given_a_Slide_object_having_background_as_slide(context, def_or_ovr):
    slide_idx = {"the default": 0, "an overridden": 1}[def_or_ovr]
    context.slide = Presentation(test_pptx("sld-slide")).slides[slide_idx]


@given("a SlideLayout object as slide")
@given("a SlideLayout object as slide_layout")
def given_a_SlideLayout_object_as_slide(context):
    prs = Presentation(test_pptx("sld-slide"))
    context.slide = context.slide_layout = prs.slide_layouts[0]


@given("a SlideLayout object having name {name} as slide")
def given_a_SlideLayout_object_having_name_as_slide(context, name):
    slide_layout_idx = 0 if name == "of no explicit value" else 1
    prs = Presentation(test_pptx("sld-slide"))
    context.slide = prs.slide_layouts[slide_layout_idx]


@given("a SlideLayout object used by {which_slides} as slide_layout")
def given_a_SlideLayout_object_used_by_slides_as_slide_layout(context, which_slides):
    slide_layout_idx = {"a slide": 0, "no slides": 1}[which_slides]
    context.prs = Presentation(test_pptx("sld-slide"))
    context.slide_layout = context.prs.slide_layouts[slide_layout_idx]


@given("a SlideMaster object as slide")
@given("a SlideMaster object as slide_master")
def given_a_SlideMaster_object_as_slide(context):
    prs = Presentation(test_pptx("sld-slide"))
    context.slide = context.slide_master = prs.slide_masters[0]


# then ====================================================


@then("len(notes_slide.shapes) is {count}")
def then_len_notes_slide_shapes_is_count(context, count):
    shapes = context.notes_slide.shapes
    assert len(shapes) == int(count)


@then("notes_slide.notes_placeholder is a NotesSlidePlaceholder object")
def then_notes_slide_notes_placeholder_is_a_NotesSlidePlacehldr_obj(context):
    notes_slide = context.notes_slide
    cls_name = type(notes_slide.notes_placeholder).__name__
    assert cls_name == "NotesSlidePlaceholder", "got %s" % cls_name


@then("notes_slide.notes_text_frame is a TextFrame object")
def then_notes_slide_notes_text_frame_is_a_TextFrame_object(context):
    notes_slide = context.notes_slide
    cls_name = type(notes_slide.notes_text_frame).__name__
    assert cls_name == "TextFrame", "got %s" % cls_name


@then("notes_slide.placeholders is a NotesSlidePlaceholders object")
def then_notes_slide_placeholders_is_a_NotesSlidePlaceholders_object(context):
    notes_slide = context.notes_slide
    assert type(notes_slide.placeholders).__name__ == "NotesSlidePlaceholders"


@then("slide in slide_layout.used_by_slides is True")
def then_slide_in_slide_layout_used_by_slides_is_True(context):
    assert context.slide in context.slide_layout.used_by_slides


@then("slide.background is a _Background object")
def then_slide_background_is_a_Background_object(context):
    cls_name = context.slide.background.__class__.__name__
    assert cls_name == "_Background", "slide.background is a %s object" % cls_name


@then("slide.follow_master_background is {value}")
def then_slide_follow_master_background_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    actual_value = context.slide.follow_master_background
    assert actual_value is expected_value, (
        "slide.follow_master_background is %s" % actual_value
    )


@then("slide.has_notes_slide is {value}")
def then_slide_has_notes_slide_is_value(context, value):
    expected_value = {"True": True, "False": False}[value]
    slide = context.slide
    assert slide.has_notes_slide is expected_value


@then("slide.name is {value}")
def then_slide_name_is_value(context, value):
    expected_name = "" if value == "the empty string" else value
    actual_name = context.slide.name
    assert actual_name == expected_name, "slide.name == %s" % actual_name


@then("slide.notes_slide is a NotesSlide object")
def then_slide_notes_slide_is_a_NotesSlide_object(context):
    notes_slide = context.notes_slide = context.slide.notes_slide
    assert type(notes_slide).__name__ == "NotesSlide"


@then("slide.placeholders is a {clsname} object")
def then_slide_placeholders_is_a_clsname_object(context, clsname):
    actual_clsname = context.slide.placeholders.__class__.__name__
    expected_clsname = clsname
    assert actual_clsname == expected_clsname, (
        "slide.placeholders is a %s object" % actual_clsname
    )


@then("slide.shapes is a {clsname} object")
def then_slide_shapes_is_a_clsname_object(context, clsname):
    actual_clsname = context.slide.shapes.__class__.__name__
    expected_clsname = clsname
    assert actual_clsname == expected_clsname, (
        "slide.shapes is a %s object" % actual_clsname
    )


@then("slide.slide_id is 256")
def then_slide_slide_id_is_256(context):
    slide = context.slide
    assert slide.slide_id == 256


@then("slide.slide_layout is the one passed in the call")
def then_slide_slide_layout_is_the_one_passed_in_the_call(context):
    slide = context.prs.slides[3]
    assert slide.slide_layout == context.slide_layout


@then("slide_layout.slide_master is a SlideMaster object")
def then_slide_layout_slide_master_is_a_SlideMaster_object(context):
    slide_layout = context.slide_layout
    assert type(slide_layout.slide_master).__name__ == "SlideMaster"


@then("slide_master.slide_layouts is a SlideLayouts object")
def then_slide_master_slide_layouts_is_a_SlideLayouts_object(context):
    slide_master = context.slide_master
    assert type(slide_master.slide_layouts).__name__ == "SlideLayouts"


@then("slide_layout.used_by_slides == ()")
def then_slide_layout_used_by_slides_eq_empty_tuple(context):
    assert context.slide_layout.used_by_slides == ()
