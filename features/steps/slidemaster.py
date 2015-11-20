# encoding: utf-8

"""
Step implementations for slide master-related features
"""

from __future__ import absolute_import

from behave import given, then

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.parts.slidelayout import SlideLayout
from pptx.parts.slidemaster import (
    _MasterPlaceholders, _MasterShapeTree, _SlideLayouts
)
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import MasterPlaceholder

from helpers import test_pptx


# given ===================================================

@given('a master placeholder collection')
def given_master_placeholder_collection(context):
    prs = Presentation(test_pptx('mst-placeholders'))
    context.master_placeholders = prs.slide_master.placeholders


@given('a master shape collection containing two shapes')
def given_master_shape_collection_containing_two_shapes(context):
    prs = Presentation(test_pptx('mst-shapes'))
    context.master_shapes = prs.slide_master.shapes


@given('a slide master having two placeholders')
def given_master_having_two_placeholders(context):
    prs = Presentation(test_pptx('mst-placeholders'))
    context.slide_master = prs.slide_master


@given('a slide master having two shapes')
def given_slide_master_having_two_shapes(context):
    prs = Presentation(test_pptx('mst-shapes'))
    context.slide_master = prs.slide_master


@given('a slide master having two slide layouts')
def given_slide_master_having_two_layouts(context):
    prs = Presentation(test_pptx('mst-slide-layouts'))
    context.slide_master = prs.slide_master


@given('a slide layout collection containing two layouts')
def given_slide_layout_collection_containing_two_layouts(context):
    prs = Presentation(test_pptx('mst-slide-layouts'))
    context.slide_layouts = prs.slide_master.slide_layouts


# then ====================================================

@then('I can access a master placeholder by index')
def then_can_access_master_placeholder_by_index(context):
    master_placeholders = context.master_placeholders
    for idx in range(2):
        master_placeholder = master_placeholders[idx]
        assert isinstance(master_placeholder, MasterPlaceholder)


@then('I can access a master placeholder by type')
def then_can_access_master_placeholder_by_type(context):
    master_placeholders = context.master_placeholders
    title_placeholder = master_placeholders.get(PP_PLACEHOLDER.TITLE)
    body_placeholder = master_placeholders.get(PP_PLACEHOLDER.BODY)
    assert title_placeholder._element is master_placeholders[0]._element
    assert body_placeholder._element is master_placeholders[1]._element


@then('I can access a master shape by index')
def then_can_access_master_shape_by_index(context):
    master_shapes = context.master_shapes
    for idx in range(2):
        master_shape = master_shapes[idx]
        assert isinstance(master_shape, BaseShape)


@then('I can access a slide layout by index')
def then_can_access_slide_layout_by_index(context):
    slide_layouts = context.slide_layouts
    for idx in range(2):
        slide_layout = slide_layouts[idx]
        assert isinstance(slide_layout, SlideLayout)


@then('I can access the placeholder collection of the slide master')
def then_can_access_placeholder_collection_of_slide_master(context):
    slide_master = context.slide_master
    master_placeholders = slide_master.placeholders
    msg = 'SlideMaster.placeholders not instance of _MasterPlaceholders'
    assert isinstance(master_placeholders, _MasterPlaceholders), msg


@then('I can access the shape collection of the slide master')
def then_can_access_shape_collection_of_slide_master(context):
    slide_master = context.slide_master
    master_shapes = slide_master.shapes
    msg = 'SlideMaster.shapes not instance of _MasterShapeTree'
    assert isinstance(master_shapes, _MasterShapeTree), msg


@then('I can access the slide layouts of the slide master')
def then_can_access_slide_layouts_of_slide_master(context):
    slide_master = context.slide_master
    slide_layouts = slide_master.slide_layouts
    msg = 'SlideMaster.slide_layouts not instance of _SlideLayouts'
    assert isinstance(slide_layouts, _SlideLayouts), msg


@then('I can iterate over the master placeholders')
def then_can_iterate_over_the_master_placeholders(context):
    master_placeholders = context.master_placeholders
    actual_count = 0
    for master_placeholder in master_placeholders:
        actual_count += 1
        assert isinstance(master_placeholder, MasterPlaceholder)
    assert actual_count == 2


@then('I can iterate over the master shapes')
def then_can_iterate_over_the_master_shapes(context):
    master_shapes = context.master_shapes
    actual_count = 0
    for master_shape in master_shapes:
        actual_count += 1
        assert isinstance(master_shape, BaseShape)
    assert actual_count == 2


@then('I can iterate over the slide layouts')
def then_can_iterate_over_the_slide_layouts(context):
    slide_layouts = context.slide_layouts
    actual_count = 0
    for slide_layout in slide_layouts:
        actual_count += 1
        assert isinstance(slide_layout, SlideLayout)
    assert actual_count == 2


@then('the length of the master shape collection is 2')
def then_len_of_master_shape_collection_is_2(context):
    slide_master = context.slide_master
    master_shapes = slide_master.shapes
    assert len(master_shapes) == 2, (
        'expected len(master_shapes) of 2, got %s' % len(master_shapes)
    )


@then('the length of the master placeholder collection is 2')
def then_len_of_placeholder_collection_is_2(context):
    slide_master = context.slide_master
    master_placeholders = slide_master.placeholders
    assert len(master_placeholders) == 2, (
        'expected len(master_placeholders) of 2, got %s' %
        len(master_placeholders)
    )


@then('the length of the slide layout collection is 2')
def then_len_of_slide_layout_collection_is_2(context):
    slide_master = context.slide_master
    slide_layouts = slide_master.slide_layouts
    assert len(slide_layouts) == 2, (
        'expected len(slide_layouts) of 2, got %s' % len(slide_layouts)
    )
