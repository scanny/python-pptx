# encoding: utf-8

"""Gherkin step implementations for slide collection-related features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, when, then

from pptx import Presentation

from helpers import test_pptx


# given ===================================================


@given("a SlideLayouts object containing 2 layouts as slide_layouts")
def given_a_SlideLayouts_object_containing_2_layouts(context):
    prs = Presentation(test_pptx("mst-slide-layouts"))
    context.slide_layouts = prs.slide_master.slide_layouts


@given("a SlideMasters object containing 2 masters")
def given_a_SlideMasters_object_containing_2_masters(context):
    prs = Presentation(test_pptx("prs-slide-masters"))
    context.slide_masters = prs.slide_masters


@given("a Slides object containing 3 slides")
def given_a_Slides_object_containing_3_slides(context):
    prs = Presentation(test_pptx("sld-slides"))
    context.prs = prs
    context.slides = prs.slides


# when ====================================================


@when("I call slides.add_slide()")
def when_I_call_slides_add_slide(context):
    context.slide_layout = context.prs.slide_masters[0].slide_layouts[0]
    context.slides.add_slide(context.slide_layout)


@when("I call slide_layouts.remove(slide_layouts[1])")
def when_I_call_slide_layouts_remove(context):
    slide_layouts = context.slide_layouts
    slide_layouts.remove(slide_layouts[1])


# then ====================================================


@then("iterating produces 3 NotesSlidePlaceholder objects")
def then_iterating_produces_3_NotesSlidePlaceholder_objects(context):
    idx = -1
    for idx, placeholder in enumerate(context.notes_slide.placeholders):
        typename = type(placeholder).__name__
        assert typename == "NotesSlidePlaceholder", "got %s" % typename
    assert idx == 2


@then("iterating slide_layouts produces 2 SlideLayout objects")
def then_iterating_slide_layouts_produces_2_SlideLayout_objects(context):
    slide_layouts = context.slide_layouts
    idx = -1
    for idx, slide_layout in enumerate(slide_layouts):
        assert type(slide_layout).__name__ == "SlideLayout"
    assert idx == 1


@then("iterating slide_masters produces 2 SlideMaster objects")
def then_iterating_slide_masters_produces_2_SlideMaster_objects(context):
    slide_masters = context.slide_masters
    idx = -1
    for idx, slide_master in enumerate(slide_masters):
        assert type(slide_master).__name__ == "SlideMaster"
    assert idx == 1


@then("iterating slides produces 3 Slide objects")
def then_iterating_slides_produces_3_Slide_objects(context):
    slides = context.slides
    idx = -1
    for idx, slide in enumerate(slides):
        assert type(slide).__name__ == "Slide"
    assert idx == 2


@then("len(slides) is {count}")
def then_len_slides_is_count(context, count):
    slides = context.slides
    assert len(slides) == int(count)


@then("len(slide_layouts) is {n}")
def then_len_slide_layouts_is_2(context, n):
    assert len(context.slide_layouts) == int(n)


@then("len(slide_masters) is 2")
def then_len_slide_masters_is_2(context):
    slide_masters = context.slide_masters
    assert len(slide_masters) == 2


@then("slide_layouts[1] is a SlideLayout object")
def then_slide_layouts_1_is_a_SlideLayout_object(context):
    slide_layouts = context.slide_layouts
    assert type(slide_layouts[1]).__name__ == "SlideLayout"


@then("slide_layouts.get_by_name(slide_layouts[1].name) is slide_layouts[1]")
def then_slide_layouts_get_by_name_is_slide_layout(context):
    slide_layouts = context.slide_layouts
    assert slide_layouts.get_by_name(slide_layouts[1].name) is slide_layouts[1]


@then("slide_layouts.index(slide_layouts[1]) == 1")
def then_slide_layouts_index_is_1(context):
    slide_layouts = context.slide_layouts
    assert slide_layouts.index(slide_layouts[1]) == 1


@then("slide_masters[1] is a SlideMaster object")
def then_slide_masters_1_is_a_SlideMaster_object(context):
    slide_masters = context.slide_masters
    assert type(slide_masters[1]).__name__ == "SlideMaster"


@then("slides.get(256) is slides[0]")
def then_slides_get_256_is_slides_0(context):
    slides = context.slides
    assert slides.get(256) is slides[0]


@then("slides.get(666, default=slides[2]) is slides[2]")
def then_slides_get_666_default_slides_2_is_slides_2(context):
    slides = context.slides
    assert slides.get(666, default=slides[2]) is slides[2]


@then("slides[2] is a Slide object")
def then_slides_2_is_a_Slide_object(context):
    slides = context.slides
    assert type(slides[2]).__name__ == "Slide"
