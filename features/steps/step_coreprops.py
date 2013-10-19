# encoding: utf-8

"""
Gherkin step implementations for core properties-related features.
"""

import os

from datetime import datetime, timedelta

from behave import given, when, then
from hamcrest import assert_that, is_, less_than

from pptx import Presentation


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
scratch_dir = absjoin(thisdir, '../_scratch')
test_file_dir = absjoin(thisdir, '../../tests/test_files')
no_core_props_pptx_path = absjoin(test_file_dir, 'no-core-props.pptx')
saved_pptx_path = absjoin(scratch_dir, 'test_out.pptx')

test_text = "python-pptx was here!"


# given ===================================================

@given('I have a reference to the core properties of a presentation')
def step_given_ref_to_core_doc_props(context):
    context.prs = Presentation()
    context.core_properties = context.prs.core_properties


# when ====================================================

@when('I open a presentation having no core properties part')
def step_when_open_presentation_with_no_core_props_part(context):
    context.prs = Presentation(no_core_props_pptx_path)


@when("I set the core properties to valid values")
def step_when_set_core_doc_props_to_valid_values(context):
    context.propvals = (
        ('author', 'Creator'),
        ('category', 'Category'),
        ('comments', 'Description'),
        ('content_status', 'Content Status'),
        ('created', datetime(2013, 6, 15, 12, 34, 56)),
        ('identifier', 'Identifier'),
        ('keywords', 'key; word; keyword'),
        ('language', 'Language'),
        ('last_modified_by', 'Last Modified By'),
        ('last_printed', datetime(2013, 6, 15, 12, 34, 56)),
        ('modified', datetime(2013, 6, 15, 12, 34, 56)),
        ('revision', 9),
        ('subject', 'Subject'),
        ('title', 'Title'),
        ('version', 'Version'),
    )
    for name, value in context.propvals:
        setattr(context.prs.core_properties, name, value)


# then ====================================================

@then('a core properties part with default values is added')
def step_then_a_core_props_part_with_def_vals_is_added(context):
    core_props = context.prs.core_properties
    assert_that(core_props.title, is_('PowerPoint Presentation'))
    assert_that(core_props.last_modified_by, is_('python-pptx'))
    assert_that(core_props.revision, is_(1))
    # core_props.modified only stores time with seconds resolution, so
    # comparison needs to be a little loose (within two seconds)
    modified_timedelta = datetime.utcnow() - core_props.modified
    max_expected_timedelta = timedelta(seconds=2)
    assert_that(modified_timedelta, less_than(max_expected_timedelta))


@then('the core properties of the presentation have the values I set')
def step_then_core_props_have_values_previously_set(context):
    core_props = Presentation(saved_pptx_path).core_properties
    for name, value in context.propvals:
        reason = "for core property '%s'" % name
        assert_that(getattr(core_props, name), is_(value), reason)
