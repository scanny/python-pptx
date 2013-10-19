# encoding: utf-8

"""
Gherkin step implementations for picture-related features.
"""

import os

from StringIO import StringIO

from behave import when, then
from hamcrest import assert_that, has_item

from pptx import packaging
from pptx import Presentation
from pptx.util import Inches


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
scratch_dir = absjoin(thisdir, '../_scratch')
test_file_dir = absjoin(thisdir, '../../tests/test_files')
saved_pptx_path = absjoin(scratch_dir, 'test_out.pptx')
test_image_path = absjoin(test_file_dir, 'python-powered.png')

test_text = "python-pptx was here!"


# when ====================================================

@when("I add a picture stream to the slide's shape collection")
def step_when_add_picture_stream(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.25), Inches(1.25))
    with open(test_image_path) as f:
        stream = StringIO(f.read())
    shapes.add_picture(stream, x, y)


@when("I add a picture to the slide's shape collection")
def step_when_add_picture(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.25), Inches(1.25))
    shapes.add_picture(test_image_path, x, y)


# then ====================================================

@then('the image is saved in the pptx file')
def step_then_img_saved_in_pptx_file(context):
    pkgng_pkg = packaging.Package().open(saved_pptx_path)
    partnames = [part.partname for part in pkgng_pkg.parts
                 if part.partname.startswith('/ppt/media/')]
    assert_that(partnames, has_item('/ppt/media/image1.png'))


@then('the picture appears in the slide')
def step_then_picture_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    shapes = sld.shapes
    classnames = [sp.__class__.__name__ for sp in shapes]
    assert_that(classnames, has_item('_Picture'))
