# encoding: utf-8

"""
Gherkin step implementations for placeholder-related features.
"""

from __future__ import absolute_import

from behave import given, when, then
from hamcrest import assert_that, equal_to, is_

from pptx import Presentation

from .helpers import saved_pptx_path, test_text


# given ===================================================

@given('a bullet body placeholder')
def given_a_bullet_body_placeholder(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[1]
    context.sld = context.prs.slides.add_slide(slidelayout)
    context.body = context.sld.shapes.placeholders[1]


# when ====================================================

@when('I indent the first paragraph')
def step_when_indent_first_paragraph(context):
    context.body.textframe.paragraphs[0].level = 1


@when("I set the title text of the slide")
def step_when_set_slide_title_text(context):
    context.sld.shapes.title.text = test_text


# then ====================================================

@then('the paragraph is indented')
def then_paragraph_is_indented(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    body = sld.shapes.placeholders[1]
    p = body.textframe.paragraphs[0]
    assert_that(p.level, is_(equal_to(1)))


@then('the text appears in the title placeholder')
def step_then_text_appears_in_title_placeholder(context):
    prs = Presentation(saved_pptx_path)
    title_shape = prs.slides[0].shapes.title
    title_text = title_shape.textframe.paragraphs[0].runs[0].text
    assert_that(title_text, is_(equal_to(test_text)))
