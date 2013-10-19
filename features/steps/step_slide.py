# encoding: utf-8

"""
Gherkin step implementations for slide-related features.
"""

import os

from behave import given, when, then
from hamcrest import assert_that, equal_to, is_

from pptx import Presentation


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
scratch_dir = absjoin(thisdir, '../_scratch')
saved_pptx_path = absjoin(scratch_dir, 'test_out.pptx')


# given ===================================================

@given('I have a reference to a blank slide')
def step_given_ref_to_blank_slide(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[6]
    context.sld = context.prs.slides.add_slide(slidelayout)


@given('I have a reference to a slide')
def step_given_ref_to_slide(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[0]
    context.sld = context.prs.slides.add_slide(slidelayout)


# when ====================================================

@when('I add a new slide')
def step_when_add_slide(context):
    slidelayout = context.prs.slidemasters[0].slidelayouts[0]
    context.prs.slides.add_slide(slidelayout)


# then ====================================================

@then('the pptx file contains a single slide')
def step_then_pptx_file_contains_single_slide(context):
    prs = Presentation(saved_pptx_path)
    assert_that(len(prs.slides), is_(equal_to(1)))
