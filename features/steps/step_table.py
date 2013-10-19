# encoding: utf-8

"""
Gherkin step implementations for table-related features.
"""

from __future__ import absolute_import

from behave import given, when, then
from hamcrest import assert_that, equal_to, has_item, is_

from pptx import Presentation
from pptx.constants import MSO
from pptx.util import Inches

from .helpers import saved_pptx_path


# given ===================================================


@given('I have a reference to a table')
def step_given_ref_to_table(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[6]
    sld = context.prs.slides.add_slide(slidelayout)
    shapes = sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    context.tbl = shapes.add_table(2, 2, x, y, cx, cy)


@given('I have a reference to a table cell')
def step_given_ref_to_table_cell(context):
    context.prs = Presentation()
    slidelayout = context.prs.slidelayouts[6]
    sld = context.prs.slides.add_slide(slidelayout)
    length = 1000
    tbl = sld.shapes.add_table(2, 2, length, length, length, length)
    context.cell = tbl.cell(0, 0)


# when ====================================================


@when("I add a table to the slide's shape collection")
def step_when_add_table(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    shapes.add_table(2, 2, x, y, cx, cy)


@when("I set the cell margins")
def step_when_set_cell_margins(context):
    context.cell.margin_top = 1000
    context.cell.margin_right = 2000
    context.cell.margin_bottom = 3000
    context.cell.margin_left = 4000


@when("I set the cell vertical anchor to middle")
def step_when_set_cell_vertical_anchor_to_middle(context):
    context.cell.vertical_anchor = MSO.ANCHOR_MIDDLE


@when("I set the first_col property to True")
def step_when_set_first_col_property_to_true(context):
    context.tbl.first_col = True


@when("I set the first_row property to True")
def step_when_set_first_row_property_to_true(context):
    context.tbl.first_row = True


@when("I set the horz_banding property to True")
def step_when_set_horz_banding_property_to_true(context):
    context.tbl.horz_banding = True


@when("I set the last_col property to True")
def step_when_set_last_col_property_to_true(context):
    context.tbl.last_col = True


@when("I set the last_row property to True")
def step_when_set_last_row_property_to_true(context):
    context.tbl.last_row = True


@when("I set the text of the first cell")
def step_when_set_text_of_first_cell(context):
    context.tbl.cell(0, 0).text = 'test text'


@when("I set the vert_banding property to True")
def step_when_set_vert_banding_property_to_true(context):
    context.tbl.vert_banding = True


@when("I set the width of the table's columns")
def step_when_set_table_column_widths(context):
    context.tbl.columns[0].width = Inches(1.50)
    context.tbl.columns[1].width = Inches(3.00)


# then ====================================================


@then('the cell contents are inset by the margins')
def step_then_cell_contents_are_inset_by_the_margins(context):
    prs = Presentation(saved_pptx_path)
    table = prs.slides[0].shapes[0]
    cell = table.cell(0, 0)
    assert_that(cell.margin_top, is_(equal_to(1000)))
    assert_that(cell.margin_right, is_(equal_to(2000)))
    assert_that(cell.margin_bottom, is_(equal_to(3000)))
    assert_that(cell.margin_left, is_(equal_to(4000)))


@then('the cell contents are vertically centered')
def step_then_cell_contents_are_vertically_centered(context):
    prs = Presentation(saved_pptx_path)
    table = prs.slides[0].shapes[0]
    cell = table.cell(0, 0)
    assert_that(cell.vertical_anchor, is_(equal_to(MSO.ANCHOR_MIDDLE)))


@then('the columns of the table have alternating shading')
def step_then_columns_of_table_have_alternating_shading(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.vert_banding, is_(True))


@then('the first column of the table has special formatting')
def step_then_first_column_of_table_has_special_formatting(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.first_col, is_(True))


@then('the first row of the table has special formatting')
def step_then_first_row_of_table_has_special_formatting(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.first_row, is_(True))


@then('the last column of the table has special formatting')
def step_then_last_column_of_table_has_special_formatting(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.last_col, is_(True))


@then('the last row of the table has special formatting')
def step_then_last_row_of_table_has_special_formatting(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.last_row, is_(True))


@then('the rows of the table have alternating shading')
def step_then_rows_of_table_have_alternating_shading(context):
    tbl = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(tbl.horz_banding, is_(True))


@then('the table appears in the slide')
def step_then_table_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    shapes = sld.shapes
    classnames = [sp.__class__.__name__ for sp in shapes]
    assert_that(classnames, has_item('_Table'))


@then('the table appears with the new column widths')
def step_then_table_appears_with_new_col_widths(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    tbl = sld.shapes[0]
    assert_that(tbl.columns[0].width, is_(equal_to(Inches(1.50))))
    assert_that(tbl.columns[1].width, is_(equal_to(Inches(3.00))))


@then('the text appears in the first cell of the table')
def step_then_text_appears_in_first_cell_of_table(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    tbl = sld.shapes[0]
    text = tbl.cell(0, 0).textframe.paragraphs[0].runs[0].text
    assert_that(text, is_(equal_to('test text')))
