# encoding: utf-8

"""
Gherkin step implementations for text-related features.
"""

from __future__ import absolute_import

import os

from behave import given, when, then

from hamcrest import assert_that, equal_to, is_

from pptx import Presentation
from pptx.constants import PP
from pptx.util import Inches

from .helpers import italics_pptx_path, saved_pptx_path


# given ===================================================

@given('a paragraph')
def given_a_paragraph(context):
    context.prs = Presentation()
    blank_slidelayout = context.prs.slidelayouts[6]
    slide = context.prs.slides.add_slide(blank_slidelayout)
    length = Inches(2.00)
    textbox = slide.shapes.add_textbox(length, length, length, length)
    context.p = textbox.textframe.paragraphs[0]


@given('a run with italics set {setting}')
def step_given_run_with_italics_set_to_setting(context, setting):
    run_idx = {'on': 0, 'off': 1, 'to None': 2}[setting]
    context.prs = Presentation(italics_pptx_path)
    runs = context.prs.slides[0].shapes[0].textframe.paragraphs[0].runs
    context.run = runs[run_idx]


@given('a textframe')
def step_given_a_textframe(context):
    context.prs = Presentation()
    blank_slidelayout = context.prs.slidelayouts[6]
    slide = context.prs.slides.add_slide(blank_slidelayout)
    length = Inches(2.00)
    textbox = slide.shapes.add_textbox(length, length, length, length)
    context.textframe = textbox.textframe


# when ====================================================


@when('I reload the presentation')
def when_reload_presentation(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)
    context.prs = Presentation(saved_pptx_path)


@when('I set the {side} margin to {inches}"')
def when_set_margin_to_value(context, side, inches):
    emu = Inches(float(inches))
    if side == 'left':
        context.textframe.margin_left = emu
    elif side == 'top':
        context.textframe.margin_top = emu
    elif side == 'right':
        context.textframe.margin_right = emu
    elif side == 'bottom':
        context.textframe.margin_bottom = emu


@when('I indent the paragraph')
def step_when_indent_first_paragraph(context):
    context.p.level = 1


@when("I set italics {setting}")
def when_set_italics_to_setting(context, setting):
    new_italics_value = {'on': True, 'off': False, 'to None': None}[setting]
    context.run.font.italic = new_italics_value


@when("I set the paragraph alignment to centered")
def step_when_set_paragraph_alignment_to_centered(context):
    context.p.alignment = PP.ALIGN_CENTER


@when("I set the textframe word wrap {setting}")
def step_when_set_textframe_word_wrap(context, setting):
    bool_val = {'on': True, 'off': False, 'to None': None}
    context.textframe.word_wrap = bool_val[setting]


# then ====================================================

@then('the paragraph is aligned centered')
def step_then_paragraph_is_aligned_centered(context):
    p = context.prs.slides[0].shapes[0].textframe.paragraphs[0]
    assert_that(p.alignment, is_(equal_to(PP.ALIGN_CENTER)))


@then('the paragraph is indented to the second level')
def step_then_paragraph_indented_to_second_level(context):
    p = context.prs.slides[0].shapes[0].textframe.paragraphs[0]
    assert_that(p.level, is_(equal_to(1)))


@then("the run that had italics set {initial} now has it set {setting}")
def step_then_run_now_has_italics_set_to_setting(context, initial, setting):
    run_idx = {'on': 0, 'off': 1, 'to None': 2}[initial]
    run = (
        context.prs
               .slides[0]
               .shapes[0]
               .textframe
               .paragraphs[0]
               .runs[run_idx]
    )
    expected_val = {'on': True, 'off': False, 'to None': None}[setting]
    assert run.font.italic == expected_val


@then('the textframe\'s {side} margin is {inches}"')
def then_textframe_margin_is_value(context, side, inches):
    textframe = context.prs.slides[0].shapes[0].textframe
    emu = Inches(float(inches))
    if side == 'left':
        assert textframe.margin_left == emu
    elif side == 'top':
        assert textframe.margin_top == emu
    elif side == 'right':
        assert textframe.margin_right == emu
    elif side == 'bottom':
        assert textframe.margin_bottom == emu


@then('the textframe word wrap is set {setting}')
def step_then_textframe_word_wrap_is_setting(context, setting):
    bool_val = {'on': True, 'off': False, 'to None': None}
    textframe = context.prs.slides[0].shapes[0].textframe
    assert_that(textframe.word_wrap, is_(bool_val[setting]))
