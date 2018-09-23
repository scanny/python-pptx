# encoding: utf-8

"""Gherkin step implementations for table-related features"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

from behave import given, when, then

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR  # noqa
from pptx.util import Inches

from helpers import test_pptx


# given ===================================================

@given('a Table object as table')
@given('a 2x2 Table object as table')
def given_a_2x2_Table_object_as_table(context):
    prs = Presentation(test_pptx('shp-shapes'))
    context.table_ = prs.slides[0].shapes[3].table


@given('a 2x3 _MergeOriginCell object as cell')
def given_a_2x3_MergeOriginCell_object_as_cell(context):
    prs = Presentation(test_pptx('tbl-cell'))
    context.cell = prs.slides[1].shapes[1].table.cell(0, 0)


@given('a 3x3 Table object as table')
@given('a 3x3 Table object with cells a to i as table')
def given_a_3x3_table_with_cells_a_to_i_as_table(context):
    prs = Presentation(test_pptx('tbl-cell'))
    # ---context.table is used by Behave for some odd reason---
    context.table_ = prs.slides[2].shapes[0].table


@given('a _Cell object as cell')
def given_a_Cell_object_as_cell(context):
    prs = Presentation(test_pptx('shp-shapes'))
    context.cell = prs.slides[0].shapes[3].table.cell(0, 0)


@given('a _Cell object containing "unladen swallows" as cell')
def given_a_Cell_object_containing_unladen_swallows_as_cell(context):
    prs = Presentation(test_pptx('tbl-cell'))
    context.cell = prs.slides[0].shapes[0].table.cell(1, 0)


@given('a _Cell object with known margins as cell')
def given_a_Cell_object_with_known_margins_as_cell(context):
    prs = Presentation(test_pptx('tbl-cell'))
    context.cell = prs.slides[0].shapes[0].table.cell(0, 0)


@given('a _Cell object with {setting} vertical alignment as cell')
def given_a_Cell_object_with_setting_vertical_alignment(context, setting):
    cell_coordinates = {
        'inherited': (0, 1),
        'middle': (0, 2),
        'bottom': (0, 3),
    }[setting]
    prs = Presentation(test_pptx('tbl-cell'))
    context.cell = prs.slides[0].shapes[0].table.cell(*cell_coordinates)


@given('a {role} _Cell object as cell')
def given_a_role_Cell_object_as_cell(context, role):
    coordinates = {
        'merge-origin': (0, 0),
        'spanned': (0, 1),
        'unmerged': (2, 2),
    }[role]
    table = Presentation(test_pptx('tbl-cell')).slides[1].shapes[0].table
    # ---create other_cell here where we know the coordinates---
    context.cell = table.cell(*coordinates)
    context.other_cell = table.cell(*coordinates)


@given('a second proxy instance for that cell as other_cell')
def given_a_second_proxy_instance_for_that_cell_as_other_cell(context):
    # ---other_cell is actually produced by prior step---
    assert context.other_cell


@given('a _Column object as column')
def given_a_Column_object_as_column(context):
    prs = Presentation(test_pptx('shp-shapes'))
    context.column = prs.slides[0].shapes[3].table.columns[0]


# when ====================================================

@when('I assign cell.margin_{side} = {value}')
def when_I_assign_cell_margin_side_eq_value(context, value, side):
    setattr(context.cell, 'margin_%s' % side, eval(value))


@when('I assign cell.text = "test text"')
def when_I_assign_cell_text(context):
    context.cell.text = 'test text'


@when('I assign cell.vertical_anchor = {value}')
def when_I_assign_cell_vertical_anchor_eq_value(context, value):
    context.cell.vertical_anchor = eval(value)


@when('I assign column.width = {value}')
def when_I_assign_column_width_eq_value(context, value):
    context.column.width = eval(value)


@when('I assign origin_cell = table.cell(0, 0)')
def when_I_assign_origin_cell_eq_table_cell_0_0(context):
    context.origin_cell = context.table_.cell(0, 0)


@when('I assign other_cell = table.cell(1, 1)')
def when_I_assign_other_cell_eq_table_cell_1_1(context):
    context.other_cell = context.table_.cell(1, 1)


@when("I assign table.first_col = True")
def when_I_assign_table_first_col_eq_True(context):
    context.table_.first_col = True


@when("I assign table.first_row = True")
def when_I_assign_table_first_row_eq_True(context):
    context.table_.first_row = True


@when("I assign table.horz_banding = True")
def when_I_assign_table_horz_banding_eq_True(context):
    context.table_.horz_banding = True


@when("I assign table.last_col = True")
def when_I_assign_table_last_col_eq_True(context):
    context.table_.last_col = True


@when("I assign table.last_row = True")
def when_I_assign_table_last_row_eq_True(context):
    context.table_.last_row = True


@when("I assign table.vert_banding = True")
def when_I_assign_table_vert_banding_eq_True(context):
    context.table_.vert_banding = True


@when('I call cell.split()')
def when_I_call_cell_split_other_cell(context):
    context.cell.split()


@when('I call origin_cell.merge(other_cell)')
def when_I_call_origin_cell_merge_other_cell(context):
    context.origin_cell.merge(context.other_cell)


# then ====================================================

@then('cell == other_cell')
def then_cell_eq_other_cell(context):
    cell, other_cell = context.cell, context.other_cell
    assert cell == other_cell, 'cell != other_cell'


@then('cell.fill is a FillFormat object')
def then_cell_fill_is_a_FillFormat_object(context):
    actual = type(context.cell.fill).__name__
    expected = 'FillFormat'
    assert actual == expected, 'cell.fill is a %s object' % actual


@then('cell.margin_{side} == Inches({num_lit})')
def then_cell_margin_side_eq_Inches_num(context, side, num_lit):
    actual = getattr(context.cell, 'margin_%s' % side)
    expected = Inches(float(num_lit))
    assert actual == expected, 'cell.margin_%s == %s' % (side, actual.inches)


@then('{cell_ref}.is_merge_origin is {bool_lit}')
def then_cell_ref_is_merge_origin_is(context, cell_ref, bool_lit):
    expected = eval(bool_lit)
    actual = getattr(context, cell_ref).is_merge_origin
    assert actual is expected, (
        '%s.is_merge_origin is %s' % (cell_ref, actual)
    )


@then('{cell_ref}.is_spanned is {bool_lit}')
def then_cell_is_spanned_is(context, cell_ref, bool_lit):
    expected = eval(bool_lit)
    actual = getattr(context, cell_ref).is_spanned
    assert actual is expected, (
        '%s.is_spanned is %s' % (cell_ref, actual)
    )


@then('{cell_ref}.text == {value}')
def then_cell_ref_text_eq_value(context, cell_ref, value):
    actual = getattr(context, cell_ref).text
    expected = eval(value)
    assert actual == expected, '%s.text == %s' % (cell_ref, actual)


@then('cell.span_height == {int_lit}')
def then_cell_span_height_eq(context, int_lit):
    expected = int(int_lit)
    actual = context.cell.span_height
    assert actual is expected, 'cell.span_height == %s' % actual


@then('cell.span_width == {int_lit}')
def then_cell_span_width_eq(context, int_lit):
    expected = int(int_lit)
    actual = context.cell.span_width
    assert actual is expected, 'cell.span_width == %s' % actual


@then('cell.vertical_anchor == {value}')
def then_cell_vertical_anchor_eq_value(context, value):
    actual = context.cell.vertical_anchor
    expected = eval(value)
    assert actual == expected, 'cell.vertical_anchor == %s' % actual


@then('column.width.inches == {float_lit}')
def then_column_width_inches_eq(context, float_lit):
    actual = context.column.width.inches
    expected = float(float_lit)
    assert actual == expected, 'column.width.inches == %s' % actual


@then('len(list(table.iter_cells())) == {int_lit}')
def then_len_list_table_iter_cells_eq(context, int_lit):
    actual = len(list(context.table_.iter_cells()))
    expected = int(int_lit)
    assert actual is expected, 'len(list(table.iter_cells())) == %s' % actual


@then('table.cell(0, 0) is a {type_name} object')
def then_table_cell_0_0_is_a_type_object(context, type_name):
    actual = type(context.table_.cell(0, 0)).__name__
    expected = type_name
    assert actual == expected, 'table.cell(0, 0) is a %s object' % actual


@then('table.columns is a {type_name} object')
def then_table_columns_is_a_type_object(context, type_name):
    actual = type(context.table_.columns).__name__
    expected = type_name
    assert actual == expected, 'table.columns is a %s object' % actual


@then('table.first_col is {bool_lit}')
def then_table_first_col_is_value(context, bool_lit):
    actual = context.table_.first_col
    expected = eval(bool_lit)
    assert actual is expected, 'table.first_col is %s' % actual


@then('table.first_row is {bool_lit}')
def then_table_first_row_is_value(context, bool_lit):
    actual = context.table_.first_row
    expected = eval(bool_lit)
    assert actual is expected, 'table.first_row is %s' % actual


@then('table.horz_banding is {bool_lit}')
def then_table_horz_banding_is_value(context, bool_lit):
    actual = context.table_.horz_banding
    expected = eval(bool_lit)
    assert actual is expected, 'table.horz_banding is %s' % actual


@then('table.last_col is {bool_lit}')
def then_table_last_col_is_value(context, bool_lit):
    actual = context.table_.last_col
    expected = eval(bool_lit)
    assert actual is expected, 'table.last_col is %s' % actual


@then('table.last_row is {bool_lit}')
def then_table_last_row_is_value(context, bool_lit):
    actual = context.table_.last_row
    expected = eval(bool_lit)
    assert actual is expected, 'table.last_row is %s' % actual


@then('table.rows is a {type_name} object')
def then_table_rows_is_a_type_object(context, type_name):
    actual = type(context.table_.rows).__name__
    expected = type_name
    assert actual == expected, 'table.rows is a %s object' % actual


@then('table.vert_banding is {bool_lit}')
def then_table_vert_banding_is_value(context, bool_lit):
    actual = context.table_.vert_banding
    expected = eval(bool_lit)
    assert actual is expected, 'table.vert_banding is %s' % actual
