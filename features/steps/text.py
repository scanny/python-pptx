# encoding: utf-8

"""Gherkin step implementations for text-related features."""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, when, then

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

from helpers import test_pptx


# given ===================================================


@given("a _Paragraph object as paragraph")
def given_a_Paragraph_object_as_paragraph(context):
    prs = Presentation(test_pptx("txt-text"))
    context.paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]


@given("a _Paragraph object containing {value} as paragraph")
def given_a_Paragraph_object_containing_value_as_paragraph(context, value):
    prs = Presentation(test_pptx("txt-text"))
    paragraph_idx = {"abc": 0, "a\vb\vc": 1}[eval(value)]
    context.paragraph = prs.slides[0].shapes[1].text_frame.paragraphs[paragraph_idx]


@given("a paragraph having line spacing of {setting}")
def given_a_paragraph_having_line_spacing_of_setting(context, setting):
    paragraph_idx = {"no explicit setting": 0, "1.5 lines": 1, "20 pt": 2}[setting]
    prs = Presentation(test_pptx("txt-paragraph-spacing"))
    text_frame = prs.slides[2].shapes[0].text_frame
    context.paragraph = text_frame.paragraphs[paragraph_idx]


@given("a paragraph having space {before_after} of {setting}")
def given_a_paragraph_having_space_before_after_of_setting(
    context, before_after, setting
):
    slide_idx = {"before": 0, "after": 1}[before_after]
    paragraph_idx = {"no explicit setting": 0, "6 pt": 1}[setting]
    prs = Presentation(test_pptx("txt-paragraph-spacing"))
    text_frame = prs.slides[slide_idx].shapes[0].text_frame
    context.paragraph = text_frame.paragraphs[paragraph_idx]


@given("a _Run object as run")
def given_a_Run_object_as_run(context):
    prs = Presentation(test_pptx("txt-text"))
    context.run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]


@given("a _Run object containing text as run")
def given_a_Run_object_containing_text_as_run(context):
    prs = Presentation(test_pptx("txt-text"))
    context.run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]


@given("a text run")
def given_a_text_run(context):
    prs = Presentation(test_pptx("txt-text"))
    context.run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]


@given("a text run in a table cell")
def given_a_text_run_in_a_table_cell(context):
    prs = Presentation(test_pptx("txt-text"))
    cell = prs.slides[1].shapes[0].table.cell(0, 0)
    context.run = cell.text_frame.paragraphs[0].runs[0]


@given("a text run having a hyperlink")
def given_a_text_run_having_a_hyperlink(context):
    prs = Presentation(test_pptx("txt-text"))
    context.run = prs.slides[0].shapes[0].text_frame.paragraphs[1].runs[0]


# when ====================================================


@when("I assign a typeface name to the font")
def when_assign_typeface_name_to_font(context):
    context.font.name = "Verdana"


@when("I assign None to hyperlink.address")
def when_assign_None_to_hyperlink_address(context):
    context.run.hyperlink.address = None


@when("I assign paragraph.alignment = PP_ALIGN.CENTER")
def when_I_assign_paragraph_alignment_eq_center(context):
    context.paragraph.alignment = PP_ALIGN.CENTER


@when("I assign paragraph.level = 1")
def when_I_assign_paragraph_leve_eq_1(context):
    context.paragraph.level = 1


@when("I assign paragraph.text = {value}")
def when_I_assign_paragraph_text_eq_value(context, value):
    context.paragraph.text = eval(value)


@when("I assign run.text = {value}")
def when_I_assign_run_text_eq_value(context, value):
    context.run.text = eval(value)


@when("I assign {value_str} to paragraph.line_spacing")
def when_I_assign_value_to_paragraph_line_spacing(context, value_str):
    value = {
        "1.5": 1.5,
        "2.0": 2.0,
        "254000": Emu(254000),
        "304800": Emu(304800),
        "None": None,
    }[value_str]
    paragraph = context.paragraph
    paragraph.line_spacing = value


@when("I assign {value_str} to paragraph.space_{before_after}")
def when_I_assign_value_to_paragraph_space_before_after(
    context, value_str, before_after
):
    value = {"76200": 76200, "38100": 38100, "None": None}[value_str]
    attr_name = {"before": "space_before", "after": "space_after"}[before_after]
    paragraph = context.paragraph
    setattr(paragraph, attr_name, value)


@when("I set the hyperlink address")
def when_set_hyperlink_address(context):
    context.run_text = "python-pptx @ GitHub"
    context.address = "https://github.com/scanny/python-pptx"

    run = context.run
    run.text = context.run_text
    hlink = run.hyperlink
    hlink.address = context.address


# then ====================================================


@then("paragraph.alignment == PP_ALIGN.CENTER")
def then_paragraph_alignment_eq_center(context):
    actual, expected = context.paragraph.alignment, PP_ALIGN.CENTER
    assert actual == expected, "paragraph.alignment == %s" % actual


@then("paragraph.level == 1")
def then_paragraph_level_eq_1(context):
    actual, expected = context.paragraph.level, 1
    assert actual == expected, "paragraph.level == %s" % actual


@then("paragraph.line_spacing is {value_str}")
def then_paragraph_line_spacing_is_value(context, value_str):
    value = {
        "None": None,
        "1.0": 1.0,
        "1.5": 1.5,
        "2.0": 2.0,
        "254000": 254000,
        "304800": 304800,
    }[value_str]
    paragraph = context.paragraph
    assert paragraph.line_spacing == value


@then("paragraph.line_spacing.pt {result}")
def then_paragraph_line_spacing_pt_result(context, result):
    value, exception = {
        "raises AttributeError": (None, AttributeError),
        "is 20.0": (20.0, None),
        "is 24.0": (24.0, None),
    }[result]
    line_spacing = context.paragraph.line_spacing
    if value is not None:
        assert line_spacing.pt == value
    if exception is not None:
        try:
            line_spacing.pt
            raise AssertionError("did not raise")
        except exception:
            pass


@then("paragraph.space_{before_after} is {value_str}")
def then_paragraph_space_before_is_value(context, before_after, value_str):
    attr_name = {"before": "space_before", "after": "space_after"}[before_after]
    value = None if value_str == "None" else int(value_str)
    paragraph = context.paragraph
    assert getattr(paragraph, attr_name) == value


@then("paragraph.space_{before_after}.pt {result}")
def then_paragraph_space_before_pt_is_value(context, before_after, result):
    attr_name = {"before": "space_before", "after": "space_after"}[before_after]
    value, exception = {
        "raises AttributeError": (None, AttributeError),
        "== 6.0": (6.0, None),
    }[result]
    space_before_after = getattr(context.paragraph, attr_name)
    if value is not None:
        assert space_before_after.pt == value
    if exception is not None:
        try:
            space_before_after.pt
            raise AssertionError("did not raise")
        except exception:
            pass


@then("paragraph.text == {value}")
def then_paragraph_text_eq_value(context, value):
    actual, expected = context.paragraph.text, eval(value)
    assert actual == expected, 'paragraph.text == "%s"' % actual


@then("paragraph.text matches the assigned string")
def then_paragraph_text_matches_the_assigned_string(context):
    paragraph = context.paragraph
    assert paragraph.text == " Boo Far \n Faz Foo "


@then("run.text == {value}")
def then_run_text_is_value(context, value):
    actual, expected = context.run.text, eval(value)
    assert actual == expected, "run.text == %s" % (actual,)


@then("run.text is a hyperlink")
def then_run_text_is_a_hyperlink(context):
    run = context.run
    hlink = run.hyperlink
    assert run.text == context.run_text
    assert hlink.address == context.address


@then("run.text is not a hyperlink")
def then_run_text_is_not_a_hyperlink(context):
    hlink = context.run.hyperlink
    assert hlink.address is None


@then("the font name matches the typeface I set")
def then_font_name_matches_typeface_I_set(context):
    assert context.font.name == "Verdana"
