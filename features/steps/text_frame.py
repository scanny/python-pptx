# encoding: utf-8

"""Step implementations for text frame-related features"""

from __future__ import absolute_import, division, print_function, unicode_literals

from behave import given, then, when

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from helpers import test_pptx


# given ===================================================


@given("a TextFrame object as text_frame")
def given_a_text_frame(context):
    context.text_frame = (
        Presentation(test_pptx("txt-text")).slides[0].shapes[0].text_frame
    )


@given("a TextFrame object containing {value} as text_frame")
def given_a_TextFrame_object_containing_value_as_text_frame(context, value):
    shape_idx = {"abc": 0, "a\nb\nc": 1}[eval(value)]
    prs = Presentation(test_pptx("txt-text-frame"))
    context.text_frame = prs.slides[1].shapes[shape_idx].text_frame


@given("a TextFrame object having auto-size of {setting} as text_frame")
def given_a_TextFrame_object_having_auto_size_of_setting(context, setting):
    shape_idx = {
        "None": 0,
        "no auto-size": 1,
        "fit shape to text": 2,
        "fit text to shape": 3,
    }[setting]
    prs = Presentation(test_pptx("txt-text-frame"))
    shape = prs.slides[0].shapes[shape_idx]
    context.text_frame = shape.text_frame


@given("a text frame with more text than will fit")
def given_a_text_frame_with_more_text_than_will_fit(context):
    prs = Presentation(test_pptx("txt-fit-text"))
    shape = prs.slides[0].shapes[0]
    context.text_frame = shape.text_frame


# when ====================================================


@when("I assign {value} to text_frame.auto_size")
def when_I_assign_value_to_text_frame_auto_size(context, value):
    text_frame = context.text_frame
    text_frame.auto_size = {
        "None": None,
        "MSO_AUTO_SIZE.NONE": MSO_AUTO_SIZE.NONE,
        "MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        "MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    }[value]


@when("I assign text_frame.margin_{side} = Inches({inches})")
def when_I_assign_text_frame_margin_side_eq_inches(context, side, inches):
    attr_name = "margin_%s" % side
    setattr(context.text_frame, attr_name, Inches(float(inches)))


@when("I assign text_frame.text = {value}")
def when_I_assign_text_frame_text_eq_value(context, value):
    context.text_frame.text = eval(value)


@when("I assign {value} to text_frame.word_wrap")
def when_I_assign_value_to_text_frame_word_wrap(context, value):
    new_value = {"True": True, "False": False, "None": None}[value]
    context.text_frame.word_wrap = new_value


@when("I call TextFrame.fit_text()")
def when_I_call_TextFrame_fit_text(context):
    from helpers import test_file

    font_file = test_file("calibriz.ttf")
    context.text_frame.fit_text(bold=True, italic=True, font_file=font_file)
    # context.text_frame.fit_text(font_family='Arial', bold=True, italic=True)


# then ====================================================


@then("text_frame.auto_size is {value}")
def then_text_frame_autosize_is_value(context, value):
    expected_value = {
        "None": None,
        "MSO_AUTO_SIZE.NONE": MSO_AUTO_SIZE.NONE,
        "MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        "MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    }[value]
    text_frame = context.text_frame
    assert text_frame.auto_size == expected_value, "got %s" % text_frame.auto_size


@then("text_frame.margin_{side}.inches == {inches}")
def then_text_frame_margin_side_inches_eq_inches(context, side, inches):
    attr_name = "margin_%s" % side
    actual = getattr(context.text_frame, attr_name).inches
    expected = float(inches)
    assert actual == expected, "text_frame.margin_%s.inches == %s" % (side, actual)


@then("text_frame.text == {value}")
def then_text_frame_text_eq_value(context, value):
    actual, expected = context.text_frame.text, eval(value)
    assert actual == expected, 'text_frame.text == "%s"' % actual


@then("text_frame.word_wrap is {value}")
def then_text_frame_word_wrap_is_value(context, value):
    expected_value = {"True": True, "False": False, "None": None}[value]
    text_frame = context.text_frame
    assert text_frame.word_wrap is expected_value


@then("the size of the text is 10pt")
def then_the_size_of_the_text_is_10pt(context):
    text_frame = context.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.size == Pt(10.0), "got %s" % run.font.size.pt
