# encoding: utf-8

"""
Step implementations for textframe-related features
"""

from __future__ import absolute_import

from behave import given, then, when

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from .helpers import test_pptx


# given ===================================================

@given('a textframe')
def given_a_textframe(context):
    context.prs = Presentation()
    blank_slide_layout = context.prs.slide_layouts[6]
    slide = context.prs.slides.add_slide(blank_slide_layout)
    length = Inches(2.00)
    textbox = slide.shapes.add_textbox(length, length, length, length)
    context.textframe = textbox.textframe


@given('a textframe containing text')
def given_a_textframe_containing_text(context):
    prs = Presentation(test_pptx('txt-text'))
    context.textframe = prs.slides[0].shapes[0].textframe


@given('a textframe having auto-size set to {setting}')
def given_a_textframe_having_auto_size_set_to_setting(context, setting):
    shape_idx = {
        'None':              0,
        'no auto-size':      1,
        'fit shape to text': 2,
        'fit text to shape': 3,
    }[setting]
    prs = Presentation(test_pptx('txt-textframe-props'))
    shape = prs.slides[0].shapes[shape_idx]
    context.textframe = shape.textframe


# when ====================================================

@when('I assign a string to textframe.text')
def when_I_assign_a_string_to_textframe_text(context):
    context.textframe.text = ' Foo Bar \n Baz Zoo '


@when("I set textframe.auto_size to {setting}")
def when_set_textframe_auto_size(context, setting):
    textframe = context.textframe
    textframe.auto_size = {
        'None': None,
        'MSO_AUTO_SIZE.NONE': MSO_AUTO_SIZE.NONE,
        'MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT': MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        'MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE': MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    }[setting]


@when("I set the textframe word wrap {setting}")
def when_set_textframe_word_wrap(context, setting):
    bool_val = {'on': True, 'off': False, 'to None': None}
    context.textframe.word_wrap = bool_val[setting]


# then ====================================================

@then('textframe.auto_size is {value}')
def then_textframe_autosize_is_value(context, value):
    expected_value = {
        'None': None,
        'MSO_AUTO_SIZE.NONE': MSO_AUTO_SIZE.NONE,
        'MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT': MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        'MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE': MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    }[value]
    textframe = context.textframe
    assert textframe.auto_size == expected_value, (
        'got %s' % textframe.auto_size
    )


@then('the textframe\'s {side} margin is {inches}"')
def then_textframe_margin_is_value(context, side, inches):
    textframe = context.prs.slides[0].shapes[0].textframe
    emu = Inches(float(inches))
    if side == 'left':
        assert textframe.margin_left == emu
    elif side == 'top':
        assert textframe.margin_top == emu
    elif side == 'right':
        assert textframe.margin_right == emu
    elif side == 'bottom':
        assert textframe.margin_bottom == emu


@then('textframe.text is the text in the shape')
def then_textframe_text_is_the_text_in_the_shape(context):
    textframe = context.textframe
    assert textframe.text == ' Foo Bar \n Baz Zoo \n1'


@then('textframe.text matches the assigned string')
def then_textframe_text_matches_the_assigned_string(context):
    textframe = context.textframe
    assert textframe.text == ' Foo Bar \n Baz Zoo '


@then('textframe.word_wrap is {value}')
def then_textframe_word_wrap_is_value(context, value):
    expected_value = {
        'on': True, 'off': False, 'to None': None
    }[value]
    textframe = context.prs.slides[0].shapes[0].textframe
    assert textframe.word_wrap is expected_value
