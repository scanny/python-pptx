#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

"""Evaluate formulae in guides."""

#===============================================================================

import math

import pptx.shapes.connector
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

from .presets import Shapes

#===============================================================================

STANGLE_PER_DEGREE = 60000

def radians(x):
    return x*math.pi/(STANGLE_PER_DEGREE*180)

def st_angle(x):
    return x*STANGLE_PER_DEGREE*180/math.pi

#===============================================================================

PRESET_VARIABLES = {
    "3cd4": 16200000.0,  # 3/4 of a Circle
    "3cd8": 8100000.0,  # 3/8 of a Circle
    "5cd8": 13500000.0,  # 5/8 of a Circle
    "7cd8": 18900000.0,  # 7/8 of a Circle
    "cd2": 10800000.0,  # 1/2 of a Circle
    "cd4": 5400000.0,  # 1/4 of a Circle
    "cd8": 2700000.0,  # 1/8 of a Circle
    "t": 0,  # Shape Top Edge
    "b": "val h",  # Shape Bottom Edge
    "vc": "*/ h 1.0 2.0",  # Vertical Center of Shape
    "hd2": "*/ h 1.0 2.0",  # 1/2 of Shape Height
    "hd3": "*/ h 1.0 3.0",  # 1/3 of Shape Height
    "hd4": "*/ h 1.0 4.0",  # 1/4 of Shape Height
    "hd5": "*/ h 1.0 5.0",  # 1/5 of Shape Height
    "hd6": "*/ h 1.0 6.0",  # 1/6 of Shape Height
    "hd8": "*/ h 1.0 8.0",  # 1/8 of Shape Height
    "l": 0,  # Shape Left Edge
    "r": "val w",  # Shape Right Edge
    "hc": "*/ w 1.0 2.0",  # Horizontal Center
    "wd2": "*/ w 1.0 2.0",  # 1/2 of Shape Width
    "wd3": "*/ w 1.0 3.0",  # 1/3 of Shape Width
    "wd4": "*/ w 1.0 4.0",  # 1/4 of Shape Width
    "wd5": "*/ w 1.0 5.0",  # 1/5 of Shape Width
    "wd6": "*/ w 1.0 6.0",  # 1/6 of Shape Width
    "wd8": "*/ w 1.0 8.0",  # 1/8 of Shape Width
    "wd10": "*/ w 1.0 10.0",  # 1/10 of Shape Width
    "ls": "max w h",  # Longest Side of Shape
    "ss": "min w h",  # Shortest Side of Shape
    "ssd2": "*/ ss 1.0 2.0",  # 1/2 Shortest Side of Shape
    "ssd4": "*/ ss 1.0 4.0",  # 1/4 Shortest Side of Shape
    "ssd6": "*/ ss 1.0 6.0",  # 1/6 Shortest Side of Shape
    "ssd8": "*/ ss 1.0 8.0",  # 1/8 Shortest Side of Shape
    "ssd16": "*/ ss 1.0 16.0",  # 1/16 Shortest Side of Shape
    "ssd32": "*/ ss 1.0 32.0",  # 1/32 Shortest Side of Shape
}

#===============================================================================

class Evaluator(object):

    formulae = {
        "*/": lambda v, x, y, z: v(x) * v(y) / v(z),  # Multiply Divide Formula
        "+-": lambda v, x, y, z: v(x) + v(y) - v(z),  # Add Subtract Formula
        "+/": lambda v, x, y, z: (v(x) + v(y)) / v(z),  # Add Divide Formula
        "?:": lambda v, x, y, z: v(y) if v(x) > 0 else v(z),  # If Else Formula
        "at2": lambda v, x, y: (  # ArcTan Formula
            st_angle(math.atan(v(y) / v(x)))
            if v(x) != 0.0
            else v("cd4" if v(y) >= 0 else "3cd4")
        ),
        "tan": lambda v, x, y: v(x) * math.tan(radians(v(y))),  # Tangent Formula
        "cat2": lambda v, x, y, z: (  # Cosine ArcTan Formula
            v(x)
            * math.cos(math.atan(v(z) / v(y)))
            if v(y) != 0.0
            else 0.0
        ),
        "cos": lambda v, x, y: v(x) * math.cos(radians(v(y))),  # Cosine Formula
        "sat2": lambda v, x, y, z: (  # Sine ArcTan Formula
            v(x)
            * math.sin(math.atan(v(z) / v(y)))
            if v(y) != 0.0
            else v(x)
            if v(z) >= 0
            else -v(x)
        ),
        "sin": lambda v, x, y: v(x) * math.sin(radians(v(y))),  # Sine Formula
        "mod": lambda v, x, y, z: math.sqrt(
            v(x) ** 2 + v(y) ** 2 + v(z) ** 2
        ),  # Modulo Formula
        "sqrt": lambda v, x: math.sqrt(v(x)),  # Square Root Formula
        "val": lambda v, x: v(x),  # Literal Value Formula
        "abs": lambda v, x: abs(v(x)),  # Absolute Value Formula
        "max": lambda v, x, y: max(v(x), v(y)),  # Maximum Value Formula
        "min": lambda v, x, y: min(v(x), v(y)),  # Minimum Value Formula
        "pin": lambda v, x, y, z: (
            v(x)
            if v(y) < v(x)  # Pin To Formula
            else v(z)
            if v(y) > v(z)
            else v(y)
        ),
    }

    @staticmethod
    def evaluate(expr, context):
        args = expr.split()
        return Evaluator.formulae[args[0]](context.evaluate, *args[1:])

#===============================================================================

class Geometry(object):
    def __init__(self, shape):
        self._xfrm = shape.element.xfrm

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            self._geometry = Shapes.lookup(shape.element.prstGeom.attrib['prst'])
            adjustments = shape.element.prstGeom.avLst

        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            self._geometry = shape.element.spPr.custGeom
            adjustments = None

        elif isinstance(shape, pptx.shapes.connector.Connector):
            self._geometry = Shapes.lookup(shape.element.spPr.prstGeom.attrib['prst'])
            adjustments = None

        else:
            print('Unknown geometry for', shape.shape_type)

        self._variables = {
            'w': shape.width,
            'h': shape.height
        }

        if self._geometry.gdLst is not None:
            for gd in self._geometry.gdLst:
                self._variables[gd.name] = gd.fmla

        if self._geometry.avLst is not None:
            for gd in self._geometry.avLst:
                self._variables[gd.name] = gd.fmla

        if adjustments is not None:
            for gd in adjustments:
                self._variables[gd.name] = gd.fmla

    @property
    def path_list(self):
        return self._geometry.pathLst

    @property
    def xfrm(self):
        return self._xfrm

    def evaluate(self, x):
        try: return float(x)
        except ValueError: pass
        try: return self.evaluate(PRESET_VARIABLES[x])
        except KeyError: pass
        try: return self.evaluate(self._variables[x])
        except KeyError: pass
        return Evaluator.evaluate(x, self)

    def point(self, pt):
        return (self.evaluate(pt.attrib['x']), self.evaluate(pt.attrib['y']))

    def attrib_value(self, element, attrib):
        return self.evaluate(element.attrib[attrib])

#===============================================================================

if __name__ == '__main__':
    variables['x'] = '*/ 1 2 3'
    print(value('val x'))
    print(value('wd10'))
    print(value('cos 1 0'))
    print(value('at2 0 1'))
    print(value('cat2 1 0 1'))
    print(value('cat2 1 1 1'))
    print(value('sat2 1 1 1'))
    print(value('sat2 1 0 1'))

#===============================================================================
