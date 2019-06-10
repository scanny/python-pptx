#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import sqrt, sin, cos, pi as PI
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.spec import autoshape_types

import pptx.shapes.connector

import svgwrite

import numpy as np

#===============================================================================

from drawml.formula import Geometry, radians

#===============================================================================

def DML(tag):
    return '{{http://schemas.openxmlformats.org/drawingml/2006/main}}{}'.format(tag)

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def svg_coords(x, y):
#====================
    return (x/EMU_PER_DOT, y/EMU_PER_DOT)

def svg_units(emu):
#===================
    return emu/EMU_PER_DOT

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = (svg_coords(xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = (svg_coords(xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    svg_coords(*bbox))
        (Bx_, By_) = svg_coords(xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = svg_coords(xfrm.ext.cx, xfrm.ext.cy)

        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1

        T_st = np.matrix([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                          [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                          [     0,      0,                 1]])
        U = np.matrix([[1, 0, -(Bx_ + Dx_/2.0)],
                       [0, 1, -(By_ + Dy_/2.0)],
                       [0, 0,                1]])

        R = np.matrix([[cos(theta), -sin(theta), 0],
                       [sin(theta),  cos(theta), 0],
                       [0,                    0, 1]])
        Flip = np.matrix([[Fx,  0, 0],
                          [ 0, Fy, 0],
                          [ 0,  0, 1]])
        T_rf = U.I*R*Flip*U

        self._T = T_rf*T_st

    def svg_matrix(self):
        return (self._T[0, 0], self._T[1, 0],
                self._T[0, 1], self._T[1, 1],
                self._T[0, 2], self._T[1, 2])

#===============================================================================

class SvgMaker(object):
    def __init__(self, slide, slide_number, slide_size, args):
        if args.debug_xml:
            xml = open(os.path.join(args.output_dir, 'slide{:02d}.xml'.format(slide_number)), 'w')
            xml.write(slide.element.xml)
            xml.close()
        self._dwg = svgwrite.Drawing(filename=os.path.join(args.output_dir, 'slide{:02d}.svg'.format(slide_number)),
                                     size=svg_coords(slide_size[0], slide_size[1]))
        self._dwg.defs.add(self._dwg.style('.non-scaling-stroke { vector-effect: non-scaling-stroke; }'))
        self.svg_from_shapes(slide.shapes, self._dwg)
        self._dwg.save()

    def svg_from_shapes(self, shapes, svg_parent):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.shape_to_svg(shape, svg_parent)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                svg_group = self._dwg.g()
                transform = Transform(shape)
                svg_group.matrix(*transform.svg_matrix())
                svg_parent.add(svg_group)
                self.svg_from_shapes(shape.shapes, svg_group)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass  # or recognise name of '#layer-id' and get layer name...

            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

    def shape_to_svg(self, shape, svg_parent):
        geometry = Geometry(shape)
        for path in geometry.path_list:
            bbox = (shape.width, shape.height) if path.w is None else (path.w, path.h)
            transform = Transform(shape, bbox)
            svg_path = self._dwg.path(fill='none', stroke_width=3, class_='non-scaling-stroke') # id='sss'
            svg_path.matrix(*transform.svg_matrix())
            first_point = None
            current_point = None
            closed = False
            for c in path.getchildren():
                if   c.tag == DML('arcTo'):
                    wR = geometry.attrib_value(c, 'wR')
                    hR = geometry.attrib_value(c, 'hR')
                    stAng = radians(geometry.attrib_value(c, 'stAng'))
                    swAng = radians(geometry.attrib_value(c, 'swAng'))
                    p1 = ellipse_point(wR, hR, stAng)
                    p2 = ellipse_point(wR, hR, stAng + swAng)
                    pt = (current_point[0] - p1[0] + p2[0],
                          current_point[1] - p1[1] + p2[1])
                    large_arc_flag = 1 if swAng >= PI else 0
                    svg_path.push('A', svg_units(wR), svg_units(hR),
                                       0, large_arc_flag, 1,
                                       svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt

                elif c.tag == DML('close'):
                    if first_point is not None and first_point == current_point:
                        closed = True
                    svg_path.push('Z')
                    first_point = None
                elif c.tag == DML('cubicBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('C', *coords)
                elif c.tag == DML('lnTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('L', svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt
                elif c.tag == DML('moveTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('M', svg_units(pt[0]), svg_units(pt[1]))
                    if first_point is None:
                        first_point = pt
                    current_point = pt
                elif c.tag == DML('quadBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('Q', *coords)
                else:
                    print('Unknown path element: {}'.format(c.tag))
            if closed:
                svg_path.attribs['fill'] = '#808080'
                svg_path.attribs['opacity'] = 0.3
                svg_path.attribs['stroke'] = 'red'
            else:
                svg_path.attribs['stroke'] = 'blue'

            svg_parent.add(svg_path)

#===============================================================================

class SvgExtract(object):
    def __init__(self, args):
        self._args = args
        self._ppt = Presentation(args.powerpoint)
        self._slides = self._ppt.slides
        self._slide_size = [self._ppt.slide_width, self._ppt.slide_height]

    def slide_to_svg(self, slide_number):
        svg_maker = SvgMaker(self._slides[slide_number-1], slide_number, self._slide_size, self._args)

    def slides_to_svg(self):
        for n in range(len(slides)):
            self.slide_to_svg(n)

#===============================================================================

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Extract geometries from Powerpoint slides.')
    parser.add_argument('--version', action='version', version='0.2.1')
    parser.add_argument('--debug-xml', action='store_true',
                        help="save a slide's DrawML for debugging")
    parser.add_argument('--slide', type=int, metavar='N',
                        help='only process this slide number (1-origin)')
    parser.add_argument('output_dir', metavar='OUTPUT_DIRECTORY',
                        help='directory in which to save geometries')
    parser.add_argument('powerpoint', metavar='POWERPOINT_FILE',
                        help='the name of a Powerpoint file')


    args = parser.parse_args()

    if not os.path.exists(args.output_dir):
        os.makedirs(args.output_dir)

    svg_extract = SvgExtract(args)
    if args.slide is None:
        svg_extract.slides_to_svg()
    else:
        svg_extract.slide_to_svg(args.slide)

#===============================================================================
