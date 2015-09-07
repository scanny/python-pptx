import pandas as pd
import xlrd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation('PP.pptx')

chart_data = ChartData()
chart_data.add_series('SingeSerie1', ((10,  5), (20,  6), (30,  7)))
chart_data.add_series('SingeSerie2', ((20, 10), (30, 11), (40, 12)))

chart_data2 = ChartData()
#chart_data2.add_series('SerieName' , ((x1, y1), (x2, y2), (x1, y2) ))
chart_data2.add_series('SingeSerie1', ((10,  5), (20,  6), (30,  7)))
chart_data2.add_series('SingeSerie2', ((20, 10), (30, 11), (40, 12)))
chart_data2.add_series('SingeSerie3', ((35, 15), (45, 16), (55, 17)))
chart_data2.add_series('SingeSerie4', ((50, 20), (60, 21), (70, 22)))

chart_data3 = ChartData()
chart_data3.add_series('SingeSerie1', ((10,  5), (20,  6), (30,  7)))
chart_data3.add_series('SingeSerie2', ((20, 10), (30, 11), (40, 12)))
chart_data3.add_series('SingeSerie3', ((35, 15), (45, 16), (55, 17)))
chart_data3.add_series('SingeSerie4', ((50, 20), (60, 21), (70, 22)))
chart_data3.add_series('SingeSerie5', ((50, 50), (60, 60), (70, 50)))

chart_data4 = ChartData()
chart_data4.categories = ['East', 'West', 'Midwest']
chart_data4.add_series('Series 1', (40, 40, 20))

chart_data5 = ChartData()
chart_data5.categories = 'Foobar', 'Barbaz', 'Bazfoo'
chart_data5.add_series('New Series 1', (5.6, 6.7, 7.8))
chart_data5.add_series('New Series 2', (2.3, 3.4, 4.5))
chart_data5.add_series('New Series 3', (8.9, 9.1, 1.2))

chart_datas = [chart_data, chart_data2, chart_data3, chart_data5, chart_data4, chart_data5]

for index, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_chart:
            shape.chart.replace_data(chart_datas[index])

prs.save('pp2.pptx')
