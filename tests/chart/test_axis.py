# encoding: utf-8

"""
Test suite for pptx.chart module
"""

from __future__ import absolute_import, print_function

import pytest

from pptx.chart.axis import (
    AxisTitle, _BaseAxis, CategoryAxis, DateAxis, MajorGridlines, TickLabels,
    ValueAxis
)
from pptx.dml.chtfmt import ChartFormat
from pptx.enum.chart import (
    XL_AXIS_CROSSES, XL_CATEGORY_TYPE,
    XL_TICK_LABEL_POSITION as XL_TICK_LBL_POS, XL_TICK_MARK
)
from pptx.text.text import Font

from ..unitutil.cxml import element, xml
from ..unitutil.mock import class_mock, instance_mock


class Describe_BaseAxis(object):

    def it_knows_whether_it_has_major_gridlines(
            self, major_gridlines_get_fixture):
        base_axis, expected_value = major_gridlines_get_fixture
        assert base_axis.has_major_gridlines is expected_value

    def it_can_change_whether_it_has_major_gridlines(
            self, major_gridlines_set_fixture):
        base_axis, new_value, expected_xml = major_gridlines_set_fixture
        base_axis.has_major_gridlines = new_value
        assert base_axis._element.xml == expected_xml

    def it_knows_whether_it_has_minor_gridlines(
            self, minor_gridlines_get_fixture):
        base_axis, expected_value = minor_gridlines_get_fixture
        assert base_axis.has_minor_gridlines is expected_value

    def it_can_change_whether_it_has_minor_gridlines(
            self, minor_gridlines_set_fixture):
        base_axis, new_value, expected_xml = minor_gridlines_set_fixture
        base_axis.has_minor_gridlines = new_value
        assert base_axis._element.xml == expected_xml

    def it_knows_whether_it_has_a_title(self, has_title_get_fixture):
        axis, expected_value = has_title_get_fixture
        assert axis.has_title is expected_value

    def it_can_change_whether_it_has_a_title(self, has_title_set_fixture):
        axis, new_value, expected_xml = has_title_set_fixture
        axis.has_title = new_value
        assert axis._element.xml == expected_xml

    def it_knows_whether_it_is_visible(self, visible_get_fixture):
        axis, expected_bool_value = visible_get_fixture
        assert axis.visible is expected_bool_value

    def it_can_change_whether_it_is_visible(self, visible_set_fixture):
        axis, new_value, expected_xml = visible_set_fixture
        axis.visible = new_value
        assert axis._element.xml == expected_xml

    def it_raises_on_assign_non_bool_to_visible(self):
        axis = _BaseAxis(None)
        with pytest.raises(ValueError):
            axis.visible = 'foobar'

    def it_knows_the_scale_maximum(self, maximum_scale_get_fixture):
        axis, expected_value = maximum_scale_get_fixture
        assert axis.maximum_scale == expected_value

    def it_can_change_the_scale_maximum(self, maximum_scale_set_fixture):
        axis, new_value, expected_xml = maximum_scale_set_fixture
        axis.maximum_scale = new_value
        assert axis._element.xml == expected_xml

    def it_knows_the_scale_minimum(self, minimum_scale_get_fixture):
        axis, expected_value = minimum_scale_get_fixture
        assert axis.minimum_scale == expected_value

    def it_can_change_the_scale_minimum(self, minimum_scale_set_fixture):
        axis, new_value, expected_xml = minimum_scale_set_fixture
        axis.minimum_scale = new_value
        assert axis._element.xml == expected_xml

    def it_knows_its_major_tick_setting(self, major_tick_get_fixture):
        axis, expected_value = major_tick_get_fixture
        assert axis.major_tick_mark == expected_value

    def it_can_change_its_major_tick_mark(self, major_tick_set_fixture):
        axis, new_value, expected_xml = major_tick_set_fixture
        axis.major_tick_mark = new_value
        assert axis._element.xml == expected_xml

    def it_knows_its_minor_tick_setting(self, minor_tick_get_fixture):
        axis, expected_value = minor_tick_get_fixture
        assert axis.minor_tick_mark == expected_value

    def it_can_change_its_minor_tick_mark(self, minor_tick_set_fixture):
        axis, new_value, expected_xml = minor_tick_set_fixture
        axis.minor_tick_mark = new_value
        assert axis._element.xml == expected_xml

    def it_knows_its_tick_label_position(self, tick_lbl_pos_get_fixture):
        axis, expected_value = tick_lbl_pos_get_fixture
        assert axis.tick_label_position == expected_value

    def it_can_change_its_tick_label_position(self, tick_lbl_pos_set_fixture):
        axis, new_value, expected_xml = tick_lbl_pos_set_fixture
        axis.tick_label_position = new_value
        assert axis._element.xml == expected_xml

    def it_provides_access_to_its_title(self, title_fixture):
        axis, AxisTitle_, axis_title_ = title_fixture
        axis_title = axis.axis_title
        AxisTitle_.assert_called_once_with(axis._element.title)
        assert axis_title is axis_title_

    def it_provides_access_to_its_format(self, format_fixture):
        axis, ChartFormat_, format_ = format_fixture
        format = axis.format
        ChartFormat_.assert_called_once_with(axis._xAx)
        assert format is format_

    def it_provides_access_to_its_major_gridlines(self, maj_grdlns_fixture):
        axis, MajorGridlines_, xAx, major_gridlines_ = maj_grdlns_fixture
        major_gridlines = axis.major_gridlines
        MajorGridlines_.assert_called_once_with(xAx)
        assert major_gridlines is major_gridlines_

    def it_provides_access_to_the_tick_labels(self, tick_labels_fixture):
        axis, tick_labels_, TickLabels_, xAx = tick_labels_fixture
        tick_labels = axis.tick_labels
        TickLabels_.assert_called_once_with(xAx)
        assert tick_labels is tick_labels_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['c:catAx', 'c:dateAx', 'c:valAx'])
    def format_fixture(self, request, ChartFormat_, format_):
        xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, ChartFormat_, format_

    @pytest.fixture(params=[
        ('c:catAx',          False),
        ('c:catAx/c:title',  True),
        ('c:dateAx',         False),
        ('c:dateAx/c:title', True),
        ('c:valAx',          False),
        ('c:valAx/c:title',  True),
    ])
    def has_title_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',  True, 'c:catAx/c:title/(c:layout,c:overlay{val=0})'),
        ('c:catAx/c:title',  True,  'c:catAx/c:title'),
        ('c:catAx/c:title',  False, 'c:catAx'),
        ('c:catAx',          False, 'c:catAx'),
        ('c:dateAx', True, 'c:dateAx/c:title/(c:layout,c:overlay{val=0})'),
        ('c:dateAx/c:title', True,  'c:dateAx/c:title'),
        ('c:dateAx/c:title', False, 'c:dateAx'),
        ('c:dateAx',         False, 'c:dateAx'),
        ('c:valAx',  True, 'c:valAx/c:title/(c:layout,c:overlay{val=0})'),
        ('c:valAx/c:title',  True,  'c:valAx/c:title'),
        ('c:valAx/c:title',  False, 'c:valAx'),
        ('c:valAx',          False, 'c:valAx'),
    ])
    def has_title_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=['c:catAx', 'c:dateAx', 'c:valAx'])
    def maj_grdlns_fixture(self, request, MajorGridlines_, major_gridlines_):
        xAx_cxml = request.param
        xAx = element(xAx_cxml)
        axis = _BaseAxis(xAx)
        return axis, MajorGridlines_, xAx, major_gridlines_

    @pytest.fixture(params=[
        ('c:catAx',                   False),
        ('c:catAx/c:majorGridlines',  True),
        ('c:dateAx',                  False),
        ('c:dateAx/c:majorGridlines', True),
        ('c:valAx',                   False),
        ('c:valAx/c:majorGridlines',  True),
    ])
    def major_gridlines_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        base_axis = _BaseAxis(element(xAx_cxml))
        return base_axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',                   True,  'c:catAx/c:majorGridlines'),
        ('c:catAx/c:majorGridlines',  True,  'c:catAx/c:majorGridlines'),
        ('c:catAx/c:majorGridlines',  False, 'c:catAx'),
        ('c:catAx',                   False, 'c:catAx'),
        ('c:dateAx',                  True,  'c:dateAx/c:majorGridlines'),
        ('c:dateAx/c:majorGridlines', True,  'c:dateAx/c:majorGridlines'),
        ('c:dateAx/c:majorGridlines', False, 'c:dateAx'),
        ('c:dateAx',                  False, 'c:dateAx'),
        ('c:valAx',                   True,  'c:valAx/c:majorGridlines'),
        ('c:valAx/c:majorGridlines',  True,  'c:valAx/c:majorGridlines'),
        ('c:valAx/c:majorGridlines',  False, 'c:valAx'),
        ('c:valAx',                   False, 'c:valAx'),
    ])
    def major_gridlines_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        base_axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return base_axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx',                           XL_TICK_MARK.CROSS),
        ('c:catAx/c:majorTickMark',           XL_TICK_MARK.CROSS),
        ('c:catAx/c:majorTickMark{val=out}',  XL_TICK_MARK.OUTSIDE),
        ('c:dateAx',                          XL_TICK_MARK.CROSS),
        ('c:dateAx/c:majorTickMark',          XL_TICK_MARK.CROSS),
        ('c:dateAx/c:majorTickMark{val=out}', XL_TICK_MARK.OUTSIDE),
        ('c:valAx',                           XL_TICK_MARK.CROSS),
        ('c:valAx/c:majorTickMark',           XL_TICK_MARK.CROSS),
        ('c:valAx/c:majorTickMark{val=in}',   XL_TICK_MARK.INSIDE),
    ])
    def major_tick_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',                             XL_TICK_MARK.INSIDE,
         'c:catAx/c:majorTickMark{val=in}'),
        ('c:catAx/c:majorTickMark{val=in}',     XL_TICK_MARK.OUTSIDE,
         'c:catAx/c:majorTickMark{val=out}'),
        ('c:catAx/c:majorTickMark{val=out}',    XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:catAx',                             XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:catAx/c:majorTickMark{val=cross}',  XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:dateAx',                            XL_TICK_MARK.INSIDE,
         'c:dateAx/c:majorTickMark{val=in}'),
        ('c:dateAx/c:majorTickMark{val=in}',    XL_TICK_MARK.OUTSIDE,
         'c:dateAx/c:majorTickMark{val=out}'),
        ('c:dateAx/c:majorTickMark{val=out}',   XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:dateAx',                            XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:dateAx/c:majorTickMark{val=cross}', XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:valAx',                             XL_TICK_MARK.INSIDE,
         'c:valAx/c:majorTickMark{val=in}'),
        ('c:valAx/c:majorTickMark{val=in}',     XL_TICK_MARK.OUTSIDE,
         'c:valAx/c:majorTickMark{val=out}'),
        ('c:valAx/c:majorTickMark{val=out}',    XL_TICK_MARK.CROSS,
         'c:valAx'),
        ('c:valAx',                             XL_TICK_MARK.CROSS,
         'c:valAx'),
        ('c:valAx/c:majorTickMark{val=cross}',  XL_TICK_MARK.CROSS,
         'c:valAx'),
    ])
    def major_tick_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx/c:scaling',                   None),
        ('c:catAx/c:scaling/c:max{val=12.34}',  12.34),
        ('c:dateAx/c:scaling',                  None),
        ('c:dateAx/c:scaling/c:max{val=42.24}', 42.24),
        ('c:valAx/c:scaling',                   None),
        ('c:valAx/c:scaling/c:max{val=23.45}',  23.45),
    ])
    def maximum_scale_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx/c:scaling', 34.56, 'c:catAx/c:scaling/c:max{val=34.56}'),
        ('c:catAx/c:scaling/c:max{val=34.56}', 42.42,
         'c:catAx/c:scaling/c:max{val=42.42}'),
        ('c:catAx/c:scaling/c:max{val=42.42}', None, 'c:catAx/c:scaling'),
        ('c:catAx/c:scaling', None, 'c:catAx/c:scaling'),
        ('c:dateAx/c:scaling', 45.67, 'c:dateAx/c:scaling/c:max{val=45.67}'),
        ('c:dateAx/c:scaling/c:max{val=45.67}', 42.42,
         'c:dateAx/c:scaling/c:max{val=42.42}'),
        ('c:dateAx/c:scaling/c:max{val=42.42}', None, 'c:dateAx/c:scaling'),
        ('c:dateAx/c:scaling', None, 'c:dateAx/c:scaling'),
        ('c:valAx/c:scaling', 56.78, 'c:valAx/c:scaling/c:max{val=56.78}'),
        ('c:valAx/c:scaling/c:max{val=56.78}', 42.42,
         'c:valAx/c:scaling/c:max{val=42.42}'),
        ('c:valAx/c:scaling/c:max{val=42.42}', None, 'c:valAx/c:scaling'),
        ('c:valAx/c:scaling', None, 'c:valAx/c:scaling'),
    ])
    def maximum_scale_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx/c:scaling',                   None),
        ('c:catAx/c:scaling/c:min{val=12.34}',  12.34),
        ('c:dateAx/c:scaling',                  None),
        ('c:dateAx/c:scaling/c:min{val=42.24}', 42.24),
        ('c:valAx/c:scaling',                   None),
        ('c:valAx/c:scaling/c:min{val=23.45}',  23.45),
    ])
    def minimum_scale_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx/c:scaling', 34.56, 'c:catAx/c:scaling/c:min{val=34.56}'),
        ('c:catAx/c:scaling/c:min{val=34.56}', 42.42,
         'c:catAx/c:scaling/c:min{val=42.42}'),
        ('c:catAx/c:scaling/c:min{val=42.42}', None, 'c:catAx/c:scaling'),
        ('c:catAx/c:scaling', None, 'c:catAx/c:scaling'),
        ('c:dateAx/c:scaling', 45.67, 'c:dateAx/c:scaling/c:min{val=45.67}'),
        ('c:dateAx/c:scaling/c:min{val=45.67}', 42.42,
         'c:dateAx/c:scaling/c:min{val=42.42}'),
        ('c:dateAx/c:scaling/c:min{val=42.42}', None, 'c:dateAx/c:scaling'),
        ('c:dateAx/c:scaling', None, 'c:dateAx/c:scaling'),
        ('c:valAx/c:scaling', 56.78, 'c:valAx/c:scaling/c:min{val=56.78}'),
        ('c:valAx/c:scaling/c:min{val=56.78}', 42.42,
         'c:valAx/c:scaling/c:min{val=42.42}'),
        ('c:valAx/c:scaling/c:min{val=42.42}', None, 'c:valAx/c:scaling'),
        ('c:valAx/c:scaling', None, 'c:valAx/c:scaling'),
    ])
    def minimum_scale_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx',                   False),
        ('c:catAx/c:minorGridlines',  True),
        ('c:dateAx',                  False),
        ('c:dateAx/c:minorGridlines', True),
        ('c:valAx',                   False),
        ('c:valAx/c:minorGridlines',  True),
    ])
    def minor_gridlines_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        base_axis = _BaseAxis(element(xAx_cxml))
        return base_axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',                   True,  'c:catAx/c:minorGridlines'),
        ('c:catAx/c:minorGridlines',  True,  'c:catAx/c:minorGridlines'),
        ('c:catAx/c:minorGridlines',  False, 'c:catAx'),
        ('c:catAx',                   False, 'c:catAx'),
        ('c:dateAx',                  True,  'c:dateAx/c:minorGridlines'),
        ('c:dateAx/c:minorGridlines', True,  'c:dateAx/c:minorGridlines'),
        ('c:dateAx/c:minorGridlines', False, 'c:dateAx'),
        ('c:dateAx',                  False, 'c:dateAx'),
        ('c:valAx',                   True,  'c:valAx/c:minorGridlines'),
        ('c:valAx/c:minorGridlines',  True,  'c:valAx/c:minorGridlines'),
        ('c:valAx/c:minorGridlines',  False, 'c:valAx'),
        ('c:valAx',                   False, 'c:valAx'),
    ])
    def minor_gridlines_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        base_axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return base_axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx',                           XL_TICK_MARK.CROSS),
        ('c:catAx/c:minorTickMark',           XL_TICK_MARK.CROSS),
        ('c:catAx/c:minorTickMark{val=out}',  XL_TICK_MARK.OUTSIDE),
        ('c:dateAx',                          XL_TICK_MARK.CROSS),
        ('c:dateAx/c:minorTickMark',          XL_TICK_MARK.CROSS),
        ('c:dateAx/c:minorTickMark{val=out}', XL_TICK_MARK.OUTSIDE),
        ('c:valAx',                           XL_TICK_MARK.CROSS),
        ('c:valAx/c:minorTickMark',           XL_TICK_MARK.CROSS),
        ('c:valAx/c:minorTickMark{val=in}',   XL_TICK_MARK.INSIDE),
    ])
    def minor_tick_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',                             XL_TICK_MARK.INSIDE,
         'c:catAx/c:minorTickMark{val=in}'),
        ('c:catAx/c:minorTickMark{val=in}',     XL_TICK_MARK.OUTSIDE,
         'c:catAx/c:minorTickMark{val=out}'),
        ('c:catAx/c:minorTickMark{val=out}',    XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:catAx',                             XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:catAx/c:minorTickMark{val=cross}',  XL_TICK_MARK.CROSS,
         'c:catAx'),
        ('c:dateAx',                            XL_TICK_MARK.INSIDE,
         'c:dateAx/c:minorTickMark{val=in}'),
        ('c:dateAx/c:minorTickMark{val=in}',    XL_TICK_MARK.OUTSIDE,
         'c:dateAx/c:minorTickMark{val=out}'),
        ('c:dateAx/c:minorTickMark{val=out}',   XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:dateAx',                            XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:dateAx/c:minorTickMark{val=cross}', XL_TICK_MARK.CROSS,
         'c:dateAx'),
        ('c:valAx',                             XL_TICK_MARK.INSIDE,
         'c:valAx/c:minorTickMark{val=in}'),
        ('c:valAx/c:minorTickMark{val=in}',     XL_TICK_MARK.OUTSIDE,
         'c:valAx/c:minorTickMark{val=out}'),
        ('c:valAx/c:minorTickMark{val=out}',    XL_TICK_MARK.CROSS,
         'c:valAx'),
        ('c:valAx',                             XL_TICK_MARK.CROSS,
         'c:valAx'),
        ('c:valAx/c:minorTickMark{val=cross}',  XL_TICK_MARK.CROSS,
         'c:valAx'),
    ])
    def minor_tick_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=['c:catAx', 'c:dateAx', 'c:valAx'])
    def tick_labels_fixture(self, request, TickLabels_, tick_labels_):
        xAx_cxml = request.param
        xAx = element(xAx_cxml)
        axis = _BaseAxis(xAx)
        return axis, tick_labels_, TickLabels_, xAx

    @pytest.fixture(params=[
        ('c:catAx',                        XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:catAx/c:tickLblPos',           XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:catAx/c:tickLblPos{val=high}', XL_TICK_LBL_POS.HIGH),
        ('c:dateAx',                       XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:dateAx/c:tickLblPos',          XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:dateAx/c:tickLblPos{val=low}', XL_TICK_LBL_POS.LOW),
        ('c:valAx',                        XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:valAx/c:tickLblPos',           XL_TICK_LBL_POS.NEXT_TO_AXIS),
        ('c:valAx/c:tickLblPos{val=none}', XL_TICK_LBL_POS.NONE),
    ])
    def tick_lbl_pos_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_value

    @pytest.fixture(params=[
        ('c:catAx',                           XL_TICK_LBL_POS.HIGH,
         'c:catAx/c:tickLblPos{val=high}'),
        ('c:catAx/c:tickLblPos{val=high}',    XL_TICK_LBL_POS.LOW,
         'c:catAx/c:tickLblPos{val=low}'),
        ('c:catAx/c:tickLblPos{val=low}',     None,
         'c:catAx/c:tickLblPos'),
        ('c:catAx',                           None,
         'c:catAx/c:tickLblPos'),
        ('c:dateAx',                          XL_TICK_LBL_POS.NEXT_TO_AXIS,
         'c:dateAx/c:tickLblPos{val=nextTo}'),
        ('c:dateAx/c:tickLblPos{val=nextTo}', XL_TICK_LBL_POS.NONE,
         'c:dateAx/c:tickLblPos{val=none}'),
        ('c:dateAx/c:tickLblPos{val=none}',   None,
         'c:dateAx/c:tickLblPos'),
        ('c:valAx',                           XL_TICK_LBL_POS.HIGH,
         'c:valAx/c:tickLblPos{val=high}'),
        ('c:valAx/c:tickLblPos{val=high}',    XL_TICK_LBL_POS.LOW,
         'c:valAx/c:tickLblPos{val=low}'),
        ('c:valAx/c:tickLblPos{val=low}',     None,
         'c:valAx/c:tickLblPos'),
    ])
    def tick_lbl_pos_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    @pytest.fixture(params=['c:catAx', 'c:dateAx', 'c:valAx'])
    def title_fixture(self, request, AxisTitle_, axis_title_):
        xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, AxisTitle_, axis_title_

    @pytest.fixture(params=[
        ('c:catAx',                      False),
        ('c:catAx/c:delete',             False),
        ('c:catAx/c:delete{val=0}',      True),
        ('c:catAx/c:delete{val=1}',      False),
        ('c:catAx/c:delete{val=false}',  True),
        ('c:dateAx',                     False),
        ('c:dateAx/c:delete',            False),
        ('c:dateAx/c:delete{val=0}',     True),
        ('c:dateAx/c:delete{val=1}',     False),
        ('c:dateAx/c:delete{val=false}', True),
        ('c:valAx',                      False),
        ('c:valAx/c:delete',             False),
        ('c:valAx/c:delete{val=0}',      True),
        ('c:valAx/c:delete{val=1}',      False),
        ('c:valAx/c:delete{val=false}',  True),
    ])
    def visible_get_fixture(self, request):
        xAx_cxml, expected_bool_value = request.param
        axis = _BaseAxis(element(xAx_cxml))
        return axis, expected_bool_value

    @pytest.fixture(params=[
        ('c:catAx',                  False, 'c:catAx/c:delete'),
        ('c:catAx/c:delete',         True,  'c:catAx/c:delete{val=0}'),
        ('c:catAx/c:delete{val=1}',  True,  'c:catAx/c:delete{val=0}'),
        ('c:catAx/c:delete{val=0}',  False, 'c:catAx/c:delete'),
        ('c:catAx',                  True,  'c:catAx/c:delete{val=0}'),
        ('c:dateAx',                 False, 'c:dateAx/c:delete'),
        ('c:dateAx/c:delete',        True,  'c:dateAx/c:delete{val=0}'),
        ('c:dateAx/c:delete{val=0}', False, 'c:dateAx/c:delete'),
        ('c:dateAx',                 True,  'c:dateAx/c:delete{val=0}'),
        ('c:valAx/c:delete',         True,  'c:valAx/c:delete{val=0}'),
        ('c:valAx/c:delete{val=1}',  False, 'c:valAx/c:delete'),
        ('c:valAx/c:delete{val=0}',  True,  'c:valAx/c:delete{val=0}'),
    ])
    def visible_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        axis = _BaseAxis(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return axis, new_value, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def AxisTitle_(self, request, axis_title_):
        return class_mock(
            request, 'pptx.chart.axis.AxisTitle', return_value=axis_title_
        )

    @pytest.fixture
    def axis_title_(self, request):
        return instance_mock(request, AxisTitle)

    @pytest.fixture
    def ChartFormat_(self, request, format_):
        return class_mock(
            request, 'pptx.chart.axis.ChartFormat', return_value=format_
        )

    @pytest.fixture
    def format_(self, request):
        return instance_mock(request, ChartFormat)

    @pytest.fixture
    def MajorGridlines_(self, request, major_gridlines_):
        return class_mock(
            request, 'pptx.chart.axis.MajorGridlines',
            return_value=major_gridlines_
        )

    @pytest.fixture
    def major_gridlines_(self, request):
        return instance_mock(request, MajorGridlines)

    @pytest.fixture
    def TickLabels_(self, request, tick_labels_):
        return class_mock(
            request, 'pptx.chart.axis.TickLabels',
            return_value=tick_labels_
        )

    @pytest.fixture
    def tick_labels_(self, request):
        return instance_mock(request, TickLabels)


class DescribeAxisTitle(object):

    def it_knows_whether_it_has_a_text_frame(self, has_tf_get_fixture):
        axis_title, expected_value = has_tf_get_fixture
        value = axis_title.has_text_frame
        assert value is expected_value

    def it_can_change_whether_it_has_a_text_frame(self, has_tf_set_fixture):
        axis_title, value, expected_xml = has_tf_set_fixture
        axis_title.has_text_frame = value
        assert axis_title._element.xml == expected_xml

    def it_provides_access_to_its_format(self, format_fixture):
        axis_title, ChartFormat_, format_ = format_fixture
        format = axis_title.format
        ChartFormat_.assert_called_once_with(axis_title._element)
        assert format is format_

    def it_provides_access_to_its_text_frame(self, text_frame_fixture):
        axis_title, TextFrame_, text_frame_ = text_frame_fixture
        text_frame = axis_title.text_frame
        TextFrame_.assert_called_once_with(
            axis_title._element.tx.rich, axis_title
        )
        assert text_frame is text_frame_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def format_fixture(self, request, ChartFormat_, format_):
        axis_title = AxisTitle(element('c:title'))
        return axis_title, ChartFormat_, format_

    @pytest.fixture(params=[
        ('c:title',               False),
        ('c:title/c:tx',          False),
        ('c:title/c:tx/c:strRef', False),
        ('c:title/c:tx/c:rich',   True),
    ])
    def has_tf_get_fixture(self, request):
        title_cxml, expected_value = request.param
        axis_title = AxisTitle(element(title_cxml))
        return axis_title, expected_value

    @pytest.fixture(params=[
        ('c:title{a:b=c}', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title{a:b=c}/c:tx', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title{a:b=c}/c:tx/c:strRef', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title/c:tx/c:rich',   True,  'c:title/c:tx/c:rich'),
        ('c:title',               False, 'c:title'),
        ('c:title/c:tx',          False, 'c:title'),
        ('c:title/c:tx/c:rich',   False, 'c:title'),
        ('c:title/c:tx/c:strRef', False, 'c:title'),
    ])
    def has_tf_set_fixture(self, request):
        title_cxml, value, expected_cxml = request.param
        axis_title = AxisTitle(element(title_cxml))
        expected_xml = xml(expected_cxml)
        return axis_title, value, expected_xml

    @pytest.fixture
    def text_frame_fixture(self, request, TextFrame_):
        axis_title = AxisTitle(element('c:title'))
        text_frame_ = TextFrame_.return_value
        return axis_title, TextFrame_, text_frame_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def ChartFormat_(self, request, format_):
        return class_mock(
            request, 'pptx.chart.axis.ChartFormat', return_value=format_
        )

    @pytest.fixture
    def format_(self, request):
        return instance_mock(request, ChartFormat)

    @pytest.fixture
    def TextFrame_(self, request):
        return class_mock(request, 'pptx.chart.axis.TextFrame')


class DescribeCategoryAxis(object):

    def it_knows_its_category_type(self, cat_type_get_fixture):
        category_axis, expected_value = cat_type_get_fixture
        assert category_axis.category_type is expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def cat_type_get_fixture(self):
        category_axis = CategoryAxis(None)
        expected_value = XL_CATEGORY_TYPE.CATEGORY_SCALE
        return category_axis, expected_value


class DescribeDateAxis(object):

    def it_knows_its_category_type(self, cat_type_get_fixture):
        date_axis, expected_value = cat_type_get_fixture
        assert date_axis.category_type is expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def cat_type_get_fixture(self):
        date_axis = DateAxis(None)
        expected_value = XL_CATEGORY_TYPE.TIME_SCALE
        return date_axis, expected_value


class DescribeMajorGridlines(object):

    def it_provides_access_to_its_format(self, format_fixture):
        gridlines, expected_xml, ChartFormat_, format_ = format_fixture
        format = gridlines.format
        assert gridlines._xAx.xml == expected_xml
        ChartFormat_.assert_called_once_with(
            gridlines._xAx.xpath('c:majorGridlines')[0]
        )
        assert format is format_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('c:valAx',                  'c:valAx/c:majorGridlines'),
        ('c:catAx/c:majorGridlines', 'c:catAx/c:majorGridlines'),
    ])
    def format_fixture(self, request, ChartFormat_, format_):
        xAx_cxml, expected_cxml = request.param
        gridlines = MajorGridlines(element(xAx_cxml))
        expected_xml = xml(expected_cxml)
        return gridlines, expected_xml, ChartFormat_, format_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def ChartFormat_(self, request, format_):
        return class_mock(
            request, 'pptx.chart.axis.ChartFormat', return_value=format_
        )

    @pytest.fixture
    def format_(self, request):
        return instance_mock(request, ChartFormat)


class DescribeTickLabels(object):

    def it_provides_access_to_its_font(self, font_fixture):
        tick_labels, Font_, defRPr, font_ = font_fixture
        font = tick_labels.font
        Font_.assert_called_once_with(defRPr)
        assert font is font_

    def it_adds_a_txPr_to_help_font(self, txPr_fixture):
        tick_labels, expected_xml = txPr_fixture
        tick_labels.font
        assert tick_labels._element.xml == expected_xml

    def it_knows_its_number_format(self, number_format_get_fixture):
        tick_labels, expected_value = number_format_get_fixture
        assert tick_labels.number_format == expected_value

    def it_can_change_its_number_format(self, number_format_set_fixture):
        tick_labels, new_value, expected_xml = number_format_set_fixture
        tick_labels.number_format = new_value
        assert tick_labels._element.xml == expected_xml

    def it_knows_whether_its_number_format_is_linked(
            self, number_format_is_linked_get_fixture):
        tick_labels, expected_value = number_format_is_linked_get_fixture
        assert tick_labels.number_format_is_linked is expected_value

    def it_can_change_whether_its_number_format_is_linked(
            self, number_format_is_linked_set_fixture):
        tick_labels, new_value, expected_xml = (
            number_format_is_linked_set_fixture
        )
        tick_labels.number_format_is_linked = new_value
        assert tick_labels._element.xml == expected_xml

    def it_knows_its_offset(self, offset_get_fixture):
        tick_labels, expected_value = offset_get_fixture
        assert tick_labels.offset == expected_value

    def it_can_change_its_offset(self, offset_set_fixture):
        tick_labels, new_value, expected_xml = offset_set_fixture
        tick_labels.offset = new_value
        assert tick_labels._element.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def font_fixture(self, Font_, font_):
        catAx = element('c:catAx/c:txPr/a:p/a:pPr/a:defRPr')
        defRPr = catAx.xpath('.//a:defRPr')[0]
        tick_labels = TickLabels(catAx)
        return tick_labels, Font_, defRPr, font_

    @pytest.fixture(params=[
        ('c:catAx',                              'General'),
        ('c:valAx/c:numFmt{formatCode=General}', 'General'),
    ])
    def number_format_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        tick_labels = TickLabels(element(xAx_cxml))
        return tick_labels, expected_value

    @pytest.fixture(params=[
        ('c:catAx', 'General',
         'c:catAx/c:numFmt{formatCode=General,sourceLinked=0}'),
        ('c:valAx/c:numFmt{formatCode=General}', '00.00',
         'c:valAx/c:numFmt{formatCode=00.00,sourceLinked=0}'),
    ])
    def number_format_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        tick_labels = TickLabels(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return tick_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx',                          False),
        ('c:valAx/c:numFmt',                 True),
        ('c:valAx/c:numFmt{sourceLinked=0}', False),
        ('c:catAx/c:numFmt{sourceLinked=1}', True),
    ])
    def number_format_is_linked_get_fixture(self, request):
        xAx_cxml, expected_value = request.param
        tick_labels = TickLabels(element(xAx_cxml))
        return tick_labels, expected_value

    @pytest.fixture(params=[
        ('c:valAx', True,  'c:valAx/c:numFmt{sourceLinked=1}'),
        ('c:catAx', False, 'c:catAx/c:numFmt{sourceLinked=0}'),
        ('c:valAx', None,  'c:valAx/c:numFmt'),
        ('c:catAx/c:numFmt', True, 'c:catAx/c:numFmt{sourceLinked=1}'),
        ('c:valAx/c:numFmt{sourceLinked=1}', False,
         'c:valAx/c:numFmt{sourceLinked=0}'),
    ])
    def number_format_is_linked_set_fixture(self, request):
        xAx_cxml, new_value, expected_xAx_cxml = request.param
        tick_labels = TickLabels(element(xAx_cxml))
        expected_xml = xml(expected_xAx_cxml)
        return tick_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:catAx',                      100),
        ('c:catAx/c:lblOffset',          100),
        ('c:catAx/c:lblOffset{val=420}', 420),
        ('c:catAx/c:lblOffset{val=004}',   4),
        ('c:catAx/c:lblOffset{val=42%}',  42),
        ('c:catAx/c:lblOffset{val=02%}',   2),
    ])
    def offset_get_fixture(self, request):
        catAx_cxml, expected_value = request.param
        tick_labels = TickLabels(element(catAx_cxml))
        return tick_labels, expected_value

    @pytest.fixture(params=[
        ('c:catAx', 420, 'c:catAx/c:lblOffset{val=420}'),
        ('c:catAx/c:lblOffset{val=420}', 100, 'c:catAx'),
    ])
    def offset_set_fixture(self, request):
        catAx_cxml, new_value, expected_catAx_cxml = request.param
        tick_labels = TickLabels(element(catAx_cxml))
        expected_xml = xml(expected_catAx_cxml)
        return tick_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:valAx{a:b=c}',
         'c:valAx{a:b=c}/c:txPr/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr)'),
        ('c:valAx{a:b=c}/c:txPr/(a:bodyPr,a:p)',
         'c:valAx{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr/a:defRPr)'),
        ('c:valAx{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr)',
         'c:valAx{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr/a:defRPr)'),
    ])
    def txPr_fixture(self, request):
        xAx_cxml, expected_cxml = request.param
        tick_labels = TickLabels(element(xAx_cxml))
        expected_xml = xml(expected_cxml)
        return tick_labels, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def Font_(self, request, font_):
        return class_mock(
            request, 'pptx.chart.axis.Font', return_value=font_
        )

    @pytest.fixture
    def font_(self, request):
        return instance_mock(request, Font)


class DescribeValueAxis(object):

    def it_knows_the_other_axis_crossing_type(self, crosses_get_fixture):
        value_axis, expected_value = crosses_get_fixture
        assert value_axis.crosses == expected_value

    def it_can_change_the_other_axis_crossing_type(self, crosses_set_fixture):
        value_axis, new_value, plotArea, expected_xml = crosses_set_fixture
        value_axis.crosses = new_value
        assert plotArea.xml == expected_xml

    def it_knows_the_other_axis_crossing_value(self, crosses_at_get_fixture):
        value_axis, expected_value = crosses_at_get_fixture
        assert value_axis.crosses_at == expected_value

    def it_can_change_the_other_axis_crossing_value(
            self, crosses_at_set_fixture):
        value_axis, new_value, plotArea, expected_xml = crosses_at_set_fixture
        value_axis.crosses_at = new_value
        assert plotArea.xml == expected_xml

    def it_knows_its_major_unit(self, major_unit_get_fixture):
        value_axis, expected_value = major_unit_get_fixture
        assert value_axis.major_unit == expected_value

    def it_can_change_its_major_unit(self, major_unit_set_fixture):
        value_axis, new_value, expected_xml = major_unit_set_fixture
        value_axis.major_unit = new_value
        assert value_axis._element.xml == expected_xml

    def it_knows_its_minor_unit(self, minor_unit_get_fixture):
        value_axis, expected_value = minor_unit_get_fixture
        assert value_axis.minor_unit == expected_value

    def it_can_change_its_minor_unit(self, minor_unit_set_fixture):
        value_axis, new_value, expected_xml = minor_unit_set_fixture
        value_axis.minor_unit = new_value
        assert value_axis._element.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('c:plotArea/(c:valAx/c:axId{val=42},c:valAx/c:crossAx{val=42})',
         None),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx/'
         'c:crossAx{val=42})', 2.4),
        ('c:plotArea/(c:dateAx/(c:axId{val=42},c:crossesAt{val=-1.2}),c:valA'
         'x/c:crossAx{val=42})', -1.2),
    ])
    def crosses_at_get_fixture(self, request):
        plotArea_cxml, expected_value = request.param
        plotArea = element(plotArea_cxml)
        valAx = plotArea.xpath('c:valAx[c:crossAx/@val="42"]')[0]
        value_axis = ValueAxis(valAx)
        return value_axis, expected_value

    @pytest.fixture(params=[
        ('c:plotArea/(c:valAx/c:axId{val=42},c:valAx/c:crossAx{val=42})',
         2.4,
         'c:plotArea/(c:valAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx/'
         'c:crossAx{val=42})'),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crosses{val=min}),c:valAx/c:'
         'crossAx{val=42})', 1.5,
         'c:plotArea/(c:catAx/(c:axId{val=42},c:crossesAt{val=1.5}),c:valAx/'
         'c:crossAx{val=42})'),
        ('c:plotArea/(c:dateAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx'
         '/c:crossAx{val=42})', 1.5,
         'c:plotArea/(c:dateAx/(c:axId{val=42},c:crossesAt{val=1.5}),c:valAx'
         '/c:crossAx{val=42})'),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crossesAt{val=1.5}),c:valAx/'
         'c:crossAx{val=42})', None,
         'c:plotArea/(c:catAx/(c:axId{val=42}),c:valAx/c:crossAx{val=42})'),
    ])
    def crosses_at_set_fixture(self, request):
        plotArea_cxml, new_value, expected_cxml = request.param
        plotArea = element(plotArea_cxml)
        valAx = plotArea.xpath('c:valAx[c:crossAx/@val="42"]')[0]
        value_axis = ValueAxis(valAx)
        expected_xml = xml(expected_cxml)
        return value_axis, new_value, plotArea, expected_xml

    @pytest.fixture(params=[
        ('c:plotArea/(c:valAx/c:axId{val=42},c:valAx/c:crossAx{val=42})',
         'CUSTOM'),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crosses{val=autoZero}),c:val'
         'Ax/c:crossAx{val=42})',
         'AUTOMATIC'),
        ('c:plotArea/(c:valAx/(c:axId{val=42},c:crosses{val=min}),c:valAx/c:'
         'crossAx{val=42})',
         'MINIMUM'),
    ])
    def crosses_get_fixture(self, request):
        cxml, member = request.param
        valAx = element(cxml).xpath('c:valAx[c:crossAx/@val="42"]')[0]
        value_axis = ValueAxis(valAx)
        expected_value = getattr(XL_AXIS_CROSSES, member)
        return value_axis, expected_value

    @pytest.fixture(params=[
        ('c:plotArea/(c:valAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx/'
         'c:crossAx{val=42})', 'AUTOMATIC',
         'c:plotArea/(c:valAx/(c:axId{val=42},c:crosses{val=autoZero}),c:val'
         'Ax/c:crossAx{val=42})'),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crosses{val=autoZero}),c:val'
         'Ax/c:crossAx{val=42})', 'MINIMUM',
         'c:plotArea/(c:catAx/(c:axId{val=42},c:crosses{val=min}),c:valAx/c:'
         'crossAx{val=42})'),
        ('c:plotArea/(c:valAx/(c:axId{val=42},c:crosses{val=min}),c:valAx/c:'
         'crossAx{val=42})', 'CUSTOM',
         'c:plotArea/(c:valAx/(c:axId{val=42},c:crossesAt{val=0.0}),c:valAx/'
         'c:crossAx{val=42})'),
        ('c:plotArea/(c:catAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx/'
         'c:crossAx{val=42})', 'CUSTOM',
         'c:plotArea/(c:catAx/(c:axId{val=42},c:crossesAt{val=2.4}),c:valAx/'
         'c:crossAx{val=42})'),
    ])
    def crosses_set_fixture(self, request):
        plotArea_cxml, member, expected_cxml = request.param
        plotArea = element(plotArea_cxml)
        valAx = plotArea.xpath('c:valAx[c:crossAx/@val="42"]')[0]
        value_axis = ValueAxis(valAx)
        new_value = getattr(XL_AXIS_CROSSES, member)
        expected_xml = xml(expected_cxml)
        return value_axis, new_value, plotArea, expected_xml

    @pytest.fixture(params=[
        ('c:valAx', None),
        ('c:valAx/c:majorUnit{val=4.2}', 4.2),
    ])
    def major_unit_get_fixture(self, request):
        valAx_cxml, expected_value = request.param
        value_axis = ValueAxis(element(valAx_cxml))
        return value_axis, expected_value

    @pytest.fixture(params=[
        ('c:valAx',                        42,
         'c:valAx/c:majorUnit{val=42.0}'),
        ('c:valAx',                        None,
         'c:valAx'),
        ('c:valAx/c:majorUnit{val=42.0}',  24.0,
         'c:valAx/c:majorUnit{val=24.0}'),
        ('c:valAx/c:majorUnit{val=42.0}',  None,
         'c:valAx'),
    ])
    def major_unit_set_fixture(self, request):
        valAx_cxml, new_value, expected_valAx_cxml = request.param
        value_axis = ValueAxis(element(valAx_cxml))
        expected_xml = xml(expected_valAx_cxml)
        return value_axis, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:valAx', None),
        ('c:valAx/c:minorUnit{val=2.4}', 2.4),
    ])
    def minor_unit_get_fixture(self, request):
        valAx_cxml, expected_value = request.param
        value_axis = ValueAxis(element(valAx_cxml))
        return value_axis, expected_value

    @pytest.fixture(params=[
        ('c:valAx',                        36,
         'c:valAx/c:minorUnit{val=36.0}'),
        ('c:valAx',                        None,
         'c:valAx'),
        ('c:valAx/c:minorUnit{val=36.0}',  12.6,
         'c:valAx/c:minorUnit{val=12.6}'),
        ('c:valAx/c:minorUnit{val=36.0}',  None,
         'c:valAx'),
    ])
    def minor_unit_set_fixture(self, request):
        valAx_cxml, new_value, expected_valAx_cxml = request.param
        value_axis = ValueAxis(element(valAx_cxml))
        expected_xml = xml(expected_valAx_cxml)
        return value_axis, new_value, expected_xml
