# encoding: utf-8

"""Test suite for pptx.chart.chart module"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.chart.axis import CategoryAxis, DateAxis, ValueAxis
from pptx.chart.chart import Chart, ChartTitle, Legend, _Plots
from pptx.chart.data import ChartData
from pptx.chart.plot import _BasePlot
from pptx.chart.series import SeriesCollection
from pptx.chart.xmlwriter import _BaseSeriesXmlRewriter
from pptx.dml.chtfmt import ChartFormat
from pptx.enum.chart import XL_CHART_TYPE
from pptx.parts.chart import ChartWorkbook
from pptx.text.text import Font

from ..unitutil.cxml import element, xml
from ..unitutil.mock import (
    call, class_mock, function_mock, instance_mock, property_mock
)


class DescribeChart(object):

    def it_provides_access_to_its_font(self, font_fixture, Font_, font_):
        chartSpace, expected_xml = font_fixture
        Font_.return_value = font_
        chart = Chart(chartSpace, None)

        font = chart.font

        assert chartSpace.xml == expected_xml
        Font_.assert_called_once_with(
            chartSpace.xpath('./c:txPr/a:p/a:pPr/a:defRPr')[0]
        )
        assert font is font_

    def it_knows_whether_it_has_a_title(self, has_title_get_fixture):
        chart, expected_value = has_title_get_fixture
        assert chart.has_title is expected_value

    def it_can_change_whether_it_has_a_title(self, has_title_set_fixture):
        chart, new_value, expected_xml = has_title_set_fixture
        chart.has_title = new_value
        assert chart._chartSpace.chart.xml == expected_xml

    def it_provides_access_to_the_chart_title(self, title_fixture):
        chart, expected_xml, ChartTitle_, chart_title_ = title_fixture

        chart_title = chart.chart_title

        assert chart.element.xpath('c:chart/c:title')[0].xml == expected_xml
        ChartTitle_.assert_called_once_with(chart.element.chart.title)
        assert chart_title is chart_title_

    def it_provides_access_to_the_category_axis(self, category_axis_fixture):
        chart, category_axis_, AxisCls_, xAx = category_axis_fixture
        category_axis = chart.category_axis
        AxisCls_.assert_called_once_with(xAx)
        assert category_axis is category_axis_

    def it_raises_when_no_category_axis(self, cat_ax_raise_fixture):
        chart = cat_ax_raise_fixture
        with pytest.raises(ValueError):
            chart.category_axis

    def it_provides_access_to_the_value_axis(self, val_ax_fixture):
        chart, ValueAxis_, valAx, value_axis_ = val_ax_fixture
        value_axis = chart.value_axis
        ValueAxis_.assert_called_once_with(valAx)
        assert value_axis is value_axis_

    def it_raises_when_no_value_axis(self, val_ax_raise_fixture):
        chart = val_ax_raise_fixture
        with pytest.raises(ValueError):
            chart.value_axis

    def it_provides_access_to_its_series(self, series_fixture):
        chart, SeriesCollection_, plotArea, series_ = series_fixture
        series = chart.series
        SeriesCollection_.assert_called_once_with(plotArea)
        assert series is series_

    def it_provides_access_to_its_plots(self, plots_fixture):
        chart, plots_, _Plots_, plotArea = plots_fixture
        plots = chart.plots
        _Plots_.assert_called_once_with(plotArea, chart)
        assert plots is plots_

    def it_knows_whether_it_has_a_legend(self, has_legend_get_fixture):
        chart, expected_value = has_legend_get_fixture
        assert chart.has_legend == expected_value

    def it_can_change_whether_it_has_a_legend(self, has_legend_set_fixture):
        chart, new_value, expected_xml = has_legend_set_fixture
        chart.has_legend = new_value
        assert chart._chartSpace.xml == expected_xml

    def it_provides_access_to_its_legend(self, legend_fixture):
        chart, Legend_, expected_calls, expected_value = legend_fixture
        legend = chart.legend
        assert Legend_.call_args_list == expected_calls
        assert legend is expected_value

    def it_knows_its_chart_type(self, chart_type_fixture):
        chart, PlotTypeInspector_, plot_, chart_type = chart_type_fixture
        _chart_type = chart.chart_type
        PlotTypeInspector_.chart_type.assert_called_once_with(plot_)
        assert _chart_type is chart_type

    def it_knows_its_style(self, style_get_fixture):
        chart, expected_value = style_get_fixture
        assert chart.chart_style == expected_value

    def it_can_change_its_style(self, style_set_fixture):
        chart, new_value, expected_xml = style_set_fixture
        chart.chart_style = new_value
        assert chart._chartSpace.xml == expected_xml

    def it_can_replace_the_chart_data(self, replace_fixture):
        (chart, chart_data_, SeriesXmlRewriterFactory_, chart_type,
         rewriter_, chartSpace, workbook_, xlsx_blob) = replace_fixture

        chart.replace_data(chart_data_)

        SeriesXmlRewriterFactory_.assert_called_once_with(
            chart_type, chart_data_
        )
        rewriter_.replace_series_data.assert_called_once_with(chartSpace)
        workbook_.update_from_xlsx_blob.assert_called_once_with(xlsx_blob)

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['c:catAx', 'c:dateAx', 'c:valAx'])
    def category_axis_fixture(self, request, CategoryAxis_, DateAxis_,
                              ValueAxis_):
        ax_tag = request.param
        chartSpace_cxml = 'c:chartSpace/c:chart/c:plotArea/%s' % ax_tag
        chartSpace = element(chartSpace_cxml)
        chart = Chart(chartSpace, None)
        AxisCls_ = {
            'c:catAx':  CategoryAxis_,
            'c:dateAx': DateAxis_,
            'c:valAx':  ValueAxis_
        }[ax_tag]
        axis_ = AxisCls_.return_value
        xAx = chartSpace.xpath('.//%s' % ax_tag)[0]
        return chart, axis_, AxisCls_, xAx

    @pytest.fixture
    def cat_ax_raise_fixture(self):
        chart = Chart(element('c:chartSpace/c:chart/c:plotArea'), None)
        return chart

    @pytest.fixture
    def chart_type_fixture(self, PlotTypeInspector_, plot_):
        chart = Chart(None, None)
        chart._plots = [plot_]
        chart_type = XL_CHART_TYPE.PIE
        PlotTypeInspector_.chart_type.return_value = chart_type
        return chart, PlotTypeInspector_, plot_, chart_type

    @pytest.fixture(params=[
        ('c:chartSpace{a:b=c}',
         'c:chartSpace{a:b=c}/c:txPr/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:chartSpace/c:txPr/a:p',
         'c:chartSpace/c:txPr/a:p/a:pPr/a:defRPr'),
        ('c:chartSpace/c:txPr/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr)',
         'c:chartSpace/c:txPr/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr)'),
    ])
    def font_fixture(self, request):
        chartSpace_cxml, expected_cxml = request.param
        chartSpace = element(chartSpace_cxml)
        expected_xml = xml(expected_cxml)
        return chartSpace, expected_xml

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart', False),
        ('c:chartSpace/c:chart/c:legend', True),
    ])
    def has_legend_get_fixture(self, request):
        chartSpace_cxml, expected_value = request.param
        chart = Chart(element(chartSpace_cxml), None)
        return chart, expected_value

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart', True,
         'c:chartSpace/c:chart/c:legend'),
    ])
    def has_legend_set_fixture(self, request):
        chartSpace_cxml, new_value, expected_chartSpace_cxml = request.param
        chart = Chart(element(chartSpace_cxml), None)
        expected_xml = xml(expected_chartSpace_cxml)
        return chart, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart', False),
        ('c:chartSpace/c:chart/c:title', True),
    ])
    def has_title_get_fixture(self, request):
        chartSpace_cxml, expected_value = request.param
        chart = Chart(element(chartSpace_cxml), None)
        return chart, expected_value

    @pytest.fixture(params=[
        ('c:chart', True, 'c:chart/c:title/(c:layout,c:overlay{val=0})'),
        ('c:chart/c:title', True, 'c:chart/c:title'),
        ('c:chart/c:title', False, 'c:chart/c:autoTitleDeleted{val=1}'),
        ('c:chart', False, 'c:chart/c:autoTitleDeleted{val=1}'),
    ])
    def has_title_set_fixture(self, request):
        chart_cxml, new_value, expected_cxml = request.param
        chart = Chart(element('c:chartSpace/%s' % chart_cxml), None)
        expected_xml = xml(expected_cxml)
        return chart, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart', False),
        ('c:chartSpace/c:chart/c:legend', True),
    ])
    def legend_fixture(self, request, Legend_, legend_):
        chartSpace_cxml, has_legend = request.param
        chartSpace = element(chartSpace_cxml)
        chart = Chart(chartSpace, None)
        expected_value, expected_calls = None, []
        if has_legend:
            expected_value = legend_
            legend_elm = chartSpace.chart.legend
            expected_calls.append(call(legend_elm))
        return chart, Legend_, expected_calls, expected_value

    @pytest.fixture
    def plots_fixture(self, _Plots_, plots_):
        chartSpace = element('c:chartSpace/c:chart/c:plotArea')
        plotArea = chartSpace.xpath('./c:chart/c:plotArea')[0]
        chart = Chart(chartSpace, None)
        return chart, plots_, _Plots_, plotArea

    @pytest.fixture
    def replace_fixture(
            self, chart_data_, SeriesXmlRewriterFactory_, series_rewriter_,
            workbook_, workbook_prop_):
        chartSpace = element('c:chartSpace/c:chart/c:plotArea/c:pieChart')
        chart = Chart(chartSpace, None)
        chart_type = XL_CHART_TYPE.PIE
        xlsx_blob = 'fooblob'
        chart_data_.xlsx_blob = xlsx_blob
        return (
            chart, chart_data_, SeriesXmlRewriterFactory_, chart_type,
            series_rewriter_, chartSpace, workbook_, xlsx_blob
        )

    @pytest.fixture
    def series_fixture(self, SeriesCollection_, series_collection_):
        chartSpace = element('c:chartSpace/c:chart/c:plotArea')
        plotArea = chartSpace.xpath('.//c:plotArea')[0]
        chart = Chart(chartSpace, None)
        return chart, SeriesCollection_, plotArea, series_collection_

    @pytest.fixture(params=[
        ('c:chartSpace/c:style{val=42}', 42),
        ('c:chartSpace',                 None),
    ])
    def style_get_fixture(self, request):
        chartSpace_cxml, expected_value = request.param
        chart = Chart(element(chartSpace_cxml), None)
        return chart, expected_value

    @pytest.fixture(params=[
        ('c:chartSpace',                4,    'c:chartSpace/c:style{val=4}'),
        ('c:chartSpace',                None, 'c:chartSpace'),
        ('c:chartSpace/c:style{val=4}', 2,    'c:chartSpace/c:style{val=2}'),
        ('c:chartSpace/c:style{val=4}', None, 'c:chartSpace'),
    ])
    def style_set_fixture(self, request):
        chartSpace_cxml, new_value, expected_chartSpace_cxml = request.param
        chart = Chart(element(chartSpace_cxml), None)
        expected_xml = xml(expected_chartSpace_cxml)
        return chart, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart',
         'c:title/(c:layout,c:overlay{val=0})'),
        ('c:chartSpace/c:chart/c:title/c:layout',
         'c:title/c:layout'),
    ])
    def title_fixture(self, request, ChartTitle_, chart_title_):
        chartSpace_cxml, expected_cxml = request.param
        chart = Chart(element(chartSpace_cxml), None)
        expected_xml = xml(expected_cxml)
        return chart, expected_xml, ChartTitle_, chart_title_

    @pytest.fixture(params=[
        ('c:chartSpace/c:chart/c:plotArea/(c:catAx,c:valAx)', 0),
        ('c:chartSpace/c:chart/c:plotArea/(c:valAx,c:valAx)', 1),
    ])
    def val_ax_fixture(self, request, ValueAxis_, value_axis_):
        chartSpace_xml, idx = request.param
        chartSpace = element(chartSpace_xml)
        chart = Chart(chartSpace, None)
        valAx = chartSpace.xpath('.//c:valAx')[idx]
        return chart, ValueAxis_, valAx, value_axis_

    @pytest.fixture
    def val_ax_raise_fixture(self):
        chart = Chart(element('c:chartSpace/c:chart/c:plotArea'), None)
        return chart

    # fixture components ---------------------------------------------

    @pytest.fixture
    def CategoryAxis_(self, request, category_axis_):
        return class_mock(
            request, 'pptx.chart.chart.CategoryAxis',
            return_value=category_axis_
        )

    @pytest.fixture
    def category_axis_(self, request):
        return instance_mock(request, CategoryAxis)

    @pytest.fixture
    def chart_data_(self, request):
        return instance_mock(request, ChartData)

    @pytest.fixture
    def ChartTitle_(self, request, chart_title_):
        return class_mock(
            request, 'pptx.chart.chart.ChartTitle', return_value=chart_title_
        )

    @pytest.fixture
    def chart_title_(self, request):
        return instance_mock(request, ChartTitle)

    @pytest.fixture
    def DateAxis_(self, request, date_axis_):
        return class_mock(
            request, 'pptx.chart.chart.DateAxis', return_value=date_axis_
        )

    @pytest.fixture
    def date_axis_(self, request):
        return instance_mock(request, DateAxis)

    @pytest.fixture
    def Font_(self, request):
        return class_mock(request, 'pptx.chart.chart.Font')

    @pytest.fixture
    def font_(self, request):
        return instance_mock(request, Font)

    @pytest.fixture
    def Legend_(self, request, legend_):
        return class_mock(
            request, 'pptx.chart.chart.Legend', return_value=legend_
        )

    @pytest.fixture
    def legend_(self, request):
        return instance_mock(request, Legend)

    @pytest.fixture
    def PlotTypeInspector_(self, request):
        return class_mock(request, 'pptx.chart.chart.PlotTypeInspector')

    @pytest.fixture
    def _Plots_(self, request, plots_):
        return class_mock(
            request, 'pptx.chart.chart._Plots', return_value=plots_
        )

    @pytest.fixture
    def plot_(self, request):
        return instance_mock(request, _BasePlot)

    @pytest.fixture
    def plots_(self, request):
        return instance_mock(request, _Plots)

    @pytest.fixture
    def SeriesCollection_(self, request, series_collection_):
        return class_mock(
            request, 'pptx.chart.chart.SeriesCollection',
            return_value=series_collection_
        )

    @pytest.fixture
    def SeriesXmlRewriterFactory_(self, request, series_rewriter_):
        return function_mock(
            request, 'pptx.chart.chart.SeriesXmlRewriterFactory',
            return_value=series_rewriter_, autospec=True
        )

    @pytest.fixture
    def series_collection_(self, request):
        return instance_mock(request, SeriesCollection)

    @pytest.fixture
    def series_rewriter_(self, request):
        return instance_mock(request, _BaseSeriesXmlRewriter)

    @pytest.fixture
    def ValueAxis_(self, request, value_axis_):
        return class_mock(
            request, 'pptx.chart.chart.ValueAxis',
            return_value=value_axis_
        )

    @pytest.fixture
    def value_axis_(self, request):
        return instance_mock(request, ValueAxis)

    @pytest.fixture
    def workbook_(self, request):
        return instance_mock(request, ChartWorkbook)

    @pytest.fixture
    def workbook_prop_(self, request, workbook_):
        return property_mock(
            request, Chart, '_workbook', return_value=workbook_
        )


class DescribeChartTitle(object):

    def it_provides_access_to_its_format(self, format_fixture):
        chart_title, ChartFormat_, format_ = format_fixture
        format = chart_title.format
        ChartFormat_.assert_called_once_with(chart_title.element)
        assert format is format_

    def it_knows_whether_it_has_a_text_frame(self, has_tf_get_fixture):
        chart_title, expected_value = has_tf_get_fixture
        value = chart_title.has_text_frame
        assert value is expected_value

    def it_can_change_whether_it_has_a_text_frame(self, has_tf_set_fixture):
        chart_title, value, expected_xml = has_tf_set_fixture
        chart_title.has_text_frame = value
        assert chart_title._element.xml == expected_xml

    def it_provides_access_to_its_text_frame(self, text_frame_fixture):
        chart_title, TextFrame_, text_frame_ = text_frame_fixture
        text_frame = chart_title.text_frame
        TextFrame_.assert_called_once_with(
            chart_title._element.tx.rich, chart_title
        )
        assert text_frame is text_frame_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def format_fixture(self, request, ChartFormat_, format_):
        chart_title = ChartTitle(element('c:title'))
        return chart_title, ChartFormat_, format_

    @pytest.fixture(params=[
        ('c:title',               False),
        ('c:title/c:tx',          False),
        ('c:title/c:tx/c:strRef', False),
        ('c:title/c:tx/c:rich',   True),
    ])
    def has_tf_get_fixture(self, request):
        title_cxml, expected_value = request.param
        chart_title = ChartTitle(element(title_cxml))
        return chart_title, expected_value

    @pytest.fixture(params=[
        ('c:title{a:b=c}', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title{a:b=c}/c:tx', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title{a:b=c}/c:tx/c:strRef', True,
         'c:title{a:b=c}/c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr'
         ')'),
        ('c:title/c:tx/c:rich',   True,  'c:title/c:tx/c:rich'),
        ('c:title',               False, 'c:title'),
        ('c:title/c:tx',          False, 'c:title'),
        ('c:title/c:tx/c:rich',   False, 'c:title'),
        ('c:title/c:tx/c:strRef', False, 'c:title'),
    ])
    def has_tf_set_fixture(self, request):
        title_cxml, value, expected_cxml = request.param
        chart_title = ChartTitle(element(title_cxml))
        expected_xml = xml(expected_cxml)
        return chart_title, value, expected_xml

    @pytest.fixture
    def text_frame_fixture(self, request, TextFrame_):
        chart_title = ChartTitle(element('c:title'))
        text_frame_ = TextFrame_.return_value
        return chart_title, TextFrame_, text_frame_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def ChartFormat_(self, request, format_):
        return class_mock(
            request, 'pptx.chart.chart.ChartFormat', return_value=format_
        )

    @pytest.fixture
    def format_(self, request):
        return instance_mock(request, ChartFormat)

    @pytest.fixture
    def TextFrame_(self, request):
        return class_mock(request, 'pptx.chart.chart.TextFrame')


class Describe_Plots(object):

    def it_supports_indexed_access(self, getitem_fixture):
        plots, idx, PlotFactory_, plot_elm, chart_, plot_ = getitem_fixture
        plot = plots[idx]
        PlotFactory_.assert_called_once_with(plot_elm, chart_)
        assert plot is plot_

    def it_supports_len(self, len_fixture):
        plots, expected_len = len_fixture
        assert len(plots) == expected_len

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('c:plotArea/c:barChart', 0),
        ('c:plotArea/(c:radarChart,c:barChart)', 1),
    ])
    def getitem_fixture(self, request, PlotFactory_, chart_, plot_):
        plotArea_cxml, idx = request.param
        plotArea = element(plotArea_cxml)
        plot_elm = plotArea[idx]
        plots = _Plots(plotArea, chart_)
        return plots, idx, PlotFactory_, plot_elm, chart_, plot_

    @pytest.fixture(params=[
        ('c:plotArea',                          0),
        ('c:plotArea/c:barChart',               1),
        ('c:plotArea/(c:barChart,c:lineChart)', 2),
    ])
    def len_fixture(self, request):
        plotArea_cxml, expected_len = request.param
        plots = _Plots(element(plotArea_cxml), None)
        return plots, expected_len

    # fixture components ---------------------------------------------

    @pytest.fixture
    def chart_(self, request):
        return instance_mock(request, Chart)

    @pytest.fixture
    def PlotFactory_(self, request, plot_):
        return function_mock(
            request, 'pptx.chart.chart.PlotFactory', return_value=plot_
        )

    @pytest.fixture
    def plot_(self, request):
        return instance_mock(request, _BasePlot)
