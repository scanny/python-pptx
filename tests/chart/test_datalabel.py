# encoding: utf-8

"""Unit test suite for the pptx.chart.datalabel module"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.chart.datalabel import DataLabel, DataLabels
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.text.text import Font

from ..unitutil.cxml import element, xml
from ..unitutil.mock import class_mock, instance_mock, method_mock


class DescribeDataLabel(object):

    def it_has_a_font(self, font_fixture):
        data_label, expected_xml = font_fixture
        font = data_label.font
        assert data_label._ser.xml == expected_xml
        assert isinstance(font, Font)

    def it_knows_its_position(self, position_get_fixture):
        data_label, expected_value = position_get_fixture
        position = data_label.position
        assert position == expected_value

    def it_can_change_its_position(self, position_set_fixture):
        data_label, value, expected_xml = position_set_fixture
        data_label.position = value
        assert data_label._element.xml == expected_xml

    def it_knows_whether_it_has_a_text_frame(self, has_tf_get_fixture):
        data_label, expected_value = has_tf_get_fixture
        value = data_label.has_text_frame
        assert value is expected_value

    def it_can_change_whether_it_has_a_text_frame(self, has_tf_set_fixture):
        data_label, value, expected_xml = has_tf_set_fixture
        data_label.has_text_frame = value
        assert data_label._element.xml == expected_xml

    def it_provides_access_to_its_text_frame(self, text_frame_fixture):
        data_label, TextFrame_, rich_, text_frame_ = text_frame_fixture
        text_frame = data_label.text_frame
        TextFrame_.assert_called_once_with(rich_, data_label)
        assert text_frame is text_frame_

    def it_gets_or_adds_rich_element_to_help(self, rich_fixture):
        data_label, expected_xml = rich_fixture
        rich = data_label._get_or_add_rich()
        assert data_label._ser.xml == expected_xml
        assert rich is data_label._ser.xpath('c:dLbls/c:dLbl/c:tx/c:rich')[0]

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('c:ser{a:b=c}',
         'c:ser{a:b=c}/c:dLbls/(c:dLbl/(c:idx{val=9},c:spPr,c:txPr/(a:bodyPr'
         ',a:lstStyle,a:p/a:pPr/a:defRPr),c:showLegendKey{val=0},c:showVal{v'
         'al=1},c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val='
         '0},c:showBubbleSize{val=0}),c:showLegendKey{val=0},c:showVal{val=0'
         '},c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0},c'
         ':showBubbleSize{val=0})'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=9},c:txPr/(a:bodyPr,a:p))',
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=9},c:txPr/(a:bodyPr,a:p/a:p'
         'Pr/a:defRPr))'),
    ])
    def font_fixture(self, request):
        ser_cxml, expected_cxml = request.param
        data_label = DataLabel(element(ser_cxml), 9)
        expected_xml = xml(expected_cxml)
        return data_label, expected_xml

    @pytest.fixture(params=[
        ('c:ser',                                              False),
        ('c:ser/c:dLbls',                                      False),
        ('c:ser/c:dLbls/c:dLbl/c:idx{val=42}',                 False),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:strRef)', False),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=24},c:tx/c:rich)',   False),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich)',   True),
    ])
    def has_tf_get_fixture(self, request):
        ser_cxml, expected_value = request.param
        data_label = DataLabel(element(ser_cxml), 42)
        return data_label, expected_value

    @pytest.fixture(params=[
        ('c:ser', False, 'c:ser'),
        ('c:ser/c:dLbls', False, 'c:ser/c:dLbls'),
        ('c:ser/c:dLbls/c:dLbl/c:idx{val=42}', False,
         'c:ser/c:dLbls/c:dLbl/c:idx{val=42}'),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:strRef)', False,
         'c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:strRef)'),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=24},c:tx/c:rich)', False,
         'c:ser/c:dLbls/c:dLbl/(c:idx{val=24},c:tx/c:rich)'),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich)', False,
         'c:ser/c:dLbls/c:dLbl/c:idx{val=42}'),
        ('c:ser{a:b=c}', True,
         'c:ser{a:b=c}/c:dLbls/(c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,'
         'a:lstStyle,a:p/a:pPr/a:defRPr),c:showLegendKey{val=0},c:showVal{va'
         'l=1},c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0'
         '},c:showBubbleSize{val=0}),c:showLegendKey{val=0},c:showVal{val=0}'
         ',c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0},c:'
         'showBubbleSize{val=0})'),
        ('c:ser{a:b=c}/c:dLbls', True,
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,a'
         ':lstStyle,a:p/a:pPr/a:defRPr),c:showLegendKey{val=0},c:showVal{val'
         '=1},c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0}'
         ',c:showBubbleSize{val=0})'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/c:idx{val=42}', True,
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,a'
         ':lstStyle,a:p/a:pPr/a:defRPr))'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:strRef)', True,
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,a'
         ':lstStyle,a:p/a:pPr/a:defRPr))'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=24},c:tx/c:rich)', True,
         'c:ser{a:b=c}/c:dLbls/(c:dLbl/(c:idx{val=24},c:tx/c:rich),c:dLbl/(c'
         ':idx{val=42},c:tx/c:rich/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr),'
         'c:showLegendKey{val=0},c:showVal{val=1},c:showCatName{val=0},c:sho'
         'wSerName{val=0},c:showPercent{val=0},c:showBubbleSize{val=0}))'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/c:idx{val=42}', True,
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,a'
         ':lstStyle,a:p/a:pPr/a:defRPr))'),
    ])
    def has_tf_set_fixture(self, request):
        ser_cxml, value, expected_cxml = request.param
        data_label = DataLabel(element(ser_cxml), 42)
        expected_xml = xml(expected_cxml)
        return data_label, value, expected_xml

    @pytest.fixture(params=[
        ('c:ser',                                                 None),
        ('c:ser/c:dLbls/c:dLbl/c:idx{val=42}',                    None),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:dLblPos{val=b})', 'BELOW'),
    ])
    def position_get_fixture(self, request):
        ser_cxml, value = request.param
        data_label = DataLabel(element(ser_cxml), 42)
        expected_value = (
            None if value is None else getattr(XL_LABEL_POSITION, value)
        )
        return data_label, expected_value

    @pytest.fixture(params=[
        ('c:ser{a:b=c}', 'CENTER',
         'c:ser{a:b=c}/c:dLbls/(c:dLbl/(c:idx{val=42},c:spPr,c:txPr/(a:bodyP'
         'r,a:lstStyle,a:p/a:pPr/a:defRPr),c:dLblPos{val=ctr},c:showLegendKe'
         'y{val=0},c:showVal{val=1},c:showCatName{val=0},c:showSerName{val=0'
         '},c:showPercent{val=0},c:showBubbleSize{val=0}),c:showLegendKey{va'
         'l=0},c:showVal{val=0},c:showCatName{val=0},c:showSerName{val=0},c:'
         'showPercent{val=0},c:showBubbleSize{val=0})'),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:dLblPos{val=ctr})', 'BELOW',
         'c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:dLblPos{val=b})'),
        ('c:ser/c:dLbls/c:dLbl/(c:idx{val=42},c:dLblPos{val=b})', None,
         'c:ser/c:dLbls/c:dLbl/c:idx{val=42}'),
        ('c:ser', None, 'c:ser'),
    ])
    def position_set_fixture(self, request):
        ser_cxml, value, expected_cxml = request.param
        data_label = DataLabel(element(ser_cxml), 42)
        new_value = (
            None if value is None else getattr(XL_LABEL_POSITION, value)
        )
        expected_xml = xml(expected_cxml)
        return data_label, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:ser{a:b=c}',
         'c:ser{a:b=c}/c:dLbls/(c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,'
         'a:lstStyle,a:p/a:pPr/a:defRPr),c:showLegendKey{val=0},c:showVal{va'
         'l=1},c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0'
         '},c:showBubbleSize{val=0}),c:showLegendKey{val=0},c:showVal{val=0}'
         ',c:showCatName{val=0},c:showSerName{val=0},c:showPercent{val=0},c:'
         'showBubbleSize{val=0})'),
        ('c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:strRef)',
         'c:ser{a:b=c}/c:dLbls/c:dLbl/(c:idx{val=42},c:tx/c:rich/(a:bodyPr,a'
         ':lstStyle,a:p/a:pPr/a:defRPr))'),
    ])
    def rich_fixture(self, request):
        ser_cxml, expected_cxml = request.param
        data_label = DataLabel(element(ser_cxml), 42)
        expected_xml = xml(expected_cxml)
        return data_label, expected_xml

    @pytest.fixture
    def text_frame_fixture(self, request, _get_or_add_rich_, TextFrame_):
        data_label = DataLabel(None, None)
        rich_ = _get_or_add_rich_.return_value
        text_frame_ = TextFrame_.return_value
        return data_label, TextFrame_, rich_, text_frame_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _get_or_add_rich_(self, request):
        return method_mock(request, DataLabel, '_get_or_add_rich')

    @pytest.fixture
    def TextFrame_(self, request):
        return class_mock(request, 'pptx.chart.datalabel.TextFrame')


class DescribeDataLabels(object):

    def it_provides_access_to_its_font(self, font_fixture):
        data_labels, Font_, defRPr, font_ = font_fixture
        font = data_labels.font
        Font_.assert_called_once_with(defRPr)
        assert font is font_

    def it_adds_a_txPr_to_help_font(self, txPr_fixture):
        data_labels, expected_xml = txPr_fixture
        data_labels.font
        assert data_labels._element.xml == expected_xml

    def it_knows_its_number_format(self, number_format_get_fixture):
        data_labels, expected_value = number_format_get_fixture
        assert data_labels.number_format == expected_value

    def it_can_change_its_number_format(self, number_format_set_fixture):
        data_labels, new_value, expected_xml = number_format_set_fixture
        data_labels.number_format = new_value
        assert data_labels._element.xml == expected_xml

    def it_knows_whether_its_number_format_is_linked(
            self, number_format_is_linked_get_fixture):
        data_labels, expected_value = number_format_is_linked_get_fixture
        assert data_labels.number_format_is_linked is expected_value

    def it_can_change_whether_its_number_format_is_linked(
            self, number_format_is_linked_set_fixture):
        data_labels, new_value, expected_xml = (
            number_format_is_linked_set_fixture
        )
        data_labels.number_format_is_linked = new_value
        assert data_labels._element.xml == expected_xml

    def it_knows_its_position(self, position_get_fixture):
        data_labels, expected_value = position_get_fixture
        assert data_labels.position == expected_value

    def it_can_change_its_position(self, position_set_fixture):
        data_labels, new_value, expected_xml = position_set_fixture
        data_labels.position = new_value
        assert data_labels._element.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def font_fixture(self, Font_, font_):
        dLbls = element('c:dLbls/c:txPr/a:p/a:pPr/a:defRPr')
        defRPr = dLbls.xpath('.//a:defRPr')[0]
        data_labels = DataLabels(dLbls)
        return data_labels, Font_, defRPr, font_

    @pytest.fixture(params=[
        ('c:dLbls',                             'General'),
        ('c:dLbls/c:numFmt{formatCode=foobar}', 'foobar'),
    ])
    def number_format_get_fixture(self, request):
        dLbls_cxml, expected_value = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        return data_labels, expected_value

    @pytest.fixture(params=[
        ('c:dLbls', 'General',
         'c:dLbls/c:numFmt{formatCode=General,sourceLinked=0}'),
        ('c:dLbls/c:numFmt{formatCode=General}', '00.00',
         'c:dLbls/c:numFmt{formatCode=00.00,sourceLinked=0}'),
    ])
    def number_format_set_fixture(self, request):
        dLbls_cxml, new_value, expected_dLbls_cxml = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        expected_xml = xml(expected_dLbls_cxml)
        return data_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:dLbls',                          True),
        ('c:dLbls/c:numFmt',                 True),
        ('c:dLbls/c:numFmt{sourceLinked=0}', False),
        ('c:dLbls/c:numFmt{sourceLinked=1}', True),
    ])
    def number_format_is_linked_get_fixture(self, request):
        dLbls_cxml, expected_value = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        return data_labels, expected_value

    @pytest.fixture(params=[
        ('c:dLbls', True,  'c:dLbls/c:numFmt{sourceLinked=1}'),
        ('c:dLbls', False, 'c:dLbls/c:numFmt{sourceLinked=0}'),
        ('c:dLbls', None,  'c:dLbls/c:numFmt'),
        ('c:dLbls/c:numFmt', True, 'c:dLbls/c:numFmt{sourceLinked=1}'),
        ('c:dLbls/c:numFmt{sourceLinked=1}', False,
         'c:dLbls/c:numFmt{sourceLinked=0}'),
    ])
    def number_format_is_linked_set_fixture(self, request):
        dLbls_cxml, new_value, expected_dLbls_cxml = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        expected_xml = xml(expected_dLbls_cxml)
        return data_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:dLbls',                       None),
        ('c:dLbls/c:dLblPos{val=inBase}', XL_LABEL_POSITION.INSIDE_BASE),
    ])
    def position_get_fixture(self, request):
        dLbls_cxml, expected_value = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        return data_labels, expected_value

    @pytest.fixture(params=[
        ('c:dLbls',                       XL_LABEL_POSITION.INSIDE_BASE,
         'c:dLbls/c:dLblPos{val=inBase}'),
        ('c:dLbls/c:dLblPos{val=inBase}', XL_LABEL_POSITION.OUTSIDE_END,
         'c:dLbls/c:dLblPos{val=outEnd}'),
        ('c:dLbls/c:dLblPos{val=inBase}', None, 'c:dLbls'),
        ('c:dLbls',                       None, 'c:dLbls'),
    ])
    def position_set_fixture(self, request):
        dLbls_cxml, new_value, expected_dLbls_cxml = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        expected_xml = xml(expected_dLbls_cxml)
        return data_labels, new_value, expected_xml

    @pytest.fixture(params=[
        ('c:dLbls{a:b=c}',
         'c:dLbls{a:b=c}/c:txPr/(a:bodyPr,a:lstStyle,a:p/a:pPr/a:defRPr)'),
        ('c:dLbls{a:b=c}/c:txPr/(a:bodyPr,a:p)',
         'c:dLbls{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr/a:defRPr)'),
        ('c:dLbls{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr)',
         'c:dLbls{a:b=c}/c:txPr/(a:bodyPr,a:p/a:pPr/a:defRPr)'),
    ])
    def txPr_fixture(self, request):
        dLbls_cxml, expected_cxml = request.param
        data_labels = DataLabels(element(dLbls_cxml))
        expected_xml = xml(expected_cxml)
        return data_labels, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def Font_(self, request, font_):
        return class_mock(
            request, 'pptx.chart.datalabel.Font', return_value=font_
        )

    @pytest.fixture
    def font_(self, request):
        return instance_mock(request, Font)
