# encoding: utf-8

"""
Test suite for pptx.chart.title module
"""

from __future__ import absolute_import, print_function

import pytest

from pptx.chart.title import ChartTitle
from pptx.text.text import TextFrame

from ..unitutil.cxml import element, xml


class DescribeChartTitle(object):

    def it_provides_access_to_its_text_frame(self, text_frame_get_fixture):
        title, expected_text_frame_xml = text_frame_get_fixture
        text_frame = title.text_frame
        assert isinstance(text_frame, TextFrame)
        assert text_frame._element.xml == xml(expected_text_frame_xml)
    
    # fixtures ---------------------------------------------

    @pytest.fixture(params=[
        ('c:title/c:tx/c:rich/(a:bodyPr,a:p)', 'c:rich/(a:bodyPr,a:p)'),
        ('c:title', 'c:rich'),
    ])
    def text_frame_get_fixture(self, request):
        title_cxml, expected_value = request.param
        title = ChartTitle(element(title_cxml))
        return title, expected_value
