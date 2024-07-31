"""Unit-test suite for `pptx.enum.shapes`."""

from __future__ import annotations

import pytest

from pptx.enum.shapes import PROG_ID


class DescribeProgId:
    """Unit-test suite for `pptx.enum.shapes.ProgId`."""

    def it_has_members_for_the_OLE_embeddings_known_to_work_on_Windows(self):
        assert PROG_ID.DOCX
        assert PROG_ID.PPTX
        assert PROG_ID.XLSX

    @pytest.mark.parametrize(
        ("member", "expected_value"),
        [(PROG_ID.DOCX, 609600), (PROG_ID.PPTX, 609600), (PROG_ID.XLSX, 609600)],
    )
    def it_knows_its_height(self, member: PROG_ID, expected_value: int):
        assert member.height == expected_value

    def it_knows_its_icon_filename(self):
        assert PROG_ID.DOCX.icon_filename == "docx-icon.emf"

    def it_knows_its_progId(self):
        assert PROG_ID.PPTX.progId == "PowerPoint.Show.12"

    def it_knows_its_width(self):
        assert PROG_ID.XLSX.width == 965200

    @pytest.mark.parametrize(
        ("value", "expected_value"),
        [
            # -DELETEME---------------------------------------------------------------
            (PROG_ID.DOCX, True),
            (PROG_ID.PPTX, True),
            (PROG_ID.XLSX, True),
            (17, False),
            ("XLSX", False),
        ],
    )
    def it_knows_each_of_its_members_is_an_instance(self, value: object, expected_value: bool):
        assert isinstance(value, PROG_ID) is expected_value
