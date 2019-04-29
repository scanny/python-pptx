# encoding: utf-8

"""
Test suite for pptx.oxml.autoshape module.
"""

from __future__ import absolute_import, print_function

import pytest

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import nsdecls
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.shapes.shared import ST_Direction, ST_PlaceholderSize

from ..unitdata.shape import a_gd, a_prstGeom, an_avLst
from ...unitutil.cxml import element


class DescribeCT_PresetGeometry2D(object):
    def it_can_get_the_gd_elms_as_a_sequence(self, gd_lst_fixture):
        prstGeom, expected_vals = gd_lst_fixture
        actual_vals = [(gd.name, gd.fmla) for gd in prstGeom.gd_lst]
        assert actual_vals == expected_vals

    def it_can_rewrite_the_gd_elms(self, rewrite_guides_fixture_):
        prstGeom, guides, expected_xml = rewrite_guides_fixture_
        print(prstGeom.xml)
        prstGeom.rewrite_guides(guides)
        assert prstGeom.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            [],
            [("adj1", "val 111")],
            [
                ("adj2", "val 222"),
                ("adj3", "val 333"),
                ("adj4", "val 444"),
                ("adj5", "val 555"),
            ],
        ]
    )
    def gd_lst_fixture(self, request):
        expected_vals = request.param
        prstGeom = self.prstGeom_bldr("foobar", expected_vals).element
        return prstGeom, expected_vals

    @pytest.fixture(
        params=[
            ("circularArrow", 5),  # all five guides for five adj shape
            ("chevron", 0),  # empty guides for single adj shape
            ("chevron", 1),  # one guide for single adj shape
        ]
    )
    def rewrite_guides_fixture_(self, request):
        prst, gd_count = request.param

        names = ("adj1", "adj2", "adj3", "adj4", "adj5")
        vals = (111, 222, 333, 444, 555)
        fmlas = [("val %d" % v) for v in vals]
        before_vals = [("adj6", "val 666")] * 3
        after_vals = zip(names[:gd_count], fmlas[:gd_count])

        prstGeom = self.prstGeom_bldr(prst, before_vals).element
        guides = zip(names[:gd_count], vals[:gd_count])
        expected_xml = self.prstGeom_bldr(prst, after_vals).xml()

        return prstGeom, guides, expected_xml

    # fixture components ---------------------------------------------

    def prstGeom_bldr(self, prst, gd_vals):
        avLst_bldr = an_avLst()
        for name, fmla in gd_vals:
            gd_bldr = a_gd().with_name(name).with_fmla(fmla)
            avLst_bldr.with_child(gd_bldr)
        prstGeom_bldr = (
            a_prstGeom().with_nsdecls().with_prst(prst).with_child(avLst_bldr)
        )
        return prstGeom_bldr


class DescribeCT_Shape(object):
    def it_knows_how_to_create_a_new_autoshape_sp(self):
        # setup ------------------------
        id_ = 9
        name = "Rounded Rectangle 8"
        prst = "roundRect"
        left, top, width, height = 111, 222, 333, 444
        xml = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%d" name="%s"/>\n    <'
            "p:cNvSpPr/>\n    <p:nvPr/>\n  </p:nvSpPr>\n  <p:spPr>\n    <a:xf"
            'rm>\n      <a:off x="%d" y="%d"/>\n      <a:ext cx="%d" cy="%d"/'
            '>\n    </a:xfrm>\n    <a:prstGeom prst="%s">\n      <a:avLst/>\n'
            '    </a:prstGeom>\n  </p:spPr>\n  <p:style>\n    <a:lnRef idx="1'
            '">\n      <a:schemeClr val="accent1"/>\n    </a:lnRef>\n    <a:f'
            'illRef idx="3">\n      <a:schemeClr val="accent1"/>\n    </a:fil'
            'lRef>\n    <a:effectRef idx="2">\n      <a:schemeClr val="accent'
            '1"/>\n    </a:effectRef>\n    <a:fontRef idx="minor">\n      <a:'
            'schemeClr val="lt1"/>\n    </a:fontRef>\n  </p:style>\n  <p:txBo'
            'dy>\n    <a:bodyPr rtlCol="0" anchor="ctr"/>\n    <a:lstStyle/>'
            '\n    <a:p>\n      <a:pPr algn="ctr"/>\n    </a:p>\n  </p:txBody'
            ">\n</p:sp>\n"
            % (nsdecls("a", "p"), id_, name, left, top, width, height, prst)
        )
        # exercise ---------------------
        sp = CT_Shape.new_autoshape_sp(id_, name, prst, left, top, width, height)
        # verify -----------------------
        assert sp.xml == xml

    def it_knows_how_to_create_a_new_placeholder_sp(self, new_ph_sp_fixture):
        id_, name, ph_type, orient, sz, idx, expected_xml = new_ph_sp_fixture
        sp = CT_Shape.new_placeholder_sp(id_, name, ph_type, orient, sz, idx)
        assert sp.xml == expected_xml

    def it_knows_how_to_create_a_new_textbox_sp(self):
        # setup ------------------------
        id_ = 9
        name = "TextBox 8"
        left, top, width, height = 111, 222, 333, 444
        xml = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%d" name="%s"/>\n    '
            '<p:cNvSpPr txBox="1"/>\n    <p:nvPr/>\n  </p:nvSpPr>\n  <p:spPr'
            '>\n    <a:xfrm>\n      <a:off x="%d" y="%d"/>\n      <a:ext cx='
            '"%d" cy="%d"/>\n    </a:xfrm>\n    <a:prstGeom prst="rect">\n  '
            "    <a:avLst/>\n    </a:prstGeom>\n    <a:noFill/>\n  </p:spPr>"
            '\n  <p:txBody>\n    <a:bodyPr wrap="none">\n      <a:spAutoFit/'
            ">\n    </a:bodyPr>\n    <a:lstStyle/>\n    <a:p/>\n  </p:txBody"
            ">\n</p:sp>\n" % (nsdecls("a", "p"), id_, name, left, top, width, height)
        )
        # exercise ---------------------
        sp = CT_Shape.new_textbox_sp(id_, name, left, top, width, height)
        # verify -----------------------
        assert sp.xml == xml

    def it_knows_whether_it_is_an_autoshape(self, is_autoshape_fixture):
        sp, expected_value = is_autoshape_fixture
        assert sp.is_autoshape is expected_value

    def it_knows_whether_it_is_an_textbox(self, is_textbox_fixture):
        sp, expected_value = is_textbox_fixture
        assert sp.is_textbox is expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            # autoshape
            ("p:sp/(p:nvSpPr/p:cNvSpPr, p:spPr/a:prstGeom)", True),
            # placeholder
            ("p:sp/(p:nvSpPr/p:nvPr/p:ph, p:spPr)", False),
            # textbox
            ("p:sp/(p:nvSpPr/p:cNvSpPr{txBox=1}, p:spPr/a:prstGeom)", False),
        ]
    )
    def is_autoshape_fixture(self, request):
        sp_cxml, expected_value = request.param
        sp = element(sp_cxml)
        return sp, expected_value

    @pytest.fixture(
        params=[
            # autoshape
            ("p:sp/(p:nvSpPr/p:cNvSpPr, p:spPr/a:prstGeom)", False),
            # placeholder
            ("p:sp/(p:nvSpPr/(p:nvPr/p:ph, p:cNvSpPr), p:spPr)", False),
            # textbox
            ("p:sp/(p:nvSpPr/p:cNvSpPr{txBox=1}, p:spPr/a:prstGeom)", True),
        ]
    )
    def is_textbox_fixture(self, request):
        sp_cxml, expected_value = request.param
        sp = element(sp_cxml)
        return sp, expected_value

    @pytest.fixture(
        params=[
            (
                2,
                "Title 1",
                PP_PLACEHOLDER.CENTER_TITLE,
                ST_Direction.HORZ,
                ST_PlaceholderSize.FULL,
                0,
                ' type="ctrTitle"',
            ),
            (
                3,
                "Date Placeholder 2",
                PP_PLACEHOLDER.DATE,
                ST_Direction.HORZ,
                ST_PlaceholderSize.HALF,
                10,
                ' type="dt" sz="half" idx="10"',
            ),
            (
                4,
                "Vertical Subtitle 3",
                PP_PLACEHOLDER.SUBTITLE,
                ST_Direction.VERT,
                ST_PlaceholderSize.FULL,
                1,
                ' type="subTitle" orient="vert" idx="1"',
            ),
            (
                5,
                "Table Placeholder 4",
                PP_PLACEHOLDER.TABLE,
                ST_Direction.HORZ,
                ST_PlaceholderSize.QUARTER,
                14,
                ' type="tbl" sz="quarter" idx="14"',
            ),
            (
                6,
                "Slide Number Placeholder 5",
                PP_PLACEHOLDER.SLIDE_NUMBER,
                ST_Direction.HORZ,
                ST_PlaceholderSize.QUARTER,
                12,
                ' type="sldNum" sz="quarter" idx="12"',
            ),
            (
                7,
                "Footer Placeholder 6",
                PP_PLACEHOLDER.FOOTER,
                ST_Direction.HORZ,
                ST_PlaceholderSize.QUARTER,
                11,
                ' type="ftr" sz="quarter" idx="11"',
            ),
            (
                8,
                "Content Placeholder 7",
                PP_PLACEHOLDER.OBJECT,
                ST_Direction.HORZ,
                ST_PlaceholderSize.FULL,
                15,
                ' idx="15"',
            ),
        ]
    )
    def new_ph_sp_fixture(self, request):
        id_, name, ph_type, orient, sz, idx, expected_attrs = request.param
        expected_xml_tmpl = (
            "<p:sp %s>\n"
            "  <p:nvSpPr>\n"
            '    <p:cNvPr id="%s" name="%s"/>\n'
            "    <p:cNvSpPr>\n"
            '      <a:spLocks noGrp="1"/>\n'
            "    </p:cNvSpPr>\n"
            "    <p:nvPr>\n"
            "      <p:ph%s/>\n"
            "    </p:nvPr>\n"
            "  </p:nvSpPr>\n"
            "  <p:spPr/>\n"
            "%s"
            "</p:sp>\n" % (nsdecls("a", "p"), "%d", "%s", "%s", "%s")
        )
        txBody_snippet = (
            "  <p:txBody>\n"
            "    <a:bodyPr/>\n"
            "    <a:lstStyle/>\n"
            "    <a:p/>\n"
            "  </p:txBody>\n"
        )
        txBody_str = txBody_snippet if id_ in (2, 4, 8) else ""
        expected_values = (id_, name, expected_attrs, txBody_str)
        expected_xml = expected_xml_tmpl % expected_values
        return id_, name, ph_type, orient, sz, idx, expected_xml
