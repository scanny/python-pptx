# encoding: utf-8

"""Test suite for pptx.oxml.shapes.shapetree module."""

from __future__ import absolute_import, division, print_function, unicode_literals

import pytest

from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.oxml.shapes.picture import CT_Picture

from ...unitutil.cxml import element, xml
from ...unitutil.mock import call, class_mock, instance_mock, method_mock, property_mock


class DescribeCT_GroupShape(object):
    def it_can_add_a_graphicFrame_element_containing_a_table(self, add_table_fixt):
        spTree, id_, name, rows, cols, x, y, cx, cy = add_table_fixt[:9]
        new_table_graphicFrame_ = add_table_fixt[9]
        insert_element_before_, graphicFrame_ = add_table_fixt[10:]

        graphicFrame = spTree.add_table(id_, name, rows, cols, x, y, cx, cy)

        new_table_graphicFrame_.assert_called_once_with(
            id_, name, rows, cols, x, y, cx, cy
        )
        insert_element_before_.assert_called_once_with(graphicFrame_, "p:extLst")
        assert graphicFrame is graphicFrame_

    def it_can_add_a_grpSp_element(self, add_grpSp_fixture):
        spTree, expected_grpSp_xml, expected_xml = add_grpSp_fixture

        grpSp = spTree.add_grpSp()

        assert grpSp.xml == expected_grpSp_xml
        assert spTree.xml == expected_xml

    def it_can_add_a_pic_element_representing_a_picture(self, add_pic_fixt):
        spTree, id_, name, desc, rId, x, y, cx, cy = add_pic_fixt[:9]
        CT_Picture_, insert_element_before_, pic_ = add_pic_fixt[9:]
        pic = spTree.add_pic(id_, name, desc, rId, x, y, cx, cy)
        CT_Picture_.new_pic.assert_called_once_with(id_, name, desc, rId, x, y, cx, cy)
        insert_element_before_.assert_called_once_with(pic_, "p:extLst")
        assert pic is pic_

    def it_can_add_an_sp_element_for_a_placeholder(self, add_placeholder_fixt):
        spTree, id_, name, ph_type, orient, sz, idx = add_placeholder_fixt[:7]
        CT_Shape_, insert_element_before_, sp_ = add_placeholder_fixt[7:]
        sp = spTree.add_placeholder(id_, name, ph_type, orient, sz, idx)
        CT_Shape_.new_placeholder_sp.assert_called_once_with(
            id_, name, ph_type, orient, sz, idx
        )
        insert_element_before_.assert_called_once_with(sp_, "p:extLst")
        assert sp is sp_

    def it_can_add_an_sp_element_for_an_autoshape(self, add_autoshape_fixt):
        spTree, id_, name, prst, x, y, cx, cy = add_autoshape_fixt[:8]
        CT_Shape_, insert_element_before_, sp_ = add_autoshape_fixt[8:]
        sp = spTree.add_autoshape(id_, name, prst, x, y, cx, cy)
        CT_Shape_.new_autoshape_sp.assert_called_once_with(
            id_, name, prst, x, y, cx, cy
        )
        insert_element_before_.assert_called_once_with(sp_, "p:extLst")
        assert sp is sp_

    def it_can_add_a_textbox_sp_element(self, add_textbox_fixt):
        spTree, id_, name, x, y, cx, cy, CT_Shape_ = add_textbox_fixt[:8]
        insert_element_before_, sp_ = add_textbox_fixt[8:]
        sp = spTree.add_textbox(id_, name, x, y, cx, cy)
        CT_Shape_.new_textbox_sp.assert_called_once_with(id_, name, x, y, cx, cy)
        insert_element_before_.assert_called_once_with(sp_, "p:extLst")
        assert sp is sp_

    def it_can_recalculate_its_pos_and_size(self, recalc_fixture):
        xSp, expected_xml, parent_sp, calls = recalc_fixture

        xSp.recalculate_extents()

        assert xSp.xml == expected_xml
        assert parent_sp.recalculate_extents.call_args_list == calls

    def it_calculates_its_child_extents_to_help(self, child_exts_fixture):
        xSp, expected_values = child_exts_fixture
        x, y, cx, cy = xSp._child_extents
        assert (x, y, cx, cy) == expected_values

    # fixtures ---------------------------------------------

    @pytest.fixture
    def add_autoshape_fixt(self, spTree, CT_Shape_, insert_element_before_, sp_):
        id_, name, prst = 42, "name", "prst"
        x, y, cx, cy = 9, 8, 7, 6
        return (
            spTree,
            id_,
            name,
            prst,
            x,
            y,
            cx,
            cy,
            CT_Shape_,
            insert_element_before_,
            sp_,
        )

    @pytest.fixture
    def add_grpSp_fixture(self):
        spTree = element("p:spTree{a:b=c,r:s=t}")
        expected_grpSp_xml = (
            # ---can't get the namespaces right with cxml, using full text--
            '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentatio'
            'nml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawi'
            'ngml/2006/main" xmlns:r="http://schemas.openxmlformats.org/offi'
            'ceDocument/2006/relationships">\n  <p:nvGrpSpPr>\n    <p:cNvPr '
            'id="1" name="Group 0"/>\n    <p:cNvGrpSpPr/>\n    <p:nvPr/>\n  '
            '</p:nvGrpSpPr>\n  <p:grpSpPr>\n    <a:xfrm>\n      <a:off x="0"'
            ' y="0"/>\n      <a:ext cx="0" cy="0"/>\n      <a:chOff x="0" y='
            '"0"/>\n      <a:chExt cx="0" cy="0"/>\n    </a:xfrm>\n  </p:grp'
            "SpPr>\n</p:grpSp>\n"
        )
        expected_xml = xml(
            "p:spTree{a:b=c,r:s=t}/p:grpSp/(p:nvGrpSpPr/(p:cNvPr{id=1,name=G"
            "roup 0},p:cNvGrpSpPr,p:nvPr),p:grpSpPr/a:xfrm/(a:off{x=0,y=0},a"
            ":ext{cx=0,cy=0},a:chOff{x=0,y=0},a:chExt{cx=0,cy=0}))"
        )
        return spTree, expected_grpSp_xml, expected_xml

    @pytest.fixture
    def add_pic_fixt(self, spTree, CT_Picture_, insert_element_before_, pic_):
        id_, name, desc, rId = 42, "name", "desc", "rId6"
        x, y, cx, cy = 6, 7, 8, 9
        return (
            spTree,
            id_,
            name,
            desc,
            rId,
            x,
            y,
            cx,
            cy,
            CT_Picture_,
            insert_element_before_,
            pic_,
        )

    @pytest.fixture
    def add_placeholder_fixt(self, spTree, CT_Shape_, insert_element_before_, sp_):
        id_, name, ph_type = 42, "name", "type"
        orient, sz, idx = "orient", "sz", 24
        return (
            spTree,
            id_,
            name,
            ph_type,
            orient,
            sz,
            idx,
            CT_Shape_,
            insert_element_before_,
            sp_,
        )

    @pytest.fixture
    def add_table_fixt(
        self, spTree, CT_GraphicalObjectFrame_, insert_element_before_, graphicFrame_
    ):
        id_, name, rows, cols = 42, "name", 12, 23
        x, y, cx, cy = 5, 4, 3, 2
        new_table_graphicFrame_ = CT_GraphicalObjectFrame_.new_table_graphicFrame
        return (
            spTree,
            id_,
            name,
            rows,
            cols,
            x,
            y,
            cx,
            cy,
            new_table_graphicFrame_,
            insert_element_before_,
            graphicFrame_,
        )

    @pytest.fixture
    def add_textbox_fixt(self, spTree, CT_Shape_, insert_element_before_, sp_):
        id_, name = 42, "name"
        x, y, cx, cy = 3, 4, 5, 6
        return (spTree, id_, name, x, y, cx, cy, CT_Shape_, insert_element_before_, sp_)

    @pytest.fixture(
        params=[
            ("p:grpSp", (0, 0, 0, 0)),
            (
                "p:grpSp/p:sp/p:spPr/a:xfrm/(a:off{x=1,y=2},a:ext{cx=3,cy=4})",
                (1, 2, 3, 4),
            ),
            (
                "p:grpSp/(p:sp/p:spPr/a:xfrm/(a:off{x=1,y=2},a:ext{cx=3,cy=4}),p:sp"
                "/p:spPr/a:xfrm/(a:off{x=10,y=20},a:ext{cx=30,cy=40}))",
                (1, 2, 39, 58),
            ),
            (
                "p:grpSp/(p:sp/p:spPr/a:xfrm/(a:off{x=100,y=50},a:ext{cx=25,cy=25})"
                ",p:sp/p:spPr/a:xfrm/(a:off{x=50,y=100},a:ext{cx=25,cy=25}),p:sp/p:"
                "spPr/a:xfrm/(a:off{x=150,y=75},a:ext{cx=25,cy=25}))",
                (50, 50, 125, 75),
            ),
        ]
    )
    def child_exts_fixture(self, request):
        xSp_cxml, expected_values = request.param
        xSp = element(xSp_cxml)
        return xSp, expected_values

    @pytest.fixture(
        params=[
            ("p:spTree", None, [], "p:spTree"),
            (
                "p:grpSp/p:grpSpPr/a:xfrm",
                (1, 2, 3, 4),
                [call()],
                "p:grpSp/p:grpSpPr/a:xfrm/(a:off{x=1,y=2},a:ext{cx=3,cy=4},a:chOff{"
                "x=1,y=2},a:chExt{cx=3,cy=4})",
            ),
        ]
    )
    def recalc_fixture(self, request, _child_extents_prop_, getparent_, grpSp_):
        xSp_cxml, extents, calls, expected_cxml = request.param
        xSp = element(xSp_cxml)

        _child_extents_prop_.return_value = extents
        expected_xml = xml(expected_cxml)
        getparent_.return_value = parent_sp = grpSp_
        return xSp, expected_xml, parent_sp, calls

    # fixture components -----------------------------------

    @pytest.fixture
    def _child_extents_prop_(self, request):
        return property_mock(request, CT_GroupShape, "_child_extents")

    @pytest.fixture
    def CT_GraphicalObjectFrame_(self, request, graphicFrame_):
        CT_GraphicalObjectFrame_ = class_mock(
            request, "pptx.oxml.shapes.groupshape.CT_GraphicalObjectFrame"
        )
        CT_GraphicalObjectFrame_.new_table_graphicFrame.return_value = graphicFrame_
        return CT_GraphicalObjectFrame_

    @pytest.fixture
    def CT_Picture_(self, request, pic_):
        CT_Picture_ = class_mock(request, "pptx.oxml.shapes.groupshape.CT_Picture")
        CT_Picture_.new_pic.return_value = pic_
        return CT_Picture_

    @pytest.fixture
    def CT_Shape_(self, request, sp_):
        CT_Shape_ = class_mock(request, "pptx.oxml.shapes.groupshape.CT_Shape")
        CT_Shape_.new_autoshape_sp.return_value = sp_
        CT_Shape_.new_placeholder_sp.return_value = sp_
        CT_Shape_.new_textbox_sp.return_value = sp_
        return CT_Shape_

    @pytest.fixture
    def getparent_(self, request):
        return method_mock(request, CT_GroupShape, "getparent")

    @pytest.fixture
    def graphicFrame_(self, request):
        return instance_mock(request, CT_GraphicalObjectFrame)

    @pytest.fixture
    def grpSp_(self, request):
        return instance_mock(request, CT_GroupShape)

    @pytest.fixture
    def insert_element_before_(self, request):
        return method_mock(request, CT_GroupShape, "insert_element_before")

    @pytest.fixture
    def pic_(self, request):
        return instance_mock(request, CT_Picture)

    @pytest.fixture
    def sp_(self, request):
        return instance_mock(request, CT_Shape)

    @pytest.fixture
    def spTree(self):
        return element("p:spTree")
