# encoding: utf-8

"""Unit-test suite for `pptx.oxml.shapes.picture` module."""

import pytest

from pptx.oxml.ns import nsdecls
from pptx.oxml.shapes.picture import CT_Picture


class DescribeCT_Picture(object):
    """Unit-test suite for `pptx.oxml.shapes.picture.CT_Picture` objects."""

    @pytest.mark.parametrize(
        "desc, xml_desc",
        (
            ("kittens.jpg", "kittens.jpg"),
            ("bits&bobs.png", "bits&amp;bobs.png"),
            ("img&.png", "img&amp;.png"),
            ("im<ag>e.png", "im&lt;ag&gt;e.png"),
        ),
    )
    def it_can_create_a_new_pic_element(self, desc, xml_desc):
        """`desc` attr (often filename) is XML-escaped to handle special characters.

        In particular, ampersand ('&'), less/greater-than ('</>') etc.
        """
        pic = CT_Picture.new_pic(
            shape_id=9, name="Picture 8", desc=desc, rId="rId42", x=1, y=2, cx=3, cy=4
        )

        assert pic.xml == (
            "<p:pic %s>\n"
            "  <p:nvPicPr>\n"
            '    <p:cNvPr id="9" name="Picture 8" descr="%s"/>\n'
            "    <p:cNvPicPr>\n"
            '      <a:picLocks noChangeAspect="1"/>\n'
            "    </p:cNvPicPr>\n"
            "    <p:nvPr/>\n"
            "  </p:nvPicPr>\n"
            "  <p:blipFill>\n"
            '    <a:blip r:embed="rId42"/>\n'
            "    <a:stretch>\n"
            "      <a:fillRect/>\n"
            "    </a:stretch>\n"
            "  </p:blipFill>\n"
            "  <p:spPr>\n"
            "    <a:xfrm>\n"
            '      <a:off x="1" y="2"/>\n'
            '      <a:ext cx="3" cy="4"/>\n'
            "    </a:xfrm>\n"
            '    <a:prstGeom prst="rect">\n'
            "      <a:avLst/>\n"
            "    </a:prstGeom>\n"
            "  </p:spPr>\n"
            "</p:pic>\n" % (nsdecls("a", "p", "r"), xml_desc)
        )

    def it_can_create_a_new_video_pic_element(self):
        pic = CT_Picture.new_video_pic(
            shape_id=42,
            shape_name="Media 41",
            video_rId="rId1",
            media_rId="rId2",
            poster_frame_rId="rId3",
            x=1,
            y=2,
            cx=3,
            cy=4,
        )

        assert pic.xml == (
            "<p:pic %s>\n"
            "  <p:nvPicPr>\n"
            '    <p:cNvPr id="42" name="Media 41">\n'
            '      <a:hlinkClick r:id="" action="ppaction://media"/>\n'
            "    </p:cNvPr>\n"
            "    <p:cNvPicPr>\n"
            '      <a:picLocks noChangeAspect="1"/>\n'
            "    </p:cNvPicPr>\n"
            "    <p:nvPr>\n"
            '      <a:videoFile r:link="rId1"/>\n'
            "      <p:extLst>\n"
            '        <p:ext uri="{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}">\n'
            '          <p14:media xmlns:p14="http://schemas.microsoft.com/office/power'
            'point/2010/main" r:embed="rId2"/>\n'
            "        </p:ext>\n"
            "      </p:extLst>\n"
            "    </p:nvPr>\n"
            "  </p:nvPicPr>\n"
            "  <p:blipFill>\n"
            '    <a:blip r:embed="rId3"/>\n'
            "    <a:stretch>\n"
            "      <a:fillRect/>\n"
            "    </a:stretch>\n"
            "  </p:blipFill>\n"
            "  <p:spPr>\n"
            "    <a:xfrm>\n"
            '      <a:off x="1" y="2"/>\n'
            '      <a:ext cx="3" cy="4"/>\n'
            "    </a:xfrm>\n"
            '    <a:prstGeom prst="rect">\n'
            "      <a:avLst/>\n"
            "    </a:prstGeom>\n"
            "  </p:spPr>\n"
            "</p:pic>\n"
        ) % nsdecls("a", "p", "r")
