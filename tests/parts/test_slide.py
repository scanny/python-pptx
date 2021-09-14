# encoding: utf-8

"""Unit-test suite for `pptx.parts.slide` module."""

import pytest

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE as XCT
from pptx.enum.shapes import PROG_ID
from pptx.media import Video
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml.slide import CT_NotesMaster, CT_NotesSlide, CT_Slide
from pptx.oxml.theme import CT_OfficeStyleSheet
from pptx.package import Package
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedPackagePart
from pptx.parts.image import Image, ImagePart
from pptx.parts.media import MediaPart
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import (
    BaseSlidePart,
    NotesMasterPart,
    NotesSlidePart,
    SlideLayoutPart,
    SlideMasterPart,
    SlidePart,
)
from pptx.slide import NotesMaster, NotesSlide, Slide, SlideLayout, SlideMaster

from ..unitutil.cxml import element
from ..unitutil.file import absjoin, test_file_dir
from ..unitutil.mock import (
    call,
    class_mock,
    initializer_mock,
    instance_mock,
    method_mock,
)


class DescribeBaseSlidePart(object):
    """Unit-test suite for `pptx.parts.slide.BaseSlidePart` objects."""

    def it_knows_its_name(self):
        slide_part = BaseSlidePart(
            None, None, None, element("p:sld/p:cSld{name=Foobar}")
        )
        assert slide_part.name == "Foobar"

    def it_can_get_a_related_image_by_rId(self, request, image_part_):
        image_ = instance_mock(request, Image)
        image_part_.image = image_
        related_part_ = method_mock(
            request, BaseSlidePart, "related_part", return_value=image_part_
        )
        slide_part = BaseSlidePart(None, None, None, None)

        image = slide_part.get_image("rId42")

        related_part_.assert_called_once_with(slide_part, "rId42")
        assert image is image_

    def it_can_add_an_image_part(self, request, image_part_):
        package_ = instance_mock(request, Package)
        package_.get_or_add_image_part.return_value = image_part_
        relate_to_ = method_mock(
            request, BaseSlidePart, "relate_to", return_value="rId6"
        )
        slide_part = BaseSlidePart(None, None, package_, None)

        image_part, rId = slide_part.get_or_add_image_part("foobar.png")

        package_.get_or_add_image_part.assert_called_once_with("foobar.png")
        relate_to_.assert_called_once_with(slide_part, image_part_, RT.IMAGE)
        assert image_part is image_part_
        assert rId == "rId6"

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)


class DescribeNotesMasterPart(object):
    """Unit-test suite for `pptx.parts.slide.NotesMasterPart` objects."""

    def it_can_create_a_notes_master_part(
        self, request, package_, notes_master_part_, theme_part_
    ):
        method_mock(
            request,
            NotesMasterPart,
            "_new",
            autospec=False,
            return_value=notes_master_part_,
        )
        method_mock(
            request,
            NotesMasterPart,
            "_new_theme_part",
            autospec=False,
            return_value=theme_part_,
        )
        notes_master_part = NotesMasterPart.create_default(package_)

        NotesMasterPart._new.assert_called_once_with(package_)
        NotesMasterPart._new_theme_part.assert_called_once_with(package_)
        notes_master_part.relate_to.assert_called_once_with(theme_part_, RT.THEME)
        assert notes_master_part is notes_master_part_

    def it_provides_access_to_its_notes_master(self, request):
        notes_master_ = instance_mock(request, NotesMaster)
        NotesMaster_ = class_mock(
            request, "pptx.parts.slide.NotesMaster", return_value=notes_master_
        )
        notesMaster = element("p:notesMaster")
        notes_master_part = NotesMasterPart(None, None, None, notesMaster)

        notes_master = notes_master_part.notes_master

        NotesMaster_.assert_called_once_with(notesMaster, notes_master_part)
        assert notes_master is notes_master_

    def it_creates_a_new_notes_master_part_to_help(
        self, request, package_, notes_master_part_
    ):
        NotesMasterPart_ = class_mock(
            request, "pptx.parts.slide.NotesMasterPart", return_value=notes_master_part_
        )
        notesMaster = element("p:notesMaster")
        method_mock(
            request,
            CT_NotesMaster,
            "new_default",
            autospec=False,
            return_value=notesMaster,
        )

        notes_master_part = NotesMasterPart._new(package_)

        CT_NotesMaster.new_default.assert_called_once_with()
        NotesMasterPart_.assert_called_once_with(
            PackURI("/ppt/notesMasters/notesMaster1.xml"),
            CT.PML_NOTES_MASTER,
            package_,
            notesMaster,
        )
        assert notes_master_part is notes_master_part_

    def it_creates_a_new_theme_part_to_help(self, request, package_, theme_part_):
        XmlPart_ = class_mock(
            request, "pptx.parts.slide.XmlPart", return_value=theme_part_
        )
        theme_elm = element("p:theme")
        method_mock(
            request,
            CT_OfficeStyleSheet,
            "new_default",
            autospec=False,
            return_value=theme_elm,
        )
        pn_tmpl = "/ppt/theme/theme%d.xml"
        partname = PackURI("/ppt/theme/theme2.xml")
        package_.next_partname.return_value = partname

        theme_part = NotesMasterPart._new_theme_part(package_)

        package_.next_partname.assert_called_once_with(pn_tmpl)
        CT_OfficeStyleSheet.new_default.assert_called_once_with()
        XmlPart_.assert_called_once_with(partname, CT.OFC_THEME, package_, theme_elm)
        assert theme_part is theme_part_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def notes_master_part_(self, request):
        return instance_mock(request, NotesMasterPart)

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def theme_part_(self, request):
        return instance_mock(request, Part)


class DescribeNotesSlidePart(object):
    """Unit-test suite for `pptx.parts.slide.NotesSlidePart` objects."""

    def it_can_create_a_notes_slide_part(
        self,
        request,
        package_,
        slide_part_,
        notes_master_part_,
        notes_slide_,
        notes_master_,
        notes_slide_part_,
    ):
        presentation_part_ = instance_mock(request, PresentationPart)
        package_.presentation_part = presentation_part_
        presentation_part_.notes_master_part = notes_master_part_
        _add_notes_slide_part_ = method_mock(
            request,
            NotesSlidePart,
            "_add_notes_slide_part",
            autospec=False,
            return_value=notes_slide_part_,
        )
        notes_slide_part_.notes_slide = notes_slide_
        notes_master_part_.notes_master = notes_master_

        notes_slide_part = NotesSlidePart.new(package_, slide_part_)

        _add_notes_slide_part_.assert_called_once_with(
            package_, slide_part_, notes_master_part_
        )
        notes_slide_.clone_master_placeholders.assert_called_once_with(notes_master_)
        assert notes_slide_part is notes_slide_part_

    def it_provides_access_to_the_notes_master(
        self, request, notes_master_, notes_master_part_
    ):
        part_related_by_ = method_mock(
            request, NotesSlidePart, "part_related_by", return_value=notes_master_part_
        )
        notes_slide_part = NotesSlidePart(None, None, None, None)
        notes_master_part_.notes_master = notes_master_

        notes_master = notes_slide_part.notes_master

        part_related_by_.assert_called_once_with(notes_slide_part, RT.NOTES_MASTER)
        assert notes_master is notes_master_

    def it_provides_access_to_its_notes_slide(self, request, notes_slide_):
        NotesSlide_ = class_mock(
            request, "pptx.parts.slide.NotesSlide", return_value=notes_slide_
        )
        notes = element("p:notes")
        notes_slide_part = NotesSlidePart(None, None, None, notes)

        notes_slide = notes_slide_part.notes_slide

        NotesSlide_.assert_called_once_with(notes, notes_slide_part)
        assert notes_slide is notes_slide_

    def it_adds_a_notes_slide_part_to_help(
        self, request, package_, slide_part_, notes_master_part_, notes_slide_part_
    ):
        NotesSlidePart_ = class_mock(
            request, "pptx.parts.slide.NotesSlidePart", return_value=notes_slide_part_
        )
        notes = element("p:notes")
        new_ = method_mock(
            request, CT_NotesSlide, "new", autospec=False, return_value=notes
        )
        package_.next_partname.return_value = PackURI(
            "/ppt/notesSlides/notesSlide42.xml"
        )

        notes_slide_part = NotesSlidePart._add_notes_slide_part(
            package_, slide_part_, notes_master_part_
        )

        package_.next_partname.assert_called_once_with(
            "/ppt/notesSlides/notesSlide%d.xml"
        )
        new_.assert_called_once_with()
        NotesSlidePart_.assert_called_once_with(
            PackURI("/ppt/notesSlides/notesSlide42.xml"),
            CT.PML_NOTES_SLIDE,
            package_,
            notes,
        )
        assert notes_slide_part_.relate_to.call_args_list == [
            call(notes_master_part_, RT.NOTES_MASTER),
            call(slide_part_, RT.SLIDE),
        ]
        assert notes_slide_part is notes_slide_part_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def notes_master_(self, request):
        return instance_mock(request, NotesMaster)

    @pytest.fixture
    def notes_master_part_(self, request):
        return instance_mock(request, NotesMasterPart)

    @pytest.fixture
    def notes_slide_(self, request):
        return instance_mock(request, NotesSlide)

    @pytest.fixture
    def notes_slide_part_(self, request):
        return instance_mock(request, NotesSlidePart)

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribeSlidePart(object):
    """Unit-test suite for `pptx.parts.slide.SlidePart` objects."""

    def it_knows_its_slide_id(self, slide_id_fixture):
        slide_part, presentation_part_, slide_id = slide_id_fixture
        _slide_id = slide_part.slide_id
        presentation_part_.slide_id.assert_called_once_with(slide_part)
        assert _slide_id is slide_id

    def it_knows_whether_it_has_a_notes_slide(self, has_notes_slide_fixture):
        slide_part, expected_value = has_notes_slide_fixture
        value = slide_part.has_notes_slide
        slide_part.part_related_by.assert_called_once_with(slide_part, RT.NOTES_SLIDE)
        assert value is expected_value

    def it_can_add_a_chart_part(self, request, package_, relate_to_):
        chart_data_ = instance_mock(request, ChartData)
        chart_part_ = instance_mock(request, ChartPart)
        ChartPart_ = class_mock(request, "pptx.parts.slide.ChartPart")
        ChartPart_.new.return_value = chart_part_
        relate_to_.return_value = "rId42"
        slide_part = SlidePart(None, None, package_, None)

        rId = slide_part.add_chart_part(XCT.RADAR, chart_data_)

        ChartPart_.new.assert_called_once_with(XCT.RADAR, chart_data_, package_)
        relate_to_.assert_called_once_with(slide_part, chart_part_, RT.CHART)
        assert rId == "rId42"

    @pytest.mark.parametrize(
        "prog_id, rel_type",
        (
            (PROG_ID.DOCX, RT.PACKAGE),
            (PROG_ID.PPTX, RT.PACKAGE),
            (PROG_ID.XLSX, RT.PACKAGE),
            ("Foo.Bar.18", RT.OLE_OBJECT),
        ),
    )
    def it_can_add_an_embedded_ole_object_part(
        self, request, package_, relate_to_, prog_id, rel_type
    ):
        _blob_from_file_ = method_mock(
            request, SlidePart, "_blob_from_file", return_value=b"012345"
        )
        embedded_package_part_ = instance_mock(request, EmbeddedPackagePart)
        EmbeddedPackagePart_ = class_mock(
            request, "pptx.parts.slide.EmbeddedPackagePart"
        )
        EmbeddedPackagePart_.factory.return_value = embedded_package_part_
        relate_to_.return_value = "rId9"
        slide_part = SlidePart(None, None, package_, None)

        _rId = slide_part.add_embedded_ole_object_part(prog_id, "workbook.xlsx")

        _blob_from_file_.assert_called_once_with(slide_part, "workbook.xlsx")
        EmbeddedPackagePart_.factory.assert_called_once_with(
            prog_id, b"012345", package_
        )
        relate_to_.assert_called_once_with(slide_part, embedded_package_part_, rel_type)
        assert _rId == "rId9"

    def it_can_get_or_add_a_video_part(self, package_, video_, relate_to_, media_part_):
        media_rId, video_rId = "rId1", "rId2"
        package_.get_or_add_media_part.return_value = media_part_
        relate_to_.side_effect = [media_rId, video_rId]
        slide_part = SlidePart(None, None, package_, None)

        result = slide_part.get_or_add_video_media_part(video_)

        package_.get_or_add_media_part.assert_called_once_with(video_)
        assert relate_to_.call_args_list == [
            call(slide_part, media_part_, RT.MEDIA),
            call(slide_part, media_part_, RT.VIDEO),
        ]
        assert result == (media_rId, video_rId)

    def it_can_create_a_new_slide_part(self, request, package_, relate_to_):
        partname = PackURI("/foobar.xml")
        _init_ = initializer_mock(request, SlidePart)
        slide_layout_part_ = instance_mock(request, SlideLayoutPart)
        CT_Slide_ = class_mock(request, "pptx.parts.slide.CT_Slide")
        CT_Slide_.new.return_value = sld = element("c:sld")

        slide_part = SlidePart.new(partname, package_, slide_layout_part_)

        _init_.assert_called_once_with(
            slide_part, partname, CT.PML_SLIDE, package_, sld
        )
        slide_part.relate_to.assert_called_once_with(
            slide_part, slide_layout_part_, RT.SLIDE_LAYOUT
        )
        assert isinstance(slide_part, SlidePart)

    def it_provides_access_to_its_slide(self, slide_fixture):
        slide_part, Slide_, sld, slide_ = slide_fixture
        slide = slide_part.slide
        Slide_.assert_called_once_with(sld, slide_part)
        assert slide is slide_

    def it_provides_access_to_the_slide_layout(self, layout_fixture):
        slide_part, slide_layout_ = layout_fixture
        slide_layout = slide_part.slide_layout
        slide_part.part_related_by.assert_called_once_with(slide_part, RT.SLIDE_LAYOUT)
        assert slide_layout is slide_layout_

    def it_knows_the_minimal_element_xml_for_a_slide(self):
        path = absjoin(test_file_dir, "minimal_slide.xml")
        sld = CT_Slide.new()
        with open(path, "r") as f:
            expected_xml = f.read()
        assert sld.xml == expected_xml

    def it_gets_its_notes_slide_to_help(self, notes_slide_fixture):
        slide_part, calls, notes_slide_ = notes_slide_fixture

        notes_slide = slide_part.notes_slide

        slide_part.part_related_by.assert_called_once_with(slide_part, RT.NOTES_SLIDE)
        assert slide_part._add_notes_slide_part.call_args_list == calls
        assert notes_slide is notes_slide_

    def it_adds_a_notes_slide_part_to_help(
        self, package_, NotesSlidePart_, notes_slide_part_, relate_to_
    ):
        NotesSlidePart_.new.return_value = notes_slide_part_
        slide_part = SlidePart(None, None, package_, None)

        notes_slide_part = slide_part._add_notes_slide_part()

        assert notes_slide_part is notes_slide_part_
        NotesSlidePart_.new.assert_called_once_with(package_, slide_part)
        relate_to_.assert_called_once_with(slide_part, notes_slide_part, RT.NOTES_SLIDE)

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[True, False])
    def has_notes_slide_fixture(self, request, part_related_by_):
        has_notes_slide = request.param
        slide_part = SlidePart(None, None, None, None)
        part_related_by_.side_effect = None if has_notes_slide else KeyError
        expected_value = has_notes_slide
        return slide_part, expected_value

    @pytest.fixture
    def layout_fixture(self, slide_layout_, part_related_by_):
        slide_part = SlidePart(None, None, None, None)
        part_related_by_.return_value.slide_layout = slide_layout_
        return slide_part, slide_layout_

    @pytest.fixture(params=[True, False])
    def notes_slide_fixture(
        self,
        request,
        notes_slide_,
        part_related_by_,
        _add_notes_slide_part_,
        notes_slide_part_,
    ):
        has_notes_slide = request.param
        slide_part = SlidePart(None, None, None, None)
        part_related_by_.return_value = notes_slide_part_
        add_calls = []
        if not has_notes_slide:
            part_related_by_.side_effect = KeyError
            add_calls.append(call(slide_part))
        notes_slide_part_.notes_slide = notes_slide_
        return slide_part, add_calls, notes_slide_

    @pytest.fixture
    def slide_fixture(self, Slide_, slide_):
        sld = element("p:sld")
        slide_part = SlidePart(None, None, None, sld)
        return slide_part, Slide_, sld, slide_

    @pytest.fixture
    def slide_id_fixture(self, package_, presentation_part_):
        slide_part = SlidePart(None, None, package_, None)
        slide_id = 256
        package_.presentation_part = presentation_part_
        presentation_part_.slide_id.return_value = slide_id
        return slide_part, presentation_part_, slide_id

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _add_notes_slide_part_(self, request, notes_slide_part_):
        return method_mock(
            request,
            SlidePart,
            "_add_notes_slide_part",
            return_value=notes_slide_part_,
            autospec=True,
        )

    @pytest.fixture
    def media_part_(self, request):
        return instance_mock(request, MediaPart)

    @pytest.fixture
    def notes_slide_(self, request):
        return instance_mock(request, NotesSlide)

    @pytest.fixture
    def NotesSlidePart_(self, request, notes_slide_part_):
        return class_mock(
            request, "pptx.parts.slide.NotesSlidePart", return_value=notes_slide_part_
        )

    @pytest.fixture
    def notes_slide_part_(self, request):
        return instance_mock(request, NotesSlidePart)

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def part_related_by_(self, request):
        return method_mock(request, SlidePart, "part_related_by", autospec=True)

    @pytest.fixture
    def presentation_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def relate_to_(self, request):
        return method_mock(request, SlidePart, "relate_to", autospec=True)

    @pytest.fixture
    def Slide_(self, request, slide_):
        return class_mock(request, "pptx.parts.slide.Slide", return_value=slide_)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def video_(self, request):
        return instance_mock(request, Video)


class DescribeSlideLayoutPart(object):
    """Unit-test suite for `pptx.parts.slide.SlideLayoutPart` objects."""

    def it_provides_access_to_its_slide_master(self, request):
        slide_master_ = instance_mock(request, SlideMaster)
        slide_master_part_ = instance_mock(
            request, SlideMasterPart, slide_master=slide_master_
        )
        part_related_by_ = method_mock(
            request, SlideLayoutPart, "part_related_by", return_value=slide_master_part_
        )
        slide_layout_part = SlideLayoutPart(None, None, None, None)

        slide_master = slide_layout_part.slide_master

        part_related_by_.assert_called_once_with(slide_layout_part, RT.SLIDE_MASTER)
        assert slide_master is slide_master_

    def it_provides_access_to_its_slide_layout(self, request):
        slide_layout_ = instance_mock(request, SlideLayout)
        SlideLayout_ = class_mock(
            request, "pptx.parts.slide.SlideLayout", return_value=slide_layout_
        )
        sldLayout = element("p:sldLayout")
        slide_layout_part = SlideLayoutPart(None, None, None, sldLayout)

        slide_layout = slide_layout_part.slide_layout

        SlideLayout_.assert_called_once_with(sldLayout, slide_layout_part)
        assert slide_layout is slide_layout_


class DescribeSlideMasterPart(object):
    """Unit-test suite for `pptx.parts.slide.SlideMasterPart` objects."""

    def it_provides_access_to_its_slide_master(self, request):
        slide_master_ = instance_mock(request, SlideMaster)
        SlideMaster_ = class_mock(
            request, "pptx.parts.slide.SlideMaster", return_value=slide_master_
        )
        sldMaster = element("p:sldMaster")
        slide_master_part = SlideMasterPart(None, None, None, sldMaster)

        slide_master = slide_master_part.slide_master

        SlideMaster_.assert_called_once_with(sldMaster, slide_master_part)
        assert slide_master is slide_master_

    def it_provides_access_to_a_related_slide_layout(self, request):
        slide_layout_ = instance_mock(request, SlideLayout)
        slide_layout_part_ = instance_mock(
            request, SlideLayoutPart, slide_layout=slide_layout_
        )
        related_part_ = method_mock(
            request, SlideMasterPart, "related_part", return_value=slide_layout_part_
        )
        slide_master_part = SlideMasterPart(None, None, None, None)

        slide_layout = slide_master_part.related_slide_layout("rId42")

        related_part_.assert_called_once_with(slide_master_part, "rId42")
        assert slide_layout is slide_layout_
