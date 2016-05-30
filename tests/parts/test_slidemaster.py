# encoding: utf-8

"""
Test suite for pptx.parts.slidemaster module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.parts.slidemaster import SlideMasterPart
from pptx.shapes.shapetree import MasterPlaceholders, MasterShapes
from pptx.slide import SlideLayout, SlideLayouts

from ..unitutil.cxml import element
from ..unitutil.mock import class_mock, instance_mock, property_mock


class DescribeSlideMasterPart(object):

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide_master, MasterPlaceholders_, spTree, placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide_master.placeholders
        MasterPlaceholders_.assert_called_once_with(spTree, slide_master)
        assert placeholders is placeholders_

    def it_provides_access_to_its_shapes(self, shapes_fixture):
        slide_master, MasterShapes_, spTree, shapes_ = shapes_fixture
        shapes = slide_master.shapes
        MasterShapes_.assert_called_once_with(spTree, slide_master)
        assert shapes is shapes_

    def it_provides_access_to_its_slide_layouts(self, layouts_fixture):
        slide_master, SlideLayouts_, sldLayoutIdLst, slide_layouts_ = (
            layouts_fixture
        )
        slide_layouts = slide_master.slide_layouts
        SlideLayouts_.assert_called_once_with(sldLayoutIdLst, slide_master)
        assert slide_layouts is slide_layouts_

    def it_provides_access_to_a_related_slide_layout(self, related_fixture):
        slide_master_part, rId, getitem_, slide_layout_ = related_fixture
        slide_layout = slide_master_part.related_slide_layout(rId)
        getitem_.assert_called_once_with(rId)
        assert slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layouts_fixture(self, SlideLayouts_, slide_layouts_):
        sldMaster = element('p:sldMaster/p:sldLayoutIdLst')
        slide_master = SlideMasterPart(None, None, sldMaster)
        sldMasterIdLst = sldMaster.sldLayoutIdLst
        return slide_master, SlideLayouts_, sldMasterIdLst, slide_layouts_

    @pytest.fixture
    def related_fixture(self, slide_layout_, related_parts_prop_):
        slide_master_part = SlideMasterPart(None, None, None, None)
        rId = 'rId42'
        getitem_ = related_parts_prop_.return_value.__getitem__
        getitem_.return_value.slide_layout = slide_layout_
        return slide_master_part, rId, getitem_, slide_layout_

    @pytest.fixture
    def placeholders_fixture(self, MasterPlaceholders_, placeholders_):
        sldMaster = element('p:sldMaster/p:cSld/p:spTree')
        slide_master = SlideMasterPart(None, None, sldMaster)
        spTree = sldMaster.xpath('//p:spTree')[0]
        return slide_master, MasterPlaceholders_, spTree, placeholders_

    @pytest.fixture
    def shapes_fixture(self, MasterShapes_, shapes_):
        sldMaster = element('p:sldMaster/p:cSld/p:spTree')
        slide_master = SlideMasterPart(None, None, sldMaster)
        spTree = sldMaster.xpath('//p:spTree')[0]
        return slide_master, MasterShapes_, spTree, shapes_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def MasterPlaceholders_(self, request, placeholders_):
        return class_mock(
            request, 'pptx.parts.slidemaster.MasterPlaceholders',
            return_value=placeholders_
        )

    @pytest.fixture
    def MasterShapes_(self, request, shapes_):
        return class_mock(
            request, 'pptx.parts.slidemaster.MasterShapes',
            return_value=shapes_
        )

    @pytest.fixture
    def placeholders_(self, request):
        return instance_mock(request, MasterPlaceholders)

    @pytest.fixture
    def related_parts_prop_(self, request):
        return property_mock(request, SlideMasterPart, 'related_parts')

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, MasterShapes)

    @pytest.fixture
    def SlideLayouts_(self, request, slide_layouts_):
        return class_mock(
            request, 'pptx.parts.slidemaster.SlideLayouts',
            return_value=slide_layouts_
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_layouts_(self, request):
        return instance_mock(request, SlideLayouts)
