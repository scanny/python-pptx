# encoding: utf-8

"""Test suite for pptx.slides module."""

from __future__ import absolute_import

import pytest

from lxml import objectify
from mock import ANY, call, Mock

from pptx.opc import package
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.opc.rels import _Relationship, RelationshipCollection
from pptx.oxml.ns import namespaces
from pptx.oxml.presentation import CT_SlideId, CT_SlideIdList
from pptx.parts.slides import (
    _BaseSlide, Slide, SlideCollection, SlideLayout, SlideMaster
)
from pptx.presentation import Package, Presentation
from pptx.shapes.shapetree import ShapeCollection

from ..unitutil import (
    absjoin, class_mock, instance_mock, loose_mock, method_mock,
    parse_xml_file, property_mock, serialize_xml, test_file_dir
)


test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')

nsmap = namespaces('a', 'r', 'p')


def actual_xml(elm):
    objectify.deannotate(elm, cleanup_namespaces=True)
    return serialize_xml(elm, pretty_print=True)


def _sldLayout1():
    path = absjoin(test_file_dir, 'slideLayout1.xml')
    sldLayout = parse_xml_file(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = ShapeCollection(spTree)
    return shapes


class Describe_BaseSlide(object):

    def it_knows_the_name_of_the_slide(self, base_slide):
        """_BaseSlide.name value is correct"""
        # setup ------------------------
        base_slide._element = _sldLayout1()
        # exercise ---------------------
        name = base_slide.name
        # verify -----------------------
        expected = 'Title Slide'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        assert actual == expected, msg

    def it_provides_access_to_the_shapes_on_the_slide(self):
        """_BaseSlide.shapes is expected size after _load()"""
        # setup ------------------------
        path = absjoin(test_file_dir, 'slide1.xml')
        with open(path, 'r') as f:
            blob = f.read()
        base_slide = _BaseSlide.load(None, None, blob)
        # exercise ---------------------
        shapes = base_slide.shapes
        # verify -----------------------
        assert len(shapes) == 9

    def it_can_add_an_image_part_to_the_slide(self, base_slide_fixture):
        base_slide, image_, rel_ = base_slide_fixture
        image, rel = base_slide._add_image(file)
        base_slide._package._images.add_image.assert_called_once_with(file)
        base_slide._relationships.get_or_add.assert_called_once_with(
            RT.IMAGE, image_
        )
        assert image is image_
        assert rel is rel_

    def it_knows_what_to_do_after_the_slide_is_unmarshaled(self):
        pass

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def base_slide_fixture(self, request, base_slide):
        # mock _BaseSlide._package._images.add_image() train wreck
        image_ = loose_mock(request, name='image_')
        pkg_ = loose_mock(request, name='_package', spec=Package)
        pkg_._images.add_image.return_value = image_
        _package = property_mock(  # noqa
            request, 'pptx.parts.slides._BaseSlide._package',
            return_value=pkg_
        )
        # mock _BaseSlide._relationships.get_or_add()
        rel_ = loose_mock(request, name='rel_')
        rels_ = loose_mock(request, name='rels_')
        rels_.get_or_add.return_value = rel_
        _relationships = property_mock(  # noqa
            request, 'pptx.parts.slides._BaseSlide._relationships',
            return_value=rels_
        )
        return base_slide, image_, rel_

    @pytest.fixture
    def base_slide(self):
        partname = PackURI('/foo/bar.xml')
        return _BaseSlide(partname, None, None)


class DescribeSlide(object):

    def it_establishes_a_relationship_to_its_slide_layout_on_construction(
            self):
        """Slide(slidelayout) adds relationship slide->slidelayout"""
        # setup ------------------------
        slidelayout = SlideLayout(None, None, _sldLayout1())
        partname = PackURI('/ppt/slides/slide1.xml')
        # exercise ---------------------
        slide = Slide.new(slidelayout, partname)
        # verify length ---------------
        assert len(slide._relationships) == 1
        # verify values ---------------
        rel = slide._relationships[0]
        expected = ('rId1', RT.SLIDE_LAYOUT, slidelayout)
        actual = (rel.rId, rel.reltype, rel.target_part)
        assert actual == expected

    # def it_creates_a_minimal_sld_element_on_construction(self, slide):
    #     """Slide._element is minimal sld on construction"""
    #     # setup ------------------------
    #     slidelayout = SlideLayout(None, None, _sldLayout1())
    #     partname = PackURI('/ppt/slides/slide1.xml')
    #     slide = Slide.new(slidelayout, partname)
    #     path = absjoin(test_file_dir, 'minimal_slide.xml')
    #     # exercise ---------------------
    #     elm = slide._element
    #     # verify -----------------------
    #     with open(path, 'r') as f:
    #         expected_xml = f.read()
    #     assert actual_xml(elm) == expected_xml

    # def it_has_slidelayout_property_of_none_on_construction(self, slide):
    #     """Slide.slidelayout property None on construction"""
    #     assert slide.slidelayout is None

    # def it_sets_slidelayout_on_load(self, slide):
    #     """Slide._load() sets slidelayout"""
    #     # setup ------------------------
    #     path = absjoin(test_file_dir, 'slide1.xml')
    #     slidelayout = Mock(name='slideLayout')
    #     slidelayout.partname = '/ppt/slideLayouts/slideLayout1.xml'
    #     rel = Mock(name='pptx.package.Relationship')
    #     rel.rId = 'rId1'
    #     rel.reltype = RT.SLIDE_LAYOUT
    #     rel.target = slidelayout
    #     pkgpart = Mock(name='pptx.package.Part')
    #     with open(path, 'rb') as f:
    #         pkgpart.blob = f.read()
    #     pkgpart.relationships = [rel]
    #     part_dict = {slidelayout.partname: slidelayout}
    #     slide_ = slide.load(pkgpart, part_dict)
    #     # exercise ---------------------
    #     retval = slide_.slidelayout
    #     # verify -----------------------
    #     expected = slidelayout
    #     actual = retval
    #     msg = "expected: %s, got %s" % (expected, actual)
    #     assert actual == expected, msg

    def it_knows_the_minimal_element_xml_for_a_slide(self, slide):
        """Slide._minimal_element generates correct XML"""
        # setup ------------------------
        path = absjoin(test_file_dir, 'minimal_slide.xml')
        # exercise ---------------------
        sld = slide._minimal_element()
        # verify -----------------------
        with open(path, 'r') as f:
            expected_xml = f.read()
        assert actual_xml(sld) == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def slide(self):
        return Slide(None, None, None)


class DescribeSlideCollection(object):

    def it_supports_indexed_access(self, slides, slide_, slide_2_):
        assert (slides[0], slides[1]) == (slide_, slide_2_)
        slides._sldIdLst.__getitem__.assert_has_calls([call(0), call(1)])

    def it_raises_on_slide_index_out_of_range(self, slides):
        with pytest.raises(IndexError):
            slides[2]

    def it_is_iterable(self, slides, slide_, slide_2_):
        assert [s for s in slides] == [slide_, slide_2_]

    def it_supports_len(self, slides):
        assert len(slides) == 2

    def it_can_add_a_new_slide(
            self, slides, Slide_, slidelayout_, prs_, slide_, sldIdLst_,
            rId_, _rename_slides_):
        slide = slides.add_slide(slidelayout_)
        # verify -----------------------
        Slide_.new.assert_called_once_with(slidelayout_, ANY)
        prs_._add_relationship.assert_called_once_with(RT.SLIDE, slide_, ANY)
        sldIdLst_.add_sldId.assert_called_once_with(rId_)
        _rename_slides_.assert_called_once_with()
        assert slide is slide_

    def it_can_assign_partnames_to_the_slides(
            self, slides, slide_, slide_2_):
        slides._rename_slides()
        assert slide_.partname == '/ppt/slides/slide1.xml'
        assert slide_2_.partname == '/ppt/slides/slide2.xml'

    def it_can_get_a_slide_based_on_a_sldId(
            self, slides, sldId_, slide_, prs_rels_, rId_):
        slide = slides._slide_from_sldId(sldId_)
        prs_rels_.part_with_rId.assert_called_once_with(rId_)
        assert slide == slide_

    def it_can_iterate_over_the_slides(self, slides, slide_, slide_2_):
        assert [s for s in slides._slides] == [slide_, slide_2_]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def prs_(self, request, rel_):
        prs_ = instance_mock(request, Presentation)
        prs_._add_relationship.return_value = rel_
        return prs_

    @pytest.fixture
    def prs_rels_(self, request, slide_, slide_2_):
        prs_rels_ = instance_mock(request, RelationshipCollection)
        prs_rels_.part_with_rId.side_effect = [slide_, slide_2_]
        return prs_rels_

    @pytest.fixture
    def rel_(self, request, rId_):
        return instance_mock(request, _Relationship, rId=rId_)

    @pytest.fixture
    def _rename_slides_(self, request):
        return method_mock(request, SlideCollection, '_rename_slides')

    @pytest.fixture
    def rId_(self, request):
        return 'rId1'

    @pytest.fixture
    def rId_2_(self, request):
        return 'rId2'

    @pytest.fixture
    def sldId_(self, request, rId_):
        return instance_mock(request, CT_SlideId, rId=rId_)

    @pytest.fixture
    def sldId_2_(self, request, rId_2_):
        return instance_mock(request, CT_SlideId, rId=rId_2_)

    @pytest.fixture
    def sldIdLst_(self, request, sldId_, sldId_2_):
        sldIdLst_ = instance_mock(request, CT_SlideIdList)
        sldIdLst_.__getitem__.side_effect = [sldId_, sldId_2_]
        sldIdLst_.__iter__.return_value = iter([sldId_, sldId_2_])
        sldIdLst_.__len__.return_value = 2
        return sldIdLst_

    @pytest.fixture
    def Slide_(self, request, slide_):
        Slide_ = class_mock(request, 'pptx.parts.slides.Slide')
        Slide_.new.return_value = slide_
        return Slide_

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_2_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slidelayout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slides(self, sldIdLst_, prs_rels_, prs_):
        return SlideCollection(sldIdLst_, prs_rels_, prs_)


class DescribeSlideLayout(object):

    def _loaded_slidelayout(self, prs_slidemaster=None):
        """
        Return SlideLayout instance loaded using mocks. *prs_slidemaster* is
        an already-loaded model-side SlideMaster instance (or mock, as
        appropriate to calling test).
        """
        # partname for related slideMaster
        sldmaster_partname = '/ppt/slideMasters/slideMaster1.xml'
        # path to test slideLayout XML
        slidelayout_path = absjoin(test_file_dir, 'slideLayout1.xml')
        # model-side slideMaster part
        if prs_slidemaster is None:
            prs_slidemaster = Mock(spec=SlideMaster)
        # a part dict containing the already-loaded model-side slideMaster
        loaded_part_dict = {sldmaster_partname: prs_slidemaster}
        # a slideMaster package part for rel target
        pkg_slidemaster_part = Mock(spec=package.Part)
        pkg_slidemaster_part.partname = sldmaster_partname
        # a package-side relationship from slideLayout to its slideMaster
        rel = Mock(name='pptx.package.Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT.SLIDE_MASTER
        rel.target = pkg_slidemaster_part
        # the slideLayout package part to send to _load()
        pkg_slidelayout_part = Mock(spec=package.Part)
        pkg_slidelayout_part.relationships = [rel]
        with open(slidelayout_path, 'rb') as f:
            pkg_slidelayout_part.blob = f.read()
        # _load and return
        slidelayout = SlideLayout()
        return slidelayout._load(pkg_slidelayout_part, loaded_part_dict)

    # def test__load_sets_slidemaster(self):
    #     """SlideLayout._load() sets slidemaster"""
    #     # setup ------------------------
    #     prs_slidemaster = Mock(spec=SlideMaster)
    #     # exercise ---------------------
    #     loaded_slidelayout = self._loaded_slidelayout(prs_slidemaster)
    #     # verify -----------------------
    #     expected = prs_slidemaster
    #     actual = loaded_slidelayout.slidemaster
    #     msg = "expected: %s, got %s" % (expected, actual)
    #     assert actual == expected, msg

    # def test_slidemaster_raises_on_ref_before_assigned(self, slidelayout):
    #     """SlideLayout.slidemaster raises on referenced before assigned"""
    #     with pytest.raises(AssertionError):
    #         slidelayout.slidemaster

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def slidelayout(self):
        return SlideLayout()


class DescribeSlideMaster(object):

    def test_slidelayouts_property_empty_on_construction(self, slidemaster):
        assert len(slidemaster.slidelayouts) == 0

    def test_slidelayouts_correct_length_after_open(self):
        slidemaster = Package.open(test_pptx_path).presentation.slidemasters[0]
        slidelayouts = slidemaster.slidelayouts
        assert len(slidelayouts) == 11

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def slidemaster(self):
        partname = PackURI('/ppt/slideMasters/slideMaster1.xml')
        return SlideMaster(partname, None, None)
