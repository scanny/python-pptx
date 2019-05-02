# encoding: utf-8

"""Test suite for pptx.shapes.autoshape module."""

from __future__ import absolute_import, division, print_function, unicode_literals

import pytest

from pptx.dml.fill import FillFormat
from pptx.dml.line import LineFormat
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.shapes.autoshape import CT_PresetGeometry2D, CT_Shape
from pptx.shapes.autoshape import Adjustment, AdjustmentCollection, AutoShapeType, Shape
from pptx.text.text import TextFrame

from ..oxml.unitdata.shape import (
    a_cNvSpPr,
    a_gd,
    a_prstGeom,
    an_avLst,
    an_nvSpPr,
    an_sp,
    an_spPr,
)
from ..unitutil.cxml import element, xml
from ..unitutil.mock import class_mock, instance_mock, property_mock


class DescribeAdjustment(object):
    def it_knows_its_effective_value(self, effective_val_fixture_):
        adjustment, expected_effective_value = effective_val_fixture_
        assert adjustment.effective_value == expected_effective_value

    # fixture --------------------------------------------------------

    def _effective_adj_val_cases():
        return [
            # no actual, effective should be determined by default value
            (50000, None, 0.5),
            # actual matches default
            (50000, 50000, 0.5),
            # actual is different than default
            (50000, 12500, 0.125),
            # actual is zero
            (50000, 0, 0.0),
            # negative default
            (-20833, None, -0.20833),
            # negative actual
            (-20833, -5678901, -56.78901),
        ]

    @pytest.fixture(params=_effective_adj_val_cases())
    def effective_val_fixture_(self, request):
        name = None
        def_val, actual, expected_effective_value = request.param
        adjustment = Adjustment(name, def_val, actual)
        return adjustment, expected_effective_value


class DescribeAdjustmentCollection(object):
    def it_should_load_default_adjustment_values(self, prstGeom_cases_):
        prstGeom, prst, expected = prstGeom_cases_
        adjustments = AdjustmentCollection(prstGeom)._adjustments
        actuals = tuple([(adj.name, adj.def_val) for adj in adjustments])
        assert len(adjustments) == len(expected)
        assert actuals == expected

    def it_should_load_adj_val_actuals_from_xml(self, load_adj_actuals_fixture_):
        prstGeom, expected_actuals, prstGeom_xml = load_adj_actuals_fixture_
        adjustments = AdjustmentCollection(prstGeom)._adjustments
        actual_actuals = dict([(a.name, a.actual) for a in adjustments])
        assert actual_actuals == expected_actuals

    def it_provides_normalized_effective_value_on_indexed_access(
        self, indexed_access_fixture_
    ):
        prstGeom, prst, expected_values = indexed_access_fixture_
        adjustments = AdjustmentCollection(prstGeom)
        actual_values = [adjustments[idx] for idx in range(len(adjustments))]
        assert actual_values == expected_values

    def it_should_update_actual_value_on_indexed_assignment(
        self, indexed_assignment_fixture_
    ):
        """
        Assignment to AdjustmentCollection[n] updates nth actual
        """
        adjs, idx, new_val, expected = indexed_assignment_fixture_
        adjs[idx] = new_val
        assert adjs._adjustments[idx].actual == expected

    def it_should_round_trip_indexed_assignment(self, adjustments):
        new_value = 0.375
        assert adjustments[0] != new_value
        # exercise ---------------------
        adjustments[0] = new_value
        # verify -----------------------
        assert adjustments[0] == new_value

    def it_should_raise_on_bad_index(self, adjustments):
        with pytest.raises(IndexError):
            adjustments[-6]
        with pytest.raises(IndexError):
            adjustments[6]
        with pytest.raises(TypeError):
            adjustments[0.0]
        with pytest.raises(TypeError):
            adjustments["0"]
        with pytest.raises(IndexError):
            adjustments[-6] = 1.0
        with pytest.raises(IndexError):
            adjustments[6] = 1.0
        with pytest.raises(TypeError):
            adjustments[0.0] = 1.0
        with pytest.raises(TypeError):
            adjustments["0"] = 1.0

    def it_should_raise_on_assigned_bad_value(self, adjustments):
        """
        AdjustmentCollection[n] = val raises on val is not number
        """
        with pytest.raises(ValueError):
            adjustments[0] = "1.0"

    def it_writes_adj_vals_to_xml_on_assignment(self, adjustments_with_prstGeom_):
        adjs, guides = adjustments_with_prstGeom_
        adjs[0] = 0.999
        adjs._prstGeom.rewrite_guides.assert_called_once_with(guides)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def adjustments(self):
        prstGeom = a_prstGeom().with_nsdecls().with_prst("chevron").element
        return AdjustmentCollection(prstGeom)

    @pytest.fixture
    def adjustments_with_prstGeom_(self, request):
        prstGeom = a_prstGeom().with_nsdecls().with_prst("chevron").element
        adjustments = AdjustmentCollection(prstGeom)
        prstGeom_ = instance_mock(request, CT_PresetGeometry2D, name="prstGeom_")
        adjustments._prstGeom = prstGeom_
        guides = [("adj", 99900)]
        return adjustments, guides

    def _adj_actuals_cases():
        gd_bldr = a_gd().with_name("adj2").with_fmla("val 25000")
        avLst_bldr = an_avLst().with_child(gd_bldr)
        mathDivide_bldr = (
            a_prstGeom().with_nsdecls().with_prst("mathDivide").with_child(avLst_bldr)
        )

        gd_bldr = a_gd().with_name("adj").with_fmla("val 25000")
        avLst_bldr = an_avLst().with_child(gd_bldr)
        rect_bldr = a_prstGeom().with_nsdecls().with_prst("rect").with_child(avLst_bldr)

        gd_bldr_1 = a_gd().with_name("adj1").with_fmla("val 111")
        gd_bldr_2 = a_gd().with_name("adj2").with_fmla("val 222")
        gd_bldr_3 = a_gd().with_name("adj3").with_fmla("val 333")
        avLst_bldr = (
            an_avLst().with_child(gd_bldr_1).with_child(gd_bldr_2).with_child(gd_bldr_3)
        )
        wedgeRoundRectCallout_bldr = (
            a_prstGeom()
            .with_nsdecls()
            .with_prst("wedgeRoundRectCallout")
            .with_child(avLst_bldr)
        )

        return [
            # no adjustment values in xml or spec
            (a_prstGeom().with_nsdecls().with_prst("rect").with_child(an_avLst()), ()),
            # no adjustment values in xml, but some in spec
            (
                a_prstGeom().with_nsdecls().with_prst("circularArrow"),
                (
                    ("adj1", None),
                    ("adj2", None),
                    ("adj3", None),
                    ("adj4", None),
                    ("adj5", None),
                ),
            ),
            # adjustment value in xml but none in spec
            (rect_bldr, ()),
            # middle adjustment value in xml
            (mathDivide_bldr, (("adj1", None), ("adj2", 25000), ("adj3", None))),
            # all adjustment values in xml
            (wedgeRoundRectCallout_bldr, (("adj1", 111), ("adj2", 222), ("adj3", 333))),
        ]

    @pytest.fixture(params=_adj_actuals_cases())
    def load_adj_actuals_fixture_(self, request):
        prstGeom_bldr, adj_vals = request.param
        prstGeom = prstGeom_bldr.element
        expected = dict(adj_vals)
        prstGeom_xml = prstGeom_bldr.xml
        return prstGeom, expected, prstGeom_xml

    def _prstGeom_cases():
        return [
            # rect has no adjustments
            ("rect", ()),
            # chevron has one simple one
            ("chevron", (("adj", 50000),)),
            # one with several and some negative
            (
                "accentBorderCallout1",
                (("adj1", 18750), ("adj2", -8333), ("adj3", 112500), ("adj4", -38333)),
            ),
            # another one with some negative
            (
                "wedgeRoundRectCallout",
                (("adj1", -20833), ("adj2", 62500), ("adj3", 16667)),
            ),
            # one with values outside normal range
            (
                "circularArrow",
                (
                    ("adj1", 12500),
                    ("adj2", 1142319),
                    ("adj3", 20457681),
                    ("adj4", 10800000),
                    ("adj5", 12500),
                ),
            ),
        ]

    @pytest.fixture(params=_prstGeom_cases())
    def prstGeom_cases_(self, request):
        prst, expected_values = request.param
        prstGeom = (
            a_prstGeom().with_nsdecls().with_prst(prst).with_child(an_avLst()).element
        )
        return prstGeom, prst, expected_values

    def _effective_val_cases():
        return [
            ("rect", ()),
            ("chevron", (0.5,)),
            ("circularArrow", (0.125, 11.42319, 204.57681, 108.0, 0.125)),
        ]

    @pytest.fixture(params=_effective_val_cases())
    def indexed_access_fixture_(self, request):
        prst, effective_values = request.param
        prstGeom = a_prstGeom().with_nsdecls().with_prst(prst).element
        return prstGeom, prst, list(effective_values)

    def _indexed_assignment_cases():
        return [("chevron", 0, 0.5, 50000), ("circularArrow", 4, 99.99, 9999000)]

    @pytest.fixture(params=_indexed_assignment_cases())
    def indexed_assignment_fixture_(self, request):
        prst, idx, new_val, expected = request.param
        prstGeom = a_prstGeom().with_nsdecls().with_prst(prst).element
        adjustments = AdjustmentCollection(prstGeom)
        return adjustments, idx, new_val, expected


class DescribeAutoShapeType(object):
    def it_knows_the_details_of_the_auto_shape_type_it_represents(self):
        autoshape_type = AutoShapeType(MSO_SHAPE.ROUNDED_RECTANGLE)
        assert autoshape_type.autoshape_type_id == MSO_SHAPE.ROUNDED_RECTANGLE
        assert autoshape_type.prst == "roundRect"
        assert autoshape_type.basename == "Rounded Rectangle"

    def it_knows_the_default_adj_vals_for_its_autoshape_type(
        self, default_adj_vals_fixture_
    ):
        prst, default_adj_vals = default_adj_vals_fixture_
        _default_adj_vals = AutoShapeType.default_adjustment_values(prst)
        assert _default_adj_vals == default_adj_vals

    def it_knows_the_autoshape_type_id_for_each_prst_key(self):
        assert AutoShapeType.id_from_prst("rect") == MSO_SHAPE.RECTANGLE

    def it_raises_when_asked_for_autoshape_type_id_with_a_bad_prst(self):
        with pytest.raises(KeyError):
            AutoShapeType.id_from_prst("badPrst")

    def it_caches_autoshape_type_lookups(self):
        autoshape_type_id = MSO_SHAPE.ROUNDED_RECTANGLE
        autoshape_type_1 = AutoShapeType(autoshape_type_id)
        autoshape_type_2 = AutoShapeType(autoshape_type_id)
        assert autoshape_type_2 is autoshape_type_1

    def it_raises_on_construction_with_bad_autoshape_type_id(self):
        with pytest.raises(KeyError):
            AutoShapeType(9999)

    # fixtures -------------------------------------------------------

    def _default_adj_vals_cases():
        return [
            (MSO_SHAPE.RECTANGLE, ()),
            (MSO_SHAPE.CHEVRON, (("adj", 50000),)),
            (
                MSO_SHAPE.LEFT_CIRCULAR_ARROW,
                (
                    ("adj1", 12500),
                    ("adj2", -1142319),
                    ("adj3", 1142319),
                    ("adj4", 10800000),
                    ("adj5", 12500),
                ),
            ),
        ]

    @pytest.fixture(params=_default_adj_vals_cases())
    def default_adj_vals_fixture_(self, request):
        prst, default_adj_vals = request.param
        return prst, default_adj_vals


class DescribeShape(object):
    """Unit-test suite for `pptx.shapes.autoshape.Shape` object."""

    def it_initializes_adjustments_on_first_ref(self, init_adjs_fixture_):
        shape, adjs_, AdjustmentCollection_, sp_ = init_adjs_fixture_
        assert shape.adjustments is adjs_
        AdjustmentCollection_.assert_called_once_with(sp_.prstGeom)

    def it_knows_its_autoshape_type(self, autoshape_type_fixture_):
        shape, expected_value = autoshape_type_fixture_
        auto_shape_type = shape.auto_shape_type
        assert auto_shape_type == expected_value

    def but_it_raises_when_auto_shape_type_called_on_non_autoshape(
        self, non_autoshape_shape_
    ):
        with pytest.raises(ValueError):
            non_autoshape_shape_.auto_shape_type

    def it_has_a_fill(self, shape):
        assert isinstance(shape.fill, FillFormat)

    def it_knows_it_has_a_text_frame(self):
        shape = Shape(None, None)
        assert shape.has_text_frame is True

    def it_has_a_line(self, shape):
        assert isinstance(shape.line, LineFormat)

    def it_knows_its_shape_type_when_its_a_placeholder(self, placeholder_shape_):
        assert placeholder_shape_.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER

    def and_it_knows_its_shape_type_when_its_not_a_placeholder(
        self, non_placeholder_shapes_
    ):
        autoshape_shape_, textbox_shape_, freeform_ = non_placeholder_shapes_
        assert autoshape_shape_.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        assert textbox_shape_.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        assert freeform_.shape_type == MSO_SHAPE_TYPE.FREEFORM

    def but_it_raises_when_shape_type_called_on_unrecognized_shape(self):
        xml_ = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/'
            '2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            '2006/main"><p:nvSpPr><p:cNvPr id="9" name="Unknown Shape Type 8"'
            "/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/></p:sp>"
        )
        sp = parse_xml(xml_)
        shape = Shape(sp, None)
        # verify -----------------------
        with pytest.raises(NotImplementedError):
            shape.shape_type

    def it_knows_what_text_it_contains(self, text_frame_prop_, text_frame_):
        text_frame_prop_.return_value = text_frame_
        text_frame_.text = "foobar"
        shape = Shape(None, None)

        text = shape.text

        assert text == "foobar"

    def it_can_change_its_text(self, text_frame_prop_, text_frame_):
        text_frame_prop_.return_value = text_frame_
        shape = Shape(None, None)

        shape.text = "føøbår"

        assert text_frame_.text == "føøbår"

    def it_provides_access_to_its_text_frame(self, text_frame_fixture):
        shape, text_frame_, TextFrame_, txBody = text_frame_fixture
        text_frame = shape.text_frame
        TextFrame_.assert_called_once_with(txBody, shape)
        assert text_frame is text_frame_

    def and_it_creates_a_txBody_if_needed(self, txBody_fixture):
        shape, expected_xml = txBody_fixture
        text_frame = shape.text_frame
        assert text_frame._txBody.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def autoshape_type_fixture_(self, shape, prst):
        sp = (
            an_sp()
            .with_nsdecls()
            .with_child(an_nvSpPr().with_child(a_cNvSpPr()))
            .with_child(an_spPr().with_child(a_prstGeom().with_prst("chevron")))
        ).element
        shape = Shape(sp, None)
        return shape, MSO_SHAPE.CHEVRON

    @pytest.fixture
    def init_adjs_fixture_(
        self, request, shape, sp_, adjustments_, AdjustmentCollection_
    ):
        return shape, adjustments_, AdjustmentCollection_, sp_

    @pytest.fixture
    def text_frame_fixture(self, TextFrame_, text_frame_):
        sp = element("p:sp/p:txBody")
        txBody = sp[0]
        shape = Shape(sp, None)
        return shape, text_frame_, TextFrame_, txBody

    @pytest.fixture
    def txBody_fixture(self, request):
        sp_cxml = "p:sp"
        expected_cxml = "p:txBody/(a:bodyPr,a:p)"
        shape = Shape(element(sp_cxml), None)
        expected_xml = xml(expected_cxml)
        return shape, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def AdjustmentCollection_(self, request, adjustments_):
        return class_mock(
            request,
            "pptx.shapes.autoshape.AdjustmentCollection",
            return_value=adjustments_,
        )

    @pytest.fixture
    def adjustments_(self, request):
        return instance_mock(request, AdjustmentCollection)

    @pytest.fixture
    def non_autoshape_shape_(self, request, sp_):
        sp_.is_autoshape = False
        return Shape(sp_, None)

    @pytest.fixture
    def non_placeholder_shapes_(self, request):
        autoshape_sp_ = instance_mock(
            request,
            CT_Shape,
            name="autoshape_sp_",
            is_autoshape=True,
            is_textbox=False,
            has_custom_geometry=False,
        )
        autoshape_shape_ = Shape(autoshape_sp_, None)
        textbox_sp_ = instance_mock(
            request,
            CT_Shape,
            name="textbox_sp_",
            is_autoshape=False,
            is_textbox=True,
            has_custom_geometry=False,
        )
        textbox_shape_ = Shape(textbox_sp_, None)
        freeform_sp_ = instance_mock(
            request,
            CT_Shape,
            name="freeform_sp_",
            is_autoshape=True,
            is_textbox=False,
            has_custom_geometry=True,
        )
        freeform_ = Shape(freeform_sp_, None)
        property_mock(request, Shape, "is_placeholder", return_value=False)
        return autoshape_shape_, textbox_shape_, freeform_

    @pytest.fixture
    def placeholder_shape_(self, request, sp_):
        placeholder_shape_ = Shape(sp_, None)
        property_mock(request, Shape, "is_placeholder", return_value=True)
        return placeholder_shape_

    @pytest.fixture
    def prst(self):
        return MSO_SHAPE.CHEVRON

    @pytest.fixture
    def shape(self, request, sp_):
        return Shape(sp_, None)

    @pytest.fixture
    def sp_(self, request, prst):
        return instance_mock(request, CT_Shape, prst=prst, is_autoshape=True)

    @pytest.fixture
    def TextFrame_(self, request, text_frame_):
        return class_mock(
            request, "pptx.shapes.autoshape.TextFrame", return_value=text_frame_
        )

    @pytest.fixture
    def text_frame_(self, request):
        return instance_mock(request, TextFrame)

    @pytest.fixture
    def text_frame_prop_(self, request):
        return property_mock(request, Shape, "text_frame")
