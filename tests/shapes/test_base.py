# encoding: utf-8

"""
Test suite for pptx.shapes.shape module
"""

from __future__ import absolute_import

import pytest

from pptx.action import ActionSetting
from pptx.dml.effect import ShadowFormat
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.shapes.shared import BaseShapeElement
from pptx.oxml.text import CT_TextBody
from pptx.shapes import Subshape
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape, _PlaceholderFormat
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.shapes.shapetree import BaseShapeFactory, SlideShapes

from ..oxml.unitdata.shape import (
    a_cNvPr, a_cxnSp, a_graphicFrame, a_grpSp, a_grpSpPr, a_p_xfrm, a_pic,
    an_ext, an_nvSpPr, an_off, an_sp, an_spPr, an_xfrm
)
from ..unitutil.cxml import element, xml
from ..unitutil.mock import (
    class_mock, instance_mock, loose_mock, property_mock
)


class DescribeBaseShape(object):

    def it_provides_access_to_its_click_action(self, click_action_fixture):
        shape, ActionSetting_, cNvPr, click_action_ = click_action_fixture
        click_action = shape.click_action
        ActionSetting_.assert_called_once_with(cNvPr, shape)
        assert click_action is click_action_

    def it_knows_its_shape_id(self, id_fixture):
        shape, expected_value = id_fixture
        assert shape.shape_id == expected_value

    def it_knows_its_name(self, name_get_fixture):
        shape, name = name_get_fixture
        assert shape.name == name

    def it_can_change_its_name(self, name_set_fixture):
        shape, new_value, expected_xml = name_set_fixture
        shape.name = new_value
        assert shape._element.xml == expected_xml

    def it_has_a_position(self, position_get_fixture):
        shape, expected_left, expected_top = position_get_fixture
        assert shape.left == expected_left
        assert shape.top == expected_top

    def it_can_change_its_position(self, position_set_fixture):
        shape, left, top, expected_xml = position_set_fixture
        shape.left = left
        shape.top = top
        assert shape._element.xml == expected_xml

    def it_has_dimensions(self, dimensions_get_fixture):
        shape, expected_width, expected_height = dimensions_get_fixture
        assert shape.width == expected_width
        assert shape.height == expected_height

    def it_can_change_its_dimensions(self, dimensions_set_fixture):
        shape, width, height, expected_xml = dimensions_set_fixture
        shape.width = width
        shape.height = height
        assert shape._element.xml == expected_xml

    def it_knows_its_rotation_angle(self, rotation_get_fixture):
        shape, expected_value = rotation_get_fixture
        assert shape.rotation == expected_value

    def it_can_change_its_rotation_angle(self, rotation_set_fixture):
        shape, new_value, expected_xml = rotation_set_fixture
        shape.rotation = new_value
        assert shape._element.xml == expected_xml

    def it_provides_access_to_its_shadow(self, shadow_fixture):
        shape, ShadowFormat_, spPr, shadow_ = shadow_fixture

        shadow = shape.shadow

        ShadowFormat_.assert_called_once_with(spPr)
        assert shadow is shadow_

    def it_knows_the_part_it_belongs_to(self, part_fixture):
        shape, parent_ = part_fixture
        part = shape.part
        assert part is parent_.part

    def it_knows_it_doesnt_have_a_text_frame(self):
        shape = BaseShape(None, None)
        assert shape.has_text_frame is False

    def it_knows_whether_it_is_a_placeholder(self, is_placeholder_fixture):
        shape, is_placeholder = is_placeholder_fixture
        assert shape.is_placeholder is is_placeholder

    def it_provides_access_to_its_placeholder_format(self, phfmt_fixture):
        shape, _PlaceholderFormat_, placeholder_format_, ph = phfmt_fixture
        placeholder_format = shape.placeholder_format
        _PlaceholderFormat_.assert_called_once_with(ph)
        assert placeholder_format is placeholder_format_

    def it_raises_when_shape_is_not_a_placeholder(self, phfmt_raise_fixture):
        shape = phfmt_raise_fixture
        with pytest.raises(ValueError):
            shape.placeholder_format

    def it_knows_it_doesnt_contain_a_chart(self):
        shape = BaseShape(None, None)
        assert shape.has_chart is False

    def it_knows_it_doesnt_contain_a_table(self):
        shape = BaseShape(None, None)
        assert shape.has_table is False

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        'p:sp/p:nvSpPr/p:cNvPr',
        'p:grpSp/p:nvGrpSpPr/p:cNvPr',
        'p:graphicFrame/p:nvGraphicFramePr/p:cNvPr',
        'p:cxnSp/p:nvCxnSpPr/p:cNvPr',
        'p:pic/p:nvPicPr/p:cNvPr',
    ])
    def click_action_fixture(self, request, ActionSetting_, action_setting_):
        sp_cxml = request.param
        sp = element(sp_cxml)
        cNvPr = sp.xpath('//p:cNvPr')[0]
        shape = BaseShape(sp, None)
        return shape, ActionSetting_, cNvPr, action_setting_

    @pytest.fixture(params=[
        ('sp',           False), ('sp_with_ext',           True),
        ('pic',          False), ('pic_with_ext',          True),
        ('graphicFrame', False), ('graphicFrame_with_ext', True),
        ('grpSp',        False), ('grpSp_with_ext',        True),
        ('cxnSp',        False), ('cxnSp_with_ext',        True),
    ])
    def dimensions_get_fixture(self, request, width, height):
        shape_elm_fixt_name, expect_values = request.param
        shape_elm = request.getfuncargvalue(shape_elm_fixt_name)
        shape = BaseShape(shape_elm, None)
        if not expect_values:
            width = height = None
        return shape, width, height

    @pytest.fixture(params=[
        ('sp',           'sp_with_ext'),
        ('pic',          'pic_with_ext'),
        ('graphicFrame', 'graphicFrame_with_ext'),
        ('grpSp',        'grpSp_with_ext'),
        ('cxnSp',        'cxnSp_with_ext'),
    ])
    def dimensions_set_fixture(self, request, width, height):
        start_elm_fixt_name, expected_elm_fixt_name = request.param
        start_elm = request.getfuncargvalue(start_elm_fixt_name)
        shape = BaseShape(start_elm, None)
        expected_xml = request.getfuncargvalue(expected_elm_fixt_name).xml
        return shape, width, height, expected_xml

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:cNvPr{id=1}', 1),
        ('p:cxnSp/p:nvCxnSpPr/p:cNvPr{id=2}', 2),
        ('p:graphicFrame/p:nvGraphicFramePr/p:cNvPr{id=3}', 3),
        ('p:grpSp/p:nvGrpSpPr/p:cNvPr{id=4}', 4),
        ('p:pic/p:nvPicPr/p:cNvPr{id=5}', 5),
    ])
    def id_fixture(self, request):
        xSp_cxml, expected_value = request.param
        shape = BaseShape(element(xSp_cxml), None)
        return shape, expected_value

    @pytest.fixture(params=[True, False])
    def is_placeholder_fixture(self, request, shape_elm_, txBody_):
        is_placeholder = request.param
        shape_elm_.has_ph_elm = is_placeholder
        shape = BaseShape(shape_elm_, None)
        return shape, is_placeholder

    @pytest.fixture
    def name_get_fixture(self, shape_name):
        shape_elm = (
            an_sp().with_nsdecls().with_child(
                an_nvSpPr().with_child(
                    a_cNvPr().with_name(shape_name)))
        ).element
        shape = BaseShape(shape_elm, None)
        return shape, shape_name

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:cNvPr{id=1,name=foo}',       Shape,     'Shape1',
         'p:sp/p:nvSpPr/p:cNvPr{id=1,name=Shape1}'),
        ('p:grpSp/p:nvGrpSpPr/p:cNvPr{id=2,name=bar}', BaseShape, 'Shape2',
         'p:grpSp/p:nvGrpSpPr/p:cNvPr{id=2,name=Shape2}'),
        ('p:graphicFrame/p:nvGraphicFramePr/p:cNvPr{id=3,name=baz}',
         GraphicFrame, 'Shape3',
         'p:graphicFrame/p:nvGraphicFramePr/p:cNvPr{id=3,name=Shape3}'),
        ('p:cxnSp/p:nvCxnSpPr/p:cNvPr{id=4,name=boo}', BaseShape, 'Shape4',
         'p:cxnSp/p:nvCxnSpPr/p:cNvPr{id=4,name=Shape4}'),
        ('p:pic/p:nvPicPr/p:cNvPr{id=5,name=far}',     Picture,   'Shape5',
         'p:pic/p:nvPicPr/p:cNvPr{id=5,name=Shape5}'),
    ])
    def name_set_fixture(self, request):
        xSp_cxml, ShapeCls, new_value, expected_xSp_cxml = request.param
        shape = ShapeCls(element(xSp_cxml), None)
        expected_xml = xml(expected_xSp_cxml)
        return shape, new_value, expected_xml

    @pytest.fixture
    def part_fixture(self, shapes_):
        parent_ = shapes_
        shape = BaseShape(None, parent_)
        return shape, parent_

    @pytest.fixture
    def phfmt_fixture(self, _PlaceholderFormat_, placeholder_format_):
        sp = element('p:sp/p:nvSpPr/p:nvPr/p:ph')
        ph = sp.xpath('//p:ph')[0]
        shape = BaseShape(sp, None)
        return shape, _PlaceholderFormat_, placeholder_format_, ph

    @pytest.fixture
    def phfmt_raise_fixture(self):
        return BaseShape(element('p:sp/p:nvSpPr/p:nvPr'), None)

    @pytest.fixture(params=[
        ('sp',           False), ('sp_with_off',           True),
        ('pic',          False), ('pic_with_off',          True),
        ('graphicFrame', False), ('graphicFrame_with_off', True),
        ('grpSp',        False), ('grpSp_with_off',        True),
        ('cxnSp',        False), ('cxnSp_with_off',        True),
    ])
    def position_get_fixture(self, request, left, top):
        shape_elm_fixt_name, expect_values = request.param
        shape_elm = request.getfuncargvalue(shape_elm_fixt_name)
        shape = BaseShape(shape_elm, None)
        if not expect_values:
            left = top = None
        return shape, left, top

    @pytest.fixture(params=[
        ('sp',           'sp_with_off'),
        ('pic',          'pic_with_off'),
        ('graphicFrame', 'graphicFrame_with_off'),
        ('grpSp',        'grpSp_with_off'),
        ('cxnSp',        'cxnSp_with_off'),
    ])
    def position_set_fixture(self, request, left, top):
        start_elm_fixt_name, expected_elm_fixt_name = request.param
        start_elm = request.getfuncargvalue(start_elm_fixt_name)
        shape = BaseShape(start_elm, None)
        expected_xml = request.getfuncargvalue(expected_elm_fixt_name).xml
        return shape, left, top, expected_xml

    @pytest.fixture(params=[
        ('p:sp/p:spPr', 0.0),
        ('p:sp/p:spPr/a:xfrm{rot=60000}', 1.0),
        ('p:sp/p:spPr/a:xfrm{rot=2545200}', 42.42),
        ('p:sp/p:spPr/a:xfrm{rot=-60000}', 359.0),
        ('p:grpSp/p:grpSpPr/a:xfrm{rot=2545200}', 42.42),
    ])
    def rotation_get_fixture(self, request):
        xSp_cxml, expected_value = request.param
        shape = BaseShapeFactory(element(xSp_cxml), None)
        return shape, expected_value

    @pytest.fixture(params=[
        ('p:sp/p:spPr/a:xfrm', 1.0,
         'p:sp/p:spPr/a:xfrm{rot=60000}'),
        ('p:sp/p:spPr/a:xfrm{rot=60000}', 0.0,
         'p:sp/p:spPr/a:xfrm'),
        ('p:sp/p:spPr/a:xfrm{rot=60000}', -420.0,
         'p:sp/p:spPr/a:xfrm{rot=18000000}'),
        ('p:grpSp/p:grpSpPr/a:xfrm', 1.0,
         'p:grpSp/p:grpSpPr/a:xfrm{rot=60000}'),
    ])
    def rotation_set_fixture(self, request):
        xSp_cxml, new_value, expected_xSp_cxml = request.param
        shape = BaseShapeFactory(element(xSp_cxml), None)
        expected_xml = xml(expected_xSp_cxml)
        return shape, new_value, expected_xml

    @pytest.fixture(params=[
        'p:sp/p:spPr',
        'p:cxnSp/p:spPr',
        'p:pic/p:spPr',
        # ---group and graphic frame shapes override this property---
    ])
    def shadow_fixture(self, request, ShadowFormat_, shadow_):
        sp_cxml = request.param
        sp = element(sp_cxml)
        spPr = sp.xpath('//p:spPr')[0]
        ShadowFormat_.return_value = shadow_

        shape = BaseShape(sp, None)
        return shape, ShadowFormat_, spPr, shadow_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def ActionSetting_(self, request, action_setting_):
        return class_mock(
            request, 'pptx.shapes.base.ActionSetting',
            return_value=action_setting_
        )

    @pytest.fixture
    def action_setting_(self, request):
        return instance_mock(request, ActionSetting)

    @pytest.fixture
    def cxnSp(self):
        return a_cxnSp().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def cxnSp_with_ext(self, width, height):
        return (
            a_cxnSp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def cxnSp_with_off(self, left, top):
        return (
            a_cxnSp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def graphicFrame(self):
        # Note that <p:xfrm> element is required on graphicFrame
        return a_graphicFrame().with_nsdecls().with_child(a_p_xfrm()).element

    @pytest.fixture
    def graphicFrame_with_ext(self, width, height):
        return (
            a_graphicFrame().with_nsdecls().with_child(
                a_p_xfrm().with_child(
                    an_ext().with_cx(width).with_cy(height)))
        ).element

    @pytest.fixture
    def graphicFrame_with_off(self, left, top):
        return (
            a_graphicFrame().with_nsdecls().with_child(
                a_p_xfrm().with_child(
                    an_off().with_x(left).with_y(top)))
        ).element

    @pytest.fixture
    def grpSp(self):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr())
        ).element

    @pytest.fixture
    def grpSp_with_ext(self, width, height):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def grpSp_with_off(self, left, top):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def height(self):
        return 654

    @pytest.fixture
    def left(self):
        return 123

    @pytest.fixture
    def pic(self):
        return a_pic().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def pic_with_off(self, left, top):
        return (
            a_pic().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def pic_with_ext(self, width, height):
        return (
            a_pic().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def _PlaceholderFormat_(self, request, placeholder_format_):
        return class_mock(
            request, 'pptx.shapes.base._PlaceholderFormat',
            return_value=placeholder_format_
        )

    @pytest.fixture
    def placeholder_format_(self, request):
        return instance_mock(request, _PlaceholderFormat)

    @pytest.fixture
    def shadow_(self, request):
        return instance_mock(request, ShadowFormat)

    @pytest.fixture
    def ShadowFormat_(self, request):
        return class_mock(request, 'pptx.shapes.base.ShadowFormat')

    @pytest.fixture
    def shape_elm_(self, request, shape_id, shape_name, txBody_):
        return instance_mock(
            request, BaseShapeElement, shape_id=shape_id,
            shape_name=shape_name, txBody=txBody_
        )

    @pytest.fixture
    def shape_id(self):
        return 42

    @pytest.fixture
    def shape_name(self):
        return 'Foobar 41'

    @pytest.fixture
    def shape_text_frame_(self, request):
        return property_mock(request, BaseShape, 'text_frame')

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, SlideShapes)

    @pytest.fixture
    def sp(self):
        return an_sp().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def sp_with_ext(self, width, height):
        return (
            an_sp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def sp_with_off(self, left, top):
        return (
            an_sp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def top(self):
        return 456

    @pytest.fixture
    def txBody_(self, request):
        return instance_mock(request, CT_TextBody)

    @pytest.fixture
    def width(self):
        return 321


class DescribeSubshape(object):

    def it_knows_the_part_it_belongs_to(self, subshape_with_parent_):
        subshape, parent_ = subshape_with_parent_
        part = subshape.part
        assert part is parent_.part

    # fixtures ---------------------------------------------

    @pytest.fixture
    def subshape_with_parent_(self, request):
        parent_ = loose_mock(request, name='parent_')
        subshape = Subshape(parent_)
        return subshape, parent_


class Describe_PlaceholderFormat(object):

    def it_knows_its_idx(self, idx_get_fixture):
        placeholder_format, expected_value = idx_get_fixture
        assert placeholder_format.idx == expected_value

    def it_knows_its_type(self, type_get_fixture):
        placeholder_format, expected_value = type_get_fixture
        assert placeholder_format.type == expected_value

    # fixtures ---------------------------------------------

    @pytest.fixture(params=[
        ('p:ph',          0),
        ('p:ph{idx=42}', 42),
    ])
    def idx_get_fixture(self, request):
        ph_cxml, expected_value = request.param
        placeholder_format = _PlaceholderFormat(element(ph_cxml))
        return placeholder_format, expected_value

    @pytest.fixture(params=[
        ('p:ph',           PP_PLACEHOLDER.OBJECT),
        ('p:ph{type=pic}', PP_PLACEHOLDER.PICTURE),
    ])
    def type_get_fixture(self, request):
        ph_cxml, expected_value = request.param
        placeholder_format = _PlaceholderFormat(element(ph_cxml))
        return placeholder_format, expected_value
