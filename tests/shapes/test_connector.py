# encoding: utf-8

"""Unit test suite for pptx.shapes.connector module."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.dml.line import LineFormat
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.util import Emu

from ..unitutil.cxml import element, xml
from ..unitutil.mock import instance_mock, method_mock


class DescribeConnector(object):

    def it_knows_its_begin_point_x_location(self, begin_x_get_fixture):
        connector, expected_value = begin_x_get_fixture
        begin_x = connector.begin_x
        assert isinstance(begin_x, Emu)
        assert connector.begin_x == expected_value

    def it_can_change_its_begin_point_x_location(self, begin_x_set_fixture):
        connector, new_x, expected_xml = begin_x_set_fixture
        connector.begin_x = new_x
        assert connector._element.xml == expected_xml

    def it_knows_its_begin_point_y_location(self, begin_y_get_fixture):
        connector, expected_value = begin_y_get_fixture
        begin_y = connector.begin_y
        assert isinstance(begin_y, Emu)
        assert connector.begin_y == expected_value

    def it_can_change_its_begin_point_y_location(self, begin_y_set_fixture):
        connector, new_y, expected_xml = begin_y_set_fixture
        connector.begin_y = new_y
        assert connector._element.xml == expected_xml

    def it_knows_its_end_point_x_location(self, end_x_get_fixture):
        connector, expected_value = end_x_get_fixture
        end_x = connector.end_x
        assert isinstance(end_x, Emu)
        assert connector.end_x == expected_value

    def it_can_change_its_end_point_x_location(self, end_x_set_fixture):
        connector, new_x, expected_xml = end_x_set_fixture
        connector.end_x = new_x
        assert connector._element.xml == expected_xml

    def it_knows_its_end_point_y_location(self, end_y_get_fixture):
        connector, expected_value = end_y_get_fixture
        end_y = connector.end_y
        assert isinstance(end_y, Emu)
        assert connector.end_y == expected_value

    def it_can_change_its_end_point_y_location(self, end_y_set_fixture):
        connector, new_y, expected_xml = end_y_set_fixture
        connector.end_y = new_y
        assert connector._element.xml == expected_xml

    def it_can_connect_its_begin_point_to_a_shape(self, begin_conn_fixture):
        connector, shape, cxn_idx = begin_conn_fixture

        connector.begin_connect(shape, cxn_idx)

        connector._connect_begin_to.assert_called_once_with(
            connector, shape, cxn_idx
        )
        connector._move_begin_to_cxn.assert_called_once_with(
            connector, shape, cxn_idx
        )

    def it_connects_its_begin_point_to_help(self, connect_begin_fixture):
        connector, shape, cxn_idx, expected_xml = connect_begin_fixture
        connector._connect_begin_to(shape, cxn_idx)
        assert connector._element.xml == expected_xml

    def it_moves_its_begin_point_to_help(self, move_begin_fixture):
        connector, shape, cxn_idx, expected_xml = move_begin_fixture
        connector._move_begin_to_cxn(shape, cxn_idx)
        assert connector._element.xml == expected_xml

    def it_can_connect_its_end_point_to_a_shape(self, end_conn_fixture):
        connector, shape, cxn_idx = end_conn_fixture

        connector.end_connect(shape, cxn_idx)

        connector._connect_end_to.assert_called_once_with(
            connector, shape, cxn_idx
        )
        connector._move_end_to_cxn.assert_called_once_with(
            connector, shape, cxn_idx
        )

    def it_connects_its_end_point_to_help(self, connect_end_fixture):
        connector, shape, cxn_idx, expected_xml = connect_end_fixture
        connector._connect_end_to(shape, cxn_idx)
        assert connector._element.xml == expected_xml

    def it_moves_its_end_point_to_help(self, move_end_fixture):
        connector, shape, cxn_idx, expected_xml = move_end_fixture
        connector._move_end_to_cxn(shape, cxn_idx)
        assert connector._element.xml == expected_xml

    def it_provides_access_to_its_line_format(self):
        connector = Connector(element('p:cxnSp/p:spPr'), None)

        line = connector.line

        assert isinstance(line, LineFormat)
        # exercise line to test parent interface, .ln and .get_or_add_ln()
        line.width = 91440
        assert line.width == 91440

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def begin_conn_fixture(self, _connect_begin_to_, _move_begin_to_cxn_):
        connector = Connector(None, None)
        shape, cxn_idx = 42, 24
        return connector, shape, cxn_idx

    @pytest.fixture(params=[
        (42, 24, False, 42),
        (24, 42, True,  66),
    ])
    def begin_x_get_fixture(self, request):
        x, cx, flipH, expected_value = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm{flipH=%d}/(a:off{x=%d,y=6},a:ext{cx=%d,cy'
            '=32})' % (flipH, x, cx)
        )
        connector = Connector(cxnSp, None)
        return connector, expected_value

    @pytest.fixture(params=[
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',           5,
         'a:xfrm/(a:off{x=5,y=1},a:ext{cx=15,cy=1})'),
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',          15,
         'a:xfrm/(a:off{x=15,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',          25,
         'a:xfrm{flipH=1}/(a:off{x=20,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})', 25,
         'a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=15,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})', 15,
         'a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',  5,
         'a:xfrm/(a:off{x=5,y=1},a:ext{cx=5,cy=1})'),
    ])
    def begin_x_set_fixture(self, request):
        xfrm_cxml, new_x, expected_cxml = request.param
        tmpl = 'p:cxnSp/p:spPr/%s'
        cxnSp = element(tmpl % xfrm_cxml)
        expected_xml = xml(tmpl % expected_cxml)
        connector = Connector(cxnSp, None)
        return connector, new_x, expected_xml

    @pytest.fixture(params=[
        (40, 60, False, 40),
        (50, 42, True,  92),
    ])
    def begin_y_get_fixture(self, request):
        y, cy, flipV, expected_value = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm{flipV=%d}/(a:off{x=6,y=%d},a:ext{cx=32,cy'
            '=%d})' % (flipV, y, cy)
        )
        connector = Connector(cxnSp, None)
        return connector, expected_value

    @pytest.fixture(params=[
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',           5,
         'a:xfrm/(a:off{x=1,y=5},a:ext{cx=1,cy=15})'),
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',          15,
         'a:xfrm/(a:off{x=1,y=15},a:ext{cx=1,cy=5})'),
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',          25,
         'a:xfrm{flipV=1}/(a:off{x=1,y=20},a:ext{cx=1,cy=5})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})', 30,
         'a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=20})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})', 15,
         'a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=5})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',  5,
         'a:xfrm/(a:off{x=1,y=5},a:ext{cx=1,cy=5})'),
    ])
    def begin_y_set_fixture(self, request):
        xfrm_cxml, new_y, expected_cxml = request.param
        tmpl = 'p:cxnSp/p:spPr/%s'
        cxnSp = element(tmpl % xfrm_cxml)
        expected_xml = xml(tmpl % expected_cxml)
        connector = Connector(cxnSp, None)
        return connector, new_y, expected_xml

    @pytest.fixture(params=[
        ('p:cxnSp{a:b=c}/p:nvCxnSpPr/p:cNvCxnSpPr',
         'p:cxnSp{a:b=c}/p:nvCxnSpPr/p:cNvCxnSpPr/a:stCxn{id=42,idx=3}'),
        ('p:cxnSp/p:nvCxnSpPr/p:cNvCxnSpPr/a:stCxn{id=66,idx=6}',
         'p:cxnSp/p:nvCxnSpPr/p:cNvCxnSpPr/a:stCxn{id=42,idx=3}'),
    ])
    def connect_begin_fixture(self, request, shape_):
        cxnSp_cxml, expected_cxml = request.param
        cxnSp = element(cxnSp_cxml)
        connector = Connector(cxnSp, None)
        shape_.shape_id, cxn_idx = 42, 3
        expected_xml = xml(expected_cxml)
        return connector, shape_, cxn_idx, expected_xml

    @pytest.fixture(params=[
        ('p:cxnSp{a:b=c}/p:nvCxnSpPr/p:cNvCxnSpPr',
         'p:cxnSp{a:b=c}/p:nvCxnSpPr/p:cNvCxnSpPr/a:endCxn{id=24,idx=2}'),
        ('p:cxnSp/p:nvCxnSpPr/p:cNvCxnSpPr/a:endCxn{id=66,idx=6}',
         'p:cxnSp/p:nvCxnSpPr/p:cNvCxnSpPr/a:endCxn{id=24,idx=2}'),
    ])
    def connect_end_fixture(self, request, shape_):
        cxnSp_cxml, expected_cxml = request.param
        cxnSp = element(cxnSp_cxml)
        connector = Connector(cxnSp, None)
        shape_.shape_id, cxn_idx = 24, 2
        expected_xml = xml(expected_cxml)
        return connector, shape_, cxn_idx, expected_xml

    @pytest.fixture
    def end_conn_fixture(self, _connect_end_to_, _move_end_to_cxn_):
        connector = Connector(None, None)
        shape, cxn_idx = 42, 24
        return connector, shape, cxn_idx

    @pytest.fixture(params=[
        (21, 32, False, 53),
        (43, 54, True,  43),
    ])
    def end_x_get_fixture(self, request):
        x, cx, flipH, expected_value = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm{flipH=%d}/(a:off{x=%d,y=6},a:ext{cx=%d,cy'
            '=60})' % (flipH, x, cx)
        )
        connector = Connector(cxnSp, None)
        return connector, expected_value

    @pytest.fixture(params=[
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',          32,
         'a:xfrm/(a:off{x=10,y=1},a:ext{cx=22,cy=1})'),
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',          15,
         'a:xfrm/(a:off{x=10,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',           5,
         'a:xfrm{flipH=1}/(a:off{x=5,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})',  5,
         'a:xfrm{flipH=1}/(a:off{x=5,y=1},a:ext{cx=15,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})', 15,
         'a:xfrm{flipH=1}/(a:off{x=15,y=1},a:ext{cx=5,cy=1})'),
        ('a:xfrm{flipH=1}/(a:off{x=10,y=1},a:ext{cx=10,cy=1})', 28,
         'a:xfrm/(a:off{x=20,y=1},a:ext{cx=8,cy=1})'),
    ])
    def end_x_set_fixture(self, request):
        xfrm_cxml, new_x, expected_cxml = request.param
        tmpl = 'p:cxnSp/p:spPr/%s'
        cxnSp = element(tmpl % xfrm_cxml)
        expected_xml = xml(tmpl % expected_cxml)
        connector = Connector(cxnSp, None)
        return connector, new_x, expected_xml

    @pytest.fixture(params=[
        (31, 42, False, 73),
        (53, 14, True,  53),
    ])
    def end_y_get_fixture(self, request):
        y, cy, flipV, expected_value = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm{flipV=%d}/(a:off{x=6,y=%d},a:ext{cx=32,cy'
            '=%d})' % (flipV, y, cy)
        )
        connector = Connector(cxnSp, None)
        return connector, expected_value

    @pytest.fixture(params=[
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',          28,
         'a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=18})'),
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',          13,
         'a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=3})'),
        ('a:xfrm/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',           4,
         'a:xfrm{flipV=1}/(a:off{x=1,y=4},a:ext{cx=1,cy=6})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})',  6,
         'a:xfrm{flipV=1}/(a:off{x=1,y=6},a:ext{cx=1,cy=14})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})', 12,
         'a:xfrm{flipV=1}/(a:off{x=1,y=12},a:ext{cx=1,cy=8})'),
        ('a:xfrm{flipV=1}/(a:off{x=1,y=10},a:ext{cx=1,cy=10})', 27,
         'a:xfrm/(a:off{x=1,y=20},a:ext{cx=1,cy=7})'),
    ])
    def end_y_set_fixture(self, request):
        xfrm_cxml, new_y, expected_cxml = request.param
        tmpl = 'p:cxnSp/p:spPr/%s'
        cxnSp = element(tmpl % xfrm_cxml)
        expected_xml = xml(tmpl % expected_cxml)
        connector = Connector(cxnSp, None)
        return connector, new_y, expected_xml

    @pytest.fixture(params=[
        (0, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=25,y=15},a:ext{cx=74,cy=123})'),
        (1, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=33},a:ext{cx=89,cy=105})'),
        (2, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=25,y=51},a:ext{cx=74,cy=87})'),
        (3, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=40,y=33},a:ext{cx=59,cy=105})'),
    ])
    def move_begin_fixture(self, request, shape_):
        cxn_idx, expected_cxml = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm/(a:off{x=66,y=99},a:ext{cx=33,cy=39})'
        )
        connector = Connector(cxnSp, None)
        shape_.left, shape_.top, shape_.width, shape_.height = 10, 15, 30, 36
        expected_xml = xml(expected_cxml)
        return connector, shape_, cxn_idx, expected_xml

    @pytest.fixture(params=[
        (0, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=15},a:ext{cx=50,cy=10})'),
        (1, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=15},a:ext{cx=40,cy=19})'),
        (2, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=15},a:ext{cx=50,cy=28})'),
        (3, 'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=15},a:ext{cx=60,cy=19})'),
    ])
    def move_end_fixture(self, request, shape_):
        cxn_idx, expected_cxml = request.param
        cxnSp = element(
            'p:cxnSp/p:spPr/a:xfrm/(a:off{x=10,y=15},a:ext{cx=10,cy=5})'
        )
        connector = Connector(cxnSp, None)
        shape_.left, shape_.top, shape_.width, shape_.height = 50, 25, 20, 18
        expected_xml = xml(expected_cxml)
        return connector, shape_, cxn_idx, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _connect_begin_to_(self, request):
        return method_mock(
            request, Connector, '_connect_begin_to', autospec=True
        )

    @pytest.fixture
    def _connect_end_to_(self, request):
        return method_mock(
            request, Connector, '_connect_end_to', autospec=True
        )

    @pytest.fixture
    def _move_begin_to_cxn_(self, request):
        return method_mock(
            request, Connector, '_move_begin_to_cxn', autospec=True
        )

    @pytest.fixture
    def _move_end_to_cxn_(self, request):
        return method_mock(
            request, Connector, '_move_end_to_cxn', autospec=True
        )

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, BaseShape)
