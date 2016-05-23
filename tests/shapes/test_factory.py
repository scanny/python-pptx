# encoding: utf-8

"""
Test suite for pptx.shapes.shapetree module
"""

from __future__ import absolute_import

import pytest

from pptx.parts.slide import Slide
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.factory import (
    BaseShapeFactory, _SlidePlaceholderFactory, SlidePlaceholders,
    SlideShapeFactory
)
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import _BaseSlidePlaceholder
from pptx.shapes.shapetree import BaseShapeTree

from ..oxml.unitdata.shape import a_cNvPr, an_nvSpPr, an_spTree
from ..unitutil.cxml import element
from ..unitutil.mock import call, class_mock, function_mock, instance_mock


class DescribeBaseShapeFactory(object):

    def it_constructs_the_appropriate_shape_instance_for_a_shape_element(
            self, factory_fixture):
        shape_elm, parent_, ShapeClass_, shape_ = factory_fixture
        shape = BaseShapeFactory(shape_elm, parent_)
        ShapeClass_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    def it_finds_an_unused_shape_id_to_help_add_shape(self, next_id_fixture):
        shapes, next_available_shape_id = next_id_fixture
        shape_id = shapes._next_shape_id
        assert shape_id == next_available_shape_id

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'])
    def factory_fixture(
            self, request, slide_, Shape_, shape_, Picture_, picture_,
            GraphicFrame_, graphic_frame_, BaseShape_, base_shape_):
        shape_cxml, ShapeClass_, shape_mock = {
            'sp':           ('p:sp',           Shape_,        shape_),
            'pic':          ('p:pic',          Picture_,      picture_),
            'graphicFrame': ('p:graphicFrame', GraphicFrame_, graphic_frame_),
            'grpSp':        ('p:grpSp',        BaseShape_,    base_shape_),
            'cxnSp':        ('p:cxnSp',        BaseShape_,    base_shape_),
        }[request.param]
        shape_elm = element(shape_cxml)
        return shape_elm, slide_, ShapeClass_, shape_mock

    @pytest.fixture(params=[
        ((), 1), ((0,), 1), ((1,), 2), ((2,), 1), ((1, 3,), 2),
        (('foobar', 0, 1, 7), 2), (('1foo', 2, 2, 2), 1), ((1, 1, 1, 4), 2),
    ])
    def next_id_fixture(self, request, slide_):
        used_ids, next_available_shape_id = request.param
        nvSpPr_bldr = an_nvSpPr()
        for used_id in used_ids:
            nvSpPr_bldr.with_child(a_cNvPr().with_id(used_id))
        spTree = an_spTree().with_nsdecls().with_child(nvSpPr_bldr).element
        print(spTree.xml)
        slide_.spTree = spTree
        shapes = BaseShapeTree(slide_)
        return shapes, next_available_shape_id

    # fixture components -----------------------------------

    @pytest.fixture
    def BaseShape_(self, request, base_shape_):
        return class_mock(
            request, 'pptx.shapes.factory.BaseShape',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def GraphicFrame_(self, request, graphic_frame_):
        return class_mock(
            request, 'pptx.shapes.factory.GraphicFrame',
            return_value=graphic_frame_
        )

    @pytest.fixture
    def graphic_frame_(self, request):
        return instance_mock(request, GraphicFrame)

    @pytest.fixture
    def Picture_(self, request, picture_):
        return class_mock(
            request, 'pptx.shapes.factory.Picture', return_value=picture_
        )

    @pytest.fixture
    def picture_(self, request):
        return instance_mock(request, Picture)

    @pytest.fixture
    def Shape_(self, request, shape_):
        return class_mock(
            request, 'pptx.shapes.factory.Shape', return_value=shape_
        )

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)


class DescribeSlidePlaceholders(object):

    def it_can_get_a_placeholder_by_idx(self, getitem_fixture):
        placeholders, idx, SlideShapeFactory_ = getitem_fixture[:3]
        shape_elm, placeholder_ = getitem_fixture[3:]

        placeholder = placeholders[idx]

        SlideShapeFactory_.assert_called_once_with(shape_elm, placeholders)
        assert placeholder is placeholder_

    def it_can_iterate_over_its_placeholders(self, iter_fixture):
        placeholders, SlideShapeFactory_ = iter_fixture[:2]
        expected_calls, expected_values = iter_fixture[2:]

        ps = [p for p in placeholders]

        assert SlideShapeFactory_.call_args_list == expected_calls
        assert ps == expected_values

    def it_knows_how_many_placeholders_it_contains(self, len_fixture):
        placeholders, expected_value = len_fixture
        assert len(placeholders) == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}',    1, 0),
        ('p:spTree/p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1}',  1, 0),
        ('p:spTree/(p:sp,p:sp/p:nvSpPr/p:nvPr/p:ph{type=title})', 0, 1),
        ('p:spTree/(p:sp,p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1})',
         1, 1),
        ('p:spTree/(p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=3})', 3, 1),
    ])
    def getitem_fixture(self, request, SlideShapeFactory_, placeholder_):
        spTree_cxml, idx, offset = request.param
        spTree = element(spTree_cxml)
        placeholders = SlidePlaceholders(spTree, None)
        shape_elm = spTree[offset]
        SlideShapeFactory_.return_value = placeholder_
        return placeholders, idx, SlideShapeFactory_, shape_elm, placeholder_

    @pytest.fixture(params=[
        ('p:spTree/('
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=body,idx=1},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=3})', (1, 0, 2)),
    ])
    def iter_fixture(self, request, SlideShapeFactory_, placeholder_):
        spTree_cxml, sequence = request.param
        spTree = element(spTree_cxml)
        placeholders = SlidePlaceholders(spTree, None)
        SlideShapeFactory_.return_value = placeholder_
        calls = [call(spTree[i], placeholders) for i in sequence]
        values = [placeholder_] * len(sequence)
        return placeholders, SlideShapeFactory_, calls, values

    @pytest.fixture(params=[
        ('p:spTree',                                                    0),
        ('p:spTree/(p:sp,p:pic,p:sp)',                                  0),
        ('p:spTree/(p:sp,p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},p:pic)', 1),
        ('p:spTree/('
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=body,idx=1},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=3})',                3),
    ])
    def len_fixture(self, request):
        spTree_cxml, length = request.param
        placeholders = SlidePlaceholders(element(spTree_cxml), None)
        return placeholders, length

    # fixture components ---------------------------------------------

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def SlideShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.factory.SlideShapeFactory',
            return_value=placeholder_
        )


class DescribeSlideShapeFactory(object):

    def it_constructs_the_right_type_of_shape(self, factory_fixture):
        element, parent_, Constructor_, shape_ = factory_fixture
        shape = SlideShapeFactory(element, parent_)
        Constructor_.assert_called_once_with(element, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=title}',      'spf'),
        ('p:pic/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}', 'spf'),
        ('p:sp',                                       'bsf'),
        ('p:pic',                                      'bsf'),
    ])
    def factory_fixture(self, request, _SlidePlaceholderFactory_,
                        placeholder_, BaseShapeFactory_, base_shape_):
        shape_cxml, shape_type = request.param
        shape_elm = element(shape_cxml)
        Constructor_, shape_ = {
            'spf': (_SlidePlaceholderFactory_, placeholder_),
            'bsf': (BaseShapeFactory_,         base_shape_),
        }[shape_type]
        return shape_elm, 42, Constructor_, shape_

    # fixture components -----------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, base_shape_):
        return function_mock(
            request, 'pptx.shapes.factory.BaseShapeFactory',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def _SlidePlaceholderFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.factory._SlidePlaceholderFactory',
            return_value=placeholder_
        )


class Describe_SlidePlaceholderFactory(object):

    def it_constructs_the_right_type_of_placeholder(self, factory_fixture):
        element, parent_, Constructor_, placeholder_ = factory_fixture
        placeholder = _SlidePlaceholderFactory(element, parent_)
        Constructor_.assert_called_once_with(element, parent_)
        assert placeholder is placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=title}',                 'sp'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}',             'pph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=clipArt,idx=1}',         'pph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=tbl,idx=1}',             'tph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=chart,idx=10}',          'cph'),
        ('p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1}',           'php'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=tbl,idx=2}',   'phgf'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=chart,idx=2}', 'phgf'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=dgm,idx=2}',   'phgf'),
    ])
    def factory_fixture(
            self, request, SlidePlaceholder_, ChartPlaceholder_,
            PicturePlaceholder_, TablePlaceholder_, PlaceholderGraphicFrame_,
            PlaceholderPicture_, placeholder_):
        shape_cxml, constructor_key = request.param
        shape_elm = element(shape_cxml)
        Constructor_, shape_ = {
            'sp':   (SlidePlaceholder_,        placeholder_),
            'cph':  (ChartPlaceholder_,        placeholder_),
            'pph':  (PicturePlaceholder_,      placeholder_),
            'tph':  (TablePlaceholder_,        placeholder_),
            'phgf': (PlaceholderGraphicFrame_, placeholder_),
            'php':  (PlaceholderPicture_,      placeholder_),
        }[constructor_key]
        return shape_elm, 42, Constructor_, shape_

    # fixture components -----------------------------------

    @pytest.fixture
    def ChartPlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.ChartPlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def PicturePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.PicturePlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def PlaceholderGraphicFrame_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.PlaceholderGraphicFrame',
            return_value=placeholder_
        )

    @pytest.fixture
    def PlaceholderPicture_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.PlaceholderPicture',
            return_value=placeholder_
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def SlidePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.SlidePlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def TablePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.factory.TablePlaceholder',
            return_value=placeholder_
        )
