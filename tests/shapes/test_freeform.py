# pyright: reportPrivateUsage=false

"""Unit-test suite for `pptx.shapes.freeform` module"""

from __future__ import annotations

import pytest

from pptx.shapes.autoshape import Shape
from pptx.shapes.freeform import (
    FreeformBuilder,
    _BaseDrawingOperation,
    _Close,
    _LineSegment,
    _MoveTo,
)
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Emu, Mm

from ..unitutil.cxml import element, xml
from ..unitutil.file import snippet_seq
from ..unitutil.mock import (
    FixtureRequest,
    Mock,
    call,
    initializer_mock,
    instance_mock,
    method_mock,
    property_mock,
)


class DescribeFreeformBuilder(object):
    """Unit-test suite for `pptx.shapes.freeform.FreeformBuilder` objects."""

    def it_provides_a_constructor(self, shapes_: Mock, _init_: Mock):
        start_x, start_y, x_scale, y_scale = 99.56, 200.49, 4.2, 2.4
        start_x_int, start_y_int = 100, 200

        builder = FreeformBuilder.new(shapes_, start_x, start_y, x_scale, y_scale)

        _init_.assert_called_once_with(builder, shapes_, start_x_int, start_y_int, x_scale, y_scale)
        assert isinstance(builder, FreeformBuilder)

    @pytest.mark.parametrize("close", [True, False])
    def it_can_add_straight_line_segments(self, request: FixtureRequest, close: bool):
        _add_line_segment_ = method_mock(request, FreeformBuilder, "_add_line_segment")
        _add_close_ = method_mock(request, FreeformBuilder, "_add_close")
        builder = FreeformBuilder(None, None, None, None, None)  # type: ignore

        return_value = builder.add_line_segments(((1, 2), (3, 4), (5, 6)), close)

        assert _add_line_segment_.call_args_list == [
            call(builder, 1, 2),
            call(builder, 3, 4),
            call(builder, 5, 6),
        ]
        assert _add_close_.call_args_list == ([call(builder)] if close else [])
        assert return_value is builder

    def it_can_move_the_pen_location(self, _MoveTo_new_: Mock, move_to_: Mock):
        x, y = 42, 24
        _MoveTo_new_.return_value = move_to_
        builder = FreeformBuilder(None, None, None, None, None)  # type: ignore

        return_value = builder.move_to(x, y)

        _MoveTo_new_.assert_called_once_with(builder, x, y)
        assert builder._drawing_operations[-1] is move_to_
        assert return_value is builder

    def it_can_build_the_specified_freeform_shape(
        self,
        shapes_: Mock,
        apply_operation_to_: Mock,
        _add_freeform_sp_: Mock,
        _start_path_: Mock,
        shape_: Mock,
    ):
        origin_x, origin_y = Mm(42), Mm(24)
        sp, path = element("p:sp"), element("a:path")
        drawing_ops = (
            _LineSegment(None, None, None),  # type: ignore
            _LineSegment(None, None, None),  # type: ignore
        )
        shapes_._shape_factory.return_value = shape_
        _add_freeform_sp_.return_value = sp
        _start_path_.return_value = path
        builder = FreeformBuilder(shapes_, None, None, None, None)  # type: ignore
        builder._drawing_operations.extend(drawing_ops)
        calls = [call(drawing_ops[0], path), call(drawing_ops[1], path)]

        shape = builder.convert_to_shape(origin_x, origin_y)

        _add_freeform_sp_.assert_called_once_with(builder, origin_x, origin_y)
        _start_path_.assert_called_once_with(builder, sp)
        assert apply_operation_to_.call_args_list == calls
        shapes_._shape_factory.assert_called_once_with(sp)
        assert shape is shape_

    @pytest.mark.parametrize(
        ("start_x", "xs", "expected_value"),
        [
            (Mm(0), (1, None, 2, 3), Mm(0)),
            (Mm(6), (1, None, 2, 3), Mm(1)),
            (Mm(50), (150, -5, None, 100), Mm(-5)),
        ],
    )
    def it_knows_the_shape_x_offset(
        self, start_x: int, xs: tuple[int | None, ...], expected_value: int
    ):
        builder = FreeformBuilder(None, start_x, None, None, None)  # type: ignore
        drawing_ops = [_Close() if x is None else _LineSegment(builder, Mm(x), Mm(0)) for x in xs]
        builder._drawing_operations.extend(drawing_ops)

        assert builder.shape_offset_x == expected_value

    @pytest.mark.parametrize(
        ("start_y", "ys", "expected_value"),
        [
            (Mm(0), (2, None, 6, 8), Mm(0)),
            (Mm(4), (2, None, 6, 8), Mm(2)),
            (Mm(19), (213, -22, None, 100), Mm(-22)),
        ],
    )
    def it_knows_the_shape_y_offset(
        self, start_y: int, ys: tuple[int | None, ...], expected_value: int
    ):
        builder = FreeformBuilder(None, None, start_y, None, None)  # type: ignore
        drawing_ops = [_Close() if y is None else _LineSegment(builder, Mm(0), Mm(y)) for y in ys]
        builder._drawing_operations.extend(drawing_ops)

        assert builder.shape_offset_y == expected_value

    def it_adds_a_freeform_sp_to_help(
        self, _left_prop_: Mock, _top_prop_: Mock, _width_prop_: Mock, _height_prop_: Mock
    ):
        origin_x, origin_y = Emu(42), Emu(24)
        spTree = element("p:spTree")
        shapes = SlideShapes(spTree, None)  # type: ignore
        _left_prop_.return_value, _top_prop_.return_value = Emu(12), Emu(34)
        _width_prop_.return_value, _height_prop_.return_value = 56, 78
        builder = FreeformBuilder(shapes, None, None, None, None)  # type: ignore
        expected_xml = snippet_seq("freeform")[0]

        sp = builder._add_freeform_sp(origin_x, origin_y)

        assert spTree.xml == expected_xml
        assert sp is spTree.xpath("p:sp")[0]

    def it_adds_a_line_segment_to_help(self, _LineSegment_new_: Mock, line_segment_: Mock):
        x, y = 4, 2
        _LineSegment_new_.return_value = line_segment_

        builder = FreeformBuilder(None, None, None, None, None)  # type: ignore

        builder._add_line_segment(x, y)

        _LineSegment_new_.assert_called_once_with(builder, x, y)
        assert builder._drawing_operations == [line_segment_]

    def it_closes_a_contour_to_help(self, _Close_new_: Mock, close_: Mock):
        _Close_new_.return_value = close_
        builder = FreeformBuilder(None, None, None, None, None)  # type: ignore

        builder._add_close()

        _Close_new_.assert_called_once_with()
        assert builder._drawing_operations == [close_]

    def it_knows_the_freeform_left_extent_to_help(self, left_fixture):
        builder, expected_value = left_fixture
        left = builder._left
        assert left == expected_value

    def it_knows_the_freeform_top_extent_to_help(self, top_fixture):
        builder, expected_value = top_fixture
        top = builder._top
        assert top == expected_value

    def it_knows_the_freeform_width_to_help(self, width_fixture):
        builder, expected_value = width_fixture
        width = builder._width
        assert width == expected_value

    @pytest.mark.parametrize(
        ("dy", "y_scale", "expected_value"),
        [(0, 2.0, 0), (24, 10.0, 240), (914400, 314.1, 287213040)],
    )
    def it_knows_the_freeform_height_to_help(
        self, dy: int, y_scale: float, expected_value: int, _dy_prop_: Mock
    ):
        _dy_prop_.return_value = dy
        builder = FreeformBuilder(None, None, None, None, y_scale)  # type: ignore
        height = builder._height
        assert height == expected_value

    def it_knows_the_local_coordinate_width_to_help(self, dx_fixture):
        builder, expected_value = dx_fixture
        dx = builder._dx
        assert dx == expected_value

    def it_knows_the_local_coordinate_height_to_help(self, dy_fixture):
        builder, expected_value = dy_fixture
        dy = builder._dy
        assert dy == expected_value

    def it_can_start_a_new_path_to_help(
        self, request: FixtureRequest, _dx_prop_: Mock, _dy_prop_: Mock
    ):
        _local_to_shape_ = method_mock(
            request, FreeformBuilder, "_local_to_shape", return_value=(101, 202)
        )
        sp = element("p:sp/p:spPr/a:custGeom")
        start_x, start_y = 42, 24
        _dx_prop_.return_value, _dy_prop_.return_value = 1001, 2002
        builder = FreeformBuilder(None, start_x, start_y, None, None)

        path = builder._start_path(sp)

        _local_to_shape_.assert_called_once_with(builder, start_x, start_y)
        assert sp.xml == xml(
            "p:sp/p:spPr/a:custGeom/a:pathLst/a:path{w=1001,h=2002}/a:moveTo" "/a:pt{x=101,y=202}"
        )
        assert path is sp.xpath(".//a:path")[-1]

    def it_translates_local_to_shape_coordinates_to_help(self, local_fixture):
        builder, local_x, local_y, expected_value = local_fixture
        shape_x_y = builder._local_to_shape(local_x, local_y)
        assert shape_x_y == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            (0, (1, None, 2, 3), 3),
            (6, (1, None, 2, 3), 5),
            (50, (150, -5, None, 100), 155),
        ]
    )
    def dx_fixture(self, request: FixtureRequest):
        start_x, xs, expected_value = request.param
        drawing_ops = []
        for x in xs:
            if x is None:
                drawing_ops.append(_Close())
            else:
                drawing_ops.append(_BaseDrawingOperation(None, x, None))

        builder = FreeformBuilder(None, start_x, None, None, None)
        builder._drawing_operations.extend(drawing_ops)
        return builder, expected_value

    @pytest.fixture(
        params=[
            (0, (1, None, 2, 3), 3),
            (2, (1, None, 2, 3), 2),
            (32, (160, -8, None, 101), 168),
        ]
    )
    def dy_fixture(self, request: FixtureRequest):
        start_y, ys, expected_value = request.param
        drawing_ops = []
        for y in ys:
            if y is None:
                drawing_ops.append(_Close())
            else:
                drawing_ops.append(_BaseDrawingOperation(None, None, y))

        builder = FreeformBuilder(None, None, start_y, None, None)
        builder._drawing_operations.extend(drawing_ops)
        return builder, expected_value

    @pytest.fixture(params=[(0, 1.0, 0), (4, 10.0, 40), (914400, 914.3, 836035920)])
    def left_fixture(self, request: FixtureRequest, shape_offset_x_prop_: Mock):
        offset_x, x_scale, expected_value = request.param
        shape_offset_x_prop_.return_value = offset_x

        builder = FreeformBuilder(None, None, None, x_scale, None)
        return builder, expected_value

    @pytest.fixture
    def local_fixture(self, shape_offset_x_prop_: Mock, shape_offset_y_prop_: Mock):
        local_x, local_y = 123, 456
        shape_offset_x_prop_.return_value = 23
        shape_offset_y_prop_.return_value = 156

        builder = FreeformBuilder(None, None, None, None, None)
        expected_value = (100, 300)
        return builder, local_x, local_y, expected_value

    @pytest.fixture
    def sp_fixture(
        self, _left_prop_: Mock, _top_prop_: Mock, _width_prop_: Mock, _height_prop_: Mock
    ):
        origin_x, origin_y = 42, 24
        spTree = element("p:spTree")
        shapes = SlideShapes(spTree, None)
        _left_prop_.return_value, _top_prop_.return_value = 12, 34
        _width_prop_.return_value, _height_prop_.return_value = 56, 78

        builder = FreeformBuilder(shapes, None, None, None, None)
        expected_xml = snippet_seq("freeform")[0]
        return builder, origin_x, origin_y, spTree, expected_xml

    @pytest.fixture(params=[(0, 11.0, 0), (100, 10.36, 1036), (914242, 943.1, 862221630)])
    def top_fixture(self, request: FixtureRequest, shape_offset_y_prop_: Mock):
        offset_y, y_scale, expected_value = request.param
        shape_offset_y_prop_.return_value = offset_y

        builder = FreeformBuilder(None, None, None, None, y_scale)
        return builder, expected_value

    @pytest.fixture(params=[(0, 1.0, 0), (42, 10.0, 420), (914400, 914.4, 836127360)])
    def width_fixture(self, request: FixtureRequest, _dx_prop_: Mock):
        dx, x_scale, expected_value = request.param
        _dx_prop_.return_value = dx

        builder = FreeformBuilder(None, None, None, x_scale, None)
        return builder, expected_value

    # fixture components -----------------------------------

    @pytest.fixture
    def _add_freeform_sp_(self, request: FixtureRequest):
        return method_mock(request, FreeformBuilder, "_add_freeform_sp", autospec=True)

    @pytest.fixture
    def apply_operation_to_(self, request: FixtureRequest):
        return method_mock(request, _LineSegment, "apply_operation_to", autospec=True)

    @pytest.fixture
    def close_(self, request: FixtureRequest):
        return instance_mock(request, _Close)

    @pytest.fixture
    def _Close_new_(self, request: FixtureRequest):
        return method_mock(request, _Close, "new", autospec=False)

    @pytest.fixture
    def _dx_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_dx")

    @pytest.fixture
    def _dy_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_dy")

    @pytest.fixture
    def _height_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_height")

    @pytest.fixture
    def _init_(self, request: FixtureRequest):
        return initializer_mock(request, FreeformBuilder, autospec=True)

    @pytest.fixture
    def _left_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_left")

    @pytest.fixture
    def line_segment_(self, request: FixtureRequest):
        return instance_mock(request, _LineSegment)

    @pytest.fixture
    def _LineSegment_new_(self, request: FixtureRequest):
        return method_mock(request, _LineSegment, "new", autospec=False)

    @pytest.fixture
    def move_to_(self, request: FixtureRequest):
        return instance_mock(request, _MoveTo)

    @pytest.fixture
    def _MoveTo_new_(self, request: FixtureRequest):
        return method_mock(request, _MoveTo, "new", autospec=False)

    @pytest.fixture
    def shape_(self, request: FixtureRequest):
        return instance_mock(request, Shape)

    @pytest.fixture
    def shape_offset_x_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "shape_offset_x")

    @pytest.fixture
    def shape_offset_y_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "shape_offset_y")

    @pytest.fixture
    def shapes_(self, request: FixtureRequest):
        return instance_mock(request, SlideShapes)

    @pytest.fixture
    def _start_path_(self, request: FixtureRequest):
        return method_mock(request, FreeformBuilder, "_start_path", autospec=True)

    @pytest.fixture
    def _top_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_top")

    @pytest.fixture
    def _width_prop_(self, request: FixtureRequest):
        return property_mock(request, FreeformBuilder, "_width")


class Describe_BaseDrawingOperation(object):
    """Unit-test suite for `pptx.shapes.freeform.BaseDrawingOperation` objects."""

    def it_knows_its_x_coordinate(self, x_fixture):
        drawing_operation, expected_value = x_fixture
        x = drawing_operation.x
        assert x == expected_value

    def it_knows_its_y_coordinate(self, y_fixture):
        drawing_operation, expected_value = y_fixture
        y = drawing_operation.y
        assert y == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def x_fixture(self):
        x = 42
        drawing_operation = _BaseDrawingOperation(None, x, None)
        expected_value = x
        return drawing_operation, expected_value

    @pytest.fixture
    def y_fixture(self):
        y = 24
        drawing_operation = _BaseDrawingOperation(None, None, y)
        expected_value = y
        return drawing_operation, expected_value


class Describe_Close(object):
    """Unit-test suite for `pptx.shapes.freeform._Close` objects."""

    def it_provides_a_constructor(self, new_fixture):
        _init_ = new_fixture

        close = _Close.new()

        _init_.assert_called_once_with()
        assert isinstance(close, _Close)

    def it_can_add_close_a_contour(self, apply_fixture):
        close, path, expected_xml = apply_fixture

        close_elm = close.apply_operation_to(path)

        assert path.xml == expected_xml
        assert close_elm is path.xpath("a:close")[-1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def new_fixture(self, _init_):
        return _init_

    # fixture components -----------------------------------

    @pytest.fixture
    def apply_fixture(self):
        path = element("a:path")
        close = _Close()
        expected_xml = xml("a:path/a:close")
        return close, path, expected_xml

    @pytest.fixture
    def _init_(self, request: FixtureRequest):
        return initializer_mock(request, _Close, autospec=True)


class Describe_LineSegment(object):
    """Unit-test suite for `pptx.shapes.freeform._LineSegment` objects."""

    def it_provides_a_constructor(self, new_fixture):
        builder_, x, y, _init_, x_int, y_int = new_fixture

        line_segment = _LineSegment.new(builder_, x, y)

        _init_.assert_called_once_with(line_segment, builder_, x_int, y_int)
        assert isinstance(line_segment, _LineSegment)

    def it_can_add_its_line_segment_to_a_path(self, apply_fixture):
        line_segment, path, expected_xml = apply_fixture

        lnTo = line_segment.apply_operation_to(path)

        assert path.xml == expected_xml
        assert lnTo is path.xpath("a:lnTo")[-1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def apply_fixture(self, builder_):
        x, y = 420, 240
        path = element("a:path")
        builder_.shape_offset_x, builder_.shape_offset_y = 100, 200

        line_segment = _LineSegment(builder_, x, y)
        expected_xml = xml("a:path/a:lnTo/a:pt{x=320,y=40}")
        return line_segment, path, expected_xml

    @pytest.fixture
    def new_fixture(self, builder_, _init_):
        x, y, x_int, y_int = 99.51, 200.49, 100, 200
        return builder_, x, y, _init_, x_int, y_int

    # fixture components -----------------------------------

    @pytest.fixture
    def builder_(self, request: FixtureRequest):
        return instance_mock(request, FreeformBuilder)

    @pytest.fixture
    def _init_(self, request: FixtureRequest):
        return initializer_mock(request, _LineSegment, autospec=True)


class Describe_MoveTo(object):
    """Unit-test suite for `pptx.shapes.freeform._MoveTo` objects."""

    def it_provides_a_constructor(self, new_fixture):
        builder_, x, y, _init_, x_int, y_int = new_fixture

        move_to = _MoveTo.new(builder_, x, y)

        _init_.assert_called_once_with(move_to, builder_, x_int, y_int)
        assert isinstance(move_to, _MoveTo)

    def it_can_add_its_move_to_a_path(self, apply_fixture):
        move_to, path, expected_xml = apply_fixture

        moveTo = move_to.apply_operation_to(path)

        assert path.xml == expected_xml
        assert moveTo is path.xpath("a:moveTo")[-1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def apply_fixture(self, builder_):
        x, y = 120, 340
        path = element("a:path")
        builder_.shape_offset_x, builder_.shape_offset_y = 100, 200

        move_to = _MoveTo(builder_, x, y)
        expected_xml = xml("a:path/a:moveTo/a:pt{x=20,y=140}")
        return move_to, path, expected_xml

    @pytest.fixture
    def new_fixture(self, builder_, _init_):
        x, y, x_int, y_int = 99.51, 200.49, 100, 200
        return builder_, x, y, _init_, x_int, y_int

    # fixture components -----------------------------------

    @pytest.fixture
    def builder_(self, request: FixtureRequest):
        return instance_mock(request, FreeformBuilder)

    @pytest.fixture
    def _init_(self, request: FixtureRequest):
        return initializer_mock(request, _MoveTo, autospec=True)
