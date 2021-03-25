# encoding: utf-8

"""Unit-test suite for pptx.shapes.graphfrm module."""

import pytest

from pptx.chart.chart import Chart
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts.chart import ChartPart
from pptx.parts.slide import SlidePart
from pptx.shapes.graphfrm import GraphicFrame
from pptx.spec import GRAPHIC_DATA_URI_CHART, GRAPHIC_DATA_URI_TABLE

from ..unitutil.cxml import element
from ..unitutil.mock import instance_mock, property_mock


class DescribeGraphicFrame(object):
    """Unit-test suite for `pptx.shapes.graphfrm.GraphicFrame` object."""

    def it_provides_access_to_the_chart_it_contains(
        self, request, has_chart_prop_, chart_part_, chart_
    ):
        has_chart_prop_.return_value = True
        property_mock(request, GraphicFrame, "chart_part", return_value=chart_part_)
        chart_part_.chart = chart_

        assert GraphicFrame(None, None).chart is chart_

    def but_it_raises_on_chart_if_there_isnt_one(self, has_chart_prop_):
        has_chart_prop_.return_value = False

        with pytest.raises(ValueError) as e:
            GraphicFrame(None, None).chart
        assert str(e.value) == "shape does not contain a chart"

    def it_provides_access_to_its_chart_part(self, request, chart_part_):
        graphicFrame = element(
            "p:graphicFrame/a:graphic/a:graphicData/c:chart{r:id=rId42}"
        )
        property_mock(
            request,
            GraphicFrame,
            "part",
            return_value=instance_mock(
                request, SlidePart, related_parts={"rId42": chart_part_}
            ),
        )
        graphic_frame = GraphicFrame(graphicFrame, None)

        assert graphic_frame.chart_part is chart_part_

    @pytest.mark.parametrize(
        "graphicData_uri, expected_value",
        (
            (GRAPHIC_DATA_URI_CHART, True),
            (GRAPHIC_DATA_URI_TABLE, False),
        ),
    )
    def it_knows_whether_it_contains_a_chart(self, graphicData_uri, expected_value):
        graphicFrame = element(
            "p:graphicFrame/a:graphic/a:graphicData{uri=%s}" % graphicData_uri
        )
        assert GraphicFrame(graphicFrame, None).has_chart is expected_value

    @pytest.mark.parametrize(
        "graphicData_uri, expected_value",
        (
            (GRAPHIC_DATA_URI_CHART, False),
            (GRAPHIC_DATA_URI_TABLE, True),
        ),
    )
    def it_knows_whether_it_contains_a_table(self, graphicData_uri, expected_value):
        graphicFrame = element(
            "p:graphicFrame/a:graphic/a:graphicData{uri=%s}" % graphicData_uri
        )
        assert GraphicFrame(graphicFrame, None).has_table is expected_value

    def it_raises_on_shadow(self):
        graphic_frame = GraphicFrame(None, None)
        with pytest.raises(NotImplementedError):
            graphic_frame.shadow

    @pytest.mark.parametrize(
        "uri, expected_value",
        (
            (GRAPHIC_DATA_URI_CHART, MSO_SHAPE_TYPE.CHART),
            (GRAPHIC_DATA_URI_TABLE, MSO_SHAPE_TYPE.TABLE),
            ("foobar", None),
        ),
    )
    def it_knows_its_shape_type(self, uri, expected_value):
        graphicFrame = element("p:graphicFrame/a:graphic/a:graphicData{uri=%s}" % uri)
        assert GraphicFrame(graphicFrame, None).shape_type is expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def chart_(self, request):
        return instance_mock(request, Chart)

    @pytest.fixture
    def chart_part_(self, request, chart_):
        return instance_mock(request, ChartPart, chart=chart_)

    @pytest.fixture
    def has_chart_prop_(self, request):
        return property_mock(request, GraphicFrame, "has_chart")
