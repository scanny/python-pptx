# encoding: utf-8

"""Test suite for pptx.shapes.group module."""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.dml.effect import ShadowFormat
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import GroupShapes

from ..unitutil.cxml import element
from ..unitutil.mock import class_mock, initializer_mock, instance_mock


class DescribeGroupShape(object):

    def it_raises_on_access_click_action(self, click_action_fixture):
        group = click_action_fixture
        with pytest.raises(TypeError):
            group.click_action

    def it_provides_access_to_its_shadow(self, ShadowFormat_, shadow_):
        grpSp = element('p:grpSp/p:grpSpPr')
        grpSpPr = grpSp.xpath('//p:grpSpPr')[0]
        ShadowFormat_.return_value = shadow_
        group_shape = GroupShape(grpSp, None)

        shadow = group_shape.shadow

        ShadowFormat_.assert_called_once_with(grpSpPr)
        assert shadow is shadow_

    def it_knows_its_shape_type(self, shape_type_fixture):
        group = shape_type_fixture
        assert group.shape_type == MSO_SHAPE_TYPE.GROUP

    def it_provides_access_to_its_sub_shapes(self, shapes_fixture):
        group, GroupShapes_init_, grpSp = shapes_fixture

        shapes = group.shapes

        GroupShapes_init_.assert_called_once_with(shapes, grpSp, group)
        assert isinstance(shapes, GroupShapes)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def click_action_fixture(self):
        return GroupShape(None, None)

    @pytest.fixture
    def shape_type_fixture(self):
        return GroupShape(None, None)

    @pytest.fixture
    def shapes_fixture(self, GroupShapes_init_):
        grpSp = element('p:grpSp')
        group = GroupShape(grpSp, None)
        return group, GroupShapes_init_, grpSp

    # fixture components ---------------------------------------------

    @pytest.fixture
    def GroupShapes_init_(self, request):
        return initializer_mock(request, GroupShapes, autospec=True)

    @pytest.fixture
    def shadow_(self, request):
        return instance_mock(request, ShadowFormat)

    @pytest.fixture
    def ShadowFormat_(self, request):
        return class_mock(request, 'pptx.shapes.group.ShadowFormat')
