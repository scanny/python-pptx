# encoding: utf-8

"""Test suite for pptx.shapes.picture module."""

from __future__ import absolute_import, division, print_function, unicode_literals

import pytest

from pptx.dml.line import LineFormat
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_MEDIA_TYPE
from pptx.parts.image import Image
from pptx.parts.slide import SlidePart
from pptx.shapes.picture import _BasePicture, _MediaFormat, Movie, Picture
from pptx.util import Pt

from ..unitutil.cxml import element, xml
from ..unitutil.mock import call, class_mock, instance_mock, property_mock


class Describe_BasePicture(object):
    def it_knows_its_cropping(self, crop_get_fixture):
        picture, prop_name, expected_value = crop_get_fixture
        crop = getattr(picture, prop_name)
        assert abs(crop - expected_value) < 0.000001

    def it_can_change_its_cropping(self, crop_set_fixture):
        picture, prop_name, value, expected_xml = crop_set_fixture
        setattr(picture, prop_name, value)
        assert picture._element.xml == expected_xml

    def it_provides_access_to_its_outline(self, line_fixture):
        picture = line_fixture
        line = picture.line
        assert isinstance(line, LineFormat)
        # exercise line to test has_line interface, .ln and .get_or_add_ln()
        picture.line.width = Pt(1)
        assert picture.line.width == Pt(1)

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            ("p:pic/p:blipFill", "left", 0.0),
            ("p:pic/p:blipFill/a:srcRect", "top", 0.0),
            ("p:pic/p:blipFill/a:srcRect{l=99999}", "bottom", 0.0),
            ("p:pic/p:blipFill/a:srcRect{l=42424}", "left", 0.42424),
            ("p:pic/p:blipFill/a:srcRect{t=-10000}", "top", -0.1),
            ("p:pic/p:blipFill/a:srcRect{r=250000}", "right", 2.5),
            ("p:pic/p:blipFill/a:srcRect{b=33333}", "bottom", 0.33333),
        ]
    )
    def crop_get_fixture(self, request):
        pic_cxml, side, expected_value = request.param
        picture = Picture(element(pic_cxml), None)
        prop_name = "crop_%s" % side
        return picture, prop_name, expected_value

    @pytest.fixture(
        params=[
            (
                "p:pic{a:b=c}/p:blipFill",
                "left",
                0.11,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{l=11000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill",
                "top",
                0.21,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{t=21000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill",
                "right",
                0.31,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{r=31000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill",
                "bottom",
                0.41,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{b=41000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{l=80000}",
                "left",
                0.21,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{l=21000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{t=70000}",
                "top",
                0.22,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{t=22000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{r=60000}",
                "right",
                0.23,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{r=23000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{b=50000}",
                "bottom",
                0.24,
                "p:pic{a:b=c}/p:blipFill/a:srcRect{b=24000}",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{l=90000}",
                "left",
                0,
                "p:pic{a:b=c}/p:blipFill/a:srcRect",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{t=91000}",
                "top",
                0.0,
                "p:pic{a:b=c}/p:blipFill/a:srcRect",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{r=92000}",
                "right",
                0,
                "p:pic{a:b=c}/p:blipFill/a:srcRect",
            ),
            (
                "p:pic{a:b=c}/p:blipFill/a:srcRect{b=93000}",
                "bottom",
                0.0,
                "p:pic{a:b=c}/p:blipFill/a:srcRect",
            ),
        ]
    )
    def crop_set_fixture(self, request):
        pic_cxml, side, value, expected_cxml = request.param
        pic = element(pic_cxml)
        picture = Picture(pic, None)
        prop_name = "crop_%s" % side
        expected_xml = xml(expected_cxml)
        return picture, prop_name, value, expected_xml

    @pytest.fixture
    def line_fixture(self):
        return _BasePicture(element("p:pic/p:spPr"), None)


class DescribeMovie(object):
    def it_knows_its_shape_type(self, shape_type_fixture):
        movie = shape_type_fixture
        assert movie.shape_type == MSO_SHAPE_TYPE.MEDIA

    def it_knows_its_media_type(self, media_type_fixture):
        movie = media_type_fixture
        assert movie.media_type == PP_MEDIA_TYPE.MOVIE

    def it_provides_access_to_its_media_format(self, format_fixture):
        movie, MediaFormat_, pic, parent, media_format_ = format_fixture
        media_format = movie.media_format
        MediaFormat_.assert_called_once_with(pic, parent)
        assert media_format is media_format_

    def it_provides_access_to_its_poster_frame_image(self, pfrm_fixture):
        movie, slide_part_, calls, expected_value = pfrm_fixture
        poster_frame = movie.poster_frame
        assert slide_part_.get_image.call_args_list == calls
        assert poster_frame == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def format_fixture(self, _MediaFormat_, media_format_):
        pic = element("p:pic")
        parent = movie = Movie(pic, None)
        return movie, _MediaFormat_, pic, parent, media_format_

    @pytest.fixture
    def media_type_fixture(self):
        return Movie(None, None)

    @pytest.fixture(
        params=[
            ("p:pic/p:blipFill/a:blip{r:embed=rId42}", True),
            ("p:pic/p:blipFill/a:blip", False),
        ]
    )
    def pfrm_fixture(self, request, slide_part_, image_, part_prop_):
        cxml, is_present = request.param
        movie = Movie(element(cxml), None)
        calls, expected_value = [], None
        part_prop_.return_value = slide_part_
        slide_part_.get_image.return_value = image_
        if is_present:
            calls.append(call("rId42"))
            expected_value = image_
        return movie, slide_part_, calls, expected_value

    @pytest.fixture
    def shape_type_fixture(self):
        return Movie(None, None)

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_(self, request):
        return instance_mock(request, Image)

    @pytest.fixture
    def _MediaFormat_(self, request, media_format_):
        return class_mock(
            request, "pptx.shapes.picture._MediaFormat", return_value=media_format_
        )

    @pytest.fixture
    def media_format_(self, request):
        return instance_mock(request, _MediaFormat)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, Movie, "part")

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribePicture(object):
    def it_knows_its_masking_shape(self, autoshape_get_fixture):
        picture, expected_value = autoshape_get_fixture
        auto_shape_type = picture.auto_shape_type
        assert auto_shape_type == expected_value

    def it_can_change_its_masking_shape(self, autoshape_set_fixture):
        picture, new_value, expected_xml = autoshape_set_fixture
        picture.auto_shape_type = new_value
        assert picture._element.xml == expected_xml

    def it_knows_its_shape_type(self, shape_type_fixture):
        picture = shape_type_fixture
        assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE

    def it_provides_access_to_its_image(self, image_fixture):
        picture, slide_part_, rId, image_ = image_fixture
        image = picture.image
        slide_part_.get_image.assert_called_once_with(rId)
        assert image is image_

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            ("", None),
            ("/a:prstGeom{prst=rect}", MSO_SHAPE.RECTANGLE),
            ("/a:prstGeom{prst=hexagon}", MSO_SHAPE.HEXAGON),
        ]
    )
    def autoshape_get_fixture(self, request):
        prstGeom_cxml, expected_value = request.param
        pic_cxml = "p:pic/p:spPr%s" % prstGeom_cxml
        picture = Picture(element(pic_cxml), None)
        return picture, expected_value

    @pytest.fixture(
        params=[
            (
                "p:pic/p:spPr/a:custGeom",
                MSO_SHAPE.RECTANGLE,
                "p:pic/p:spPr/a:prstGeom{prst=rect}",
            ),
            (
                "p:pic/p:spPr/a:prstGeom{prst=rect}",
                MSO_SHAPE.OVAL,
                "p:pic/p:spPr/a:prstGeom{prst=ellipse}",
            ),
            (
                "p:pic/p:spPr/a:prstGeom{prst=ellipse}",
                MSO_SHAPE.PIE,
                "p:pic/p:spPr/a:prstGeom{prst=pie}",
            ),
        ]
    )
    def autoshape_set_fixture(self, request):
        pic_cxml, new_value, expected_cxml = request.param
        picture = Picture(element(pic_cxml), None)
        expected_xml = xml(expected_cxml)
        return picture, new_value, expected_xml

    @pytest.fixture
    def image_fixture(self, part_prop_, slide_part_, image_):
        pic_cxml, rId = "p:pic/p:blipFill/a:blip{r:embed=rId42}", "rId42"
        picture = Picture(element(pic_cxml), None)
        slide_part_.get_image.return_value = image_
        return picture, slide_part_, rId, image_

    @pytest.fixture
    def shape_type_fixture(self):
        return Picture(None, None)

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_(self, request):
        return instance_mock(request, Image)

    @pytest.fixture
    def part_prop_(self, request, slide_part_):
        return property_mock(request, Picture, "part", return_value=slide_part_)

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)
