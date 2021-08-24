# encoding: utf-8

"""Unit-test suite for `pptx.shapes.placeholder` module."""

import pytest

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE as XCT
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.shapes.shared import ST_Direction, ST_PlaceholderSize
from pptx.parts.image import ImagePart
from pptx.parts.slide import NotesSlidePart, SlidePart
from pptx.shapes.placeholder import (
    BasePlaceholder,
    _BaseSlidePlaceholder,
    ChartPlaceholder,
    _InheritsDimensions,
    LayoutPlaceholder,
    MasterPlaceholder,
    NotesSlidePlaceholder,
    PicturePlaceholder,
    PlaceholderGraphicFrame,
    PlaceholderPicture,
    TablePlaceholder,
)
from pptx.shapes.shapetree import NotesSlidePlaceholders
from pptx.slide import NotesMaster, SlideLayout, SlideMaster

from ..oxml.unitdata.shape import (
    a_graphicFrame,
    a_ph,
    an_nvGraphicFramePr,
    an_nvPicPr,
    an_nvPr,
    an_nvSpPr,
    an_sp,
)
from ..unitutil.cxml import element, xml
from ..unitutil.file import snippet_seq
from ..unitutil.mock import class_mock, instance_mock, method_mock, property_mock


class Describe_BaseSlidePlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder._BaseSlidePlaceholder` object."""

    def it_knows_its_shape_type(self):
        placeholder = _BaseSlidePlaceholder(None, None)
        assert placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER

    def it_provides_override_dimensions_when_present(self, override_fixture):
        placeholder, prop_name, expected_value = override_fixture
        assert getattr(placeholder, prop_name) == expected_value

    @pytest.mark.parametrize("prop_name", ("left", "top", "width", "height"))
    def it_provides_inherited_dims_when_no_override(self, request, prop_name):
        method_mock(request, _BaseSlidePlaceholder, "_inherited_value", return_value=42)
        placeholder = _BaseSlidePlaceholder(element("p:sp/p:spPr"), None)

        value = getattr(placeholder, prop_name)

        placeholder._inherited_value.assert_called_once_with(placeholder, prop_name)
        assert value == 42

    def it_gets_an_inherited_dim_value_to_help(self, base_val_fixture):
        placeholder, attr_name, expected_value = base_val_fixture
        value = placeholder._inherited_value(attr_name)
        assert value == expected_value

    def it_finds_its_base_placeholder_to_help(self, base_ph_fixture):
        placeholder, layout_, idx, layout_placeholder_ = base_ph_fixture
        base_placeholder = placeholder._base_placeholder
        layout_.placeholders.get.assert_called_once_with(idx=idx)
        assert base_placeholder is layout_placeholder_

    def it_can_override_inherited_dimensions(self, dim_set_fixture):
        placeholder, prop_name, value, expected_xml = dim_set_fixture
        setattr(placeholder, prop_name, value)
        assert placeholder._element.xml == expected_xml

    def it_replaces_a_placeholder_element_to_help(self, replace_fixture):
        placeholder, element, spTree, expected_xml = replace_fixture
        placeholder._replace_placeholder_with(element)
        assert spTree.xml == expected_xml
        assert placeholder._element is None

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def base_ph_fixture(self, request, part_prop_, slide_layout_, layout_placeholder_):
        sp_cxml = "p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}"
        placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        part_prop_.return_value.slide_layout = slide_layout_
        slide_layout_.placeholders.get.return_value = layout_placeholder_
        return placeholder, slide_layout_, 1, layout_placeholder_

    @pytest.fixture(params=[(True, 42), (False, None)])
    def base_val_fixture(self, request, _base_placeholder_prop_, layout_placeholder_):
        has_base_placeholder, expected_value = request.param
        placeholder = _BaseSlidePlaceholder(None, None)
        if has_base_placeholder:
            layout_placeholder_.width = expected_value
            _base_placeholder_prop_.return_value = layout_placeholder_
        else:
            _base_placeholder_prop_.return_value = None
        return placeholder, "width", expected_value

    @pytest.fixture(
        params=[
            ("p:sp/p:spPr/a:xfrm", "left", 1, "p:sp/p:spPr/a:xfrm/a:off{x=1,y=0}"),
            ("p:sp/p:spPr/a:xfrm", "top", 2, "p:sp/p:spPr/a:xfrm/a:off{x=0,y=2}"),
            ("p:sp/p:spPr/a:xfrm", "width", 3, "p:sp/p:spPr/a:xfrm/a:ext{cx=3,cy=0}"),
            ("p:sp/p:spPr/a:xfrm", "height", 4, "p:sp/p:spPr/a:xfrm/a:ext{cx=0,cy=4}"),
        ]
    )
    def dim_set_fixture(self, request):
        sp_cxml, prop_name, value, expected_cxml = request.param
        placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        expected_xml = xml(expected_cxml)
        return placeholder, prop_name, value, expected_xml

    @pytest.fixture(
        params=[
            ("left", "p:sp/p:spPr/a:xfrm/a:off{x=12}", 12),
            ("top", "p:sp/p:spPr/a:xfrm/a:off{y=34}", 34),
            ("width", "p:sp/p:spPr/a:xfrm/a:ext{cx=56}", 56),
            ("height", "p:sp/p:spPr/a:xfrm/a:ext{cy=78}", 78),
        ]
    )
    def override_fixture(self, request):
        prop_name, sp_cxml, value = request.param
        slide_placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        return slide_placeholder, prop_name, value

    @pytest.fixture
    def replace_fixture(self):
        spTree = element("p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=10}")
        pic = element("p:pic/p:nvPicPr/p:nvPr")
        placeholder = _BaseSlidePlaceholder(spTree[0], None)
        expected_xml = xml("p:spTree/p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=10}")
        return placeholder, pic, spTree, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _base_placeholder_prop_(self, request):
        return property_mock(request, _BaseSlidePlaceholder, "_base_placeholder")

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def part_prop_(self, request, slide_part_):
        return property_mock(
            request, _BaseSlidePlaceholder, "part", return_value=slide_part_
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribeBasePlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.BasePlaceholder` object."""

    def it_knows_its_idx_value(self, idx_fixture):
        placeholder, idx = idx_fixture
        assert placeholder.idx == idx

    def it_knows_its_orient_value(self, orient_fixture):
        placeholder, orient = orient_fixture
        assert placeholder.orient == orient

    def it_raises_on_ph_orient_when_not_a_placeholder(self, orient_raise_fixture):
        shape = orient_raise_fixture
        with pytest.raises(ValueError):
            shape.orient

    def it_knows_its_sz_value(self, sz_fixture):
        placeholder, sz = sz_fixture
        assert placeholder.sz == sz

    def it_knows_its_placeholder_type(self, type_fixture):
        placeholder, ph_type = type_fixture
        assert placeholder.ph_type == ph_type

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            ("sp", None, "title", 0),
            ("sp", 0, "title", 0),
            ("sp", 3, "body", 3),
            ("pic", 6, "pic", 6),
            ("graphicFrame", 9, "tbl", 9),
        ]
    )
    def idx_fixture(self, request):
        tagname, idx, ph_type, expected_idx = request.param
        shape_elm = self.shape_elm_factory(tagname, ph_type, idx)
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_idx

    @pytest.fixture(
        params=[(None, ST_Direction.HORZ), (ST_Direction.VERT, ST_Direction.VERT)]
    )
    def orient_fixture(self, request):
        orient, expected_orient = request.param
        ph_bldr = a_ph()
        if orient is not None:
            ph_bldr.with_orient(orient)
        shape_elm = (
            an_sp()
            .with_nsdecls("p")
            .with_child(an_nvSpPr().with_child(an_nvPr().with_child(ph_bldr)))
        ).element
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_orient

    @pytest.fixture
    def orient_raise_fixture(self):
        shape_elm = (
            an_sp().with_nsdecls("p").with_child(an_nvSpPr().with_child(an_nvPr()))
        ).element
        shape = BasePlaceholder(shape_elm, None)
        return shape

    @pytest.fixture(
        params=[
            (None, ST_PlaceholderSize.FULL),
            (ST_PlaceholderSize.HALF, ST_PlaceholderSize.HALF),
        ]
    )
    def sz_fixture(self, request):
        sz, expected_sz = request.param
        ph_bldr = a_ph()
        if sz is not None:
            ph_bldr.with_sz(sz)
        shape_elm = (
            an_sp()
            .with_nsdecls("p")
            .with_child(an_nvSpPr().with_child(an_nvPr().with_child(ph_bldr)))
        ).element
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_sz

    @pytest.fixture(
        params=[
            ("sp", None, 1, PP_PLACEHOLDER.OBJECT),
            ("sp", "title", 0, PP_PLACEHOLDER.TITLE),
            ("pic", "pic", 6, PP_PLACEHOLDER.PICTURE),
            ("graphicFrame", "tbl", 9, PP_PLACEHOLDER.TABLE),
        ]
    )
    def type_fixture(self, request):
        tagname, ph_type, idx, expected_ph_type = request.param
        shape_elm = self.shape_elm_factory(tagname, ph_type, idx)
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_ph_type

    # fixture components ---------------------------------------------

    @staticmethod
    def shape_elm_factory(tagname, ph_type, idx):
        root_bldr, nvXxPr_bldr = {
            "sp": (an_sp().with_nsdecls("p"), an_nvSpPr()),
            "pic": (an_sp().with_nsdecls("p"), an_nvPicPr()),
            "graphicFrame": (a_graphicFrame().with_nsdecls("p"), an_nvGraphicFramePr()),
        }[tagname]
        ph_bldr = {
            None: a_ph().with_idx(idx),
            "obj": a_ph().with_idx(idx),
            "title": a_ph().with_type("title"),
            "body": a_ph().with_type("body").with_idx(idx),
            "pic": a_ph().with_type("pic").with_idx(idx),
            "tbl": a_ph().with_type("tbl").with_idx(idx),
        }[ph_type]
        return (
            root_bldr.with_child(nvXxPr_bldr.with_child(an_nvPr().with_child(ph_bldr)))
        ).element


class DescribeChartPlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.ChartPlaceholder` object."""

    def it_can_insert_a_chart_into_itself(self, request, part_prop_):
        slide_part_ = instance_mock(request, SlidePart)
        slide_part_.add_chart_part.return_value = "rId6"
        part_prop_.return_value = slide_part_
        graphicFrame = element("p:graphicFrame")
        _new_chart_graphicFrame_ = method_mock(
            request,
            ChartPlaceholder,
            "_new_chart_graphicFrame",
            return_value=graphicFrame,
        )
        _replace_placeholder_with_ = method_mock(
            request, ChartPlaceholder, "_replace_placeholder_with"
        )
        placeholder_graphic_frame_ = instance_mock(request, PlaceholderGraphicFrame)
        PlaceholderGraphicFrame_ = class_mock(
            request,
            "pptx.shapes.placeholder.PlaceholderGraphicFrame",
            return_value=placeholder_graphic_frame_,
        )
        chart_data_ = instance_mock(request, ChartData)
        chart_ph = ChartPlaceholder(
            element("p:sp/p:spPr/a:xfrm/(a:off{x=1,y=2},a:ext{cx=3,cy=4})"), "parent"
        )

        ph_graphic_frame = chart_ph.insert_chart(XCT.PIE, chart_data_)

        slide_part_.add_chart_part.assert_called_once_with(XCT.PIE, chart_data_)
        _new_chart_graphicFrame_.assert_called_once_with(chart_ph, "rId6", 1, 2, 3, 4)
        _replace_placeholder_with_.assert_called_once_with(chart_ph, graphicFrame)
        PlaceholderGraphicFrame_.assert_called_once_with(graphicFrame, chart_ph._parent)
        assert ph_graphic_frame is placeholder_graphic_frame_

    def it_creates_a_graphicFrame_element_to_help(self, new_fixture):
        chart_ph, rId, x, y, cx, cy, expected_xml = new_fixture
        graphicFrame = chart_ph._new_chart_graphicFrame(rId, x, y, cx, cy)
        assert graphicFrame.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def new_fixture(self):
        sp_cxml = "p:sp/p:nvSpPr/p:cNvPr{id=4,name=bar}"
        chart_ph = ChartPlaceholder(element(sp_cxml), None)
        rId, x, y, cx, cy = "rId42", 1, 2, 3, 4
        expected_xml = snippet_seq("placeholders")[1]
        return chart_ph, rId, x, y, cx, cy, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_prop_(self, request, slide_):
        return property_mock(request, ChartPlaceholder, "part", return_value=slide_)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, SlidePart)


class DescribeLayoutPlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.LayoutPlaceholder` object."""

    def it_uses_InheritsDimensions_mixin(self):
        layout_placeholder = LayoutPlaceholder(None, None)
        assert isinstance(layout_placeholder, _InheritsDimensions)

    def it_finds_its_base_placeholder_to_help(self, base_ph_fixture):
        layout_placeholder, master_, mstr_ph_type, master_placeholder_ = base_ph_fixture
        master_placeholder = layout_placeholder._base_placeholder
        master_.placeholders.get.assert_called_once_with(mstr_ph_type, None)
        assert master_placeholder is master_placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            ("p:sp/p:nvSpPr/p:nvPr/p:ph{type=body}", PP_PLACEHOLDER.BODY),
            ("p:sp/p:nvSpPr/p:nvPr/p:ph{type=tbl}", PP_PLACEHOLDER.BODY),
            ("p:sp/p:nvSpPr/p:nvPr/p:ph{type=title}", PP_PLACEHOLDER.TITLE),
        ]
    )
    def base_ph_fixture(self, request, slide_master_, master_placeholder_, part_prop_):
        sp_cxml, mstr_ph_type = request.param
        sp = element(sp_cxml)
        layout_placeholder = LayoutPlaceholder(sp, None)
        part_prop_.return_value.slide_master = slide_master_
        slide_master_.placeholders.get.return_value = master_placeholder_
        return (layout_placeholder, slide_master_, mstr_ph_type, master_placeholder_)

    # fixture components ---------------------------------------------

    @pytest.fixture
    def master_placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, LayoutPlaceholder, "part")

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)


class DescribeNotesSlidePlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.NotesSlidePlaceholder` object."""

    def it_finds_its_base_placeholder_to_help(self, base_ph_fixture):
        placeholder, notes_master_, ph_type, master_placeholder_ = base_ph_fixture
        base_placeholder = placeholder._base_placeholder
        notes_master_.placeholders.get.assert_called_once_with(ph_type=ph_type)
        assert base_placeholder is master_placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def base_ph_fixture(
        self,
        request,
        notes_master_,
        master_placeholder_,
        part_prop_,
        notes_slide_part_,
        notes_slide_placeholders_,
    ):
        sp_cxml = "p:sp/p:nvSpPr/p:nvPr/p:ph{type=body}"
        placeholder = NotesSlidePlaceholder(element(sp_cxml), None)
        notes_slide_part_.notes_master = notes_master_
        notes_master_.placeholders = notes_slide_placeholders_
        notes_slide_placeholders_.get.return_value = master_placeholder_
        return (placeholder, notes_master_, PP_PLACEHOLDER.BODY, master_placeholder_)

    # fixture components ---------------------------------------------

    @pytest.fixture
    def notes_slide_placeholders_(self, request):
        return instance_mock(request, NotesSlidePlaceholders)

    @pytest.fixture
    def master_placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder)

    @pytest.fixture
    def notes_master_(self, request):
        return instance_mock(request, NotesMaster)

    @pytest.fixture
    def notes_slide_part_(self, request):
        return instance_mock(request, NotesSlidePart)

    @pytest.fixture
    def part_prop_(self, request, notes_slide_part_):
        return property_mock(
            request, NotesSlidePlaceholder, "part", return_value=notes_slide_part_
        )


class DescribePicturePlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.PicturePlaceholder` object."""

    def it_can_insert_a_picture_into_itself(self, request):
        pic = element("p:pic")
        _new_placeholder_pic_ = method_mock(
            request, PicturePlaceholder, "_new_placeholder_pic", return_value=pic
        )
        _replace_placeholder_with_ = method_mock(
            request, PicturePlaceholder, "_replace_placeholder_with"
        )
        placeholder_picture_ = instance_mock(request, PlaceholderPicture)
        PlaceholderPicture_ = class_mock(
            request,
            "pptx.shapes.placeholder.PlaceholderPicture",
            return_value=placeholder_picture_,
        )
        picture_ph = PicturePlaceholder(None, "parent")

        placeholder_picture = picture_ph.insert_picture("foobar.png")

        _new_placeholder_pic_.assert_called_once_with(picture_ph, "foobar.png")
        _replace_placeholder_with_.assert_called_once_with(picture_ph, pic)
        PlaceholderPicture_.assert_called_once_with(pic, picture_ph._parent)
        assert placeholder_picture is placeholder_picture_

    @pytest.mark.parametrize(
        "image_size, crop_attr_names",
        (((444, 333), ("l", "r")), ((333, 444), ("t", "b"))),
    )
    def it_creates_a_pic_element_to_help(self, request, image_size, crop_attr_names):
        _get_or_add_image_ = method_mock(
            request,
            PicturePlaceholder,
            "_get_or_add_image",
            return_value=(42, "bar", image_size),
        )
        picture_ph = PicturePlaceholder(
            element(
                "p:sp/(p:nvSpPr/p:cNvPr{id=2,name=foo},p:spPr/a:xfrm/a:ext{cx=99"
                ",cy=99})"
            ),
            None,
        )

        pic = picture_ph._new_placeholder_pic("foobar.png")

        _get_or_add_image_.assert_called_once_with(picture_ph, "foobar.png")
        assert pic.xml == xml(
            "p:pic/(p:nvPicPr/(p:cNvPr{id=2,name=foo,descr=bar},p:cNvPicPr/a"
            ":picLocks{noGrp=1,noChangeAspect=1},p:nvPr),p:blipFill/(a:blip{"
            "r:embed=42},a:srcRect{%s=12500,%s=12500},a:stretch/a:fillRect),"
            "p:spPr)" % crop_attr_names
        )

    def it_adds_an_image_to_help(self, get_or_add_fixture):
        placeholder, image_file, expected_value = get_or_add_fixture

        value = placeholder._get_or_add_image(image_file)

        placeholder.part.get_or_add_image_part.assert_called_once_with(image_file)
        assert value == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def get_or_add_fixture(self, part_prop_, image_part_):
        placeholder = PicturePlaceholder(None, None)
        image_file, rId, desc, image_size = "f.png", "rId6", "desc", (42, 24)
        part_prop_.return_value.get_or_add_image_part.return_value = (image_part_, rId)
        image_part_.desc, image_part_._px_size = desc, image_size
        expected_value = rId, desc, image_size
        return placeholder, image_file, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)

    @pytest.fixture
    def part_prop_(self, request, slide_):
        return property_mock(request, PicturePlaceholder, "part", return_value=slide_)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, SlidePart)


class DescribeTablePlaceholder(object):
    """Unit-test suite for `pptx.shapes.placeholder.TablePlaceholder` object."""

    def it_can_insert_a_table_into_itself(self, request):
        graphicFrame = element("p:graphicFrame")
        _new_placeholder_table_ = method_mock(
            request,
            TablePlaceholder,
            "_new_placeholder_table",
            return_value=graphicFrame,
        )
        _replace_placeholder_with_ = method_mock(
            request, TablePlaceholder, "_replace_placeholder_with"
        )
        placeholder_graphic_frame_ = instance_mock(request, PlaceholderGraphicFrame)
        PlaceholderGraphicFrame_ = class_mock(
            request,
            "pptx.shapes.placeholder.PlaceholderGraphicFrame",
            return_value=placeholder_graphic_frame_,
        )
        table_ph = TablePlaceholder(None, "parent")

        ph_graphic_frame = table_ph.insert_table(4, 2)

        _new_placeholder_table_.assert_called_once_with(table_ph, 4, 2)
        _replace_placeholder_with_.assert_called_once_with(table_ph, graphicFrame)
        PlaceholderGraphicFrame_.assert_called_once_with(graphicFrame, table_ph._parent)
        assert ph_graphic_frame is placeholder_graphic_frame_

    def it_creates_a_graphicFrame_element_to_help(self):
        table_ph = TablePlaceholder(
            element(
                "p:sp/(p:nvSpPr/p:cNvPr{id=2,name=foo},p:spPr/a:xfrm/(a:off{x=1,y=2},"
                "a:ext{cx=3,cy=4}))"
            ),
            None,
        )
        graphicFrame = table_ph._new_placeholder_table(1, 1)

        assert graphicFrame.xml == snippet_seq("placeholders")[0]
