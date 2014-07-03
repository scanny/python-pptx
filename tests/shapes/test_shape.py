# encoding: utf-8

"""
Test suite for pptx.shapes.shape module
"""

from __future__ import absolute_import

import pytest

from pptx.oxml.shapes.shared import BaseShapeElement
from pptx.oxml.text import CT_TextBody
from pptx.parts.slide import _SlideShapeTree
from pptx.shapes import Subshape
from pptx.shapes.shape import BaseShape
from pptx.text import TextFrame

from ..oxml.unitdata.shape import (
    a_cNvPr, a_cxnSp, a_graphicFrame, a_grpSp, a_grpSpPr, a_p_xfrm, a_pic,
    an_ext, an_nvSpPr, an_off, an_sp, an_spPr, an_xfrm
)
from ..unitutil.mock import (
    class_mock, instance_mock, loose_mock, property_mock
)


class DescribeBaseShape(object):

    def it_knows_its_shape_id(self, id_fixture):
        shape, shape_id = id_fixture
        assert shape.id == shape_id

    def it_knows_its_name(self, name_fixture):
        shape, name = name_fixture
        assert shape.name == name

    def it_has_a_position(self, position_get_fixture):
        shape, expected_left, expected_top = position_get_fixture
        assert shape.left == expected_left
        assert shape.top == expected_top

    def it_has_dimensions(self, dimensions_get_fixture):
        shape, expected_width, expected_height = dimensions_get_fixture
        assert shape.width == expected_width
        assert shape.height == expected_height

    def it_can_change_its_position(self, position_set_fixture):
        shape, left, top, expected_xml = position_set_fixture
        shape.left = left
        shape.top = top
        assert shape._element.xml == expected_xml

    def it_can_change_its_dimensions(self, dimensions_set_fixture):
        shape, width, height, expected_xml = dimensions_set_fixture
        shape.width = width
        shape.height = height
        assert shape._element.xml == expected_xml

    def it_knows_the_part_it_belongs_to(self, part_fixture):
        shape, parent_ = part_fixture
        part = shape.part
        assert part is parent_.part

    def it_knows_whether_it_can_contain_text(self, has_textframe_fixture):
        shape, has_textframe = has_textframe_fixture
        assert shape.has_textframe is has_textframe

    def it_knows_whether_it_is_a_placeholder(self, is_placeholder_fixture):
        shape, is_placeholder = is_placeholder_fixture
        assert shape.is_placeholder is is_placeholder

    def it_knows_it_doesnt_contain_a_chart(self):
        shape = BaseShape(None, None)
        assert shape.has_chart is False

    def it_provides_access_to_its_textframe(self, textframe_fixture):
        shape, TextFrame_, txBody_, textframe_ = textframe_fixture
        textframe = shape.textframe
        TextFrame_.assert_called_once_with(txBody_, shape)
        assert textframe is textframe_

    def it_raises_when_no_textframe(self, no_textframe_fixture):
        shape = no_textframe_fixture
        with pytest.raises(ValueError):
            shape.textframe

    def it_can_set_the_shape_text_to_a_string(self, text_set_fixture):
        shape = text_set_fixture
        shape.text = 'føøbår'
        assert shape.textframe.text == u'føøbår'

    def it_raises_on_assign_text_where_no_textframe(
            self, no_textframe_fixture):
        shape = no_textframe_fixture
        with pytest.raises(TypeError):
            shape.text = 'foobar'

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('sp',           False), ('sp_with_ext',           True),
        ('pic',          False), ('pic_with_ext',          True),
        ('graphicFrame', False), ('graphicFrame_with_ext', True),
        ('grpSp',        False), ('grpSp_with_ext',        True),
        ('cxnSp',        False), ('cxnSp_with_ext',        True),
    ])
    def dimensions_get_fixture(self, request, width, height):
        shape_elm_fixt_name, expect_values = request.param
        shape_elm = request.getfuncargvalue(shape_elm_fixt_name)
        shape = BaseShape(shape_elm, None)
        if not expect_values:
            width = height = None
        return shape, width, height

    @pytest.fixture(params=[
        ('sp',           'sp_with_ext'),
        ('pic',          'pic_with_ext'),
        ('graphicFrame', 'graphicFrame_with_ext'),
        ('grpSp',        'grpSp_with_ext'),
        ('cxnSp',        'cxnSp_with_ext'),
    ])
    def dimensions_set_fixture(self, request, width, height):
        start_elm_fixt_name, expected_elm_fixt_name = request.param
        start_elm = request.getfuncargvalue(start_elm_fixt_name)
        shape = BaseShape(start_elm, None)
        expected_xml = request.getfuncargvalue(expected_elm_fixt_name).xml
        return shape, width, height, expected_xml

    @pytest.fixture
    def id_fixture(self, shape_id):
        shape_elm = (
            an_sp().with_nsdecls().with_child(
                an_nvSpPr().with_child(
                    a_cNvPr().with_id(shape_id)))
        ).element
        shape = BaseShape(shape_elm, None)
        return shape, shape_id

    @pytest.fixture(params=[True, False])
    def has_textframe_fixture(self, request, shape_elm_, txBody_):
        has_textframe = request.param
        shape_elm_.txBody = txBody_ if has_textframe else None
        shape = BaseShape(shape_elm_, None)
        return shape, has_textframe

    @pytest.fixture(params=[True, False])
    def is_placeholder_fixture(self, request, shape_elm_, txBody_):
        is_placeholder = request.param
        shape_elm_.has_ph_elm = is_placeholder
        shape = BaseShape(shape_elm_, None)
        return shape, is_placeholder

    @pytest.fixture
    def name_fixture(self, shape_name):
        shape_elm = (
            an_sp().with_nsdecls().with_child(
                an_nvSpPr().with_child(
                    a_cNvPr().with_name(shape_name)))
        ).element
        shape = BaseShape(shape_elm, None)
        return shape, shape_name

    @pytest.fixture
    def no_textframe_fixture(self, shape_elm_):
        shape_elm_.txBody = None
        shape = BaseShape(shape_elm_, None)
        return shape

    @pytest.fixture
    def part_fixture(self, shapes_):
        parent_ = shapes_
        shape = BaseShape(None, parent_)
        return shape, parent_

    @pytest.fixture(params=[
        ('sp',           False), ('sp_with_off',           True),
        ('pic',          False), ('pic_with_off',          True),
        ('graphicFrame', False), ('graphicFrame_with_off', True),
        ('grpSp',        False), ('grpSp_with_off',        True),
        ('cxnSp',        False), ('cxnSp_with_off',        True),
    ])
    def position_get_fixture(self, request, left, top):
        shape_elm_fixt_name, expect_values = request.param
        shape_elm = request.getfuncargvalue(shape_elm_fixt_name)
        shape = BaseShape(shape_elm, None)
        if not expect_values:
            left = top = None
        return shape, left, top

    @pytest.fixture(params=[
        ('sp',           'sp_with_off'),
        ('pic',          'pic_with_off'),
        ('graphicFrame', 'graphicFrame_with_off'),
        ('grpSp',        'grpSp_with_off'),
        ('cxnSp',        'cxnSp_with_off'),
    ])
    def position_set_fixture(self, request, left, top):
        start_elm_fixt_name, expected_elm_fixt_name = request.param
        start_elm = request.getfuncargvalue(start_elm_fixt_name)
        shape = BaseShape(start_elm, None)
        expected_xml = request.getfuncargvalue(expected_elm_fixt_name).xml
        return shape, left, top, expected_xml

    @pytest.fixture
    def textframe_fixture(self, shape_elm_, TextFrame_, txBody_, textframe_):
        shape = BaseShape(shape_elm_, None)
        return shape, TextFrame_, txBody_, textframe_

    @pytest.fixture
    def text_set_fixture(self, shape_elm_, shape_textframe_):
        shape = BaseShape(shape_elm_, None)
        return shape

    # fixture components ---------------------------------------------

    @pytest.fixture
    def cxnSp(self):
        return a_cxnSp().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def cxnSp_with_ext(self, width, height):
        return (
            a_cxnSp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def cxnSp_with_off(self, left, top):
        return (
            a_cxnSp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def graphicFrame(self):
        # Note that <p:xfrm> element is required on graphicFrame
        return a_graphicFrame().with_nsdecls().with_child(a_p_xfrm()).element

    @pytest.fixture
    def graphicFrame_with_ext(self, width, height):
        return (
            a_graphicFrame().with_nsdecls().with_child(
                a_p_xfrm().with_child(
                    an_ext().with_cx(width).with_cy(height)))
        ).element

    @pytest.fixture
    def graphicFrame_with_off(self, left, top):
        return (
            a_graphicFrame().with_nsdecls().with_child(
                a_p_xfrm().with_child(
                    an_off().with_x(left).with_y(top)))
        ).element

    @pytest.fixture
    def grpSp(self):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr())
        ).element

    @pytest.fixture
    def grpSp_with_ext(self, width, height):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def grpSp_with_off(self, left, top):
        return (
            a_grpSp().with_nsdecls('p', 'a').with_child(
                a_grpSpPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def height(self):
        return 654

    @pytest.fixture
    def left(self):
        return 123

    @pytest.fixture
    def pic(self):
        return a_pic().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def pic_with_off(self, left, top):
        return (
            a_pic().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def pic_with_ext(self, width, height):
        return (
            a_pic().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def shape_elm_(self, request, shape_id, shape_name, txBody_):
        return instance_mock(
            request, BaseShapeElement, shape_id=shape_id,
            shape_name=shape_name, txBody=txBody_
        )

    @pytest.fixture
    def shape_id(self):
        return 42

    @pytest.fixture
    def shape_name(self):
        return 'Foobar 41'

    @pytest.fixture
    def shape_textframe_(self, request):
        return property_mock(request, BaseShape, 'textframe')

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, _SlideShapeTree)

    @pytest.fixture
    def sp(self):
        return an_sp().with_nsdecls().with_child(an_spPr()).element

    @pytest.fixture
    def sp_with_ext(self, width, height):
        return (
            an_sp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width).with_cy(height))))
        ).element

    @pytest.fixture
    def sp_with_off(self, left, top):
        return (
            an_sp().with_nsdecls().with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_off().with_x(left).with_y(top))))
        ).element

    @pytest.fixture
    def TextFrame_(self, request, textframe_):
        return class_mock(
            request, 'pptx.shapes.shape.TextFrame', return_value=textframe_
        )

    @pytest.fixture
    def textframe_(self, request):
        return instance_mock(request, TextFrame)

    @pytest.fixture
    def top(self):
        return 456

    @pytest.fixture
    def txBody_(self, request):
        return instance_mock(request, CT_TextBody)

    @pytest.fixture
    def width(self):
        return 321


class DescribeSubshape(object):

    def it_knows_the_part_it_belongs_to(self, subshape_with_parent_):
        subshape, parent_ = subshape_with_parent_
        part = subshape.part
        assert part is parent_.part

    # fixtures ---------------------------------------------

    @pytest.fixture
    def subshape_with_parent_(self, request):
        parent_ = loose_mock(request, name='parent_')
        subshape = Subshape(parent_)
        return subshape, parent_
