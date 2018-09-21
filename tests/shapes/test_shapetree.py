# encoding: utf-8

"""Unit test suite for pptx.shapes.shapetree module"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.compat import BytesIO
from pptx.chart.data import ChartData
from pptx.enum.shapes import (
    MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, PP_PLACEHOLDER
)
from pptx.oxml import parse_xml
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.shapes.shared import BaseShapeElement, ST_Direction
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.parts.image import ImagePart
from pptx.parts.slide import SlidePart
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import (
    _BaseSlidePlaceholder, LayoutPlaceholder, MasterPlaceholder,
    NotesSlidePlaceholder
)
from pptx.shapes.shapetree import (
    _BaseGroupShapes, BasePlaceholders, BaseShapeFactory, _BaseShapes,
    GroupShapes, LayoutPlaceholders, _LayoutShapeFactory, LayoutShapes,
    MasterPlaceholders, _MasterShapeFactory, MasterShapes,
    _MoviePicElementCreator, NotesSlidePlaceholders, _NotesSlideShapeFactory,
    NotesSlideShapes, _SlidePlaceholderFactory, SlidePlaceholders,
    SlideShapeFactory, SlideShapes
)
from pptx.slide import SlideLayout, SlideMaster
from pptx.table import Table

from ..oxml.unitdata.shape import a_ph, a_pic, an_nvPr, an_nvSpPr, an_sp
from ..unitutil.cxml import element, xml
from ..unitutil.file import snippet_seq
from ..unitutil.mock import (
    ANY, call, class_mock, function_mock, initializer_mock, instance_mock,
    method_mock, property_mock
)


class DescribeBaseShapeFactory(object):

    def it_constructs_the_right_shape_for_an_element(self, factory_fixture):
        shape_elm, parent_, ShapeClass_, shape_ = factory_fixture
        shape = BaseShapeFactory(shape_elm, parent_)
        ShapeClass_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp', Shape),
        ('p:pic', Picture),
        ('p:pic/p:nvPicPr/p:nvPr/a:videoFile', Movie),
        ('p:graphicFrame', GraphicFrame),
        ('p:grpSp', GroupShape),
        ('p:cxnSp', Connector),
    ])
    def factory_fixture(self, request, parent_):
        shape_cxml, ShapeCls = request.param
        shape_elm = element(shape_cxml)
        shape_mock = instance_mock(request, ShapeCls)
        ShapeClass_ = class_mock(
            request, 'pptx.shapes.shapetree.%s' % ShapeCls.__name__,
            return_value=shape_mock
        )
        return shape_elm, parent_, ShapeClass_, shape_mock

    # fixture components -----------------------------------

    @pytest.fixture
    def parent_(self, request):
        return instance_mock(request, SlideShapes)


class Describe_BaseShapes(object):

    def it_knows_how_many_shapes_it_contains(self, len_fixture):
        shapes, expected_count = len_fixture
        assert len(shapes) == expected_count

    def it_can_iterate_over_the_shapes_it_contains(self, iter_fixture):
        shapes, expected_shapes, BaseShapeFactory_, calls, = iter_fixture
        assert [s for s in shapes] == expected_shapes
        assert BaseShapeFactory_.call_args_list == calls

    def it_iterates_shape_elements_to_help__iter__(self, iter_elms_fixture):
        shapes, expected_elms = iter_elms_fixture
        assert [e for e in shapes._iter_member_elms()] == expected_elms

    def it_supports_indexed_access(self, getitem_fixture):
        shapes, idx, BaseShapeFactory_, sp, shape_ = getitem_fixture
        shape = shapes[idx]
        BaseShapeFactory_.assert_called_once_with(sp, shapes)
        assert shape is shape_

    def it_raises_on_shape_index_out_of_range(self, getitem_raises_fixture):
        shapes = getitem_raises_fixture
        with pytest.raises(IndexError):
            shapes[2]

    def it_can_clone_a_placeholder(self, clone_ph_fixture):
        shapes, placeholder_, expected_xml = clone_ph_fixture
        shapes.clone_placeholder(placeholder_)
        assert shapes._element.xml == expected_xml

    def it_knows_if_turbo_add_is_enabled(self, turbo_fixture):
        shapes, expected_value = turbo_fixture
        turbo_add_enabled = shapes.turbo_add_enabled
        assert turbo_add_enabled == expected_value

    def it_can_change_turbo_add_enabled(self, turbo_set_fixture):
        shapes, value, expected_value = turbo_set_fixture
        shapes.turbo_add_enabled = value
        assert shapes.turbo_add_enabled == expected_value

    def it_finds_the_next_shape_id_to_help(self, next_id_fixture):
        shapes, expected_value = next_id_fixture
        assert shapes._next_shape_id == expected_value

    def it_finds_the_next_placeholder_name_to_help(self, ph_name_fixture):
        shapes, ph_type, sp_id, orient, expected_value = ph_name_fixture
        assert shapes._next_ph_name(ph_type, sp_id, orient) == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def clone_ph_fixture(self, placeholder_):
        shapes = SlideShapes(element('p:spTree{a:b=c}'), None)
        expected_xml = xml(
            'p:spTree{a:b=c}/p:sp/(p:nvSpPr/(p:cNvPr{id=1,name=Vertical Char'
            't Placeholder 0},p:cNvSpPr/a:spLocks{noGrp=1},p:nvPr/p:ph{type='
            'chart,idx=42,orient=vert,sz=half}),p:spPr)'
        )
        placeholder_.element = element(
            'p:sp/p:nvSpPr/p:nvPr/p:ph{type=chart,idx=42,orient=vert,sz=half'
            '}'
        )
        return shapes, placeholder_, expected_xml

    @pytest.fixture
    def getitem_fixture(self, BaseShapeFactory_, shape_):
        spTree = element('p:spTree/(p:sp,p:sp)')
        shapes = _BaseShapes(spTree, None)
        idx = 1
        sp = spTree.xpath('p:sp')[idx]
        return shapes, idx, BaseShapeFactory_, sp, shape_

    @pytest.fixture
    def getitem_raises_fixture(self):
        spTree = element('p:spTree/(p:sp,p:sp)')
        shapes = _BaseShapes(spTree, None)
        return shapes

    @pytest.fixture
    def iter_fixture(self, BaseShapeFactory_):
        spTree = element('p:spTree/(p:sp,p:sp)')
        sps = spTree.xpath('p:sp')
        shapes = _BaseShapes(spTree, None)
        expected_shapes = [Shape(None, None), Shape(None, None)]
        calls = [call(sps[0], shapes), call(sps[1], shapes)]
        BaseShapeFactory_.side_effect = iter(expected_shapes)
        return shapes, expected_shapes, BaseShapeFactory_, calls

    @pytest.fixture
    def iter_elms_fixture(self):
        spTree = element('p:spTree/(p:sp,p:sp)')
        sps = spTree.xpath('p:sp')
        shapes = _BaseShapes(spTree, None)
        return shapes, sps

    @pytest.fixture
    def len_fixture(self):
        shapes = _BaseShapes(element('p:spTree/(p:spPr,p:sp,p:sp)'), None)
        expected_count = 2
        return shapes, expected_count

    @pytest.fixture(params=[
        ('p:spTree/p:nvSpPr',                                 1),
        ('p:spTree/p:nvSpPr/p:cNvPr{id=0}',                   1),
        ('p:spTree/p:nvSpPr/p:cNvPr{id=1}',                   2),
        ('p:spTree/p:nvSpPr/p:cNvPr{id=2}',                   3),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=1},p:cNvPr{id=3})',   4),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=foo},p:cNvPr{id=2})', 3),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=1fo},p:cNvPr{id=2})', 3),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=1},p:cNvPr{id=1},p:'
         'cNvPr{id=1},p:cNvPr{id=4})',                        5),
    ])
    def next_id_fixture(self, request):
        spTree_cxml, expected_value = request.param
        shapes = _BaseShapes(element(spTree_cxml), None)
        return shapes, expected_value

    @pytest.fixture(params=[
        (PP_PLACEHOLDER.OBJECT, 3, ST_Direction.HORZ,
         'Content Placeholder 2'),
        (PP_PLACEHOLDER.TABLE,  4, ST_Direction.HORZ,
         'Table Placeholder 4'),
        (PP_PLACEHOLDER.TABLE,  7, ST_Direction.VERT,
         'Vertical Table Placeholder 6'),
        (PP_PLACEHOLDER.TITLE,  2, ST_Direction.HORZ,
         'Title 2'),
    ])
    def ph_name_fixture(self, request):
        ph_type, sp_id, orient, expected_name = request.param
        spTree = element(
            'p:spTree/(p:cNvPr{name=Title 1},p:cNvPr{name=Table Placeholder '
            '3})'
        )
        shapes = SlideShapes(spTree, None)
        return shapes, ph_type, sp_id, orient, expected_name

    @pytest.fixture(params=[
        (None, False),
        (42, True),
    ])
    def turbo_fixture(self, request):
        cached_max_shape_id, expected_value = request.param
        shapes = _BaseShapes(None, None)
        if cached_max_shape_id:
            shapes._cached_max_shape_id = cached_max_shape_id
        return shapes, expected_value

    @pytest.fixture(params=[
        ('p:spTree/p:nvSpPr',                                 True),
        ('p:spTree/p:nvSpPr/p:cNvPr{id=2}',                   True),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=1},p:cNvPr{id=3})',   False),
        ('p:spTree/p:nvSpPr/(p:cNvPr{id=1},p:cNvPr{id=1},p:'
         'cNvPr{id=1},p:cNvPr{id=4})',                        True),
    ])
    def turbo_set_fixture(self, request):
        spTree_cxml, value = request.param
        shapes = _BaseShapes(element(spTree_cxml), None)
        expected_value = value
        return shapes, value, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree.BaseShapeFactory',
            return_value=shape_, autospec=True
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, BaseShape)


class Describe_BaseGroupShapes(object):

    def it_can_add_a_chart(self, add_chart_fixture):
        shapes, chart_type, x, y, cx, cy, chart_data_ = add_chart_fixture[:7]
        rId_, graphicFrame, graphic_frame_ = add_chart_fixture[7:]

        graphic_frame = shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data_
        )

        shapes.part.add_chart_part.assert_called_once_with(
            chart_type, chart_data_
        )
        shapes._add_chart_graphicFrame.assert_called_once_with(
            shapes, rId_, x, y, cx, cy
        )
        shapes._recalculate_extents.assert_called_once_with(shapes)
        shapes._shape_factory.assert_called_once_with(shapes, graphicFrame)
        assert graphic_frame is graphic_frame_

    def it_can_add_a_connector_shape(self, connector_fixture):
        shapes, connector_type, begin_x, begin_y = connector_fixture[:4]
        end_x, end_y, cxnSp_, connector_ = connector_fixture[4:]

        connector = shapes.add_connector(
            connector_type, begin_x, begin_y, end_x, end_y
        )

        shapes._add_cxnSp.assert_called_once_with(
            shapes, connector_type, begin_x, begin_y, end_x, end_y
        )
        shapes._recalculate_extents.assert_called_once_with(shapes)
        shapes._shape_factory.assert_called_once_with(shapes, cxnSp_)
        assert connector is connector_

    def it_can_provide_a_freeform_builder(self, freeform_fixture):
        shapes, start_x, start_y, scale = freeform_fixture[:4]
        FreeformBuilder_new_, x_scale, y_scale = freeform_fixture[4:7]
        builder_ = freeform_fixture[7]

        builder = shapes.build_freeform(start_x, start_y, scale)

        FreeformBuilder_new_.assert_called_once_with(
            shapes, start_x, start_y, x_scale, y_scale
        )
        assert builder is builder_

    def it_can_add_a_group_shape(self, group_fixture):
        shapes, spTree, grpSp, group_shape_ = group_fixture

        group_shape = shapes.add_group_shape()

        spTree.add_grpSp.assert_called_once_with(spTree)
        shapes._shape_factory.assert_called_once_with(shapes, grpSp)
        assert group_shape is group_shape_

    def it_can_add_a_picture(self, picture_fixture):
        shapes, image_file, x, y, cx, cy = picture_fixture[:6]
        image_part_, rId, pic, picture_ = picture_fixture[6:]

        picture = shapes.add_picture(image_file, x, y, cx, cy)

        shapes.part.get_or_add_image_part.assert_called_once_with(
            image_file
        )
        shapes._add_pic_from_image_part.assert_called_once_with(
            shapes, image_part_, rId, x, y, cx, cy
        )
        shapes._recalculate_extents.assert_called_once_with(shapes)
        shapes._shape_factory.assert_called_once_with(shapes, pic)
        assert picture is picture_

    def it_can_add_a_shape(self, shape_fixture):
        shapes, autoshape_type_id, x, y, cx, cy = shape_fixture[:6]
        AutoShapeType_, autoshape_type_, sp, shape_ = shape_fixture[6:]

        shape = shapes.add_shape(autoshape_type_id, x, y, cx, cy)

        AutoShapeType_.assert_called_once_with(autoshape_type_id)
        shapes._add_sp.assert_called_once_with(
            shapes, autoshape_type_, x, y, cx, cy
        )
        shapes._recalculate_extents.assert_called_once_with(shapes)
        shapes._shape_factory.assert_called_once_with(shapes, sp)
        assert shape is shape_

    def it_can_add_a_textbox(self, textbox_fixture):
        shapes, x, y, cx, cy, sp, shape_ = textbox_fixture

        shape = shapes.add_textbox(x, y, cx, cy)

        shapes._add_textbox_sp.assert_called_once_with(shapes, x, y, cx, cy)
        shapes._recalculate_extents.assert_called_once_with(shapes)
        shapes._shape_factory.assert_called_once_with(shapes, sp)
        assert shape is shape_

    def it_knows_the_index_of_each_of_its_shapes(self, index_fixture):
        shapes, shape_, expected_value = index_fixture
        assert shapes.index(shape_) == expected_value

    def it_raises_on_index_where_shape_not_found(self, index_raises_fixture):
        shapes, shape_ = index_raises_fixture
        with pytest.raises(ValueError):
            shapes.index(shape_)

    def it_adds_a_chart_graphicFrame_to_help(self, add_cht_gr_frm_fixture):
        shapes, rId, x, y, cx, cy, expected_xml = add_cht_gr_frm_fixture

        graphicFrame = shapes._add_chart_graphicFrame(rId, x, y, cx, cy)

        assert shapes._element.xml == expected_xml
        assert graphicFrame is shapes._element.xpath('p:graphicFrame')[0]

    def it_adds_a_cxnSp_to_help(self, add_cxnSp_fixture):
        shapes, connector_type, begin_x, begin_y = add_cxnSp_fixture[:4]
        end_x, end_y, expected_xml = add_cxnSp_fixture[4:]

        cxnSp = shapes._add_cxnSp(
            connector_type, begin_x, begin_y, end_x, end_y
        )

        assert cxnSp is shapes._element.xpath('p:cxnSp')[0]
        assert cxnSp.xml == expected_xml

    def it_adds_a_pic_element_to_help(self, add_pic_fixture):
        shapes, image_part_, rId, x, y, cx, cy = add_pic_fixture[:7]
        expected_xml = add_pic_fixture[7]

        pic = shapes._add_pic_from_image_part(image_part_, rId, x, y, cx, cy)

        image_part_.scale.assert_called_once_with(cx, cy)
        assert shapes._element.xml == expected_xml
        assert pic is shapes._element.xpath('p:pic')[0]

    def it_adds_an_sp_element_to_help(self, add_sp_fixture):
        shapes, autoshape_type_, x, y, cx, cy, expected_xml = add_sp_fixture

        sp = shapes._add_sp(autoshape_type_, x, y, cx, cy)

        assert shapes._element.xml == expected_xml
        assert sp is shapes._element.xpath('p:sp')[0]

    def it_adds_a_textbox_sp_element_to_help(self, add_textbox_sp_fixture):
        shapes, x, y, cx, cy, expected_xml = add_textbox_sp_fixture

        sp = shapes._add_textbox_sp(x, y, cx, cy)

        assert shapes._element.xml == expected_xml
        assert sp is shapes._element.xpath('p:sp')[0]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_chart_fixture(
            self, chart_data_, _add_chart_graphicFrame_, graphic_frame_,
            part_prop_, slide_part_, _recalculate_extents_, _shape_factory_):
        shapes = _BaseGroupShapes(None, None)
        chart_type = 0
        rId, x, y, cx, cy = 'rId42', 1, 2, 3, 4
        graphicFrame = element('p:graphicFrame')

        part_prop_.return_value = slide_part_
        slide_part_.add_chart_part.return_value = rId
        _add_chart_graphicFrame_.return_value = graphicFrame
        _shape_factory_.return_value = graphic_frame_

        return (
            shapes, chart_type, x, y, cx, cy, chart_data_, rId,
            graphicFrame, graphic_frame_
        )

    @pytest.fixture
    def add_cht_gr_frm_fixture(self):
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        rId, x, y, cx, cy = 'rId42', 1, 2, 3, 4
        expected_xml = (
            '<p:spTree xmlns:p="http://schemas.openxmlformats.org/presentati'
            'onml/2006/main">\n  <p:graphicFrame xmlns:a="http://schemas.ope'
            'nxmlformats.org/drawingml/2006/main">\n    <p:nvGraphicFramePr>'
            '\n      <p:cNvPr id="1" name="Chart 0"/>\n      <p:cNvGraphicFr'
            'amePr>\n        <a:graphicFrameLocks noGrp="1"/>\n      </p:cNv'
            'GraphicFramePr>\n      <p:nvPr/>\n    </p:nvGraphicFramePr>\n  '
            '  <p:xfrm>\n      <a:off x="1" y="2"/>\n      <a:ext cx="3" cy='
            '"4"/>\n    </p:xfrm>\n    <a:graphic>\n      <a:graphicData uri'
            '="http://schemas.openxmlformats.org/drawingml/2006/chart">\n   '
            '     <c:chart xmlns:c="http://schemas.openxmlformats.org/drawin'
            'gml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/offi'
            'ceDocument/2006/relationships" r:id="rId42"/>\n      </a:graphi'
            'cData>\n    </a:graphic>\n  </p:graphicFrame>\n</p:spTree>'
        )
        return shapes, rId, x, y, cx, cy, expected_xml

    @pytest.fixture(params=[
        (1, 2, 3, 5,
         'p:spPr/(a:xfrm/(a:off{x=1,y=2},a:ext{cx=2,cy=3})'),
        (8, 3, 4, 9,
         'p:spPr/(a:xfrm{flipH=1}/(a:off{x=4,y=3},a:ext{cx=4,cy=6})'),
        (1, 6, 5, 2,
         'p:spPr/(a:xfrm{flipV=1}/(a:off{x=1,y=2},a:ext{cx=4,cy=4})'),
        (9, 8, 2, 3,
         'p:spPr/(a:xfrm{flipH=1,flipV=1}/(a:off{x=2,y=3},a:ext{cx=7,cy=5})'),
    ])
    def add_cxnSp_fixture(self, request):
        begin_x, begin_y, end_x, end_y, spPr_cxml = request.param
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        connector_type = MSO_CONNECTOR.STRAIGHT
        tmpl_cxml = (
            'p:cxnSp/(p:nvCxnSpPr/(p:cNvPr{id=1,name=Connector 0},p:cNvCxnSp'
            'Pr,p:nvPr),%s,a:prstGeom{prst=line}/a:avLst),p:style/(a:lnRef{i'
            'dx=2}/a:schemeClr{val=accent1},a:fillRef{idx=0}/a:schemeClr{val'
            '=accent1},a:effectRef{idx=1}/a:schemeClr{val=accent1},a:fontRef'
            '{idx=minor}/a:schemeClr{val=tx1}))'
        )
        expected_xml = xml(tmpl_cxml % spPr_cxml)
        return (
            shapes, connector_type, begin_x, begin_y, end_x, end_y,
            expected_xml
        )

    @pytest.fixture
    def add_pic_fixture(self, image_part_, _next_shape_id_prop_):
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        rId, x, y, cx, cy = 'rId24', 10, 11, 12, 13

        _next_shape_id_prop_.return_value = 42
        image_part_.scale.return_value = (101, 102)
        image_part_.desc = 'sprocket.jpg'
        expected_xml = (
            '<p:spTree xmlns:p="http://schemas.openxmlformats.org/presentati'
            'onml/2006/main">\n  <p:pic xmlns:a="http://schemas.openxmlforma'
            'ts.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlform'
            'ats.org/officeDocument/2006/relationships">\n    <p:nvPicPr>\n '
            '     <p:cNvPr id="42" name="Picture 41" descr="sprocket.jpg"/>'
            '\n      <p:cNvPicPr>\n        <a:picLocks noChangeAspect="1"/>'
            '\n      </p:cNvPicPr>\n      <p:nvPr/>\n    </p:nvPicPr>\n    <'
            'p:blipFill>\n      <a:blip r:embed="rId24"/>\n      <a:stretch>'
            '\n        <a:fillRect/>\n      </a:stretch>\n    </p:blipFill>'
            '\n    <p:spPr>\n      <a:xfrm>\n        <a:off x="10" y="11"/>'
            '\n        <a:ext cx="101" cy="102"/>\n      </a:xfrm>\n      <a'
            ':prstGeom prst="rect">\n        <a:avLst/>\n      </a:prstGeom>'
            '\n    </p:spPr>\n  </p:pic>\n</p:spTree>'
        )
        return shapes, image_part_, rId, x, y, cx, cy, expected_xml

    @pytest.fixture
    def add_sp_fixture(self, autoshape_type_, _next_shape_id_prop_):
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        x, y, cx, cy = 8, 7, 6, 5

        _next_shape_id_prop_.return_value = 7
        autoshape_type_.basename = 'Rounded Rectangle'
        autoshape_type_.prst = 'roundRect'

        expected_xml = (
            '<p:spTree xmlns:p="http://schemas.openxmlformats.org/presentati'
            'onml/2006/main">\n  <p:sp xmlns:a="http://schemas.openxmlformat'
            's.org/drawingml/2006/main">\n    <p:nvSpPr>\n      <p:cNvPr id='
            '"7" name="Rounded Rectangle 6"/>\n      <p:cNvSpPr/>\n      <p:'
            'nvPr/>\n    </p:nvSpPr>\n    <p:spPr>\n      <a:xfrm>\n        '
            '<a:off x="8" y="7"/>\n        <a:ext cx="6" cy="5"/>\n      </a'
            ':xfrm>\n      <a:prstGeom prst="roundRect">\n        <a:avLst/>'
            '\n      </a:prstGeom>\n    </p:spPr>\n    <p:style>\n      <a:l'
            'nRef idx="1">\n        <a:schemeClr val="accent1"/>\n      </a:'
            'lnRef>\n      <a:fillRef idx="3">\n        <a:schemeClr val="ac'
            'cent1"/>\n      </a:fillRef>\n      <a:effectRef idx="2">\n    '
            '    <a:schemeClr val="accent1"/>\n      </a:effectRef>\n      <'
            'a:fontRef idx="minor">\n        <a:schemeClr val="lt1"/>\n     '
            ' </a:fontRef>\n    </p:style>\n    <p:txBody>\n      <a:bodyPr '
            'rtlCol="0" anchor="ctr"/>\n      <a:lstStyle/>\n      <a:p>\n  '
            '      <a:pPr algn="ctr"/>\n      </a:p>\n    </p:txBody>\n  </p'
            ':sp>\n</p:spTree>'
        )
        return shapes, autoshape_type_, x, y, cx, cy, expected_xml

    @pytest.fixture
    def add_textbox_sp_fixture(self, _next_shape_id_prop_):
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        x, y, cx, cy = 1, 2, 3, 4

        _next_shape_id_prop_.return_value = 6

        expected_xml = (
            '<p:spTree xmlns:p="http://schemas.openxmlformats.org/presentati'
            'onml/2006/main">\n  <p:sp xmlns:a="http://schemas.openxmlformat'
            's.org/drawingml/2006/main">\n    <p:nvSpPr>\n      <p:cNvPr id='
            '"6" name="TextBox 5"/>\n      <p:cNvSpPr txBox="1"/>\n      <p:'
            'nvPr/>\n    </p:nvSpPr>\n    <p:spPr>\n      <a:xfrm>\n        '
            '<a:off x="1" y="2"/>\n        <a:ext cx="3" cy="4"/>\n      </a'
            ':xfrm>\n      <a:prstGeom prst="rect">\n        <a:avLst/>\n   '
            '   </a:prstGeom>\n      <a:noFill/>\n    </p:spPr>\n    <p:txBo'
            'dy>\n      <a:bodyPr wrap="none">\n        <a:spAutoFit/>\n    '
            '  </a:bodyPr>\n      <a:lstStyle/>\n      <a:p/>\n    </p:txBod'
            'y>\n  </p:sp>\n</p:spTree>'
        )
        return shapes, x, y, cx, cy, expected_xml

    @pytest.fixture
    def connector_fixture(self, _add_cxnSp_, _shape_factory_,
                          _recalculate_extents_, connector_):
        shapes = _BaseGroupShapes(element('p:spTree'), None)
        connector_type = MSO_CONNECTOR.STRAIGHT
        begin_x, begin_y, end_x, end_y = 1, 2, 3, 4
        cxnSp = element('p:cxnSp')

        _add_cxnSp_.return_value = cxnSp
        _shape_factory_.return_value = connector_

        return (
            shapes, connector_type, begin_x, begin_y, end_x, end_y, cxnSp,
            connector_
        )

    @pytest.fixture(params=[
        (0,   0,   1.0,        1.0, 1.0),
        (100, 200, 2.0,        2.0, 2.0),
        (100, 200, (4.2, 2.4), 4.2, 2.4),
    ])
    def freeform_fixture(self, request, FreeformBuilder_new_, builder_):
        start_x, start_y, scale, x_scale, y_scale = request.param
        shapes = _BaseGroupShapes(None, None)
        FreeformBuilder_new_.return_value = builder_
        return (
            shapes, start_x, start_y, scale, FreeformBuilder_new_, x_scale,
            y_scale, builder_
        )

    @pytest.fixture
    def group_fixture(self, CT_GroupShape_add_grpSp_, _shape_factory_,
                      group_shape_):
        spTree = element('p:spTree{id=2e838acdc755e83113ed03904d2fe081f}')
        grpSp = element('p:grpSp{id=052874e154b48f9bec4266f80913cae38f}')
        shapes = _BaseGroupShapes(spTree, None)

        CT_GroupShape_add_grpSp_.return_value = grpSp
        _shape_factory_.return_value = group_shape_

        return shapes, spTree, grpSp, group_shape_

    @pytest.fixture(params=[
        ('p:spTree/(p:sp,p:sp,p:sp)', SlideShapes, 0),
        ('p:spTree/(p:sp,p:sp,p:sp)', SlideShapes, 1),
        ('p:spTree/(p:sp,p:sp,p:sp)', SlideShapes, 2),
        ('p:grpSp/(p:sp,p:sp,p:sp)', GroupShapes, 0),
        ('p:grpSp/(p:sp,p:sp,p:sp)', GroupShapes, 1),
        ('p:grpSp/(p:sp,p:sp,p:sp)', GroupShapes, 2),
    ])
    def index_fixture(self, request, shape_):
        grpSp_cxml, ShapesCls, idx = request.param
        grpSp = element(grpSp_cxml)
        sps = grpSp.xpath('p:sp')
        shapes = ShapesCls(grpSp, None)
        shape_.element = sps[idx]
        expected_value = idx
        return shapes, shape_, expected_value

    @pytest.fixture(params=[
        ('p:spTree/(p:sp,p:sp,p:sp)', SlideShapes),
        ('p:grpSp/(p:sp,p:sp,p:sp)',  GroupShapes),
    ])
    def index_raises_fixture(self, request, shape_):
        grpSp_cxml, ShapesCls = request.param
        shapes = SlideShapes(element(grpSp_cxml), None)
        shape_.element = element('p:sp')
        return shapes, shape_

    @pytest.fixture
    def _next_shape_id_prop_(self, request):
        return property_mock(request, _BaseGroupShapes, '_next_shape_id')

    @pytest.fixture
    def picture_fixture(self, part_prop_, slide_part_, image_part_,
                        _add_pic_from_image_part_, _recalculate_extents_,
                        _shape_factory_, picture_):
        shapes = _BaseGroupShapes(None, None)
        image_file, x, y, cx, cy, rId = 'foobar.png', 1, 2, 3, 4, 'rId42'
        pic = element('p:pic')

        part_prop_.return_value = slide_part_
        slide_part_.get_or_add_image_part.return_value = image_part_, rId
        _add_pic_from_image_part_.return_value = pic
        _shape_factory_.return_value = picture_
        return (
            shapes, image_file, x, y, cx, cy, image_part_, rId, pic, picture_
        )

    @pytest.fixture
    def shape_fixture(self, AutoShapeType_, autoshape_type_, _add_sp_,
                      _recalculate_extents_, _shape_factory_, shape_):
        shapes = _BaseGroupShapes(None, None)
        autoshape_type_id = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        x, y, cx, cy = 21, 22, 23, 24
        sp = element('p:sp')

        AutoShapeType_.return_value = autoshape_type_
        _add_sp_.return_value = sp
        _shape_factory_.return_value = shape_

        return (
            shapes, autoshape_type_id, x, y, cx, cy, AutoShapeType_,
            autoshape_type_, sp, shape_
        )

    @pytest.fixture
    def textbox_fixture(self, _add_textbox_sp_, _recalculate_extents_,
                        _shape_factory_, shape_):
        shapes = _BaseGroupShapes(None, None)
        x, y, cx, cy = 31, 32, 33, 34
        sp = element('p:sp')

        _add_textbox_sp_.return_value = sp
        _shape_factory_.return_value = shape_

        return shapes, x, y, cx, cy, sp, shape_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _add_chart_graphicFrame_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_add_chart_graphicFrame',
            autospec=True
        )

    @pytest.fixture
    def _add_cxnSp_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_add_cxnSp', autospec=True
        )

    @pytest.fixture
    def _add_pic_from_image_part_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_add_pic_from_image_part',
            autospec=True
        )

    @pytest.fixture
    def _add_sp_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_add_sp', autospec=True
        )

    @pytest.fixture
    def _add_textbox_sp_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_add_textbox_sp', autospec=True
        )

    @pytest.fixture
    def autoshape_type_(self, request):
        return instance_mock(request, AutoShapeType)

    @pytest.fixture
    def AutoShapeType_(self, request):
        return class_mock(request, 'pptx.shapes.shapetree.AutoShapeType')

    @pytest.fixture
    def builder_(self, request):
        return instance_mock(request, FreeformBuilder)

    @pytest.fixture
    def chart_data_(self, request):
        return instance_mock(request, ChartData)

    @pytest.fixture
    def connector_(self, request):
        return instance_mock(request, Connector)

    @pytest.fixture
    def CT_GroupShape_add_grpSp_(self, request):
        return method_mock(
            request, CT_GroupShape, 'add_grpSp', autospec=True
        )

    @pytest.fixture
    def FreeformBuilder_new_(self, request):
        return method_mock(request, FreeformBuilder, 'new')

    @pytest.fixture
    def graphic_frame_(self, request):
        return instance_mock(request, GraphicFrame)

    @pytest.fixture
    def group_shape_(self, request):
        return instance_mock(request, GroupShape)

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, _BaseGroupShapes, 'part')

    @pytest.fixture
    def picture_(self, request):
        return instance_mock(request, Picture)

    @pytest.fixture
    def _recalculate_extents_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_recalculate_extents', autospec=True
        )

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def _shape_factory_(self, request):
        return method_mock(
            request, _BaseGroupShapes, '_shape_factory', autospec=True
        )

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribeGroupShapes(object):

    def it_recalculates_its_extents_to_help(self, recalc_fixture):
        shapes = recalc_fixture
        shapes._recalculate_extents()
        shapes._grpSp.recalculate_extents.assert_called_once_with()

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def recalc_fixture(self, grpSp_):
        return GroupShapes(grpSp_, None)

    # fixture components ---------------------------------------------

    @pytest.fixture
    def grpSp_(self, request):
        return instance_mock(request, CT_GroupShape)


class DescribeBasePlaceholders(object):

    def it_contains_only_placeholder_shapes(self, member_fixture):
        shape_elm_, is_ph_shape = member_fixture
        _is_ph_shape = BasePlaceholders._is_member_elm(shape_elm_)
        assert _is_ph_shape == is_ph_shape

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[True, False])
    def member_fixture(self, request, shape_elm_):
        is_ph_shape = request.param
        shape_elm_.has_ph_elm = is_ph_shape
        return shape_elm_, is_ph_shape

    # fixture components ---------------------------------------------

    @pytest.fixture
    def shape_elm_(self, request):
        return instance_mock(request, BaseShapeElement)


class DescribeNotesSlidePlaceholders(object):

    def it_brokers_access_to_its_shape_factory(self, factory_fixture):
        placeholders, sp, _NotesSlideShapeFactory_, placeholder_ = (
            factory_fixture
        )
        placeholder = placeholders._shape_factory(sp)
        _NotesSlideShapeFactory_.assert_called_once_with(sp, placeholders)
        assert placeholder is placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def factory_fixture(self, _NotesSlideShapeFactory_, placeholder_):
        placeholders = NotesSlidePlaceholders(None, None)
        sp = element('p:sp')
        return placeholders, sp, _NotesSlideShapeFactory_, placeholder_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _NotesSlideShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._NotesSlideShapeFactory',
            return_value=placeholder_, autospec=True
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, NotesSlidePlaceholder)


class DescribeNotesSlideShapes(object):

    def it_knows_notes_slide_placeholder_basenames(self, basename_fixture):
        notes_slide_shapes, ph_type, expected_value = basename_fixture
        assert notes_slide_shapes.ph_basename(ph_type) == expected_value

    def it_brokers_access_to_its_shape_factory(self, factory_fixture):
        shapes, sp, _NotesSlideShapeFactory_, shape_ = factory_fixture
        shape = shapes._shape_factory(sp)
        _NotesSlideShapeFactory_.assert_called_once_with(sp, shapes)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        (PP_PLACEHOLDER.BODY,   'Notes Placeholder'),
        (PP_PLACEHOLDER.HEADER, 'Header Placeholder'),
    ])
    def basename_fixture(self, request):
        ph_type, expected_value = request.param
        notes_slide_shapes = NotesSlideShapes(None, None)
        return notes_slide_shapes, ph_type, expected_value

    @pytest.fixture
    def factory_fixture(self, _NotesSlideShapeFactory_, shape_):
        shapes = NotesSlideShapes(None, None)
        sp = element('p:sp')
        return shapes, sp, _NotesSlideShapeFactory_, shape_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _NotesSlideShapeFactory_(self, request, shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree._NotesSlideShapeFactory',
            return_value=shape_, autospec=True
        )

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, Shape)


class DescribeSlidePlaceholders(object):

    def it_can_get_a_placeholder_by_idx(self, getitem_fixture):
        placeholders, idx, SlideShapeFactory_ = getitem_fixture[:3]
        shape_elm, placeholder_ = getitem_fixture[3:]

        placeholder = placeholders[idx]

        SlideShapeFactory_.assert_called_once_with(shape_elm, placeholders)
        assert placeholder is placeholder_

    def it_can_iterate_over_its_placeholders(self, iter_fixture):
        placeholders, SlideShapeFactory_ = iter_fixture[:2]
        expected_calls, expected_values = iter_fixture[2:]

        ps = [p for p in placeholders]

        assert SlideShapeFactory_.call_args_list == expected_calls
        assert ps == expected_values

    def it_knows_how_many_placeholders_it_contains(self, len_fixture):
        placeholders, expected_value = len_fixture
        assert len(placeholders) == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}',    1, 0),
        ('p:spTree/p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1}',  1, 0),
        ('p:spTree/(p:sp,p:sp/p:nvSpPr/p:nvPr/p:ph{type=title})', 0, 1),
        ('p:spTree/(p:sp,p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1})',
         1, 1),
        ('p:spTree/(p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=3})', 3, 1),
    ])
    def getitem_fixture(self, request, SlideShapeFactory_, placeholder_):
        spTree_cxml, idx, offset = request.param
        spTree = element(spTree_cxml)
        placeholders = SlidePlaceholders(spTree, None)
        shape_elm = spTree[offset]
        SlideShapeFactory_.return_value = placeholder_
        return placeholders, idx, SlideShapeFactory_, shape_elm, placeholder_

    @pytest.fixture(params=[
        ('p:spTree/('
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=body,idx=1},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=3})', (1, 0, 2)),
    ])
    def iter_fixture(self, request, SlideShapeFactory_, placeholder_):
        spTree_cxml, sequence = request.param
        spTree = element(spTree_cxml)
        placeholders = SlidePlaceholders(spTree, None)
        SlideShapeFactory_.return_value = placeholder_
        calls = [call(spTree[i], placeholders) for i in sequence]
        values = [placeholder_] * len(sequence)
        return placeholders, SlideShapeFactory_, calls, values

    @pytest.fixture(params=[
        ('p:spTree',                                                    0),
        ('p:spTree/(p:sp,p:pic,p:sp)',                                  0),
        ('p:spTree/(p:sp,p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},p:pic)', 1),
        ('p:spTree/('
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=body,idx=1},'
         'p:sp/p:nvSpPr/p:nvPr/p:ph{type=title},'
         'p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=3})',                3),
    ])
    def len_fixture(self, request):
        spTree_cxml, length = request.param
        placeholders = SlidePlaceholders(element(spTree_cxml), None)
        return placeholders, length

    # fixture components ---------------------------------------------

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def SlideShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree.SlideShapeFactory',
            return_value=placeholder_
        )


class Describe_SlidePlaceholderFactory(object):

    def it_constructs_the_right_type_of_placeholder(self, factory_fixture):
        element, parent_, Constructor_, placeholder_ = factory_fixture
        placeholder = _SlidePlaceholderFactory(element, parent_)
        Constructor_.assert_called_once_with(element, parent_)
        assert placeholder is placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=title}',                 'sp'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}',             'pph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=clipArt,idx=1}',         'pph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=tbl,idx=1}',             'tph'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=chart,idx=10}',          'cph'),
        ('p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=1}',           'php'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=tbl,idx=2}',   'phgf'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=chart,idx=2}', 'phgf'),
        ('p:graphicFrame/p:nvSpPr/p:nvPr/p:ph{type=dgm,idx=2}',   'phgf'),
    ])
    def factory_fixture(
            self, request, SlidePlaceholder_, ChartPlaceholder_,
            PicturePlaceholder_, TablePlaceholder_, PlaceholderGraphicFrame_,
            PlaceholderPicture_, placeholder_):
        shape_cxml, constructor_key = request.param
        shape_elm = element(shape_cxml)
        Constructor_, shape_ = {
            'sp':   (SlidePlaceholder_,        placeholder_),
            'cph':  (ChartPlaceholder_,        placeholder_),
            'pph':  (PicturePlaceholder_,      placeholder_),
            'tph':  (TablePlaceholder_,        placeholder_),
            'phgf': (PlaceholderGraphicFrame_, placeholder_),
            'php':  (PlaceholderPicture_,      placeholder_),
        }[constructor_key]
        return shape_elm, 42, Constructor_, shape_

    # fixture components -----------------------------------

    @pytest.fixture
    def ChartPlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.ChartPlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def PicturePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.PicturePlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def PlaceholderGraphicFrame_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.PlaceholderGraphicFrame',
            return_value=placeholder_
        )

    @pytest.fixture
    def PlaceholderPicture_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.PlaceholderPicture',
            return_value=placeholder_
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def SlidePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.SlidePlaceholder',
            return_value=placeholder_
        )

    @pytest.fixture
    def TablePlaceholder_(self, request, placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.TablePlaceholder',
            return_value=placeholder_
        )


class DescribeSlideShapeFactory(object):

    def it_constructs_the_right_type_of_shape(self, factory_fixture):
        element, parent_, Constructor_, shape_ = factory_fixture
        shape = SlideShapeFactory(element, parent_)
        Constructor_.assert_called_once_with(element, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp/p:nvSpPr/p:nvPr/p:ph{type=title}',      'spf'),
        ('p:pic/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}', 'spf'),
        ('p:sp',                                       'bsf'),
        ('p:pic',                                      'bsf'),
    ])
    def factory_fixture(self, request, _SlidePlaceholderFactory_,
                        placeholder_, BaseShapeFactory_, base_shape_):
        shape_cxml, shape_type = request.param
        shape_elm = element(shape_cxml)
        Constructor_, shape_ = {
            'spf': (_SlidePlaceholderFactory_, placeholder_),
            'bsf': (BaseShapeFactory_,         base_shape_),
        }[shape_type]
        return shape_elm, 42, Constructor_, shape_

    # fixture components -----------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, base_shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree.BaseShapeFactory',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, _BaseSlidePlaceholder)

    @pytest.fixture
    def _SlidePlaceholderFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._SlidePlaceholderFactory',
            return_value=placeholder_
        )


class DescribeSlideShapes(object):

    def it_provides_access_to_its_shape_factory(self, factory_fixture):
        shapes, sp, SlideShapeFactory_, shape_ = factory_fixture
        shape = shapes._shape_factory(sp)
        SlideShapeFactory_.assert_called_once_with(sp, shapes)
        assert shape is shape_

    def it_can_find_the_title_placeholder(self, title_fixture):
        shapes, _shape_factory_, calls, shape_ = title_fixture
        title_placeholder = shapes.title
        assert _shape_factory_.call_args_list == calls
        assert title_placeholder is shape_

    def it_can_add_a_movie(self, movie_fixture):
        shapes, movie_file, x, y, cx, cy = movie_fixture[:6]
        poster_frame_image, mime_type, shape_id_ = movie_fixture[6:9]
        _MoviePicElementCreator_, movie_pic = movie_fixture[9:11]
        _add_video_timing_, _shape_factory_, movie_ = movie_fixture[11:]

        movie = shapes.add_movie(
            movie_file, x, y, cx, cy, poster_frame_image, mime_type
        )

        _MoviePicElementCreator_.new_movie_pic.assert_called_once_with(
            shapes, shape_id_, movie_file, x, y, cx, cy, poster_frame_image,
            mime_type
        )
        shapes._spTree[-1] is movie_pic
        _add_video_timing_.assert_called_once_with(shapes, movie_pic)
        _shape_factory_.assert_called_once_with(shapes, movie_pic)
        assert movie is movie_

    def it_can_add_a_table(self, table_fixture):
        shapes, rows, cols, x, y, cx, cy, table_, expected_xml = table_fixture

        table = shapes.add_table(rows, cols, x, y, cx, cy)

        graphicFrame = shapes._element.xpath('p:graphicFrame')[0]
        shapes._shape_factory.assert_called_once_with(shapes, graphicFrame)
        assert table is table_
        assert shapes._element.xml == expected_xml

    def it_can_clone_placeholder_shapes_from_a_layout(self, clone_fixture):
        shapes, slide_layout_, calls = clone_fixture
        shapes.clone_layout_placeholders(slide_layout_)
        assert shapes.clone_placeholder.call_args_list == calls

    def it_adds_a_video_timing_to_help(self, add_timing_fixture):
        shapes, pic, sld, expected_xml = add_timing_fixture
        shapes._add_video_timing(pic)
        assert sld.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        (0, 1),  # no timing gets timing with one video
        (1, 2),  # timing with one video gets a second video
        (3, 1),  # timing without p:childTnLst parent gets replaced
    ])
    def add_timing_fixture(self, request):
        before_idx, after_idx = request.param
        snippets = snippet_seq('timing')
        sld = parse_xml(snippets[before_idx])
        spTree = sld.xpath('.//p:spTree')[0]
        shapes = SlideShapes(spTree, None)
        pic = element('p:pic/p:nvPicPr/p:cNvPr{id=42}')
        expected_xml = snippets[after_idx]
        return shapes, pic, sld, expected_xml

    @pytest.fixture
    def clone_fixture(self, slide_layout_, clone_placeholder_, placeholder_):
        shapes = SlideShapes(None, None)
        calls = [call(shapes, placeholder_)]
        slide_layout_.iter_cloneable_placeholders.return_value = (
            iter([placeholder_])
        )
        return shapes, slide_layout_, calls

    @pytest.fixture
    def factory_fixture(self, SlideShapeFactory_, shape_):
        shapes = SlideShapes(None, None)
        sp = element('p:sp')
        return shapes, sp, SlideShapeFactory_, shape_

    @pytest.fixture
    def movie_fixture(self, _MoviePicElementCreator_, _add_video_timing_,
                      _shape_factory_, movie_, _next_shape_id_prop_):
        shapes = SlideShapes(element('p:spTree'), None)
        movie_file, x, y, cx, cy = 'foobar.mp4', 1, 2, 3, 4
        poster_frame_image, mime_type = 'foobar.png', 'video/mp4'
        movie_pic = element('p:pic')
        _MoviePicElementCreator_.new_movie_pic.return_value = movie_pic
        _shape_factory_.return_value = movie_
        shape_id_ = _next_shape_id_prop_.return_value
        return (
            shapes, movie_file, x, y, cx, cy, poster_frame_image, mime_type,
            shape_id_, _MoviePicElementCreator_, movie_pic,
            _add_video_timing_, _shape_factory_, movie_
        )

    @pytest.fixture
    def table_fixture(self, table_, _shape_factory_):
        shapes = SlideShapes(element('p:spTree'), None)
        rows, cols, x, y, cx, cy = 1, 2, 10, 11, 12, 13
        _shape_factory_.return_value = table_
        expected_xml = (
            '<p:spTree xmlns:p="http://schemas.openxmlformats.org/presentati'
            'onml/2006/main">\n  <p:graphicFrame xmlns:a="http://schemas.ope'
            'nxmlformats.org/drawingml/2006/main">\n    <p:nvGraphicFramePr>'
            '\n      <p:cNvPr id="1" name="Table 0"/>\n      <p:cNvGraphicFr'
            'amePr>\n        <a:graphicFrameLocks noGrp="1"/>\n      </p:cNv'
            'GraphicFramePr>\n      <p:nvPr/>\n    </p:nvGraphicFramePr>\n  '
            '  <p:xfrm>\n      <a:off x="10" y="11"/>\n      <a:ext cx="12" '
            'cy="13"/>\n    </p:xfrm>\n    <a:graphic>\n      <a:graphicData'
            ' uri="http://schemas.openxmlformats.org/drawingml/2006/table">'
            '\n        <a:tbl>\n          <a:tblPr firstRow="1" bandRow="1">'
            '\n            <a:tableStyleId>{5C22544A-7EE6-4342-B048-85BDC9FD'
            '1C3A}</a:tableStyleId>\n          </a:tblPr>\n          <a:tblG'
            'rid>\n            <a:gridCol w="6"/>\n            <a:gridCol w='
            '"6"/>\n          </a:tblGrid>\n          <a:tr h="13">\n       '
            '     <a:tc>\n              <a:txBody>\n                <a:bodyP'
            'r/>\n                <a:lstStyle/>\n                <a:p/>\n   '
            '           </a:txBody>\n              <a:tcPr/>\n            </'
            'a:tc>\n            <a:tc>\n              <a:txBody>\n          '
            '      <a:bodyPr/>\n                <a:lstStyle/>\n             '
            '   <a:p/>\n              </a:txBody>\n              <a:tcPr/>\n'
            '            </a:tc>\n          </a:tr>\n        </a:tbl>\n     '
            ' </a:graphicData>\n    </a:graphic>\n  </p:graphicFrame>\n</p:s'
            'pTree>'
        )
        return shapes, rows, cols, x, y, cx, cy, table_, expected_xml

    @pytest.fixture(params=[
        ('p:spTree/(p:sp,p:sp/p:nvSpPr/p:nvPr/p:ph{idx=0})', True),
        ('p:spTree/(p:sp,p:sp)',                             False),
    ])
    def title_fixture(self, request, _shape_factory_, shape_):
        spTree_cxml, found = request.param
        spTree = element(spTree_cxml)
        shapes = SlideShapes(spTree, None)
        calls = [call(shapes, spTree.xpath('p:sp')[1])] if found else []
        _shape_ = shape_ if found else None
        return shapes, _shape_factory_, calls, _shape_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _add_video_timing_(self, request):
        return method_mock(
            request, SlideShapes, '_add_video_timing', autospec=True
        )

    @pytest.fixture
    def clone_placeholder_(self, request):
        return method_mock(
            request, SlideShapes, 'clone_placeholder', autospec=True
        )

    @pytest.fixture
    def movie_(self, request):
        return instance_mock(request, Movie)

    @pytest.fixture
    def _MoviePicElementCreator_(self, request):
        return class_mock(
            request, 'pptx.shapes.shapetree._MoviePicElementCreator',
            autospec=True
        )

    @pytest.fixture
    def _next_shape_id_prop_(self, request, shape_id_):
        return property_mock(
            request, SlideShapes, '_next_shape_id', return_value=shape_id_
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, Shape)

    @pytest.fixture
    def shape_id_(self):
        return 42

    @pytest.fixture
    def _shape_factory_(self, request, shape_):
        return method_mock(
            request, SlideShapes, '_shape_factory', return_value=shape_,
            autospec=True
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def SlideShapeFactory_(self, request, shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree.SlideShapeFactory',
            return_value=shape_, autospec=True
        )

    @pytest.fixture
    def table_(self, request):
        return instance_mock(request, Table)


class DescribeLayoutShapes(object):

    def it_provides_access_to_its_shape_factory(self, factory_fixture):
        shapes, sp, _LayoutShapeFactory_, placeholder_ = factory_fixture
        placeholder = shapes._shape_factory(sp)
        _LayoutShapeFactory_.assert_called_once_with(sp, shapes)
        assert placeholder is placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def factory_fixture(self, _LayoutShapeFactory_, placeholder_):
        shapes = LayoutShapes(None, None)
        sp = element('p:sp')
        return shapes, sp, _LayoutShapeFactory_, placeholder_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _LayoutShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._LayoutShapeFactory',
            return_value=placeholder_, autospec=True
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)


class Describe_LayoutShapeFactory(object):

    def it_constructs_the_right_shape_for_an_element(self, factory_fixture):
        shape_elm, parent_, ShapeConstructor_, shape_ = factory_fixture
        shape = _LayoutShapeFactory(shape_elm, parent_)
        ShapeConstructor_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['ph', 'sp', 'pic'])
    def factory_fixture(
            self, request, ph_bldr, parent_, _LayoutPlaceholder_,
            layout_placeholder_, BaseShapeFactory_, base_shape_):
        shape_bldr, ShapeConstructor_, shape_mock = {
            'ph':  (ph_bldr, _LayoutPlaceholder_, layout_placeholder_),
            'sp':  (an_sp(), BaseShapeFactory_,   base_shape_),
            'pic': (a_pic(), BaseShapeFactory_,   base_shape_),
        }[request.param]
        shape_elm = shape_bldr.with_nsdecls().element
        return shape_elm, parent_, ShapeConstructor_, shape_mock

    # fixture components ---------------------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, base_shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree.BaseShapeFactory',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def _LayoutPlaceholder_(self, request, layout_placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.LayoutPlaceholder',
            return_value=layout_placeholder_
        )

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def parent_(self, request):
        return instance_mock(request, BasePlaceholders)

    @pytest.fixture
    def ph_bldr(self):
        return (
            an_sp().with_child(
                an_nvSpPr().with_child(
                    an_nvPr().with_child(
                        a_ph().with_idx(1))))
        )


class DescribeLayoutPlaceholders(object):

    def it_provides_access_to_its_shape_factory(self, factory_fixture):
        placeholders, sp, _LayoutShapeFactory_, placeholder_ = factory_fixture
        placeholder = placeholders._shape_factory(sp)
        _LayoutShapeFactory_.assert_called_once_with(sp, placeholders)
        assert placeholder is placeholder_

    def it_can_find_a_placeholder_by_idx_value(self, get_fixture):
        placeholders, idx, placeholder_ = get_fixture
        assert placeholders.get(idx) is placeholder_

    def it_returns_default_on_ph_idx_not_found(self, default_fixture):
        placeholders, default = default_fixture
        assert placeholders.get(42, default) is default

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def default_fixture(self, _iter_):
        placeholders = LayoutPlaceholders(None, None)
        default = 'barfoo'
        return placeholders, default

    @pytest.fixture
    def factory_fixture(self, _LayoutShapeFactory_, placeholder_):
        placeholders = LayoutPlaceholders(None, None)
        sp = element('p:sp')
        return placeholders, sp, _LayoutShapeFactory_, placeholder_

    @pytest.fixture(params=[0, 1])
    def get_fixture(self, request, _iter_, placeholder_, placeholder_2_):
        idx = request.param
        layout_placeholders = LayoutPlaceholders(None, None)
        _placeholder_ = (placeholder_, placeholder_2_)[idx]
        placeholder_.element.ph_idx, placeholder_2_.element.ph_idx = 0, 1
        return layout_placeholders, idx, _placeholder_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _iter_(self, request, placeholder_, placeholder_2_):
        return method_mock(
            request, LayoutPlaceholders, '__iter__',
            return_value=iter([placeholder_, placeholder_2_])
        )

    @pytest.fixture
    def _LayoutShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._LayoutShapeFactory',
            return_value=placeholder_, autospec=True
        )

    @pytest.fixture
    def ph_elm_(self, request):
        return instance_mock(request, CT_Shape)

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def placeholder_2_(self, request):
        return instance_mock(request, LayoutPlaceholder)


class Describe_MasterShapeFactory(object):

    def it_constructs_a_master_placeholder_for_a_shape_element(
            self, factory_fixture):
        shape_elm, parent_, ShapeConstructor_, shape_ = factory_fixture
        shape = _MasterShapeFactory(shape_elm, parent_)
        ShapeConstructor_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['ph', 'sp', 'pic'])
    def factory_fixture(
            self, request, ph_bldr, slide_master_, _MasterPlaceholder_,
            master_placeholder_, BaseShapeFactory_, base_shape_):
        shape_bldr, ShapeConstructor_, shape_mock = {
            'ph':  (ph_bldr, _MasterPlaceholder_, master_placeholder_),
            'sp':  (an_sp(), BaseShapeFactory_,   base_shape_),
            'pic': (a_pic(), BaseShapeFactory_,   base_shape_),
        }[request.param]
        shape_elm = shape_bldr.with_nsdecls().element
        return shape_elm, slide_master_, ShapeConstructor_, shape_mock

    # fixture components -----------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, base_shape_):
        return function_mock(
            request, 'pptx.shapes.shapetree.BaseShapeFactory',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def _MasterPlaceholder_(self, request, master_placeholder_):
        return class_mock(
            request, 'pptx.shapes.shapetree.MasterPlaceholder',
            return_value=master_placeholder_
        )

    @pytest.fixture
    def master_placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder)

    @pytest.fixture
    def ph_bldr(self):
        return (
            an_sp().with_child(
                an_nvSpPr().with_child(
                    an_nvPr().with_child(
                        a_ph().with_idx(1))))
        )

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)


class DescribeMasterShapes(object):

    def it_provides_access_to_its_shape_factory(self, factory_fixture):
        shapes, sp, _MasterShapeFactory_, placeholder_ = factory_fixture
        placeholder = shapes._shape_factory(sp)
        _MasterShapeFactory_.assert_called_once_with(sp, shapes)
        assert placeholder is placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def factory_fixture(self, _MasterShapeFactory_, placeholder_):
        shapes = MasterShapes(None, None)
        sp = element('p:sp')
        return shapes, sp, _MasterShapeFactory_, placeholder_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _MasterShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._MasterShapeFactory',
            return_value=placeholder_, autospec=True
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder)


class DescribeMasterPlaceholders(object):

    def it_provides_access_to_its_shape_factory(self, factory_fixture):
        placeholders, sp, _MasterShapeFactory_, placeholder_ = factory_fixture
        placeholder = placeholders._shape_factory(sp)
        _MasterShapeFactory_.assert_called_once_with(sp, placeholders)
        assert placeholder is placeholder_

    def it_can_find_a_placeholder_by_type(self, get_fixture):
        placeholders, ph_type, placeholder_ = get_fixture
        assert placeholders.get(ph_type) is placeholder_

    def it_returns_default_on_ph_type_not_found(self, default_fixture):
        placeholders, default = default_fixture
        assert placeholders.get(42, default) is default

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def default_fixture(self, _iter_):
        placeholders = MasterPlaceholders(None, None)
        default = 'barfoo'
        return placeholders, default

    @pytest.fixture
    def factory_fixture(self, _MasterShapeFactory_, placeholder_):
        placeholders = MasterPlaceholders(None, None)
        sp = element('p:sp')
        return placeholders, sp, _MasterShapeFactory_, placeholder_

    @pytest.fixture(params=['title', 'body'])
    def get_fixture(self, request, _iter_, placeholder_, placeholder_2_):
        ph_type = request.param
        placeholders = MasterPlaceholders(None, None)
        _placeholder_ = {
            'title': placeholder_, 'body': placeholder_2_
        }[ph_type]
        return placeholders, ph_type, _placeholder_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _iter_(self, request, placeholder_, placeholder_2_):
        return method_mock(
            request, MasterPlaceholders, '__iter__',
            return_value=iter([placeholder_, placeholder_2_])
        )

    @pytest.fixture
    def _MasterShapeFactory_(self, request, placeholder_):
        return function_mock(
            request, 'pptx.shapes.shapetree._MasterShapeFactory',
            return_value=placeholder_, autospec=True
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder, ph_type='title')

    @pytest.fixture
    def placeholder_2_(self, request):
        return instance_mock(request, MasterPlaceholder, ph_type='body')


class Describe_MoviePicElementCreator(object):

    def it_creates_a_new_movie_pic_element(self, movie_pic_fixture):
        shapes_, shape_id, movie_file, x, y, cx, cy = movie_pic_fixture[:7]
        poster_frame_image, mime_type = movie_pic_fixture[7:9]
        _MoviePicElementCreator_init_, _pic_prop_ = movie_pic_fixture[9:11]
        pic_ = movie_pic_fixture[11]

        pic = _MoviePicElementCreator.new_movie_pic(
            shapes_, shape_id, movie_file, x, y, cx, cy, poster_frame_image,
            mime_type
        )

        _MoviePicElementCreator_init_.assert_called_once_with(
            ANY, shapes_, shape_id, movie_file, x, y, cx, cy,
            poster_frame_image, mime_type
        )
        _pic_prop_.assert_called_once_with()
        assert pic is pic_

    def it_creates_a_pic_element(self, pic_fixture):
        movie_pic_element_creator, new_video_pic_, shape_id = pic_fixture[:3]
        shape_name, video_rId, media_rId, poster_frame_rId = pic_fixture[3:7]
        x, y, cx, cy, pic_ = pic_fixture[7:]

        pic = movie_pic_element_creator._pic

        new_video_pic_.assert_called_once_with(
            shape_id, shape_name, video_rId, media_rId, poster_frame_rId, x,
            y, cx, cy
        )
        assert pic is pic_

    def it_knows_the_shape_name_to_help(self, shape_name_fixture):
        movie_pic_element_creator, filename = shape_name_fixture
        shape_name = movie_pic_element_creator._shape_name
        assert shape_name == filename

    def it_constructs_the_video_to_help(self, video_fixture):
        movie_pic_element_creator, movie_file = video_fixture[:2]
        mime_type, video_ = video_fixture[2:]
        video = movie_pic_element_creator._video
        Video.from_path_or_file_like.assert_called_once_with(
            movie_file, mime_type
        )
        assert video is video_

    def it_knows_the_media_rId_to_help(self, media_rId_fixture):
        movie_pic_element_creator, expected_value = media_rId_fixture
        assert movie_pic_element_creator._media_rId == expected_value

    def it_knows_the_video_rId_to_help(self, video_rId_fixture):
        movie_pic_element_creator, expected_value = video_rId_fixture
        assert movie_pic_element_creator._video_rId == expected_value

    def it_adds_the_poster_frame_image_to_help(self, pfrm_rId_fixture):
        movie_pic_element_creator, slide_part_ = pfrm_rId_fixture[:2]
        poster_frame_image_file, expected_value = pfrm_rId_fixture[2:]

        poster_frame_rId = movie_pic_element_creator._poster_frame_rId

        slide_part_.get_or_add_image_part.assert_called_once_with(
            poster_frame_image_file
        )
        assert poster_frame_rId == expected_value

    def it_gets_the_poster_frame_image_file_to_help(self, pfrm_img_fixture):
        movie_pic_element_creator, BytesIO_ = pfrm_img_fixture[:2]
        calls, expected_value = pfrm_img_fixture[2:]
        image_file = movie_pic_element_creator._poster_frame_image_file
        assert BytesIO_.call_args_list == calls
        assert image_file == expected_value

    def it_gets_the_video_part_rIds_to_help(self, part_rIds_fixture):
        movie_pic_element_creator, slide_part_ = part_rIds_fixture[:2]
        video_, media_rId, video_rId = part_rIds_fixture[2:]

        result = movie_pic_element_creator._video_part_rIds

        slide_part_.get_or_add_video_media_part.assert_called_once_with(
            video_
        )
        assert result == (media_rId, video_rId)

    def it_gets_the_slide_part_to_help(self, slide_part_fixture):
        movie_pic_element_creator, slide_part_ = slide_part_fixture
        slide_part = movie_pic_element_creator._slide_part
        assert slide_part is slide_part_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def media_rId_fixture(self, _video_part_rIds_prop_):
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, None, None
        )
        expected_value = 'rId24'
        _video_part_rIds_prop_.return_value = (expected_value, 'rId666')
        return movie_pic_element_creator, expected_value

    @pytest.fixture
    def movie_pic_fixture(self, shapes_, _MoviePicElementCreator_init_,
                          _pic_prop_, pic_):
        shape_id, movie_file, x, y, cx, cy = 42, 'movie.mp4', 1, 2, 3, 4
        poster_frame_image, mime_type = 'image.png', 'video/mp4'
        return (
            shapes_, shape_id, movie_file, x, y, cx, cy, poster_frame_image,
            mime_type, _MoviePicElementCreator_init_, _pic_prop_, pic_
        )

    @pytest.fixture
    def part_rIds_fixture(self, slide_part_, video_, _slide_part_prop_,
                          _video_prop_):
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, None, None
        )
        media_rId, video_rId = 'rId42', 'rId24'
        _slide_part_prop_.return_value = slide_part_
        slide_part_.get_or_add_video_media_part.return_value = (
            media_rId, video_rId
        )
        _video_prop_.return_value = video_
        return (
            movie_pic_element_creator, slide_part_, video_, media_rId,
            video_rId
        )

    @pytest.fixture(params=[
        'image.png',
        None,
    ])
    def pfrm_img_fixture(self, request, BytesIO_, stream_):
        poster_frame_file = request.param
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, poster_frame_file, None
        )
        if poster_frame_file is None:
            calls = [call(SPEAKER_IMAGE_BYTES)]
            BytesIO_.return_value = stream_
            expected_value = stream_
        else:
            calls = []
            expected_value = poster_frame_file
        return movie_pic_element_creator, BytesIO_, calls, expected_value

    @pytest.fixture
    def pfrm_rId_fixture(self, _slide_part_prop_, slide_part_,
                         _poster_frame_image_file_prop_):
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, None, None
        )
        poster_frame_image_file, expected_value = 'image.png', 'rId42'
        _slide_part_prop_.return_value = slide_part_
        _poster_frame_image_file_prop_.return_value = poster_frame_image_file
        slide_part_.get_or_add_image_part.return_value = (
            None, expected_value
        )
        return (
            movie_pic_element_creator, slide_part_, poster_frame_image_file,
            expected_value
        )

    @pytest.fixture
    def pic_fixture(self, new_video_pic_, pic_, _shape_name_prop_,
                    _video_rId_prop_, _media_rId_prop_,
                    _poster_frame_rId_prop_):
        shape_id, x, y, cx, cy = 42, 1, 2, 3, 4
        movie_pic_element_creator = _MoviePicElementCreator(
            None, shape_id, None, x, y, cx, cy, None, None
        )
        _shape_name_prop_.return_value = shape_name = 'movie.mp4'
        _video_rId_prop_.return_value = video_rId = 'rId1'
        _media_rId_prop_.return_value = media_rId = 'rId2',
        _poster_frame_rId_prop_.return_value = poster_frame_rId = 'rId3'
        new_video_pic_.return_value = pic_
        return (
            movie_pic_element_creator, new_video_pic_, shape_id, shape_name,
            video_rId, media_rId, poster_frame_rId, x, y, cx, cy, pic_
        )

    @pytest.fixture
    def shape_name_fixture(self, _video_prop_, video_):
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, None, None
        )
        _video_prop_.return_value = video_
        video_.filename = filename = 'movie.mp4'
        return movie_pic_element_creator, filename

    @pytest.fixture
    def slide_part_fixture(self, shapes_, slide_part_):
        movie_pic_element_creator = _MoviePicElementCreator(
            shapes_, None, None, None, None, None, None, None, None
        )
        shapes_.part = slide_part_
        return movie_pic_element_creator, slide_part_

    @pytest.fixture
    def video_rId_fixture(self, _video_part_rIds_prop_):
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, None, None, None, None, None, None, None
        )
        expected_value = 'rId42'
        _video_part_rIds_prop_.return_value = ('rId666', expected_value)
        return movie_pic_element_creator, expected_value

    @pytest.fixture
    def video_fixture(self, video_, from_path_or_file_like_):
        movie_file, mime_type = 'movie.mp4', 'video/mp4'
        movie_pic_element_creator = _MoviePicElementCreator(
            None, None, movie_file, None, None, None, None, None, mime_type
        )
        from_path_or_file_like_.return_value = video_
        return movie_pic_element_creator, movie_file, mime_type, video_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def BytesIO_(self, request):
        return class_mock(request, 'pptx.shapes.shapetree.BytesIO')

    @pytest.fixture
    def from_path_or_file_like_(self, request):
        return method_mock(request, Video, 'from_path_or_file_like')

    @pytest.fixture
    def _media_rId_prop_(self, request):
        return property_mock(request, _MoviePicElementCreator, '_media_rId')

    @pytest.fixture
    def _MoviePicElementCreator_init_(self, request):
        return initializer_mock(
            request, _MoviePicElementCreator, autospec=True
        )

    @pytest.fixture
    def new_video_pic_(self, request):
        return method_mock(request, CT_Picture, 'new_video_pic')

    @pytest.fixture
    def pic_(self):
        return element('p:pic')

    @pytest.fixture
    def _pic_prop_(self, request, pic_):
        return property_mock(
            request, _MoviePicElementCreator, '_pic', return_value=pic_
        )

    @pytest.fixture
    def _poster_frame_image_file_prop_(self, request):
        return property_mock(
            request, _MoviePicElementCreator, '_poster_frame_image_file'
        )

    @pytest.fixture
    def _poster_frame_rId_prop_(self, request):
        return property_mock(
            request, _MoviePicElementCreator, '_poster_frame_rId'
        )

    @pytest.fixture
    def _shape_name_prop_(self, request):
        return property_mock(request, _MoviePicElementCreator, '_shape_name')

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, _BaseShapes)

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)

    @pytest.fixture
    def _slide_part_prop_(self, request):
        return property_mock(request, _MoviePicElementCreator, '_slide_part')

    @pytest.fixture
    def stream_(self, request):
        return instance_mock(request, BytesIO)

    @pytest.fixture
    def video_(self, request):
        return instance_mock(request, Video)

    @pytest.fixture
    def _video_prop_(self, request):
        return property_mock(request, _MoviePicElementCreator, '_video')

    @pytest.fixture
    def _video_rId_prop_(self, request):
        return property_mock(request, _MoviePicElementCreator, '_video_rId')

    @pytest.fixture
    def _video_part_rIds_prop_(self, request):
        return property_mock(
            request, _MoviePicElementCreator, '_video_part_rIds'
        )


class Describe_NotesSlideShapeFactory(object):

    def it_constructs_the_right_shape_for_an_element(self, factory_fixture):
        shape_elm, parent_, ShapeConstructor_, shape_ = factory_fixture
        shape = _NotesSlideShapeFactory(shape_elm, parent_)
        ShapeConstructor_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp',                      'BaseShapeFactory'),
        ('p:sp/p:nvSpPr/p:nvPr/p:ph', 'NotesSlidePlaceholder'),
    ])
    def factory_fixture(self, request, shapes_, shape_):
        sp_cxml, shape_cls_name = request.param
        shape_elm = element(sp_cxml)
        ShapeConstructor_ = class_mock(
            request, 'pptx.shapes.shapetree.%s' % shape_cls_name,
            return_value=shape_
        )
        return shape_elm, shapes_, ShapeConstructor_, shape_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, _BaseShapes)
