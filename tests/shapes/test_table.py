# encoding: utf-8

"""
Test suite for pptx.table module.
"""

from __future__ import absolute_import, print_function

import pytest

from pptx.dml.fill import FillFormat
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.table import (
    _Cell, _CellCollection, _Column, _ColumnCollection, _Row, _RowCollection,
    Table
)
from pptx.util import Inches, Length, Pt

from ..unitutil.cxml import element, xml
from ..unitutil.mock import class_mock, instance_mock


class DescribeTable(object):

    def it_provides_access_to_its_cells(self, cell_fixture):
        table, row_idx, col_idx, expected_cell_ = cell_fixture
        cell = table.cell(row_idx, col_idx)
        assert cell is expected_cell_

    def it_provides_access_to_its_rows(self, rows_fixture):
        table, expected_rows_ = rows_fixture
        assert table.rows is expected_rows_

    def it_provides_access_to_its_columns(self, columns_fixture):
        table, expected_columns_ = columns_fixture
        assert table.columns is expected_columns_

    def it_updates_graphic_frame_width_on_width_change(self, dx_fixture):
        table, expected_width = dx_fixture
        table.notify_width_changed()
        assert table._graphic_frame.width == expected_width

    def it_updates_graphic_frame_height_on_height_change(self, dy_fixture):
        table, expected_height = dy_fixture
        table.notify_height_changed()
        assert table._graphic_frame.height == expected_height

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def cell_fixture(self, table, row_, cell_):
        table._rows = [row_]
        row_.cells = [cell_]
        row_idx = col_idx = 0
        return table, row_idx, col_idx, cell_

    @pytest.fixture
    def columns_fixture(self, table, columns_):
        table._columns = columns_
        return table, columns_

    @pytest.fixture
    def dx_fixture(self, graphic_frame_):
        tbl_cxml = 'a:tbl/a:tblGrid/(a:gridCol{w=111},a:gridCol{w=222})'
        table = Table(element(tbl_cxml), graphic_frame_)
        expected_width = 333
        return table, expected_width

    @pytest.fixture
    def dy_fixture(self, graphic_frame_):
        tbl_cxml = 'a:tbl/(a:tr{h=100},a:tr{h=200})'
        table = Table(element(tbl_cxml), graphic_frame_)
        expected_height = 300
        return table, expected_height

    @pytest.fixture
    def rows_fixture(self, table, rows_):
        table._rows = rows_
        return table, rows_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def cell_(self, request):
        return instance_mock(request, _Cell)

    @pytest.fixture
    def columns_(self, request):
        return instance_mock(request, _ColumnCollection)

    @pytest.fixture
    def graphic_frame_(self, request):
        return instance_mock(request, GraphicFrame)

    @pytest.fixture
    def row_(self, request):
        return instance_mock(request, _Row)

    @pytest.fixture
    def rows_(self, request):
        return instance_mock(request, _RowCollection)

    @pytest.fixture
    def table(self):
        return Table(element('a:tbl'), None)


class DescribeTableBooleanProperties(object):

    def it_knows_its_boolean_property_settings(self, boolprop_get_fixture):
        table, boolprop_name, expected_value = boolprop_get_fixture
        boolprop_value = getattr(table, boolprop_name)
        assert boolprop_value is expected_value

    def it_can_change_its_boolean_property_settings(
            self, boolprop_set_fixture):
        table, boolprop_name, new_value, expected_xml = boolprop_set_fixture
        setattr(table, boolprop_name, new_value)
        assert table._tbl.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('a:tbl',                         'first_row',    False),
        ('a:tbl/a:tblPr',                 'first_row',    False),
        ('a:tbl/a:tblPr{firstRow=1}',     'first_row',    True),
        ('a:tbl/a:tblPr{firstRow=0}',     'first_row',    False),
        ('a:tbl/a:tblPr{firstRow=true}',  'first_row',    True),
        ('a:tbl/a:tblPr{firstRow=false}', 'first_row',    False),
        ('a:tbl/a:tblPr{firstCol=1}',     'first_col',    True),
        ('a:tbl/a:tblPr{lastRow=0}',      'last_row',     False),
        ('a:tbl/a:tblPr{lastCol=true}',   'last_col',     True),
        ('a:tbl/a:tblPr{bandRow=false}',  'horz_banding', False),
        ('a:tbl/a:tblPr',                 'vert_banding', False),
    ])
    def boolprop_get_fixture(self, request):
        tbl_cxml, boolprop_name, expected_value = request.param
        table = Table(element(tbl_cxml), None)
        return table, boolprop_name, expected_value

    @pytest.fixture(params=[
        ('a:tbl',         'first_row', True,  'a:tbl/a:tblPr{firstRow=1}'),
        ('a:tbl',         'first_row', False, 'a:tbl/a:tblPr'),
        ('a:tbl/a:tblPr', 'first_row', True,  'a:tbl/a:tblPr{firstRow=1}'),
        ('a:tbl/a:tblPr', 'first_row', False, 'a:tbl/a:tblPr'),
        ('a:tbl/a:tblPr{firstRow=true}',  'first_row', True,
         'a:tbl/a:tblPr{firstRow=1}'),
        ('a:tbl/a:tblPr{firstRow=false}', 'first_row', False,
         'a:tbl/a:tblPr'),
        ('a:tbl/a:tblPr{bandRow=1}',      'first_row', True,
         'a:tbl/a:tblPr{bandRow=1,firstRow=1}'),
        ('a:tbl', 'first_col',    True, 'a:tbl/a:tblPr{firstCol=1}'),
        ('a:tbl', 'last_row',     True, 'a:tbl/a:tblPr{lastRow=1}'),
        ('a:tbl', 'last_col',     True, 'a:tbl/a:tblPr{lastCol=1}'),
        ('a:tbl', 'horz_banding', True, 'a:tbl/a:tblPr{bandRow=1}'),
        ('a:tbl', 'vert_banding', True, 'a:tbl/a:tblPr{bandCol=1}'),
    ])
    def boolprop_set_fixture(self, request):
        tbl_cxml, boolprop_name, new_value, expected_tbl_cxml = request.param
        table = Table(element(tbl_cxml), None)
        expected_xml = xml(expected_tbl_cxml)
        return table, boolprop_name, new_value, expected_xml


class Describe_Cell(object):

    def it_has_a_fill(self, fill_fixture):
        cell = fill_fixture
        assert isinstance(cell.fill, FillFormat)

    def it_knows_its_margin_settings(self, margin_get_fixture):
        cell, margin_prop_name, expected_value = margin_get_fixture
        margin_value = getattr(cell, margin_prop_name)
        assert margin_value == expected_value

    def it_can_change_its_margin_settings(self, margin_set_fixture):
        cell, margin_prop_name, new_value, expected_xml = margin_set_fixture
        setattr(cell, margin_prop_name, new_value)
        assert cell._tc.xml == expected_xml

    def it_raises_on_margin_assigned_other_than_int_or_None(
            self, margin_raises_fixture):
        cell, margin_attr_name, val_of_invalid_type = margin_raises_fixture
        with pytest.raises(TypeError):
            setattr(cell, margin_attr_name, val_of_invalid_type)

    def it_knows_its_vertical_anchor_setting(self, anchor_get_fixture):
        cell, expected_value = anchor_get_fixture
        assert cell.vertical_anchor == expected_value

    def it_can_change_its_vertical_anchor(self, anchor_set_fixture):
        cell, new_value, expected_xml = anchor_set_fixture
        cell.vertical_anchor = new_value
        assert cell._tc.xml == expected_xml

    def it_can_replace_its_contents_with_a_string(self, text_set_fixture):
        cell, text, expected_xml = text_set_fixture
        cell.text = text
        assert cell._tc.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('a:tc',                    None),
        ('a:tc/a:tcPr',             None),
        ('a:tc/a:tcPr{anchor=t}',   MSO_ANCHOR.TOP),
        ('a:tc/a:tcPr{anchor=ctr}', MSO_ANCHOR.MIDDLE),
        ('a:tc/a:tcPr{anchor=b}',   MSO_ANCHOR.BOTTOM),
    ])
    def anchor_get_fixture(self, request):
        tc_cxml, expected_value = request.param
        cell = _Cell(element(tc_cxml), None)
        return cell, expected_value

    @pytest.fixture(params=[
        ('a:tc', None,              'a:tc'),
        ('a:tc', MSO_ANCHOR.TOP,    'a:tc/a:tcPr{anchor=t}'),
        ('a:tc', MSO_ANCHOR.MIDDLE, 'a:tc/a:tcPr{anchor=ctr}'),
        ('a:tc', MSO_ANCHOR.BOTTOM, 'a:tc/a:tcPr{anchor=b}'),
        ('a:tc/a:tcPr{anchor=t}',   MSO_ANCHOR.MIDDLE,
         'a:tc/a:tcPr{anchor=ctr}'),
        ('a:tc/a:tcPr{anchor=ctr}', None,
         'a:tc/a:tcPr'),
    ])
    def anchor_set_fixture(self, request):
        tc_cxml, new_value, expected_tc_cxml = request.param
        cell = _Cell(element(tc_cxml), None)
        expected_xml = xml(expected_tc_cxml)
        return cell, new_value, expected_xml

    @pytest.fixture
    def fill_fixture(self, cell):
        return cell

    @pytest.fixture(params=[
        ('a:tc/a:tcPr{marL=82296}', 'margin_left',   Inches(0.09)),
        ('a:tc/a:tcPr{marR=73152}', 'margin_right',  Inches(0.08)),
        ('a:tc/a:tcPr{marT=64008}', 'margin_top',    Inches(0.07)),
        ('a:tc/a:tcPr{marB=54864}', 'margin_bottom', Inches(0.06)),
        ('a:tc',                    'margin_left',   Inches(0.1)),
        ('a:tc/a:tcPr',             'margin_right',  Inches(0.1)),
        ('a:tc',                    'margin_top',    Inches(0.05)),
        ('a:tc/a:tcPr',             'margin_bottom', Inches(0.05)),
    ])
    def margin_get_fixture(self, request):
        tc_cxml, margin_prop_name, expected_value = request.param
        cell = _Cell(element(tc_cxml), None)
        return cell, margin_prop_name, expected_value

    @pytest.fixture(params=[
        ('a:tc', 'margin_left',   Inches(0.08), 'a:tc/a:tcPr{marL=73152}'),
        ('a:tc', 'margin_right',  Inches(0.08), 'a:tc/a:tcPr{marR=73152}'),
        ('a:tc', 'margin_top',    Inches(0.08), 'a:tc/a:tcPr{marT=73152}'),
        ('a:tc', 'margin_bottom', Inches(0.08), 'a:tc/a:tcPr{marB=73152}'),
        ('a:tc', 'margin_left',   None,         'a:tc'),
        ('a:tc/a:tcPr{marL=42}', 'margin_left', None,
         'a:tc/a:tcPr'),
    ])
    def margin_set_fixture(self, request):
        tc_cxml, margin_prop_name, new_value, expected_tc_cxml = request.param
        cell = _Cell(element(tc_cxml), None)
        expected_xml = xml(expected_tc_cxml)
        return cell, margin_prop_name, new_value, expected_xml

    @pytest.fixture(params=[
        'margin_left', 'margin_right', 'margin_top', 'margin_bottom'
    ])
    def margin_raises_fixture(self, request):
        margin_prop_name = request.param
        cell = _Cell(element('a:tc'), None)
        val_of_invalid_type = 'foobar'
        return cell, margin_prop_name, val_of_invalid_type

    @pytest.fixture(params=[
        ('a:tc', 'foobar',
         'a:tc/a:txBody/(a:bodyPr, a:p/a:r/a:t"foobar")'),
        ('a:tc/a:txBody/(a:bodyPr, a:p/a:r/(a:br,a:t"bar"))', 'foobar',
         'a:tc/a:txBody/(a:bodyPr, a:p/a:r/a:t"foobar")'),
    ])
    def text_set_fixture(self, request):
        tc_cxml, new_text, expected_cxml = request.param
        cell = _Cell(element(tc_cxml), None)
        expected_xml = xml(expected_cxml)
        return cell, new_text, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def cell(self):
        return _Cell(element('a:tc'), None)


class Describe_CellCollection(object):

    def it_knows_how_many_cells_it_contains(self, len_fixture):
        cells, expected_count = len_fixture
        assert len(cells) == expected_count

    def it_can_iterate_over_the_cells_it_contains(self, iter_fixture):
        cells, expected_tc_lst = iter_fixture
        count = 0
        for idx, cell in enumerate(cells):
            assert isinstance(cell, _Cell)
            assert cell._tc is expected_tc_lst[idx]
            count += 1
        assert count == len(expected_tc_lst)

    def it_supports_indexed_access(self, getitem_fixture):
        cells, expected_tc_lst = getitem_fixture
        for idx, tc in enumerate(expected_tc_lst):
            cell = cells[idx]
            assert isinstance(cell, _Cell)
            assert cell._tc is tc

    def it_raises_on_indexed_access_out_of_range(self):
        cells = _CellCollection(element('a:tr/a:tc'), None)
        with pytest.raises(IndexError):
            cells[-1]
        with pytest.raises(IndexError):
            cells[9]

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        'a:tr', 'a:tr/a:tc', 'a:tr/(a:tc, a:tc, a:tc)',
    ])
    def getitem_fixture(self, request):
        tr_cxml = request.param
        tr = element(tr_cxml)
        cells = _CellCollection(tr, None)
        expected_cell_lst = tr.findall(qn('a:tc'))
        return cells, expected_cell_lst

    @pytest.fixture(params=[
        'a:tr', 'a:tr/a:tc', 'a:tr/(a:tc, a:tc, a:tc)',
    ])
    def iter_fixture(self, request):
        tr_cxml = request.param
        tr = element(tr_cxml)
        cells = _CellCollection(tr, None)
        expected_cell_lst = tr.findall(qn('a:tc'))
        return cells, expected_cell_lst

    @pytest.fixture(params=[
        ('a:tr', 0), ('a:tr/a:tc', 1), ('a:tr/(a:tc, a:tc)', 2),
    ])
    def len_fixture(self, request):
        tr_cxml, expected_len = request.param
        cells = _CellCollection(element(tr_cxml), None)
        return cells, expected_len


class Describe_Column(object):

    def it_knows_its_width(self, width_get_fixture):
        column, expected_value = width_get_fixture
        width = column.width
        assert width == expected_value
        assert isinstance(width, Length)

    def it_can_change_its_width(self, width_set_fixture):
        column, new_width, expected_xml, parent_ = width_set_fixture
        column.width = new_width
        assert column._gridCol.xml == expected_xml
        parent_.notify_width_changed.assert_called_once_with()

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('a:gridCol{w=914400}', Inches(1)),
        ('a:gridCol{w=10pt}',   Pt(10)),
    ])
    def width_get_fixture(self, request):
        gridCol_cxml, expected_value = request.param
        column = _Column(element(gridCol_cxml), None)
        return column, expected_value

    @pytest.fixture(params=[
        ('a:gridCol{w=12pt}', Inches(1), 'a:gridCol{w=914400}'),
        ('a:gridCol{w=1234}', Inches(1), 'a:gridCol{w=914400}'),
    ])
    def width_set_fixture(self, request, parent_):
        gridCol_cxml, new_width, expected_gridCol_cxml = request.param
        column = _Column(element(gridCol_cxml), parent_)
        expected_xml = xml(expected_gridCol_cxml)
        return column, new_width, expected_xml, parent_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def parent_(self, request):
        return instance_mock(request, _ColumnCollection)


class Describe_ColumnCollection(object):

    def it_knows_how_many_columns_it_contains(self, len_fixture):
        columns, expected_count = len_fixture
        assert len(columns) == expected_count

    def it_can_iterate_over_the_columns_it_contains(self, iter_fixture):
        columns, expected_gridCol_lst = iter_fixture
        count = 0
        for idx, column in enumerate(columns):
            assert isinstance(column, _Column)
            assert column._gridCol is expected_gridCol_lst[idx]
            count += 1
        assert count == len(expected_gridCol_lst)

    def it_supports_indexed_access(self, getitem_fixture):
        columns, expected_gridCol_lst = getitem_fixture
        for idx, gridCol in enumerate(expected_gridCol_lst):
            column = columns[idx]
            assert isinstance(column, _Column)
            assert column._gridCol is gridCol

    def it_raises_on_indexed_access_out_of_range(self):
        columns = _ColumnCollection(
            element('a:tbl/a:tblGrid/a:gridCol'), None
        )
        with pytest.raises(IndexError):
            columns[-1]
        with pytest.raises(IndexError):
            columns[9]

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        'a:tbl/a:tblGrid',
        'a:tbl/a:tblGrid/a:gridCol',
        'a:tbl/a:tblGrid/(a:gridCol, a:gridCol, a:gridCol)',
    ])
    def getitem_fixture(self, request):
        tbl_cxml = request.param
        tbl = element(tbl_cxml)
        columns = _ColumnCollection(tbl, None)
        expected_column_lst = tbl.xpath('//a:gridCol')
        return columns, expected_column_lst

    @pytest.fixture(params=[
        'a:tbl/a:tblGrid',
        'a:tbl/a:tblGrid/a:gridCol',
        'a:tbl/a:tblGrid/(a:gridCol, a:gridCol, a:gridCol)',
    ])
    def iter_fixture(self, request):
        tbl_cxml = request.param
        tbl = element(tbl_cxml)
        columns = _ColumnCollection(tbl, None)
        expected_column_lst = tbl.xpath('//a:gridCol')
        print(expected_column_lst)
        return columns, expected_column_lst

    @pytest.fixture(params=[
        ('a:tbl/a:tblGrid', 0),
        ('a:tbl/a:tblGrid/a:gridCol', 1),
        ('a:tbl/a:tblGrid/(a:gridCol,a:gridCol)', 2),
    ])
    def len_fixture(self, request):
        tbl_cxml, expected_len = request.param
        columns = _ColumnCollection(element(tbl_cxml), None)
        return columns, expected_len


class Describe_Row(object):

    def it_knows_its_height(self, height_get_fixture):
        row, expected_value = height_get_fixture
        height = row.height
        assert height == expected_value
        assert isinstance(height, Length)

    def it_can_change_its_height(self, height_set_fixture):
        row, new_height, expected_xml, parent_ = height_set_fixture
        row.height = new_height
        assert row._tr.xml == expected_xml
        parent_.notify_height_changed.assert_called_once_with()

    def it_provides_access_to_its_cells(self, cells_fixture):
        row, _CellCollection_, cells_ = cells_fixture
        cells = row.cells
        _CellCollection_.assert_called_once_with(row._tr, row)
        assert cells is cells_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def cells_fixture(self, _CellCollection_, cells_):
        row = _Row(element('a:tr'), None)
        return row, _CellCollection_, cells_

    @pytest.fixture(params=[
        ('a:tr{h=914400}', Inches(1)),
        ('a:tr{h=10pt}',   Pt(10)),
    ])
    def height_get_fixture(self, request):
        tr_cxml, expected_value = request.param
        row = _Row(element(tr_cxml), None)
        return row, expected_value

    @pytest.fixture(params=[
        ('a:tr{h=12pt}', Inches(1), 'a:tr{h=914400}'),
        ('a:tr{h=1234}', Inches(1), 'a:tr{h=914400}'),
    ])
    def height_set_fixture(self, request, parent_):
        tr_cxml, new_height, expected_tr_cxml = request.param
        row = _Row(element(tr_cxml), parent_)
        expected_xml = xml(expected_tr_cxml)
        return row, new_height, expected_xml, parent_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _CellCollection_(self, request, cells_):
        return class_mock(
            request, 'pptx.shapes.table._CellCollection', return_value=cells_
        )

    @pytest.fixture
    def cells_(self, request):
        return instance_mock(request, _CellCollection)

    @pytest.fixture
    def parent_(self, request):
        return instance_mock(request, _RowCollection)


class Describe_RowCollection(object):

    def it_knows_how_many_rows_it_contains(self, len_fixture):
        rows, expected_count = len_fixture
        assert len(rows) == expected_count

    def it_can_iterate_over_the_rows_it_contains(self, iter_fixture):
        rows, expected_tr_lst = iter_fixture
        count = 0
        for idx, row in enumerate(rows):
            assert isinstance(row, _Row)
            assert row._tr is expected_tr_lst[idx]
            count += 1
        assert count == len(expected_tr_lst)

    def it_supports_indexed_access(self, getitem_fixture):
        rows, expected_tr_lst = getitem_fixture
        for idx, tr in enumerate(expected_tr_lst):
            row = rows[idx]
            assert isinstance(row, _Row)
            assert row._tr is tr

    def it_raises_on_indexed_access_out_of_range(self):
        rows = _RowCollection(element('a:tbl/a:tr'), None)
        with pytest.raises(IndexError):
            rows[-1]
        with pytest.raises(IndexError):
            rows[9]

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        'a:tbl', 'a:tbl/a:tr', 'a:tbl/(a:tr, a:tr, a:tr)',
    ])
    def getitem_fixture(self, request):
        tbl_cxml = request.param
        tbl = element(tbl_cxml)
        rows = _RowCollection(tbl, None)
        expected_row_lst = tbl.findall(qn('a:tr'))
        return rows, expected_row_lst

    @pytest.fixture(params=[
        'a:tbl', 'a:tbl/a:tr', 'a:tbl/(a:tr, a:tr, a:tr)',
    ])
    def iter_fixture(self, request):
        tbl_cxml = request.param
        tbl = element(tbl_cxml)
        rows = _RowCollection(tbl, None)
        expected_row_lst = tbl.findall(qn('a:tr'))
        return rows, expected_row_lst

    @pytest.fixture(params=[
        ('a:tbl', 0), ('a:tbl/a:tr', 1), ('a:tbl/(a:tr, a:tr)', 2),
    ])
    def len_fixture(self, request):
        tbl_cxml, expected_len = request.param
        rows = _RowCollection(element(tbl_cxml), None)
        return rows, expected_len
