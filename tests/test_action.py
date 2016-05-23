# encoding: utf-8

"""
Test suite for pptx.action module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.action import ActionSetting, Hyperlink
from pptx.enum.action import PP_ACTION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from .unitutil.cxml import element, xml
from .unitutil.mock import call, class_mock, instance_mock, property_mock


class DescribeActionSetting(object):

    def it_knows_its_action_type(self, action_fixture):
        action_setting, expected_action = action_fixture
        action = action_setting.action
        assert action is expected_action

    def it_provides_access_to_its_hyperlink(self, hyperlink_fixture):
        action_setting, hyperlink_, Hyperlink_ = hyperlink_fixture[:3]
        xPr, parent = hyperlink_fixture[3:]
        hyperlink = action_setting.hyperlink
        Hyperlink_.assert_called_once_with(xPr, parent, False)
        assert hyperlink is hyperlink_

    def it_can_find_its_slide_jump_target(self, target_slide_fixture):
        action_setting, expected_value = target_slide_fixture
        target_slide = action_setting.target_slide
        assert target_slide == expected_value

    def it_raises_on_no_next_prev_slide(self, tgt_sld_raise_fixture):
        action_setting = tgt_sld_raise_fixture
        with pytest.raises(ValueError):
            action_setting.target_slide

    def it_knows_its_slide_index_to_help(self, _slide_index_fixture):
        action_setting, slides, slide, expected_value = _slide_index_fixture
        slide_index = action_setting._slide_index
        slides.index.assert_called_once_with(slide)
        assert slide_index == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:cNvPr',              None,
         PP_ACTION.NONE),
        ('p:cNvPr/a:hlinkClick', None,
         PP_ACTION.HYPERLINK),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=firstslide',
         PP_ACTION.FIRST_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=lastslide',
         PP_ACTION.LAST_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=nextslide',
         PP_ACTION.NEXT_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=previousslide',
         PP_ACTION.PREVIOUS_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=endshow',
         PP_ACTION.END_SHOW),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinksldjump',
         PP_ACTION.NAMED_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkfile',
         PP_ACTION.OPEN_FILE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkpres',
         PP_ACTION.PLAY),
        ('p:cNvPr/a:hlinkClick', 'ppaction://customshow',
         PP_ACTION.NAMED_SLIDE_SHOW),
        ('p:cNvPr/a:hlinkClick', 'ppaction://ole',
         PP_ACTION.OLE_VERB),
        ('p:cNvPr/a:hlinkClick', 'ppaction://macro',
         PP_ACTION.RUN_MACRO),
        ('p:cNvPr/a:hlinkClick', 'ppaction://program',
         PP_ACTION.RUN_PROGRAM),
        ('p:cNvPr/a:hlinkClick',
         'ppaction://hlinkshowjump?jump=lastslideviewed',
         PP_ACTION.LAST_SLIDE_VIEWED),
    ])
    def action_fixture(self, request):
        cNvPr_cxml, action_text, expected_action = request.param
        cNvPr = element(cNvPr_cxml)
        if action_text is not None:
            cNvPr.hlinkClick.action = action_text
        action_setting = ActionSetting(cNvPr, None)
        return action_setting, expected_action

    @pytest.fixture
    def hyperlink_fixture(self, Hyperlink_, hyperlink_):
        xPr, parent = 'xPr', 'parent'
        action_setting = ActionSetting(xPr, parent)
        return action_setting, hyperlink_, Hyperlink_, xPr, parent

    @pytest.fixture
    def _slide_index_fixture(self, request, part_prop_):
        action_setting = ActionSetting(None, None)
        slide_part = part_prop_.return_value
        slides = slide_part.package.presentation_part.presentation.slides
        slide = slide_part.slide
        expected_value = 123
        slides.index.return_value = expected_value
        return action_setting, slides, slide, expected_value

    @pytest.fixture(params=[
        (PP_ACTION.NONE,                None),
        (PP_ACTION.HYPERLINK,           None),
        (PP_ACTION.FIRST_SLIDE,         0),
        (PP_ACTION.LAST_SLIDE,          5),
        (PP_ACTION.NEXT_SLIDE,          3),
        (PP_ACTION.PREVIOUS_SLIDE,      1),
        (PP_ACTION.END_SHOW,            None),
        (PP_ACTION.NAMED_SLIDE,         4),
        (PP_ACTION.OPEN_FILE,           None),
        (PP_ACTION.PLAY,                None),
        (PP_ACTION.NAMED_SLIDE_SHOW,    None),
        (PP_ACTION.OLE_VERB,            None),
        (PP_ACTION.RUN_MACRO,           None),
        (PP_ACTION.RUN_PROGRAM,         None),
        (PP_ACTION.LAST_SLIDE_VIEWED,   None),
    ])
    def target_slide_fixture(
            self, request, action_prop_, _slide_index_prop_, part_prop_):
        action_type, expected_value = request.param
        cNvPr = element('p:cNvPr/a:hlinkClick{r:id=rId6}')
        action_setting = ActionSetting(cNvPr, None)
        action_prop_.return_value = action_type
        _slide_index_prop_.return_value = 2
        # this becomes the return value of ActionSetting._slides
        prs_part_ = part_prop_.return_value.package.presentation_part
        prs_part_.presentation.slides = [0, 1, 2, 3, 4, 5]
        related_parts_ = part_prop_.return_value.related_parts
        related_parts_.__getitem__.return_value.slide = 4
        return action_setting, expected_value

    @pytest.fixture(params=[
        (PP_ACTION.NEXT_SLIDE,     2),
        (PP_ACTION.PREVIOUS_SLIDE, 0),
    ])
    def tgt_sld_raise_fixture(self, request, action_prop_, part_prop_,
                              _slide_index_prop_):
        action_type, slide_idx = request.param
        action_setting = ActionSetting(None, None)
        action_prop_.return_value = action_type
        # this becomes the return value of ActionSetting._slides
        part_prop_.return_value.package.presentation.slides = [0, 1, 2]
        _slide_index_prop_.return_value = slide_idx
        return action_setting

    # fixture components ---------------------------------------------

    @pytest.fixture
    def action_prop_(self, request):
        return property_mock(request, ActionSetting, 'action')

    @pytest.fixture
    def Hyperlink_(self, request, hyperlink_):
        return class_mock(
            request, 'pptx.action.Hyperlink', return_value=hyperlink_
        )

    @pytest.fixture
    def hyperlink_(self, request):
        return instance_mock(request, Hyperlink)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, ActionSetting, 'part')

    @pytest.fixture
    def _slide_index_prop_(self, request):
        return property_mock(request, ActionSetting, '_slide_index')


class DescribeHyperlink(object):

    def it_knows_the_target_url_of_the_hyperlink(self, address_fixture):
        hyperlink, rId, expected_address = address_fixture
        address = hyperlink.address
        hyperlink.part.target_ref.assert_called_once_with(rId)
        assert address == expected_address

    def it_knows_when_theres_no_url(self, no_address_fixture):
        hyperlink = no_address_fixture
        assert hyperlink.address is None

    def it_can_remove_its_url(self, remove_fixture):
        hyperlink, calls, expected_xml = remove_fixture

        hyperlink.address = None

        assert hyperlink.part.drop_rel.call_args_list == calls
        assert hyperlink._element.xml == expected_xml

    def it_can_set_its_target_url(self, update_fixture):
        hyperlink, url, calls, expected_xml = update_fixture

        hyperlink.address = url

        assert hyperlink.part.relate_to.call_args_list == calls
        assert hyperlink._element.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def address_fixture(self, part_prop_):
        cNvPr_cxml, rId = 'p:cNvPr/a:hlinkClick{r:id=rId1}', 'rId1'
        expected_address = 'http://foobar.com'
        cNvPr = element(cNvPr_cxml)
        hyperlink = Hyperlink(cNvPr, None)
        part_prop_.return_value.target_ref.return_value = expected_address
        return hyperlink, rId, expected_address

    @pytest.fixture(params=[
        'p:cNvPr',
        'p:cNvPr/a:hlinkClick',
    ])
    def no_address_fixture(self, request):
        cNvPr_cxml = request.param
        cNvPr = element(cNvPr_cxml)
        hyperlink = Hyperlink(cNvPr, None)
        return hyperlink

    @pytest.fixture(params=[
        ('p:cNvPr{a:a=a,r:r=r}',                         []),
        ('p:cNvPr{a:a=a,r:r=r}/a:hlinkClick',            []),
        ('p:cNvPr{a:a=a,r:r=r}/a:hlinkClick{r:id=rId3}', [call('rId3')]),
    ])
    def remove_fixture(self, request, part_prop_):
        cNvPr_cxml, calls = request.param
        cNvPr = element(cNvPr_cxml)
        expected_xml = xml('p:cNvPr{a:a=a,r:r=r}')
        hyperlink = Hyperlink(cNvPr, None)
        return hyperlink, calls, expected_xml

    @pytest.fixture(params=[
        ('p:cNvPr{a:a=a,r:r=r}', 'http://foo.com',              False, True,
         'p:cNvPr{a:a=a,r:r=r}/a:hlinkClick{r:id=rId3}'),
        ('p:cNvPr{a:a=a,r:r=r}', 'http://bar.com',              True,  True,
         'p:cNvPr{a:a=a,r:r=r}/a:hlinkHover{r:id=rId3}'),
        ('p:cNvPr/a:hlinkClick{r:id=rId6}', 'http://baz.com',   False, True,
         'p:cNvPr/a:hlinkClick{r:id=rId3}'),
        ('p:cNvPr/a:hlinkHover{r:id=rId6}', 'http://zab.com',   True,  True,
         'p:cNvPr/a:hlinkHover{r:id=rId3}'),
        ('p:cNvPr/a:hlinkHover{r:id=rId6}', 'http://boo.com',   False, True,
         'p:cNvPr/(a:hlinkClick{r:id=rId3},a:hlinkHover{r:id=rId6})'),
        ('p:cNvPr{a:a=a,r:r=r}/a:hlinkClick{r:id=rId6}', None,  False, False,
         'p:cNvPr{a:a=a,r:r=r}'),
    ])
    def update_fixture(self, request, part_prop_):
        cNvPr_cxml, url, hover, called, expected_cNvPr_cxml = request.param
        cNvPr = element(cNvPr_cxml)
        hyperlink = Hyperlink(cNvPr, None, hover)
        calls = [call(url, RT.HYPERLINK, is_external=True)] if called else []
        part_prop_.return_value.relate_to.return_value = 'rId3'
        expected_xml = xml(expected_cNvPr_cxml)
        return hyperlink, url, calls, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, Hyperlink, 'part')
