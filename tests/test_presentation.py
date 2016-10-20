# encoding: utf-8

"""
Test suite for pptx.presentation module.
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.parts.coreprops import CorePropertiesPart
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import NotesMasterPart
from pptx.presentation import Presentation
from pptx.slide import SlideLayouts, SlideMaster, SlideMasters, Slides

from .unitutil.cxml import element, xml
from .unitutil.mock import class_mock, instance_mock, property_mock


class DescribePresentation(object):

    def it_knows_the_height_of_its_slides(self, sld_height_get_fixture):
        prs, expected_value = sld_height_get_fixture
        assert prs.slide_height == expected_value

    def it_can_change_the_height_of_its_slides(self, sld_height_set_fixture):
        prs, slide_height, expected_xml = sld_height_set_fixture
        prs.slide_height = slide_height
        assert prs._element.xml == expected_xml

    def it_knows_the_width_of_its_slides(self, sld_width_get_fixture):
        prs, expected_value = sld_width_get_fixture
        assert prs.slide_width == expected_value

    def it_can_change_the_width_of_its_slides(self, sld_width_set_fixture):
        prs, slide_width, expected_xml = sld_width_set_fixture
        prs.slide_width = slide_width
        assert prs._element.xml == expected_xml

    def it_knows_its_part(self, part_fixture):
        prs, prs_part_ = part_fixture
        assert prs.part is prs_part_

    def it_provides_access_to_its_core_properties(self, core_props_fixture):
        prs, core_properties_ = core_props_fixture
        assert prs.core_properties is core_properties_

    def it_provides_access_to_its_notes_master(self, notes_master_fixture):
        prs, notes_master_ = notes_master_fixture
        assert prs.notes_master is notes_master_

    def it_provides_access_to_its_slides(self, slides_fixture):
        prs, rename_slide_parts_, rIds = slides_fixture[:3]
        Slides_, slides_, expected_xml = slides_fixture[3:]
        slides = prs.slides
        rename_slide_parts_.assert_called_once_with(rIds)
        Slides_.assert_called_once_with(
            prs._element.xpath('p:sldIdLst')[0], prs
        )
        assert prs._element.xml == expected_xml
        assert slides is slides_

    def it_provides_access_to_its_slide_layouts(self, layouts_fixture):
        prs, slide_layouts_ = layouts_fixture
        assert prs.slide_layouts is slide_layouts_

    def it_provides_access_to_its_slide_master(self, master_fixture):
        prs, getitem_, slide_master_ = master_fixture
        slide_master = prs.slide_master
        getitem_.assert_called_once_with(0)
        assert slide_master is slide_master_

    def it_provides_access_to_its_slide_masters(self, masters_fixture):
        prs, SlideMasters_, slide_masters_, expected_xml = masters_fixture
        slide_masters = prs.slide_masters
        SlideMasters_.assert_called_once_with(
            prs._element.xpath('p:sldMasterIdLst')[0], prs
        )
        assert slide_masters is slide_masters_
        assert prs._element.xml == expected_xml

    def it_can_save_the_presentation_to_a_file(self, save_fixture):
        prs, file_, prs_part_ = save_fixture
        prs.save(file_)
        prs_part_.save.assert_called_once_with(file_)

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def core_props_fixture(self, prs_part_, core_properties_):
        prs = Presentation(None, prs_part_)
        prs_part_.core_properties = core_properties_
        return prs, core_properties_

    @pytest.fixture
    def layouts_fixture(self, masters_prop_, slide_layouts_):
        prs = Presentation(None, None)
        masters_prop_.return_value.__getitem__.return_value.slide_layouts = (
            slide_layouts_
        )
        return prs, slide_layouts_

    @pytest.fixture
    def master_fixture(self, masters_prop_, slide_master_):
        prs = Presentation(None, None)
        getitem_ = masters_prop_.return_value.__getitem__
        getitem_.return_value = slide_master_
        return prs, getitem_, slide_master_

    @pytest.fixture(params=[
        ('p:presentation',
         'p:presentation/p:sldMasterIdLst'),
        ('p:presentation/p:sldMasterIdLst',
         'p:presentation/p:sldMasterIdLst'),
    ])
    def masters_fixture(self, request, SlideMasters_, slide_masters_):
        prs_cxml, expected_cxml = request.param
        prs = Presentation(element(prs_cxml), None)
        expected_xml = xml(expected_cxml)
        return prs, SlideMasters_, slide_masters_, expected_xml

    @pytest.fixture
    def notes_master_fixture(self, prs_part_, notes_master_):
        prs = Presentation(None, prs_part_)
        prs_part_.notes_master = notes_master_
        return prs, notes_master_

    @pytest.fixture
    def part_fixture(self, prs_part_):
        prs = Presentation(None, prs_part_)
        return prs, prs_part_

    @pytest.fixture
    def save_fixture(self, prs_part_):
        prs = Presentation(None, prs_part_)
        file_ = 'foobar.docx'
        return prs, file_, prs_part_

    @pytest.fixture(params=[
        ('p:presentation',                None),
        ('p:presentation/p:sldSz{cy=42}', 42),
    ])
    def sld_height_get_fixture(self, request):
        prs_cxml, expected_value = request.param
        prs = Presentation(element(prs_cxml), None)
        return prs, expected_value

    @pytest.fixture(params=[
        ('p:presentation',
         'p:presentation/p:sldSz{cy=914400}'),
        ('p:presentation/p:sldSz{cy=424242}',
         'p:presentation/p:sldSz{cy=914400}'),
    ])
    def sld_height_set_fixture(self, request):
        prs_cxml, expected_cxml = request.param
        prs = Presentation(element(prs_cxml), None)
        expected_xml = xml(expected_cxml)
        return prs, 914400, expected_xml

    @pytest.fixture(params=[
        ('p:presentation',                None),
        ('p:presentation/p:sldSz{cx=42}', 42),
    ])
    def sld_width_get_fixture(self, request):
        prs_cxml, expected_value = request.param
        prs = Presentation(element(prs_cxml), None)
        return prs, expected_value

    @pytest.fixture(params=[
        ('p:presentation',
         'p:presentation/p:sldSz{cx=914400}'),
        ('p:presentation/p:sldSz{cx=424242}',
         'p:presentation/p:sldSz{cx=914400}'),
    ])
    def sld_width_set_fixture(self, request):
        prs_cxml, expected_cxml = request.param
        prs = Presentation(element(prs_cxml), None)
        expected_xml = xml(expected_cxml)
        return prs, 914400, expected_xml

    @pytest.fixture(params=[
        ('p:presentation', [], 'p:presentation/p:sldIdLst'),
        ('p:presentation/p:sldIdLst/p:sldId{r:id=a}', ['a'],
         'p:presentation/p:sldIdLst/p:sldId{r:id=a}'),
        ('p:presentation/p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})',
         ['a', 'b'],
         'p:presentation/p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})'),
    ])
    def slides_fixture(self, request, part_prop_, Slides_, slides_):
        prs_cxml, rIds, expected_cxml = request.param
        prs = Presentation(element(prs_cxml), None)
        rename_slide_parts_ = part_prop_.return_value.rename_slide_parts
        expected_xml = xml(expected_cxml)
        return prs, rename_slide_parts_, rIds, Slides_, slides_, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def core_properties_(self, request):
        return instance_mock(request, CorePropertiesPart)

    @pytest.fixture
    def masters_prop_(self, request):
        return property_mock(request, Presentation, 'slide_masters')

    @pytest.fixture
    def notes_master_(self, request):
        return instance_mock(request, NotesMasterPart)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, Presentation, 'part')

    @pytest.fixture
    def prs_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def slide_layouts_(self, request):
        return instance_mock(request, SlideLayouts)

    @pytest.fixture
    def SlideMasters_(self, request, slide_masters_):
        return class_mock(
            request, 'pptx.presentation.SlideMasters',
            return_value=slide_masters_
        )

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)

    @pytest.fixture
    def slide_masters_(self, request):
        return instance_mock(request, SlideMasters)

    @pytest.fixture
    def Slides_(self, request, slides_):
        return class_mock(
            request, 'pptx.presentation.Slides', return_value=slides_
        )

    @pytest.fixture
    def slides_(self, request):
        return instance_mock(request, Slides)
