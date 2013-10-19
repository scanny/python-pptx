# -*- coding: utf-8 -*-
#
# test_shapes.py
#
# Copyright (C) 2012, 2013 Steve Canny scanny@cisco.com
#
# This module is part of python-pptx and is released under
# the MIT License: http://www.opensource.org/licenses/mit-license.php

"""Test suite for pptx.shapes module."""

import os

from hamcrest import assert_that, equal_to, is_, is_not, same_instance
from mock import MagicMock, Mock, patch, PropertyMock
# from unittest2 import skip

from pptx.constants import MSO_AUTO_SHAPE_TYPE as MAST, MSO
from pptx.oxml import _SubElement, nsdecls, oxml_fromstring, oxml_parse
from pptx.presentation import _SlideLayout
from pptx.shapes import (
    _Adjustment, _AdjustmentCollection, _AutoShapeType, _BaseShape, _Cell,
    _CellCollection, _Column, _ColumnCollection, _Placeholder, _Row,
    _RowCollection, _Shape, _ShapeCollection
)
from pptx.spec import namespaces
from pptx.spec import (
    PH_TYPE_CTRTITLE, PH_TYPE_DT, PH_TYPE_FTR, PH_TYPE_OBJ, PH_TYPE_SLDNUM,
    PH_TYPE_SUBTITLE, PH_TYPE_TBL, PH_TYPE_TITLE, PH_ORIENT_HORZ,
    PH_ORIENT_VERT, PH_SZ_FULL, PH_SZ_HALF, PH_SZ_QUARTER
)
from pptx.util import Inches
from testdata import (
    a_prstGeom, test_shape_elements, test_shapes, test_table_objects
)
from testing import TestCase


# module globals -------------------------------------------------------------
def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
test_file_dir = absjoin(thisdir, 'test_files')

test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_bmp_path = absjoin(test_file_dir, 'python.bmp')
new_image_path = absjoin(test_file_dir, 'monty-truth.png')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')
images_pptx_path = absjoin(test_file_dir, 'with_images.pptx')

nsmap = namespaces('a', 'r', 'p')


def _sldLayout1():
    path = os.path.join(thisdir, 'test_files/slideLayout1.xml')
    sldLayout = oxml_parse(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = _ShapeCollection(spTree)
    return shapes


class Test_Adjustment(TestCase):
    """Test _Adjustment"""
    def test_it_should_have_correct_effective_value(self):
        """_Adjustment.effective_value is correct"""
        # setup ------------------------
        name = "don't care"
        cases = (
            # no actual, effective should be determined by default value
            (50000, None, 0.5),
            # actual matches default
            (50000, 50000, 0.5),
            # actual is different than default
            (50000, 12500, 0.125),
            # actual is zero
            (50000, 0, 0.0),
            # negative default
            (-20833, None, -0.20833),
            # negative actual
            (-20833, -5678901, -56.78901),
        )
        # verify -----------------------
        for def_val, actual, expected in cases:
            adjustment = _Adjustment(name, def_val, actual)
            assert_that(adjustment.effective_value, is_(equal_to(expected)))


class Test_AdjustmentCollection(TestCase):
    """Test _AdjustmentCollection"""
    def test_it_should_load_default_adjustment_values(self):
        """_AdjustmentCollection() loads default adjustment values"""
        # setup ------------------------
        cases = (
            ('rect', ()),
            ('chevron', (('adj', 50000),)),
            ('accentBorderCallout1',
             (('adj1', 18750), ('adj2', -8333), ('adj3', 112500),
              ('adj4', -38333))),
            ('wedgeRoundRectCallout',
             (('adj1', -20833), ('adj2', 62500), ('adj3', 16667))),
            ('circularArrow',
             (('adj1', 12500), ('adj2', 1142319), ('adj3', 20457681),
              ('adj4', 10800000), ('adj5', 12500))),
        )
        for prst, expected_values in cases:
            prstGeom = a_prstGeom(prst).with_avLst.element
            adjustments = _AdjustmentCollection(prstGeom)._adjustments
            # verify -------------------
            reason = ("\n     failed for prst: '%s'" % prst)
            assert_that(len(adjustments),
                        is_(equal_to(len(expected_values))),
                        reason)
            actuals = tuple([(adj.name, adj.def_val) for adj in adjustments])
            assert_that(actuals, is_(equal_to(expected_values)), reason)

    def test_it_should_load_adj_val_actuals_from_xml(self):
        """_AdjustmentCollection() loads adjustment value actuals from XML"""
        # setup ------------------------
        def expected_actual(adj_vals, adjustment_name):
            for name, actual in adj_vals:
                if name == adjustment_name:
                    return actual
            return None

        cases = (
            # no adjustment values in xml or spec
            (a_prstGeom('rect').with_avLst, ()),
            # no adjustment values in xml, but some in spec
            (a_prstGeom('circularArrow'), ()),
            # adjustment value in xml but none in spec
            (a_prstGeom('rect').with_gd(), ()),
            # middle adjustment value in xml
            (a_prstGeom('mathDivide').with_gd(name='adj2'),
             (('adj2', 25000),)),
            # all adjustment values in xml
            (a_prstGeom('wedgeRoundRectCallout').with_gd(111, 'adj1')
                                                .with_gd(222, 'adj2')
                                                .with_gd(333, 'adj3'),
             (('adj1', 111), ('adj2', 222), ('adj3', 333))),
        )
        # verify -----------------------
        for prstGeom_builder, adj_vals in cases:
            prstGeom = prstGeom_builder.element
            adjustments = _AdjustmentCollection(prstGeom)._adjustments
            for adjustment in adjustments:
                expected_value = expected_actual(adj_vals, adjustment.name)
                reason = (
                    "failed for adj val name '%s', for this XML:\n\n%s" %
                    (adjustment.name, prstGeom_builder.xml))
                assert_that(adjustment.actual, is_(expected_value), reason)

    def test_it_should_return_effective_value_on_indexed_access(self):
        """_AdjustmentCollection[n] is normalized effective value of nth"""
        # setup ------------------------
        cases = (
            ('rect', ()),
            ('chevron', (0.5,)),
            ('circularArrow', (0.125, 11.42319, 204.57681, 108.0, 0.125)),
        )
        for prst, expected_values in cases:
            prstGeom = a_prstGeom(prst).element
            # exercise -----------------
            adjustments = _AdjustmentCollection(prstGeom)
            # verify -------------------
            reason = "failed on case: prst=='%s'" % prst
            assert_that(len(adjustments), is_(len(expected_values)), reason)
            retvals = tuple([adj for adj in adjustments])
            assert_that(retvals, is_(equal_to(expected_values)), reason)

    def test_it_should_update_actual_value_on_indexed_assignment(self):
        """Assignment to _AdjustmentCollection[n] updates nth actual"""
        # setup ------------------------
        cases = (
            ('chevron', 0, 0.5, 50000),
            ('circularArrow', 2, 99.99, 9999000),
        )
        prst = 'chevron'
        for prst, idx, new_value, expected in cases:
            prstGeom = a_prstGeom(prst).element
            adjustments = _AdjustmentCollection(prstGeom)
            # exercise -----------------
            adjustments[idx] = new_value
            # verify -------------------
            reason = "failed on case: prst=='%s'" % prst
            assert_that(adjustments._adjustments[idx].actual,
                        is_(equal_to(expected)),
                        reason)

    def test_it_should_round_trip_indexed_assignment(self):
        """Assignment to _AdjustmentCollection[n] round-trips"""
        # setup ------------------------
        new_value = 0.375
        prstGeom = a_prstGeom('chevron').element
        adjustments = _AdjustmentCollection(prstGeom)
        assert_that(adjustments[0], is_not(equal_to(new_value)))
        # exercise ---------------------
        adjustments[0] = new_value
        # verify -----------------------
        assert_that(adjustments[0], is_(equal_to(new_value)))

    def test_it_should_raise_on_bad_key(self):
        """_AdjustmentCollection[idx] raises on invalid idx"""
        # setup ------------------------
        prstGeom = a_prstGeom('chevron').element
        adjustments = _AdjustmentCollection(prstGeom)
        # verify -----------------------
        with self.assertRaises(IndexError):
            adjustments[-6]
        with self.assertRaises(IndexError):
            adjustments[6]
        with self.assertRaises(TypeError):
            adjustments[0.0]
        with self.assertRaises(TypeError):
            adjustments['0']
        with self.assertRaises(IndexError):
            adjustments[-6] = 1.0
        with self.assertRaises(IndexError):
            adjustments[6] = 1.0
        with self.assertRaises(TypeError):
            adjustments[0.0] = 1.0
        with self.assertRaises(TypeError):
            adjustments['0'] = 1.0

    def test_it_should_raise_on_assigned_bad_value(self):
        """_AdjustmentCollection[n] = val raises on val is not number"""
        # setup ------------------------
        prstGeom = a_prstGeom('chevron').element
        adjustments = _AdjustmentCollection(prstGeom)
        # verify -----------------------
        with self.assertRaises(ValueError):
            adjustments[0] = 'foobar'

    def test_writes_adj_vals_to_xml_on_assignment(self):
        """_AdjustmentCollection writes adj vals to XML on assignment"""
        # setup ------------------------
        prstGeom = a_prstGeom('chevron').element
        adjustments = _AdjustmentCollection(prstGeom)
        __prstGeom = Mock(name='__prstGeom')
        adjustments._AdjustmentCollection__prstGeom = __prstGeom
        # exercise ---------------------
        adjustments[0] = 0.999
        # verify -----------------------
        assert_that(__prstGeom.rewrite_guides.call_count, is_(1))


class Test_AutoShapeType(TestCase):
    """Test _AutoShapeType"""
    def test_construction_return_values(self):
        """_AutoShapeType() returns instance with correct property values"""
        # setup ------------------------
        id_ = MAST.ROUNDED_RECTANGLE
        prst = 'roundRect'
        basename = 'Rounded Rectangle'
        # exercise ---------------------
        autoshape_type = _AutoShapeType(id_)
        # verify -----------------------
        assert_that(autoshape_type.autoshape_type_id, is_(equal_to(id_)))
        assert_that(autoshape_type.prst, is_(equal_to(prst)))
        assert_that(autoshape_type.basename, is_(equal_to(basename)))

    def test_default_adjustment_values_return_value(self):
        """_AutoShapeType.default_adjustment_values() return val correct"""
        # setup ------------------------
        cases = (
            ('rect', ()),
            ('chevron', (('adj', 50000),)),
            ('leftCircularArrow',
             (('adj1', 12500), ('adj2', -1142319), ('adj3', 1142319),
              ('adj4', 10800000), ('adj5', 12500))),
        )
        # verify -----------------------
        for prst, expected_vals in cases:
            def_adj_vals = _AutoShapeType.default_adjustment_values(prst)
            assert_that(def_adj_vals, is_(equal_to(expected_vals)))

    def test__lookup_id_by_prst_return_value(self):
        """_AutoShapeType._lookup_id_by_prst() return value is correct"""
        # setup ------------------------
        autoshape_type_id = MAST.ROUNDED_RECTANGLE
        prst = 'roundRect'
        # exercise ---------------------
        retval = _AutoShapeType._lookup_id_by_prst(prst)
        # verify -----------------------
        assert_that(retval, is_(equal_to(autoshape_type_id)))

    def test__lookup_id_raises_on_bad_prst(self):
        """_AutoShapeType._lookup_id_by_prst() raises on bad prst"""
        # setup ------------------------
        prst = 'badPrst'
        # verify -----------------------
        with self.assertRaises(KeyError):
            _AutoShapeType._lookup_id_by_prst(prst)

    def test_second_construction_returns_cached_instance(self):
        """_AutoShapeType() returns cached instance on duplicate call"""
        # setup ------------------------
        id_ = MAST.ROUNDED_RECTANGLE
        ast1 = _AutoShapeType(id_)
        # exercise ---------------------
        ast2 = _AutoShapeType(id_)
        # verify -----------------------
        assert_that(ast2, is_(equal_to(ast1)))

    def test_construction_raises_on_bad_autoshape_type_id(self):
        """_AutoShapeType() raises on bad auto shape type id"""
        # setup ------------------------
        autoshape_type_id = 9999
        # verify -----------------------
        with self.assertRaises(KeyError):
            _AutoShapeType(autoshape_type_id)


class Test_BaseShape(TestCase):
    """Test _BaseShape"""
    def setUp(self):
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        self.sld = oxml_parse(path).getroot()
        xpath = './p:cSld/p:spTree/p:pic'
        pic = self.sld.xpath(xpath, namespaces=nsmap)[0]
        self.base_shape = _BaseShape(pic)

    def test_has_textframe_value(self):
        """_BaseShape.has_textframe value correct"""
        # setup ------------------------
        spTree = self.sld.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
        shapes = _ShapeCollection(spTree)
        indexes = []
        # exercise ---------------------
        for idx, shape in enumerate(shapes):
            if shape.has_textframe:
                indexes.append(idx)
        # verify -----------------------
        expected = [0, 1, 3, 5, 6]
        actual = indexes
        msg = ("expected txBody element in shapes %s, got %s" %
               (expected, actual))
        self.assertEqual(expected, actual, msg)

    def test_id_value(self):
        """_BaseShape.id value is correct"""
        # exercise ---------------------
        id = self.base_shape.id
        # verify -----------------------
        expected = 6
        actual = id
        msg = "expected %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_is_placeholder_true_for_placeholder(self):
        """_BaseShape.is_placeholder True for placeholder shape"""
        # setup ------------------------
        xpath = './p:cSld/p:spTree/p:sp'
        sp = self.sld.xpath(xpath, namespaces=nsmap)[0]
        base_shape = _BaseShape(sp)
        # verify -----------------------
        actual = base_shape.is_placeholder
        msg = "expected True, got %s" % (actual)
        self.assertTrue(actual, msg)

    def test_is_placeholder_false_for_non_placeholder(self):
        """_BaseShape.is_placeholder False for non-placeholder shape"""
        # verify -----------------------
        actual = self.base_shape.is_placeholder
        msg = "expected False, got %s" % (actual)
        self.assertFalse(actual, msg)

    def test__is_title_true_for_title_placeholder(self):
        """_BaseShape._is_title True for title placeholder shape"""
        # setup ------------------------
        xpath = './p:cSld/p:spTree/p:sp'
        title_placeholder_sp = self.sld.xpath(xpath, namespaces=nsmap)[0]
        base_shape = _BaseShape(title_placeholder_sp)
        # verify -----------------------
        actual = base_shape._is_title
        msg = "expected True, got %s" % (actual)
        self.assertTrue(actual, msg)

    def test__is_title_false_for_no_ph_element(self):
        """_BaseShape._is_title False on shape has no <p:ph> element"""
        # setup ------------------------
        self.base_shape._element = Mock(name='_element')
        self.base_shape._element.xpath.return_value = []
        # verify -----------------------
        assert_that(self.base_shape._is_title, is_(False))

    def test_name_value(self):
        """_BaseShape.name value is correct"""
        # exercise ---------------------
        name = self.base_shape.name
        # verify -----------------------
        expected = 'Picture 5'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_shape_name_returns_none_for_unimplemented_shape_types(self):
        """_BaseShape.shape_name returns None for unimplemented shape types"""
        assert_that(self.base_shape.shape_type, is_(None))

    def test_textframe_raises_on_no_textframe(self):
        """_BaseShape.textframe raises on shape with no text frame"""
        with self.assertRaises(ValueError):
            self.base_shape.textframe

    def test_text_setter_structure_and_value(self):
        """assign to _BaseShape.text yields single run para set to value"""
        # setup ------------------------
        test_text = 'python-pptx was here!!'
        xpath = './p:cSld/p:spTree/p:sp'
        textbox_sp = self.sld.xpath(xpath, namespaces=nsmap)[2]
        base_shape = _BaseShape(textbox_sp)
        # exercise ---------------------
        base_shape.text = test_text
        # verify paragraph count ------
        expected = 1
        actual = len(base_shape.textframe.paragraphs)
        msg = "expected paragraph count %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)
        # verify value ----------------
        expected = test_text
        actual = base_shape.textframe.paragraphs[0].runs[0].text
        msg = "expected text '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_text_setter_raises_on_no_textframe(self):
        """assignment to _BaseShape.text raises for shape with no text frame"""
        with self.assertRaises(TypeError):
            self.base_shape.text = 'test text'


class Test_Cell(TestCase):
    """Test _Cell"""
    def setUp(self):
        tc_xml = '<a:tc %s><a:txBody><a:p/></a:txBody></a:tc>' % nsdecls('a')
        test_tc_elm = oxml_fromstring(tc_xml)
        self.cell = _Cell(test_tc_elm)

    def test_text_round_trips_intact(self):
        """_Cell.text (setter) sets cell text"""
        # setup ------------------------
        test_text = 'test_text'
        # exercise ---------------------
        self.cell.text = test_text
        # verify -----------------------
        text = self.cell.textframe.paragraphs[0].runs[0].text
        assert_that(text, is_(equal_to(test_text)))

    def test_margin_values(self):
        """_Cell.margin_x values are calculated correctly"""
        # mockery ----------------------
        marT_val, marR_val, marB_val, marL_val = 12, 34, 56, 78
        # -- CT_TableCell
        tc = MagicMock()
        marT = type(tc).marT = PropertyMock(return_value=marT_val)
        marR = type(tc).marR = PropertyMock(return_value=marR_val)
        marB = type(tc).marB = PropertyMock(return_value=marB_val)
        marL = type(tc).marL = PropertyMock(return_value=marL_val)
        # setup ------------------------
        cell = _Cell(tc)
        # exercise ---------------------
        margin_top = cell.margin_top
        margin_right = cell.margin_right
        margin_bottom = cell.margin_bottom
        margin_left = cell.margin_left
        # verify -----------------------
        marT.assert_called_once_with()
        marR.assert_called_once_with()
        marB.assert_called_once_with()
        marL.assert_called_once_with()
        assert_that(margin_top, is_(equal_to(marT_val)))
        assert_that(margin_right, is_(equal_to(marR_val)))
        assert_that(margin_bottom, is_(equal_to(marB_val)))
        assert_that(margin_left, is_(equal_to(marL_val)))

    def test_margin_assignment(self):
        """Assignment to _Cell.margin_x sets value"""
        # mockery ----------------------
        # -- CT_TableCell
        tc = MagicMock()
        marT = type(tc).marT = PropertyMock()
        marR = type(tc).marR = PropertyMock()
        marB = type(tc).marB = PropertyMock()
        marL = type(tc).marL = PropertyMock()
        # setup ------------------------
        top, right, bottom, left = 12, 34, 56, 78
        cell = _Cell(tc)
        # exercise ---------------------
        cell.margin_top = top
        cell.margin_right = right
        cell.margin_bottom = bottom
        cell.margin_left = left
        # verify -----------------------
        marT.assert_called_once_with(top)
        marR.assert_called_once_with(right)
        marB.assert_called_once_with(bottom)
        marL.assert_called_once_with(left)

    def test_margin_assignment_raises_on_not_int_or_none(self):
        """Assignment to _Cell.margin_x raises for not (int or none)"""
        # setup ------------------------
        cell = test_table_objects.cell
        bad_margin = 'foobar'
        # verify -----------------------
        with self.assertRaises(ValueError):
            cell.margin_top = bad_margin
        with self.assertRaises(ValueError):
            cell.margin_right = bad_margin
        with self.assertRaises(ValueError):
            cell.margin_bottom = bad_margin
        with self.assertRaises(ValueError):
            cell.margin_left = bad_margin

    @patch('pptx.shapes.VerticalAnchor')
    def test_vertical_anchor_value(self, VerticalAnchor):
        """_Cell.vertical_anchor value is calculated correctly"""
        # mockery ----------------------
        # loose mocks
        anchor_val = Mock(name='anchor_val')
        vertical_anchor = Mock(name='vertical_anchor')
        # CT_TableCell
        tc = MagicMock()
        anchor_prop = type(tc).anchor = PropertyMock(return_value=anchor_val)
        # VerticalAnchor
        from_text_anchoring_type = VerticalAnchor.from_text_anchoring_type
        from_text_anchoring_type.return_value = vertical_anchor
        # setup ------------------------
        cell = _Cell(tc)
        # exercise ---------------------
        retval = cell.vertical_anchor
        # verify -----------------------
        anchor_prop.assert_called_once_with()
        from_text_anchoring_type.assert_called_once_with(anchor_val)
        assert_that(retval, is_(same_instance(vertical_anchor)))

    @patch('pptx.shapes.VerticalAnchor')
    def test_vertical_anchor_assignment(self, VerticalAnchor):
        """Assignment to _Cell.vertical_anchor sets value"""
        # mockery ----------------------
        # -- loose mocks
        vertical_anchor = Mock(name='vertical_anchor')
        anchor_val = Mock(name='anchor_val')
        # -- CT_TableCell
        tc = MagicMock()
        anchor_prop = type(tc).anchor = PropertyMock()
        # -- VerticalAnchor
        to_text_anchoring_type = VerticalAnchor.to_text_anchoring_type
        to_text_anchoring_type.return_value = anchor_val
        # setup ------------------------
        cell = _Cell(tc)
        # exercise ---------------------
        cell.vertical_anchor = vertical_anchor
        # verify -----------------------
        to_text_anchoring_type.assert_called_once_with(vertical_anchor)
        anchor_prop.assert_called_once_with(anchor_val)


class Test_CellCollection(TestCase):
    """Test _CellCollection"""
    def setUp(self):
        tr_xml = (
            '<a:tr %s h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:t'
            'c><a:txBody><a:p/></a:txBody></a:tc></a:tr>' % nsdecls('a')
        )
        test_tr_elm = oxml_fromstring(tr_xml)
        self.cells = _CellCollection(test_tr_elm)

    def test_is_indexable(self):
        """_CellCollection indexable (e.g. no TypeError on 'cells[0]')"""
        # verify -----------------------
        try:
            self.cells[0]
        except TypeError:
            msg = "'_CellCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_CellCollection is iterable (e.g. ``for cell in cells:``)"""
        # verify -----------------------
        count = 0
        try:
            for cell in self.cells:
                count += 1
        except TypeError:
            msg = "_CellCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_CellCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.cells[9]

    def test_cell_count_correct(self):
        """len(_CellCollection) returns correct cell count"""
        # verify -----------------------
        assert_that(len(self.cells), is_(equal_to(2)))


class Test_Column(TestCase):
    """Test _Column"""
    def setUp(self):
        gridCol_xml = '<a:gridCol %s w="3048000"/>' % nsdecls('a')
        test_gridCol_elm = oxml_fromstring(gridCol_xml)
        self.column = _Column(test_gridCol_elm, Mock(name='table'))

    def test_width_from_xml_correct(self):
        """_Column.width returns correct value from gridCol XML element"""
        # verify -----------------------
        assert_that(self.column.width, is_(equal_to(3048000)))

    def test_width_round_trips_intact(self):
        """_Column.width round-trips intact"""
        # setup ------------------------
        self.column.width = 999
        # verify -----------------------
        assert_that(self.column.width, is_(equal_to(999)))

    def test_set_width_raises_on_bad_value(self):
        """_Column.width raises on attempt to assign invalid value"""
        test_cases = ('abc', '1', -1)
        for value in test_cases:
            with self.assertRaises(ValueError):
                self.column.width = value


class Test_ColumnCollection(TestCase):
    """Test _ColumnCollection"""
    def setUp(self):
        tbl_xml = (
            '<a:tbl %s><a:tblGrid><a:gridCol w="3048000"/><a:gridCol w="30480'
            '00"/></a:tblGrid></a:tbl>' % nsdecls('a')
        )
        test_tbl_elm = oxml_fromstring(tbl_xml)
        self.columns = _ColumnCollection(test_tbl_elm, Mock(name='table'))

    def test_is_indexable(self):
        """_ColumnCollection indexable (e.g. no TypeError on 'columns[0]')"""
        # verify -----------------------
        try:
            self.columns[0]
        except TypeError:
            msg = "'_ColumnCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_ColumnCollection is iterable (e.g. ``for col in columns:``)"""
        # verify -----------------------
        count = 0
        try:
            for col in self.columns:
                count += 1
        except TypeError:
            msg = "_ColumnCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_ColumnCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.columns[9]

    def test_column_count_correct(self):
        """len(_ColumnCollection) returns correct column count"""
        # verify -----------------------
        assert_that(len(self.columns), is_(equal_to(2)))


class Test_Placeholder(TestCase):
    """Test _Placeholder"""
    def test_property_values(self):
        """_Placeholder property values are correct"""
        # setup ------------------------
        expected_values = (
            (PH_TYPE_CTRTITLE, PH_ORIENT_HORZ, PH_SZ_FULL,     0),
            (PH_TYPE_DT,       PH_ORIENT_HORZ, PH_SZ_HALF,    10),
            (PH_TYPE_SUBTITLE, PH_ORIENT_VERT, PH_SZ_FULL,     1),
            (PH_TYPE_TBL,      PH_ORIENT_HORZ, PH_SZ_QUARTER, 14),
            (PH_TYPE_SLDNUM,   PH_ORIENT_HORZ, PH_SZ_QUARTER, 12),
            (PH_TYPE_FTR,      PH_ORIENT_HORZ, PH_SZ_QUARTER, 11))
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        for idx, sp in enumerate(shapes):
            ph = _Placeholder(sp)
            values = (ph.type, ph.orient, ph.sz, ph.idx)
            # verify ----------------------
            expected = expected_values[idx]
            actual = values
            msg = ("expected shapes[%d] values %s, got %s"
                   % (idx, expected, actual))
            self.assertEqual(expected, actual, msg)


class Test_Picture(TestCase):
    """Test _Picture"""
    def test_shape_type_value_correct_for_picture(self):
        """_Shape.shape_type value is correct for picture"""
        # setup ------------------------
        picture = test_shapes.picture
        # verify -----------------------
        assert_that(picture.shape_type, is_(equal_to(MSO.PICTURE)))


class Test_Row(TestCase):
    """Test _Row"""
    def setUp(self):
        tr_xml = (
            '<a:tr %s h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:t'
            'c><a:txBody><a:p/></a:txBody></a:tc></a:tr>' % nsdecls('a')
        )
        test_tr_elm = oxml_fromstring(tr_xml)
        self.row = _Row(test_tr_elm, Mock(name='table'))

    def test_height_from_xml_correct(self):
        """_Row.height returns correct value from tr XML element"""
        # verify -----------------------
        assert_that(self.row.height, is_(equal_to(370840)))

    def test_height_round_trips_intact(self):
        """_Row.height round-trips intact"""
        # setup ------------------------
        self.row.height = 999
        # verify -----------------------
        assert_that(self.row.height, is_(equal_to(999)))

    def test_set_height_raises_on_bad_value(self):
        """_Row.height raises on attempt to assign invalid value"""
        test_cases = ('abc', '1', -1)
        for value in test_cases:
            with self.assertRaises(ValueError):
                self.row.height = value


class Test_RowCollection(TestCase):
    """Test _RowCollection"""
    def setUp(self):
        tbl_xml = (
            '<a:tbl %s><a:tr h="370840"><a:tc><a:txBody><a:p/></a:txBody></a:'
            'tc><a:tc><a:txBody><a:p/></a:txBody></a:tc></a:tr><a:tr h="37084'
            '0"><a:tc><a:txBody><a:p/></a:txBody></a:tc><a:tc><a:txBody><a:p/'
            '></a:txBody></a:tc></a:tr></a:tbl>' % nsdecls('a')
        )
        test_tbl_elm = oxml_fromstring(tbl_xml)
        self.rows = _RowCollection(test_tbl_elm, Mock(name='table'))

    def test_is_indexable(self):
        """_RowCollection indexable (e.g. no TypeError on 'rows[0]')"""
        # verify -----------------------
        try:
            self.rows[0]
        except TypeError:
            msg = "'_RowCollection' object does not support indexing"
            self.fail(msg)
        except IndexError:
            pass

    def test_is_iterable(self):
        """_RowCollection is iterable (e.g. ``for row in rows:``)"""
        # verify -----------------------
        count = 0
        try:
            for row in self.rows:
                count += 1
        except TypeError:
            msg = "_RowCollection object is not iterable"
            self.fail(msg)
        assert_that(count, is_(equal_to(2)))

    def test_raises_on_idx_out_of_range(self):
        """_RowCollection raises on index out of range"""
        with self.assertRaises(IndexError):
            self.rows[9]

    def test_row_count_correct(self):
        """len(_RowCollection) returns correct row count"""
        # verify -----------------------
        assert_that(len(self.rows), is_(equal_to(2)))


class Test_Shape(TestCase):
    """Test _Shape"""
    @patch('pptx.shapes._BaseShape.__init__')
    @patch('pptx.shapes._AdjustmentCollection')
    def test_it_initializes_adjustments_on_construction(
            self, _AdjustmentCollection, _BaseShape__init__):
        """_Shape() initializes adjustments on construction"""
        # setup ------------------------
        adjustments = Mock(name='adjustments')
        _AdjustmentCollection.return_value = adjustments
        sp = Mock(name='sp')
        # exercise ---------------------
        shape = _Shape(sp)
        # verify -----------------------
        _BaseShape__init__.assert_called_once_with(sp)
        _AdjustmentCollection.assert_called_once_with(sp.prstGeom)
        assert_that(shape.adjustments, is_(adjustments))

    def test_auto_shape_type_value_correct(self):
        """_Shape.auto_shape_type value is correct"""
        # setup ------------------------
        rounded_rectangle = test_shapes.rounded_rectangle
        # verify -----------------------
        assert_that(rounded_rectangle.auto_shape_type,
                    is_(equal_to(MAST.ROUNDED_RECTANGLE)))

    def test_auto_shape_type_raises_on_non_auto_shape(self):
        """_Shape.auto_shape_type raises on non auto shape"""
        # setup ------------------------
        textbox = test_shapes.textbox
        # verify -----------------------
        with self.assertRaises(ValueError):
            textbox.auto_shape_type

    def test_shape_type_value_correct(self):
        """_Shape.shape_type value is correct for all sub-types"""
        # setup ------------------------
        autoshape = test_shapes.autoshape
        placeholder = test_shapes.placeholder
        textbox = test_shapes.textbox
        # verify -----------------------
        assert_that(autoshape.shape_type, is_(equal_to(MSO.AUTO_SHAPE)))
        assert_that(placeholder.shape_type, is_(equal_to(MSO.PLACEHOLDER)))
        assert_that(textbox.shape_type, is_(equal_to(MSO.TEXT_BOX)))

    def test_shape_type_raises_on_unrecognized_shape_type(self):
        """_Shape.shape_type raises on unrecognized shape type"""
        # setup ------------------------
        xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/'
            '2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            '2006/main"><p:nvSpPr><p:cNvPr id="9" name="Unknown Shape Type 8"'
            '/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/></p:sp>'
        )
        sp = oxml_fromstring(xml)
        shape = _Shape(sp)
        # verify -----------------------
        with self.assertRaises(NotImplementedError):
            shape.shape_type


class Test_ShapeCollection(TestCase):
    """Test _ShapeCollection"""
    def setUp(self):
        path = absjoin(test_file_dir, 'slide1.xml')
        sld = oxml_parse(path).getroot()
        spTree = sld.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
        self.shapes = _ShapeCollection(spTree)

    def test_construction_size(self):
        """_ShapeCollection is expected size after construction"""
        # verify -----------------------
        self.assertLength(self.shapes, 9)

    def test_constructor_raises_on_contentPart_shape(self):
        """_ShapeCollection() raises on contentPart shape"""
        # setup ------------------------
        spTree = test_shape_elements.empty_spTree
        _SubElement(spTree, 'p:contentPart')
        # verify -----------------------
        with self.assertRaises(ValueError):
            _ShapeCollection(spTree)

    @patch('pptx.shapes.CT_Shape')
    @patch('pptx.shapes._Shape')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    @patch('pptx.shapes._AutoShapeType')
    def test_add_shape_collaboration(self, _AutoShapeType, __next_shape_id,
                                     _Shape, CT_Shape):
        """_ShapeCollection.add_shape() calls the right collaborators"""
        # constant values -------------
        basename = 'Rounded Rectangle'
        prst = 'roundRect'
        id_, name = 9, '%s 8' % basename
        left, top, width, height = 111, 222, 333, 444
        autoshape_type_id = MAST.ROUNDED_RECTANGLE
        # setup mockery ---------------
        autoshape_type = Mock(name='autoshape_type')
        autoshape_type.basename = basename
        autoshape_type.prst = prst
        _AutoShapeType.return_value = autoshape_type
        __next_shape_id.return_value = id_
        sp = Mock(name='sp')
        CT_Shape.new_autoshape_sp.return_value = sp
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        shape = Mock('shape')
        _Shape.return_value = shape
        # exercise ---------------------
        retval = shapes.add_shape(autoshape_type_id, left, top, width, height)
        # verify -----------------------
        _AutoShapeType.assert_called_once_with(autoshape_type_id)
        CT_Shape.new_autoshape_sp.assert_called_once_with(
            id_, name, prst, left, top, width, height)
        _Shape.assert_called_once_with(sp)
        __spTree.append.assert_called_once_with(sp)
        __shapes.append.assert_called_once_with(shape)
        assert_that(retval, is_(equal_to(shape)))

    @patch('pptx.shapes._Picture')
    @patch('pptx.shapes.CT_Picture')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_picture_collaboration(self, next_shape_id, CT_Picture,
                                       _Picture):
        """_ShapeCollection.add_picture() calls the right collaborators"""
        # constant values -------------
        file = test_image_path
        left, top, width, height = 1, 2, 3, 4
        id_, name, desc = 12, 'Picture 11', 'image1.jpeg'
        rId = 'rId1'
        # setup mockery ---------------
        next_shape_id.return_value = id_
        image = Mock(name='image', _desc=desc)
        image._scale.return_value = width, height
        rel = Mock(name='rel', _rId=rId)
        slide = Mock(name='slide')
        slide._add_image.return_value = image, rel
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = _ShapeCollection(test_shape_elements.empty_spTree, slide)
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        pic = Mock(name='pic')
        CT_Picture.new_pic.return_value = pic
        picture = Mock(name='picture')
        _Picture.return_value = picture
        # # exercise --------------------
        retval = shapes.add_picture(file, left, top, width, height)
        # verify -----------------------
        shapes._ShapeCollection__slide._add_image.assert_called_once_with(file)
        image._scale.assert_called_once_with(width, height)
        CT_Picture.new_pic.assert_called_once_with(
            id_, name, desc, rId, left, top, width, height)
        __spTree.append.assert_called_once_with(pic)
        _Picture.assert_called_once_with(pic)
        __shapes.append.assert_called_once_with(picture)
        assert_that(retval, is_(equal_to(picture)))

    @patch('pptx.shapes._Table')
    @patch('pptx.shapes.CT_GraphicalObjectFrame')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_table_collaboration(
            self, __next_shape_id, CT_GraphicalObjectFrame, _Table):
        """_ShapeCollection.add_table() calls the right collaborators"""
        # constant values -------------
        id_, name = 9, 'Table 8'
        rows, cols = 2, 3
        left, top, width, height = 111, 222, 333, 444
        # setup mockery ---------------
        __next_shape_id.return_value = id_
        graphicFrame = Mock(name='graphicFrame')
        CT_GraphicalObjectFrame.new_table.return_value = graphicFrame
        __spTree = Mock(name='__spTree')
        __shapes = Mock(name='__shapes')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        shapes._ShapeCollection__shapes = __shapes
        table = Mock('table')
        _Table.return_value = table
        # exercise ---------------------
        retval = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        __next_shape_id.assert_called_once_with()
        CT_GraphicalObjectFrame.new_table.assert_called_once_with(
            id_, name, rows, cols, left, top, width, height)
        __spTree.append.assert_called_once_with(graphicFrame)
        _Table.assert_called_once_with(graphicFrame)
        __shapes.append.assert_called_once_with(table)
        assert_that(retval, is_(equal_to(table)))

    @patch('pptx.shapes.CT_Shape')
    @patch('pptx.shapes._Shape')
    @patch('pptx.shapes._ShapeCollection._ShapeCollection__next_shape_id',
           new_callable=PropertyMock)
    def test_add_textbox_collaboration(self, __next_shape_id, _Shape,
                                       CT_Shape):
        """_ShapeCollection.add_textbox() calls the right collaborators"""
        # constant values -------------
        id_, name = 9, 'TextBox 8'
        left, top, width, height = 111, 222, 333, 444
        # setup mockery ---------------
        sp = Mock(name='sp')
        shape = Mock('shape')
        __spTree = Mock(name='__spTree')
        shapes = test_shapes.empty_shape_collection
        shapes._ShapeCollection__spTree = __spTree
        __next_shape_id.return_value = id_
        CT_Shape.new_textbox_sp.return_value = sp
        _Shape.return_value = shape
        # exercise ---------------------
        retval = shapes.add_textbox(left, top, width, height)
        # verify -----------------------
        CT_Shape.new_textbox_sp.assert_called_once_with(
            id_, name, left, top, width, height)
        _Shape.assert_called_once_with(sp)
        __spTree.append.assert_called_once_with(sp)
        assert_that(shapes._ShapeCollection__shapes[0], is_(equal_to(shape)))
        assert_that(retval, is_(equal_to(shape)))

    def test_title_value(self):
        """_ShapeCollection.title value is ref to correct shape"""
        # exercise ---------------------
        title_shape = self.shapes.title
        # verify -----------------------
        expected = 0
        actual = self.shapes.index(title_shape)
        msg = "expected shapes[%d], got shapes[%d]" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_title_is_none_on_no_title_placeholder(self):
        """_ShapeCollection.title value is None when no title placeholder"""
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        # verify -----------------------
        assert_that(shapes.title, is_(None))

    def test_placeholders_values(self):
        """_ShapeCollection.placeholders values are correct and sorted"""
        # setup ------------------------
        expected_values = (
            ('Title 1',                    PH_TYPE_CTRTITLE,  0),
            ('Vertical Subtitle 2',        PH_TYPE_SUBTITLE,  1),
            ('Date Placeholder 7',         PH_TYPE_DT,       10),
            ('Footer Placeholder 4',       PH_TYPE_FTR,      11),
            ('Slide Number Placeholder 5', PH_TYPE_SLDNUM,   12),
            ('Table Placeholder 3',        PH_TYPE_TBL,      14))
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        placeholders = shapes.placeholders
        # verify -----------------------
        for idx, ph in enumerate(placeholders):
            values = (ph.name, ph.type, ph.idx)
            expected = expected_values[idx]
            actual = values
            msg = ("expected placeholders[%d] values %s, got %s" %
                   (idx, expected, actual))
            self.assertEqual(expected, actual, msg)

    def test__clone_layout_placeholders_shapes(self):
        """_ShapeCollection._clone_layout_placeholders clones shapes"""
        # setup ------------------------
        expected_values = (
            [2, 'Title 1',             PH_TYPE_CTRTITLE,  0],
            [3, 'Vertical Subtitle 2', PH_TYPE_SUBTITLE,  1],
            [4, 'Table Placeholder 3', PH_TYPE_TBL,      14])
        slidelayout = _SlideLayout()
        slidelayout._shapes = _sldLayout1_shapes()
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        shapes._clone_layout_placeholders(slidelayout)
        # verify -----------------------
        for idx, sp in enumerate(shapes):
            # verify is placeholder ---
            is_placeholder = sp.is_placeholder
            msg = ("expected shapes[%d].is_placeholder == True %r"
                   % (idx, sp))
            self.assertTrue(is_placeholder, msg)
            # verify values -----------
            ph = _Placeholder(sp)
            expected = expected_values[idx]
            actual = [ph.id, ph.name, ph.type, ph.idx]
            msg = ("expected placeholder[%d] values %s, got %s"
                   % (idx, expected, actual))
            self.assertEqual(expected, actual, msg)

    def test___clone_layout_placeholder_values(self):
        """_ShapeCollection.__clone_layout_placeholder() values correct"""
        # setup ------------------------
        layout_shapes = _sldLayout1_shapes()
        layout_ph_shapes = [sp for sp in layout_shapes if sp.is_placeholder]
        shapes = test_shapes.empty_shape_collection
        expected_values = (
            [2, 'Title 1',                    PH_TYPE_CTRTITLE,  0],
            [3, 'Date Placeholder 2',         PH_TYPE_DT,       10],
            [4, 'Vertical Subtitle 3',        PH_TYPE_SUBTITLE,  1],
            [5, 'Table Placeholder 4',        PH_TYPE_TBL,      14],
            [6, 'Slide Number Placeholder 5', PH_TYPE_SLDNUM,   12],
            [7, 'Footer Placeholder 6',       PH_TYPE_FTR,      11])
        # exercise ---------------------
        for idx, layout_ph_sp in enumerate(layout_ph_shapes):
            layout_ph = _Placeholder(layout_ph_sp)
            sp = shapes._ShapeCollection__clone_layout_placeholder(layout_ph)
            # verify ------------------
            ph = _Placeholder(sp)
            expected = expected_values[idx]
            actual = [ph.id, ph.name, ph.type, ph.idx]
            msg = "expected placeholder values %s, got %s" % (expected, actual)
            self.assertEqual(expected, actual, msg)

    def test___clone_layout_placeholder_xml(self):
        """_ShapeCollection.__clone_layout_placeholder() emits correct XML"""
        # setup ------------------------
        layout_shapes = _sldLayout1_shapes()
        layout_ph_shapes = [sp for sp in layout_shapes if sp.is_placeholder]
        shapes = test_shapes.empty_shape_collection
        expected_xml_tmpl = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%s" name="%s"/>\n    <'
            'p:cNvSpPr>\n      <a:spLocks noGrp="1"/>\n    </p:cNvSpPr>\n    '
            '<p:nvPr>\n      <p:ph type="%s"%s/>\n    </p:nvPr>\n  </p:nvSpPr'
            '>\n  <p:spPr/>\n%s</p:sp>\n' %
            (nsdecls('p', 'a'), '%d', '%s', '%s', '%s', '%s')
        )
        txBody_snippet = (
            '  <p:txBody>\n    <a:bodyPr/>\n    <a:lstStyle/>\n    <a:p/>\n  '
            '</p:txBody>\n')
        expected_values = [
            (2, 'Title 1', PH_TYPE_CTRTITLE, '', txBody_snippet),
            (3, 'Date Placeholder 2', PH_TYPE_DT, ' sz="half" idx="10"', ''),
            (4, 'Vertical Subtitle 3', PH_TYPE_SUBTITLE,
                ' orient="vert" idx="1"', txBody_snippet),
            (5, 'Table Placeholder 4', PH_TYPE_TBL,
                ' sz="quarter" idx="14"', ''),
            (6, 'Slide Number Placeholder 5', PH_TYPE_SLDNUM,
                ' sz="quarter" idx="12"', ''),
            (7, 'Footer Placeholder 6', PH_TYPE_FTR,
                ' sz="quarter" idx="11"', '')]
                    # verify ----------------------
        for idx, layout_ph_sp in enumerate(layout_ph_shapes):
            layout_ph = _Placeholder(layout_ph_sp)
            sp = shapes._ShapeCollection__clone_layout_placeholder(layout_ph)
            ph = _Placeholder(sp)
            expected_xml = expected_xml_tmpl % expected_values[idx]
            self.assertEqualLineByLine(expected_xml, ph._element)

    def test___next_ph_name_return_value(self):
        """
        _ShapeCollection.__next_ph_name() returns correct value

        * basename + 'Placeholder' + num, e.g. 'Table Placeholder 8'
        * numpart of name defaults to id-1, but increments until unique
        * prefix 'Vertical' if orient="vert"

        """
        cases = (
            (PH_TYPE_OBJ,   3, PH_ORIENT_HORZ, 'Content Placeholder 2'),
            (PH_TYPE_TBL,   4, PH_ORIENT_HORZ, 'Table Placeholder 4'),
            (PH_TYPE_TBL,   7, PH_ORIENT_VERT, 'Vertical Table Placeholder 6'),
            (PH_TYPE_TITLE, 2, PH_ORIENT_HORZ, 'Title 2'))
        # setup ------------------------
        shapes = _sldLayout1_shapes()
        for ph_type, id, orient, expected_name in cases:
            # exercise --------------------
            name = shapes._ShapeCollection__next_ph_name(ph_type, id, orient)
            # verify ----------------------
            expected = expected_name
            actual = name
            msg = ("expected placeholder name '%s', got '%s'" %
                   (expected, actual))
            self.assertEqual(expected, actual, msg)

    def test___next_shape_id_value(self):
        """_ShapeCollection.__next_shape_id value is correct"""
        # setup ------------------------
        shapes = _sldLayout1_shapes()
        # exercise ---------------------
        next_id = shapes._ShapeCollection__next_shape_id
        # verify -----------------------
        expected = 4
        actual = next_id
        msg = "expected %d, got %d" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class Test_Table(TestCase):
    """Test _Table"""
    def test_initial_height_divided_evenly_between_rows(self):
        """Table creation height divided evenly between rows"""
        # constant values -------------
        rows = cols = 3
        left = top = Inches(1.0)
        width = Inches(2.0)
        height = 1000
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        table = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        assert_that(table.rows[0].height, is_(equal_to(333)))
        assert_that(table.rows[1].height, is_(equal_to(333)))
        assert_that(table.rows[2].height, is_(equal_to(334)))

    def test_initial_width_divided_evenly_between_columns(self):
        """Table creation width divided evenly between columns"""
        # constant values -------------
        rows = cols = 3
        left = top = Inches(1.0)
        width = 1000
        height = Inches(2.0)
        shapes = test_shapes.empty_shape_collection
        # exercise ---------------------
        table = shapes.add_table(rows, cols, left, top, width, height)
        # verify -----------------------
        assert_that(table.columns[0].width, is_(equal_to(333)))
        assert_that(table.columns[1].width, is_(equal_to(333)))
        assert_that(table.columns[2].width, is_(equal_to(334)))

    def test_height_sum_of_row_heights(self):
        """_Table.height is sum of row heights"""
        # constant values -------------
        rows = cols = 2
        left = top = width = height = Inches(2.0)
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        tbl = shapes.add_table(rows, cols, left, top, width, height)
        tbl.rows[0].height = 100
        tbl.rows[1].height = 200
        # verify -----------------------
        sum_of_row_heights = 300
        assert_that(tbl.height, is_(equal_to(sum_of_row_heights)))

    def test_width_sum_of_col_widths(self):
        """_Table.width is sum of column widths"""
        # constant values -------------
        rows = cols = 2
        left = top = width = height = Inches(2.0)
        # setup ------------------------
        shapes = test_shapes.empty_shape_collection
        tbl = shapes.add_table(rows, cols, left, top, width, height)
        tbl.columns[0].width = 100
        tbl.columns[1].width = 200
        # verify -----------------------
        sum_of_col_widths = tbl.columns[0].width + tbl.columns[1].width
        assert_that(tbl.width, is_(equal_to(sum_of_col_widths)))


class Test_TableBooleanProperties(TestCase):
    """Test _Table"""
    def setUp(self):
        """Test fixture for _Table boolean properties"""
        shapes = test_shapes.empty_shape_collection
        self.table = shapes.add_table(2, 2, 1000, 1000, 1000, 1000)
        self.assignment_cases = (
            (True,  True),
            (False, False),
            (0,     False),
            (1,     True),
            ('',    False),
            ('foo', True)
        )

    def mockery(self, property_name, property_return_value=None):
        """
        Return property of *property_name* on self.table with return value of
        *property_return_value*.
        """
        # mock <a:tbl> element of _Table so we can mock its properties
        tbl = MagicMock()
        self.table._Table__tbl_elm = tbl
        # create a suitable mock for the property
        property_ = PropertyMock()
        if property_return_value:
            property_.return_value = property_return_value
        # and attach it the the <a:tbl> element object (class actually)
        setattr(type(tbl), property_name, property_)
        return property_

    def test_first_col_property_value(self):
        """_Table.first_col property value is calculated correctly"""
        # mockery ----------------------
        firstCol_val = True
        firstCol = self.mockery('firstCol', firstCol_val)
        # exercise ---------------------
        retval = self.table.first_col
        # verify -----------------------
        firstCol.assert_called_once_with()
        assert_that(retval, is_(equal_to(firstCol_val)))

    def test_first_row_property_value(self):
        """_Table.first_row property value is calculated correctly"""
        # mockery ----------------------
        firstRow_val = True
        firstRow = self.mockery('firstRow', firstRow_val)
        # exercise ---------------------
        retval = self.table.first_row
        # verify -----------------------
        firstRow.assert_called_once_with()
        assert_that(retval, is_(equal_to(firstRow_val)))

    def test_horz_banding_property_value(self):
        """_Table.horz_banding property value is calculated correctly"""
        # mockery ----------------------
        bandRow_val = True
        bandRow = self.mockery('bandRow', bandRow_val)
        # exercise ---------------------
        retval = self.table.horz_banding
        # verify -----------------------
        bandRow.assert_called_once_with()
        assert_that(retval, is_(equal_to(bandRow_val)))

    def test_last_col_property_value(self):
        """_Table.last_col property value is calculated correctly"""
        # mockery ----------------------
        lastCol_val = True
        lastCol = self.mockery('lastCol', lastCol_val)
        # exercise ---------------------
        retval = self.table.last_col
        # verify -----------------------
        lastCol.assert_called_once_with()
        assert_that(retval, is_(equal_to(lastCol_val)))

    def test_last_row_property_value(self):
        """_Table.last_row property value is calculated correctly"""
        # mockery ----------------------
        lastRow_val = True
        lastRow = self.mockery('lastRow', lastRow_val)
        # exercise ---------------------
        retval = self.table.last_row
        # verify -----------------------
        lastRow.assert_called_once_with()
        assert_that(retval, is_(equal_to(lastRow_val)))

    def test_vert_banding_property_value(self):
        """_Table.vert_banding property value is calculated correctly"""
        # mockery ----------------------
        bandCol_val = True
        bandCol = self.mockery('bandCol', bandCol_val)
        # exercise ---------------------
        retval = self.table.vert_banding
        # verify -----------------------
        bandCol.assert_called_once_with()
        assert_that(retval, is_(equal_to(bandCol_val)))

    def test_first_col_assignment(self):
        """Assignment to _Table.first_col sets attribute value"""
        # mockery ----------------------
        firstCol = self.mockery('firstCol')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.first_col = assigned_value
            firstCol.assert_called_once_with(called_with_value)
            firstCol.reset_mock()

    def test_first_row_assignment(self):
        """Assignment to _Table.first_row sets attribute value"""
        # mockery ----------------------
        firstRow = self.mockery('firstRow')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.first_row = assigned_value
            firstRow.assert_called_once_with(called_with_value)
            firstRow.reset_mock()

    def test_horz_banding_assignment(self):
        """Assignment to _Table.horz_banding sets attribute value"""
        # mockery ----------------------
        bandRow = self.mockery('bandRow')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.horz_banding = assigned_value
            bandRow.assert_called_once_with(called_with_value)
            bandRow.reset_mock()

    def test_last_col_assignment(self):
        """Assignment to _Table.last_col sets attribute value"""
        # mockery ----------------------
        lastCol = self.mockery('lastCol')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.last_col = assigned_value
            lastCol.assert_called_once_with(called_with_value)
            lastCol.reset_mock()

    def test_last_row_assignment(self):
        """Assignment to _Table.last_row sets attribute value"""
        # mockery ----------------------
        lastRow = self.mockery('lastRow')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.last_row = assigned_value
            lastRow.assert_called_once_with(called_with_value)
            lastRow.reset_mock()

    def test_vert_banding_assignment(self):
        """Assignment to _Table.vert_banding sets attribute value"""
        # mockery ----------------------
        bandCol = self.mockery('bandCol')
        # verify -----------------------
        for assigned_value, called_with_value in self.assignment_cases:
            self.table.vert_banding = assigned_value
            bandCol.assert_called_once_with(called_with_value)
            bandCol.reset_mock()
