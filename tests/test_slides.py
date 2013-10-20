# encoding: utf-8

"""Test suite for pptx.slides module."""

import os

from hamcrest import assert_that, is_
from mock import Mock, patch, PropertyMock

from pptx.opc_constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.oxml import oxml_parse
from pptx.presentation import _Package, Presentation
from pptx.shapes.shapetree import _ShapeCollection
from pptx.slides import (
    _BaseSlide, _Slide, _SlideCollection, _SlideLayout, _SlideMaster
)
from pptx.spec import namespaces

from pptx import packaging

from testing import TestCase


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))


thisdir = os.path.split(__file__)[0]
test_file_dir = absjoin(thisdir, 'test_files')

test_image_path = absjoin(test_file_dir, 'python-icon.jpeg')
test_pptx_path = absjoin(test_file_dir, 'test.pptx')

nsmap = namespaces('a', 'r', 'p')


def _sldLayout1():
    path = os.path.join(thisdir, 'test_files/slideLayout1.xml')
    sldLayout = oxml_parse(path).getroot()
    return sldLayout


def _sldLayout1_shapes():
    sldLayout = _sldLayout1()
    spTree = sldLayout.xpath('./p:cSld/p:spTree', namespaces=nsmap)[0]
    shapes = _ShapeCollection(spTree)
    return shapes


class Test_BaseSlide(TestCase):
    """Test _BaseSlide"""
    def setUp(self):
        self.base_slide = _BaseSlide()

    def test_name_value(self):
        """_BaseSlide.name value is correct"""
        # setup ------------------------
        self.base_slide._element = _sldLayout1()
        # exercise ---------------------
        name = self.base_slide.name
        # verify -----------------------
        expected = 'Title Slide'
        actual = name
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_shapes_size_after__load(self):
        """_BaseSlide.shapes is expected size after _load()"""
        # setup ------------------------
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        pkgpart = Mock(name='pptx.packaging.Part')
        pkgpart.partname = '/ppt/slides/slide1.xml'
        with open(path, 'r') as f:
            pkgpart.blob = f.read()
        pkgpart.relationships = []
        part_dict = {}
        self.base_slide._load(pkgpart, part_dict)
        # exercise ---------------------
        shapes = self.base_slide.shapes
        # verify -----------------------
        self.assertLength(shapes, 9)

    @patch('pptx.slides._BaseSlide._package', new_callable=PropertyMock)
    def test__add_image_collaboration(self, _package):
        """_BaseSlide._add_image() returns (image, rel) tuple"""
        # setup ------------------------
        base_slide = self.base_slide
        image = Mock(name='image')
        rel = Mock(name='rel')
        base_slide._package._images.add_image.return_value = image
        base_slide._add_relationship = Mock('_add_relationship')
        base_slide._add_relationship.return_value = rel
        file = test_image_path
        # exercise ---------------------
        retval_image, retval_rel = base_slide._add_image(file)
        # verify -----------------------
        base_slide._package._images.add_image.assert_called_once_with(file)
        base_slide._add_relationship.assert_called_once_with(RT.IMAGE, image)
        assert_that(retval_image, is_(image))
        assert_that(retval_rel, is_(rel))


class Test_Slide(TestCase):
    """Test _Slide"""
    def setUp(self):
        self.sld = _Slide()

    def test_constructor_sets_correct_content_type(self):
        """_Slide constructor sets correct content type"""
        # exercise ---------------------
        content_type = self.sld._content_type
        # verify -----------------------
        expected = CT.PML_SLIDE
        actual = content_type
        msg = "expected '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_construction_adds_slide_layout_relationship(self):
        """_Slide(slidelayout) adds relationship slide->slidelayout"""
        # setup ------------------------
        slidelayout = _SlideLayout()
        slidelayout._shapes = _sldLayout1_shapes()
        # exercise ---------------------
        slide = _Slide(slidelayout)
        # verify length ---------------
        expected = 1
        actual = len(slide._relationships)
        msg = ("expected len(slide._relationships) of %d, got %d"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)
        # verify values ---------------
        rel = slide._relationships[0]
        expected = ('rId1', RT.SLIDE_LAYOUT, slidelayout)
        actual = (rel._rId, rel._reltype, rel._target)
        msg = "expected relationship\n%s\ngot\n%s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test__element_minimal_sld_on_construction(self):
        """_Slide._element is minimal sld on construction"""
        # setup ------------------------
        path = os.path.join(thisdir, 'test_files/minimal_slide.xml')
        # exercise ---------------------
        elm = self.sld._element
        # verify -----------------------
        with open(path, 'r') as f:
            expected_xml = f.read()
        self.assertEqualLineByLine(expected_xml, elm)

    def test_slidelayout_property_none_on_construction(self):
        """_Slide.slidelayout property None on construction"""
        # verify -----------------------
        self.assertIsProperty(self.sld, 'slidelayout', None)

    def test__load_sets_slidelayout(self):
        """_Slide._load() sets slidelayout"""
        # setup ------------------------
        path = os.path.join(thisdir, 'test_files/slide1.xml')
        slidelayout = Mock(name='slideLayout')
        slidelayout.partname = '/ppt/slideLayouts/slideLayout1.xml'
        rel = Mock(name='pptx.packaging._Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT.SLIDE_LAYOUT
        rel.target = slidelayout
        pkgpart = Mock(name='pptx.packaging.Part')
        with open(path, 'rb') as f:
            pkgpart.blob = f.read()
        pkgpart.relationships = [rel]
        part_dict = {slidelayout.partname: slidelayout}
        slide = self.sld._load(pkgpart, part_dict)
        # exercise ---------------------
        retval = slide.slidelayout
        # verify -----------------------
        expected = slidelayout
        actual = retval
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test___minimal_element_xml(self):
        """_Slide.__minimal_element generates correct XML"""
        # setup ------------------------
        path = os.path.join(thisdir, 'test_files/minimal_slide.xml')
        # exercise ---------------------
        sld = self.sld._Slide__minimal_element
        # verify -----------------------
        with open(path, 'r') as f:
            expected_xml = f.read()
        self.assertEqualLineByLine(expected_xml, sld)


class Test_SlideCollection(TestCase):
    """Test _SlideCollection"""
    def setUp(self):
        prs = Presentation()
        self.slides = _SlideCollection(prs)

    def test_add_slide_returns_slide(self):
        """_SlideCollection.add_slide() returns instance of _Slide"""
        # exercise ---------------------
        retval = self.slides.add_slide(None)
        # verify -----------------------
        self.assertIsInstance(retval, _Slide)

    def test_add_slide_sets_slidelayout(self):
        """
        _SlideCollection.add_slide() sets _Slide.slidelayout

        Kind of a throw-away test, but was helpful for initial debugging.
        """
        # setup ------------------------
        slidelayout = Mock(name='slideLayout')
        slidelayout.shapes = []
        slide = self.slides.add_slide(slidelayout)
        # exercise ---------------------
        retval = slide.slidelayout
        # verify -----------------------
        expected = slidelayout
        actual = retval
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_add_slide_adds_slide_layout_relationship(self):
        """_SlideCollection.add_slide() adds relationship prs->slide"""
        # setup ------------------------
        prs = Presentation()
        slides = prs.slides
        slidelayout = _SlideLayout()
        slidelayout._shapes = []
        # exercise ---------------------
        slide = slides.add_slide(slidelayout)
        # verify length ---------------
        expected = 1
        actual = len(prs._relationships)
        msg = ("expected len(prs._relationships) of %d, got %d"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)
        # verify values ---------------
        rel = prs._relationships[0]
        expected = ('rId1', RT.SLIDE, slide)
        actual = (rel._rId, rel._reltype, rel._target)
        msg = ("expected relationship 1:, got 2:\n1: %s\n2: %s"
               % (expected, actual))
        self.assertEqual(expected, actual, msg)

    def test_add_slide_sets_partname(self):
        """_SlideCollection.add_slide() sets partname of new slide"""
        # setup ------------------------
        prs = Presentation()
        slides = prs.slides
        slidelayout = _SlideLayout()
        slidelayout._shapes = []
        # exercise ---------------------
        slide = slides.add_slide(slidelayout)
        # verify -----------------------
        expected = '/ppt/slides/slide1.xml'
        actual = slide.partname
        msg = "expected partname '%s', got '%s'" % (expected, actual)
        self.assertEqual(expected, actual, msg)


class Test_SlideLayout(TestCase):
    """Test _SlideLayout"""
    def setUp(self):
        self.slidelayout = _SlideLayout()

    def __loaded_slidelayout(self, prs_slidemaster=None):
        """
        Return _SlideLayout instance loaded using mocks. *prs_slidemaster* is
        an already-loaded model-side _SlideMaster instance (or mock, as
        appropriate to calling test).
        """
        # partname for related slideMaster
        sldmaster_partname = '/ppt/slideMasters/slideMaster1.xml'
        # path to test slideLayout XML
        slidelayout_path = absjoin(test_file_dir, 'slideLayout1.xml')
        # model-side slideMaster part
        if prs_slidemaster is None:
            prs_slidemaster = Mock(spec=_SlideMaster)
        # a part dict containing the already-loaded model-side slideMaster
        loaded_part_dict = {sldmaster_partname: prs_slidemaster}
        # a slideMaster package part for rel target
        pkg_slidemaster_part = Mock(spec=packaging.Part)
        pkg_slidemaster_part.partname = sldmaster_partname
        # a package-side relationship from slideLayout to its slideMaster
        rel = Mock(name='pptx.packaging._Relationship')
        rel.rId = 'rId1'
        rel.reltype = RT.SLIDE_MASTER
        rel.target = pkg_slidemaster_part
        # the slideLayout package part to send to _load()
        pkg_slidelayout_part = Mock(spec=packaging.Part)
        pkg_slidelayout_part.relationships = [rel]
        with open(slidelayout_path, 'rb') as f:
            pkg_slidelayout_part.blob = f.read()
        # _load and return
        slidelayout = _SlideLayout()
        return slidelayout._load(pkg_slidelayout_part, loaded_part_dict)

    def test__load_sets_slidemaster(self):
        """_SlideLayout._load() sets slidemaster"""
        # setup ------------------------
        prs_slidemaster = Mock(spec=_SlideMaster)
        # exercise ---------------------
        loaded_slidelayout = self.__loaded_slidelayout(prs_slidemaster)
        # verify -----------------------
        expected = prs_slidemaster
        actual = loaded_slidelayout.slidemaster
        msg = "expected: %s, got %s" % (expected, actual)
        self.assertEqual(expected, actual, msg)

    def test_slidemaster_is_readonly(self):
        """_SlideLayout.slidemaster is read-only"""
        # verify -----------------------
        self.assertIsReadOnly(self.slidelayout, 'slidemaster')

    def test_slidemaster_raises_on_ref_before_assigned(self):
        """_SlideLayout.slidemaster raises on referenced before assigned"""
        with self.assertRaises(AssertionError):
            self.slidelayout.slidemaster


class Test_SlideMaster(TestCase):
    """Test _SlideMaster"""
    def setUp(self):
        self.sldmaster = _SlideMaster()

    def test_slidelayouts_property_empty_on_construction(self):
        """_SlideMaster.slidelayouts property empty on construction"""
        # verify -----------------------
        self.assertIsSizedProperty(self.sldmaster, 'slidelayouts', 0)

    def test_slidelayouts_correct_length_after_open(self):
        """_SlideMaster.slidelayouts correct length after open"""
        # setup ------------------------
        pkg = _Package(test_pptx_path)
        slidemaster = pkg.presentation.slidemasters[0]
        # exercise ---------------------
        slidelayouts = slidemaster.slidelayouts
        # verify -----------------------
        self.assertLength(slidelayouts, 11)
