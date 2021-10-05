# encoding: utf-8

"""Unit-test suite for `pptx.text.text` module."""

from __future__ import unicode_literals

import pytest

from pptx.compat import is_unicode
from pptx.dml.color import ColorFormat
from pptx.dml.fill import FillFormat
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_UNDERLINE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.shapes.autoshape import Shape
from pptx.text.text import (
    Font,
    _Hyperlink,
    _Paragraph,
    _Run,
    TextFrame,
    TextFont,
    TextListStyle,
    ParagraphProperties,
)
from pptx.text.bullets import TextBullet, TextBulletColor, TextBulletSize, TextBulletTypeface
from pptx.util import Inches, Pt

from ..oxml.unitdata.text import a_p, a_t, an_hlinkClick, an_r, an_rPr, an_pPr, an_extLst, an_ext, an_hlinkClr
from ..unitutil.cxml import element, xml
from ..unitutil.mock import (
    class_mock,
    instance_mock,
    loose_mock,
    method_mock,
    property_mock,
)


class DescribeTextFrame(object):
    """Unit-test suite for `pptx.text.text.TextFrame` object."""

    def it_can_add_a_paragraph_to_itself(self, add_paragraph_fixture):
        text_frame, expected_xml = add_paragraph_fixture
        text_frame.add_paragraph()
        assert text_frame._txBody.xml == expected_xml

    def it_knows_its_autosize_setting(self, autosize_get_fixture):
        text_frame, expected_value = autosize_get_fixture
        assert text_frame.auto_size == expected_value

    def it_can_change_its_autosize_setting(self, autosize_set_fixture):
        text_frame, value, expected_xml = autosize_set_fixture
        text_frame.auto_size = value
        assert text_frame._txBody.xml == expected_xml

    @pytest.mark.parametrize(
        "txBody_cxml",
        (
            "p:txBody/(a:p,a:p,a:p)",
            'p:txBody/a:p/a:r/a:t"foo"',
            'p:txBody/a:p/(a:br,a:r/a:t"foo")',
            'p:txBody/a:p/(a:fld,a:br,a:r/a:t"foo")',
        ),
    )
    def it_can_clear_itself_of_content(self, txBody_cxml):
        text_frame = TextFrame(element(txBody_cxml), None)
        text_frame.clear()
        assert text_frame._element.xml == xml("p:txBody/a:p")

    def it_knows_its_margin_settings(self, margin_get_fixture):
        text_frame, prop_name, unit, expected_value = margin_get_fixture
        margin_value = getattr(text_frame, prop_name)
        assert getattr(margin_value, unit) == expected_value

    def it_can_change_its_margin_settings(self, margin_set_fixture):
        text_frame, prop_name, new_value, expected_xml = margin_set_fixture
        setattr(text_frame, prop_name, new_value)
        assert text_frame._txBody.xml == expected_xml

    def it_knows_its_vertical_alignment(self, anchor_get_fixture):
        text_frame, expected_value = anchor_get_fixture
        assert text_frame.vertical_anchor == expected_value

    def it_can_change_its_vertical_alignment(self, anchor_set_fixture):
        text_frame, new_value, expected_xml = anchor_set_fixture
        text_frame.vertical_anchor = new_value
        assert text_frame._element.xml == expected_xml

    def it_knows_its_word_wrap_setting(self, wrap_get_fixture):
        text_frame, expected_value = wrap_get_fixture
        assert text_frame.word_wrap == expected_value

    def it_can_change_its_word_wrap_setting(self, wrap_set_fixture):
        text_frame, new_value, expected_xml = wrap_set_fixture
        text_frame.word_wrap = new_value
        assert text_frame._element.xml == expected_xml

    def it_provides_access_to_its_paragraphs(self, paragraphs_fixture):
        text_frame, ps = paragraphs_fixture
        paragraphs = text_frame.paragraphs
        assert len(paragraphs) == len(ps)
        for idx, paragraph in enumerate(paragraphs):
            assert isinstance(paragraph, _Paragraph)
            assert paragraph._element is ps[idx]

    def it_raises_on_attempt_to_set_margin_to_non_int(self):
        text_frame = TextFrame(element("p:txBody/a:bodyPr"), None)
        with pytest.raises(TypeError):
            text_frame.margin_bottom = "0.1"

    def it_knows_the_part_it_belongs_to(self, text_frame_with_parent_):
        text_frame, parent_ = text_frame_with_parent_
        part = text_frame.part
        assert part is parent_.part

    def it_knows_what_text_it_contains(
        self, request, text_get_fixture, paragraphs_prop_
    ):
        paragraph_texts, expected_value = text_get_fixture
        paragraphs_prop_.return_value = tuple(
            instance_mock(request, _Paragraph, text=text) for text in paragraph_texts
        )
        text_frame = TextFrame(None, None)

        text = text_frame.text

        assert text == expected_value

    def it_can_replace_the_text_it_contains(self, text_set_fixture):
        txBody, text, expected_xml = text_set_fixture
        text_frame = TextFrame(txBody, None)

        text_frame.text = text

        assert text_frame._element.xml == expected_xml

    def it_can_resize_its_text_to_best_fit(self, request, text_prop_):
        family, max_size, bold, italic, font_file, font_size = (
            "Family",
            42,
            "bold",
            "italic",
            "font_file",
            21,
        )
        text_prop_.return_value = "some text"
        _best_fit_font_size_ = method_mock(
            request, TextFrame, "_best_fit_font_size", return_value=font_size
        )
        _apply_fit_ = method_mock(request, TextFrame, "_apply_fit")
        text_frame = TextFrame(None, None)

        text_frame.fit_text(family, max_size, bold, italic, font_file)

        _best_fit_font_size_.assert_called_once_with(
            text_frame, family, max_size, bold, italic, font_file
        )
        _apply_fit_.assert_called_once_with(text_frame, family, font_size, bold, italic)

    def it_calculates_its_best_fit_font_size_to_help_fit_text(self, size_font_fixture):
        text_frame, family, max_size, bold, italic = size_font_fixture[:5]
        FontFiles_, TextFitter_, text, extents = size_font_fixture[5:9]
        font_file_, font_size_ = size_font_fixture[9:]

        font_size = text_frame._best_fit_font_size(family, max_size, bold, italic, None)

        FontFiles_.find.assert_called_once_with(family, bold, italic)
        TextFitter_.best_fit_font_size.assert_called_once_with(
            text, extents, max_size, font_file_
        )
        assert font_size is font_size_

    def it_calculates_its_effective_size_to_help_fit_text(self):
        sp_cxml = (
            "p:sp/(p:spPr/a:xfrm/(a:off{x=914400,y=914400},a:ext{cx=914400,c"
            "y=914400}),p:txBody/(a:bodyPr,a:p))"
        )
        text_frame = Shape(element(sp_cxml), None).text_frame
        assert text_frame._extents == (731520, 822960)

    def it_applies_fit_to_help_fit_text(self, request):
        family, font_size, bold, italic = "Family", 42, True, False
        _set_font_ = method_mock(request, TextFrame, "_set_font")
        text_frame = TextFrame(element("p:txBody/a:bodyPr"), None)

        text_frame._apply_fit(family, font_size, bold, italic)

        assert text_frame.auto_size is MSO_AUTO_SIZE.NONE
        assert text_frame.word_wrap is True
        _set_font_.assert_called_once_with(text_frame, family, font_size, bold, italic)

    def it_sets_its_font_to_help_fit_text(self, set_font_fixture):
        text_frame, family, size, bold, italic, expected_xml = set_font_fixture
        text_frame._set_font(family, size, bold, italic)
        assert text_frame._element.xml == expected_xml

    def it_provides_access_to_its_text_list_style(self, text_frame):
        assert isinstance(text_frame.list_style, TextListStyle)

    # fixtures ---------------------------------------------

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", "p:txBody/(a:bodyPr,a:p)"),
            ("p:txBody/(a:bodyPr,a:p)", "p:txBody/(a:bodyPr,a:p,a:p)"),
        ]
    )
    def add_paragraph_fixture(self, request):
        txBody_cxml, expected_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_cxml)
        return text_frame, expected_xml

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", None),
            ("p:txBody/a:bodyPr{anchor=t}", MSO_ANCHOR.TOP),
            ("p:txBody/a:bodyPr{anchor=b}", MSO_ANCHOR.BOTTOM),
        ]
    )
    def anchor_get_fixture(self, request):
        txBody_cxml, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        return text_frame, expected_value

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", MSO_ANCHOR.TOP, "p:txBody/a:bodyPr{anchor=t}"),
            (
                "p:txBody/a:bodyPr{anchor=t}",
                MSO_ANCHOR.MIDDLE,
                "p:txBody/a:bodyPr{anchor=ctr}",
            ),
            (
                "p:txBody/a:bodyPr{anchor=ctr}",
                MSO_ANCHOR.BOTTOM,
                "p:txBody/a:bodyPr{anchor=b}",
            ),
            ("p:txBody/a:bodyPr{anchor=b}", None, "p:txBody/a:bodyPr"),
        ]
    )
    def anchor_set_fixture(self, request):
        txBody_cxml, new_value, expected_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_cxml)
        return text_frame, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", None),
            ("p:txBody/a:bodyPr/a:noAutofit", MSO_AUTO_SIZE.NONE),
            ("p:txBody/a:bodyPr/a:spAutoFit", MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT),
            ("p:txBody/a:bodyPr/a:normAutofit", MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE),
        ]
    )
    def autosize_get_fixture(self, request):
        txBody_cxml, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        return text_frame, expected_value

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", MSO_AUTO_SIZE.NONE, "p:txBody/a:bodyPr/a:noAutofit"),
            (
                "p:txBody/a:bodyPr/a:noAutofit",
                MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
                "p:txBody/a:bodyPr/a:spAutoFit",
            ),
            (
                "p:txBody/a:bodyPr/a:spAutoFit",
                MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                "p:txBody/a:bodyPr/a:normAutofit",
            ),
            ("p:txBody/a:bodyPr/a:normAutofit", None, "p:txBody/a:bodyPr"),
        ]
    )
    def autosize_set_fixture(self, request):
        txBody_cxml, value, expected_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_cxml)
        return text_frame, value, expected_xml

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", "left", "emu", Inches(0.1)),
            ("p:txBody/a:bodyPr", "top", "emu", Inches(0.05)),
            ("p:txBody/a:bodyPr", "right", "emu", Inches(0.1)),
            ("p:txBody/a:bodyPr", "bottom", "emu", Inches(0.05)),
            ("p:txBody/a:bodyPr{lIns=9144}", "left", "cm", 0.0254),
            ("p:txBody/a:bodyPr{tIns=18288}", "top", "mm", 0.508),
            ("p:txBody/a:bodyPr{rIns=76200}", "right", "pt", 6.0),
            ("p:txBody/a:bodyPr{bIns=36576}", "bottom", "inches", 0.04),
        ]
    )
    def margin_get_fixture(self, request):
        txBody_cxml, side, unit, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        prop_name = "margin_%s" % side
        return text_frame, prop_name, unit, expected_value

    @pytest.fixture(
        params=[
            (
                "p:txBody/a:bodyPr",
                "left",
                Inches(0.11),
                "p:txBody/a:bodyPr{lIns=100584}",
            ),
            (
                "p:txBody/a:bodyPr{tIns=1234}",
                "top",
                Inches(0.12),
                "p:txBody/a:bodyPr{tIns=109728}",
            ),
            (
                "p:txBody/a:bodyPr{rIns=2345}",
                "right",
                Inches(0.13),
                "p:txBody/a:bodyPr{rIns=118872}",
            ),
            (
                "p:txBody/a:bodyPr{bIns=3456}",
                "bottom",
                Inches(0.14),
                "p:txBody/a:bodyPr{bIns=128016}",
            ),
            ("p:txBody/a:bodyPr", "left", Inches(0.1), "p:txBody/a:bodyPr"),
            ("p:txBody/a:bodyPr", "top", Inches(0.05), "p:txBody/a:bodyPr"),
            ("p:txBody/a:bodyPr", "right", Inches(0.1), "p:txBody/a:bodyPr"),
            ("p:txBody/a:bodyPr", "bottom", Inches(0.05), "p:txBody/a:bodyPr"),
        ]
    )
    def margin_set_fixture(self, request):
        txBody_cxml, side, new_value, expected_txBody_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        prop_name = "margin_%s" % side
        expected_xml = xml(expected_txBody_cxml)
        return text_frame, prop_name, new_value, expected_xml

    @pytest.fixture(params=["p:txBody", "p:txBody/a:p", "p:txBody/(a:p,a:p)"])
    def paragraphs_fixture(self, request):
        txBody_cxml = request.param
        txBody = element(txBody_cxml)
        text_frame = TextFrame(txBody, None)
        ps = txBody.xpath(".//a:p")
        return text_frame, ps

    @pytest.fixture(
        params=[
            (
                "p:txBody/(a:bodyPr,a:p/a:r)",
                True,
                False,
                "p:txBody/(a:bodyPr,a:p/(a:r/a:rPr{sz=600,b=1,i=0}/a:latin{typeface"
                "=F},a:endParaRPr{sz=600,b=1,i=0}/a:latin{typeface=F}))",
            ),
            (
                "p:txBody/a:p/a:br",
                True,
                False,
                "p:txBody/a:p/(a:br/a:rPr{sz=600,b=1,i=0}/a:latin{typeface=F},a:end"
                "ParaRPr{sz=600,b=1,i=0}/a:latin{typeface=F})",
            ),
            (
                "p:txBody/a:p/a:fld",
                True,
                False,
                "p:txBody/a:p/(a:fld/a:rPr{sz=600,b=1,i=0}/a:latin{typeface=F},a:en"
                "dParaRPr{sz=600,b=1,i=0}/a:latin{typeface=F})",
            ),
        ]
    )
    def set_font_fixture(self, request):
        txBody_cxml, bold, italic, expected_cxml = request.param
        family, size = "F", 6
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_cxml)
        return text_frame, family, size, bold, italic, expected_xml

    @pytest.fixture
    def size_font_fixture(self, FontFiles_, TextFitter_, text_prop_, _extents_prop_):
        text_frame = TextFrame(None, None)
        family, max_size, bold, italic = "Family", 42, True, False
        text, extents, font_size, font_file = "text", (111, 222), 21, "f.ttf"
        text_prop_.return_value = text
        _extents_prop_.return_value = extents
        FontFiles_.find.return_value = font_file
        TextFitter_.best_fit_font_size.return_value = font_size
        return (
            text_frame,
            family,
            max_size,
            bold,
            italic,
            FontFiles_,
            TextFitter_,
            text,
            extents,
            font_file,
            font_size,
        )

    @pytest.fixture(
        params=[(["foobar"], "foobar"), (["foo", "bar", "baz"], "foo\nbar\nbaz")]
    )
    def text_get_fixture(self, request):
        paragraph_texts, expected_value = request.param
        return paragraph_texts, expected_value

    @pytest.fixture(
        params=[
            # ---empty to something---
            ("p:txBody/a:p", "foobar", 'p:txBody/a:p/a:r/a:t"foobar"'),
            # ---something to something else---
            ('p:txBody/a:p/a:r/a:t"foobar"', "barfoo", 'p:txBody/a:p/a:r/a:t"barfoo"'),
            # ---single paragraph to multiple---
            (
                'p:txBody/a:p/a:r/a:t"barfoo"',
                "foo\nbar",
                'p:txBody/(a:p/a:r/a:t"foo",a:p/a:r/a:t"bar")',
            ),
            # ---multiple paragraphs to single---
            (
                'p:txBody/(a:p/a:r/a:t"foo",a:p/a:r/a:t"bar")',
                "barfoo",
                'p:txBody/a:p/a:r/a:t"barfoo"',
            ),
            # ---something to empty---
            ('p:txBody/a:p/a:r/a:t"foobar"', "", "p:txBody/a:p"),
            # ---vertical-tab becomes line-break---
            ("p:txBody/a:p", "a\vb", 'p:txBody/a:p/(a:r/a:t"a",a:br,a:r/a:t"b")'),
        ]
    )
    def text_set_fixture(self, request):
        txBody_cxml, text, expected_cxml = request.param
        txBody = element(txBody_cxml)
        expected_xml = xml(expected_cxml)
        return txBody, text, expected_xml

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", None),
            ("p:txBody/a:bodyPr{wrap=square}", True),
            ("p:txBody/a:bodyPr{wrap=none}", False),
        ]
    )
    def wrap_get_fixture(self, request):
        txBody_cxml, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        return text_frame, expected_value

    @pytest.fixture(
        params=[
            ("p:txBody/a:bodyPr", True, "p:txBody/a:bodyPr{wrap=square}"),
            ("p:txBody/a:bodyPr{wrap=square}", False, "p:txBody/a:bodyPr{wrap=none}"),
            ("p:txBody/a:bodyPr{wrap=none}", None, "p:txBody/a:bodyPr"),
        ]
    )
    def wrap_set_fixture(self, request):
        txBody_cxml, new_value, expected_txBody_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_txBody_cxml)
        return text_frame, new_value, expected_xml

    # fixture components -----------------------------------

    @pytest.fixture
    def _extents_prop_(self, request):
        return property_mock(request, TextFrame, "_extents")

    @pytest.fixture
    def FontFiles_(self, request):
        return class_mock(request, "pptx.text.text.FontFiles")

    @pytest.fixture
    def paragraphs_prop_(self, request):
        return property_mock(request, TextFrame, "paragraphs")

    @pytest.fixture
    def TextFitter_(self, request):
        return class_mock(request, "pptx.text.text.TextFitter")

    @pytest.fixture
    def text_frame_with_parent_(self, request):
        parent_ = loose_mock(request, name="parent_")
        text_frame = TextFrame(None, parent_)
        return text_frame, parent_

    @pytest.fixture
    def text_prop_(self, request):
        return property_mock(request, TextFrame, "text")

    @pytest.fixture
    def text_frame(self, request):
        parent_ = loose_mock(request, name="parent_")
        return TextFrame(element("p:txBody"), parent_)

class DescribeFont(object):
    """Unit-test suite for `pptx.text.text.Font` object."""

    def it_knows_its_bold_setting(self, bold_get_fixture):
        font, expected_value = bold_get_fixture
        assert font.bold == expected_value

    def it_can_change_its_bold_setting(self, bold_set_fixture):
        font, new_value, expected_xml = bold_set_fixture
        font.bold = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_italic_setting(self, italic_get_fixture):
        font, expected_value = italic_get_fixture
        assert font.italic == expected_value

    def it_can_change_its_italic_setting(self, italic_set_fixture):
        font, new_value, expected_xml = italic_set_fixture
        font.italic = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_language_id(self, language_id_get_fixture):
        font, expected_value = language_id_get_fixture
        assert font.language_id == expected_value

    def it_can_change_its_language_id_setting(self, language_id_set_fixture):
        font, new_value, expected_xml = language_id_set_fixture
        font.language_id = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_underline_setting(self, underline_get_fixture):
        font, expected_value = underline_get_fixture
        assert font.underline is expected_value, "got %s" % font.underline

    def it_can_change_its_underline_setting(self, underline_set_fixture):
        font, new_value, expected_xml = underline_set_fixture
        font.underline = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_strikethrough_setting(self, strikethrough_get_fixture):
        font, expected_value = strikethrough_get_fixture
        assert font.strikethrough == expected_value

    def it_can_change_its_strikethrough_setting(self, strikethrough_set_fixture):
        font, new_value, expected_xml = strikethrough_set_fixture
        font.strikethrough = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_baseline(self, baseline_get_fixture):
        font, expected_value = baseline_get_fixture
        assert font.baseline == expected_value

    def it_can_change_its_baseline(self, baseline_set_fixture):
        font, new_value, expected_xml = baseline_set_fixture
        font.baseline = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_size(self, size_get_fixture):
        font, expected_value = size_get_fixture
        assert font.size == expected_value

    def it_can_change_its_size(self, size_set_fixture):
        font, new_value, expected_xml = size_set_fixture
        font.size = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_latin_typeface(self, name_get_fixture):
        font, expected_value = name_get_fixture
        assert font.name == expected_value

    def it_can_change_its_latin_typeface(self, name_set_fixture):
        font, new_value, expected_xml = name_set_fixture
        font.name = new_value
        assert font._element.xml == expected_xml

    def it_provides_access_to_its_color(self, font):
        assert isinstance(font.color, ColorFormat)

    def it_provides_access_to_its_fill(self, font):
        assert isinstance(font.fill, FillFormat)

    # fixtures ---------------------------------------------

    @pytest.fixture(
        params=[("a:rPr", None), ("a:rPr{b=0}", False), ("a:rPr{b=1}", True)]
    )
    def bold_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", True, "a:rPr{b=1}"),
            ("a:rPr{b=1}", False, "a:rPr{b=0}"),
            ("a:rPr{b=0}", None, "a:rPr"),
        ]
    )
    def bold_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(
        params=[("a:rPr", None), ("a:rPr{i=0}", False), ("a:rPr{i=1}", True)]
    )
    def italic_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", True, "a:rPr{i=1}"),
            ("a:rPr{i=1}", False, "a:rPr{i=0}"),
            ("a:rPr{i=0}", None, "a:rPr"),
        ]
    )
    def italic_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:rPr", MSO_LANGUAGE_ID.NONE),
            ("a:rPr{lang=pl-PL}", MSO_LANGUAGE_ID.POLISH),
            ("a:rPr{lang=de-AT}", MSO_LANGUAGE_ID.GERMAN_AUSTRIA),
            ("a:rPr{lang=fr-FR}", MSO_LANGUAGE_ID.FRENCH),
        ]
    )
    def language_id_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", MSO_LANGUAGE_ID.ZULU, "a:rPr{lang=zu-ZA}"),
            ("a:rPr{lang=zu-ZA}", MSO_LANGUAGE_ID.URDU, "a:rPr{lang=ur-PK}"),
            ("a:rPr{lang=ur-PK}", MSO_LANGUAGE_ID.NONE, "a:rPr"),
            ("a:rPr{lang=ur-PK}", None, "a:rPr"),
            ("a:rPr", MSO_LANGUAGE_ID.NONE, "a:rPr"),
            ("a:rPr", None, "a:rPr"),
        ]
    )
    def language_id_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(
        params=[("a:rPr", None), ("a:rPr/a:latin{typeface=Foobar}", "Foobar")]
    )
    def name_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", "Foobar", "a:rPr/a:latin{typeface=Foobar}"),
            (
                "a:rPr/a:latin{typeface=Foobar}",
                "Barfoo",
                "a:rPr/a:latin{typeface=Barfoo}",
            ),
            ("a:rPr/a:latin{typeface=Barfoo}", None, "a:rPr"),
        ]
    )
    def name_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[("a:rPr", None), ("a:rPr{baseline=30000}", .3)])
    def baseline_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[("a:rPr", -.25, "a:rPr{baseline=-25000}"), ("a:rPr{baseline=2500}", None, "a:rPr")]
    )
    def baseline_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[("a:rPr", None), ("a:rPr{sz=2400}", 304800)])
    def size_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[("a:rPr", Pt(24), "a:rPr{sz=2400}"), ("a:rPr{sz=2400}", None, "a:rPr")]
    )
    def size_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:rPr", None),
            ("a:rPr{u=none}", False),
            ("a:rPr{u=sng}", True),
            ("a:rPr{u=dbl}", MSO_UNDERLINE.DOUBLE_LINE),
        ]
    )
    def underline_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", True, "a:rPr{u=sng}"),
            ("a:rPr{u=sng}", False, "a:rPr{u=none}"),
            ("a:rPr{u=none}", MSO_UNDERLINE.WAVY_LINE, "a:rPr{u=wavy}"),
            ("a:rPr{u=wavy}", MSO_UNDERLINE.NONE, "a:rPr{u=none}"),
            ("a:rPr{u=wavy}", None, "a:rPr"),
        ]
    )
    def underline_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:rPr", None),
            ("a:rPr{strike=noStrike}", False),
            ("a:rPr{strike=sngStrike}", True),
            ("a:rPr{strike=dblStrike}", 'dblStrike'),
        ]
    )
    def strikethrough_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(
        params=[
            ("a:rPr", True, "a:rPr{strike=sngStrike}"),
            ("a:rPr{strike=sngStrike}", False, "a:rPr{strike=noStrike}"),
            ("a:rPr{strike=noStrike}", 'dblStrike', "a:rPr{strike=dblStrike}"),
            ("a:rPr{strike=dblStrike}", 'sngStrike', "a:rPr{strike=sngStrike}"),
            ("a:rPr{strike=sngStrike}", None, "a:rPr"),
        ]
    )
    def strikethrough_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def font(self):
        return Font(element("a:rPr"))


class Describe_Hyperlink(object):
    """Unit-test suite for `pptx.text.text._Hyperlink` object."""

    def it_knows_the_target_url_of_the_hyperlink(self, hlink_with_url_):
        hlink, rId, url = hlink_with_url_
        assert hlink.address == url
        hlink.part.target_ref.assert_called_once_with(rId)

    def it_has_None_for_address_when_no_hyperlink_is_present(self, hlink):
        assert hlink.address is None

    def it_can_set_the_target_url(self, hlink, rPr_with_hlinkClick_xml, url):
        hlink.address = url
        # verify -----------------------
        hlink.part.relate_to.assert_called_once_with(
            url, RT.HYPERLINK, is_external=True
        )
        assert hlink._rPr.xml == rPr_with_hlinkClick_xml
        assert hlink.address == url

    def it_can_remove_the_hyperlink(self, remove_hlink_fixture_):
        hlink, rPr_xml, rId = remove_hlink_fixture_
        hlink.address = None
        assert hlink._rPr.xml == rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId)

    def it_should_remove_the_hyperlink_when_url_set_to_empty_string(
        self, remove_hlink_fixture_
    ):
        hlink, rPr_xml, rId = remove_hlink_fixture_
        hlink.address = ""
        assert hlink._rPr.xml == rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId)

    def it_can_change_the_target_url(self, change_hlink_fixture_):
        # fixture ----------------------
        hlink, rId_existing, new_url, new_rPr_xml = change_hlink_fixture_
        # exercise ---------------------
        hlink.address = new_url
        # verify -----------------------
        assert hlink._rPr.xml == new_rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId_existing)
        hlink.part.relate_to.assert_called_once_with(
            new_url, RT.HYPERLINK, is_external=True
        )

    def it_can_add_hyperlink_color_ext(self, add_hyperlink_color_fixture):
        hyperlink, expected_xml = add_hyperlink_color_fixture
        hyperlink.add_hyperlink_color()
        assert hyperlink._rPr.xml == expected_xml


    # fixtures ---------------------------------------------

    @pytest.fixture
    def change_hlink_fixture_(
        self, request, hlink_with_hlinkClick, rId, rId_2, part_, url_2
    ):
        hlinkClick_bldr = an_hlinkClick().with_rId(rId_2)
        new_rPr_xml = an_rPr().with_nsdecls("a", "r").with_child(hlinkClick_bldr).xml()
        part_.relate_to.return_value = rId_2
        property_mock(request, _Hyperlink, "part", return_value=part_)
        return hlink_with_hlinkClick, rId, url_2, new_rPr_xml

    @pytest.fixture
    def hlink(self, request, part_):
        rPr = an_rPr().with_nsdecls("a", "r").element
        hlink = _Hyperlink(rPr, None)
        property_mock(request, _Hyperlink, "part", return_value=part_)
        return hlink

    @pytest.fixture
    def hlink_with_hlinkClick(self, request, rPr_with_hlinkClick_bldr):
        rPr = rPr_with_hlinkClick_bldr.element
        return _Hyperlink(rPr, None)

    @pytest.fixture
    def hlink_with_url_(self, request, part_, hlink_with_hlinkClick, rId, url):
        property_mock(request, _Hyperlink, "part", return_value=part_)
        return hlink_with_hlinkClick, rId, url

    @pytest.fixture
    def part_(self, request, url, rId):
        """
        Mock Part instance suitable for patching into _Hyperlink.part
        property. It returns url for target_ref() and rId for relate_to().
        """
        part_ = instance_mock(request, Part)
        part_.target_ref.return_value = url
        part_.relate_to.return_value = rId
        return part_

    @pytest.fixture
    def rId(self):
        return "rId2"

    @pytest.fixture
    def rId_2(self):
        return "rId6"

    @pytest.fixture
    def remove_hlink_fixture_(self, request, hlink_with_hlinkClick, rPr_xml, rId):
        property_mock(request, _Hyperlink, "part")
        return hlink_with_hlinkClick, rPr_xml, rId

    @pytest.fixture
    def rPr_with_hlinkClick_bldr(self, rId):
        hlinkClick_bldr = an_hlinkClick().with_rId(rId)
        rPr_bldr = an_rPr().with_nsdecls("a", "r").with_child(hlinkClick_bldr)
        return rPr_bldr

    @pytest.fixture
    def rPr_with_hlinkClick_xml(self, rPr_with_hlinkClick_bldr):
        return rPr_with_hlinkClick_bldr.xml()

    @pytest.fixture
    def rPr_xml(self):
        return an_rPr().with_nsdecls("a", "r").xml()

    @pytest.fixture
    def url(self):
        return "https://github.com/scanny/python-pptx"

    @pytest.fixture
    def url_2(self):
        return "https://pypi.python.org/pypi/python-pptx"

    @pytest.fixture
    def add_hyperlink_color_fixture(
        self, request, hlink_with_hlinkClick, rId, part_, url_2
    ):
        hlinkClick_bldr = an_hlinkClick().with_rId(rId)
        extList_bldr = an_extLst()
        ext_bldr = an_ext()
        uri = "{A12FA001-AC4F-418D-AE19-62706E023703}"
        
        new_rPr_with_extList = an_rPr().with_nsdecls("a", "r") \
                                .with_child(hlinkClick_bldr \
                                .with_child(extList_bldr \
                                .with_child(an_ext().with_uri(uri) \
                                .with_child(an_hlinkClr().with_nsdecls("ahyp").with_val("tx")))))
        

        part_.relate_to.return_value = rId
        property_mock(request, _Hyperlink, "part", return_value=part_)
        return hlink_with_hlinkClick, new_rPr_with_extList.xml()

class Describe_Paragraph(object):
    """Unit test suite for pptx.text.text._Paragraph object."""

    def it_can_add_a_line_break(self, line_break_fixture):
        paragraph, expected_xml = line_break_fixture
        paragraph.add_line_break()
        assert paragraph._p.xml == expected_xml

    def it_can_add_a_run(self, paragraph, p_with_r_xml):
        run = paragraph.add_run()
        assert paragraph._p.xml == p_with_r_xml
        assert isinstance(run, _Run)

    def it_knows_its_horizontal_alignment(self, alignment_get_fixture):
        paragraph, expected_value = alignment_get_fixture
        assert paragraph.alignment == expected_value

    def it_can_change_its_horizontal_alignment(self, alignment_set_fixture):
        paragraph, new_value, expected_xml = alignment_set_fixture
        paragraph.alignment = new_value
        assert paragraph._element.xml == expected_xml

    def it_can_clear_itself_of_content(self, clear_fixture):
        paragraph, expected_xml = clear_fixture
        paragraph.clear()
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_the_default_paragraph_font(self, paragraph, Font_):
        font = paragraph.font
        Font_.assert_called_once_with(paragraph._defRPr)
        assert font == Font_.return_value

    def it_knows_its_indentation_level(self, level_get_fixture):
        paragraph, expected_value = level_get_fixture
        assert paragraph.level == expected_value

    def it_can_change_its_indentation_level(self, level_set_fixture):
        paragraph, new_value, expected_xml = level_set_fixture
        paragraph.level = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_left_margin(self, left_margin_get_fixture):
        paragraph, expected_value = left_margin_get_fixture
        assert paragraph.margin_left == expected_value

    def it_can_change_its_left_margin(self, left_margin_set_fixture):
        paragraph, new_value, expected_xml = left_margin_set_fixture
        paragraph.margin_left = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_right_margin(self, right_margin_get_fixture):
        paragraph, expected_value = right_margin_get_fixture
        assert paragraph.margin_right == expected_value

    def it_can_change_its_right_margin(self, right_margin_set_fixture):
        paragraph, new_value, expected_xml = right_margin_set_fixture
        paragraph.margin_right = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_indent(self, indent_get_fixture):
        paragraph, expected_value = indent_get_fixture
        assert paragraph.indent == expected_value

    def it_can_change_its_indent(self, indent_set_fixture):
        paragraph, new_value, expected_xml = indent_set_fixture
        paragraph.indent = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_line_spacing(self, spacing_get_fixture):
        paragraph, expected_value = spacing_get_fixture
        assert paragraph.line_spacing == expected_value

    def it_can_change_its_line_spacing(self, spacing_set_fixture):
        paragraph, new_value, expected_xml = spacing_set_fixture
        paragraph.line_spacing = new_value
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_its_runs(self, runs_fixture):
        paragraph, expected_text = runs_fixture
        runs = paragraph.runs
        assert tuple(r.text for r in runs) == expected_text
        for r in runs:
            assert isinstance(r, _Run)
            assert r._parent == paragraph

    def it_knows_its_space_after(self, after_get_fixture):
        paragraph, expected_value = after_get_fixture
        assert paragraph.space_after == expected_value

    def it_can_change_its_space_after(self, after_set_fixture):
        paragraph, new_value, expected_xml = after_set_fixture
        paragraph.space_after = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_space_before(self, before_get_fixture):
        paragraph, expected_value = before_get_fixture
        assert paragraph.space_before == expected_value

    def it_can_change_its_space_before(self, before_set_fixture):
        paragraph, new_value, expected_xml = before_set_fixture
        paragraph.space_before = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_what_text_it_contains(self, text_get_fixture):
        p, expected_value = text_get_fixture
        paragraph = _Paragraph(p, None)

        text = paragraph.text

        assert text == expected_value
        assert is_unicode(text)

    def it_can_change_its_text(self, text_set_fixture):
        p, value, expected_xml = text_set_fixture
        paragraph = _Paragraph(p, None)

        paragraph.text = value

        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_its_bullet_text(self, paragraph):
        assert isinstance(paragraph.bullet_text, TextBullet)

    def it_provides_access_to_its_bullet_color(self, paragraph):
        assert isinstance(paragraph.bullet_color, TextBulletColor)

    def it_provides_access_to_its_bullet_size(self, paragraph):
        assert isinstance(paragraph.bullet_size, TextBulletSize)

    def it_provides_access_to_its_bullet_font(self, paragraph):
        assert isinstance(paragraph.bullet_font, TextBulletTypeface)

    # fixtures ---------------------------------------------

    @pytest.fixture(
        params=[
            ("a:p", None),
            ("a:p/a:pPr", None),
            ("a:p/a:pPr/a:spcAft/a:spcPct{val=150000}", None),
            ("a:p/a:pPr/a:spcAft/a:spcPts{val=600}", 76200),
        ]
    )
    def after_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", Pt(8.333), "a:p/a:pPr/a:spcAft/a:spcPts{val=833}"),
            (
                "a:p/a:pPr/a:spcAft/a:spcPts{val=600}",
                Pt(42),
                "a:p/a:pPr/a:spcAft/a:spcPts{val=4200}",
            ),
            (
                "a:p/a:pPr/a:spcAft/a:spcPct{val=150000}",
                Pt(24),
                "a:p/a:pPr/a:spcAft/a:spcPts{val=2400}",
            ),
            ("a:p/a:pPr/a:spcAft/a:spcPts{val=600}", None, "a:p/a:pPr"),
        ]
    )
    def after_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:p", None),
            ("a:p/a:pPr{algn=ctr}", PP_ALIGN.CENTER),
            ("a:p/a:pPr{algn=r}", PP_ALIGN.RIGHT),
        ]
    )
    def alignment_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", PP_ALIGN.LEFT, "a:p/a:pPr{algn=l}"),
            ("a:p/a:pPr{algn=l}", PP_ALIGN.JUSTIFY, "a:p/a:pPr{algn=just}"),
            ("a:p/a:pPr{algn=just}", None, "a:p/a:pPr"),
        ]
    )
    def alignment_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:p", None),
            ("a:p/a:pPr", None),
            ("a:p/a:pPr/a:spcBef/a:spcPct{val=150000}", None),
            ("a:p/a:pPr/a:spcBef/a:spcPts{val=600}", 76200),
        ]
    )
    def before_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", Pt(8.333), "a:p/a:pPr/a:spcBef/a:spcPts{val=833}"),
            (
                "a:p/a:pPr/a:spcBef/a:spcPts{val=600}",
                Pt(42),
                "a:p/a:pPr/a:spcBef/a:spcPts{val=4200}",
            ),
            (
                "a:p/a:pPr/a:spcBef/a:spcPct{val=150000}",
                Pt(24),
                "a:p/a:pPr/a:spcBef/a:spcPts{val=2400}",
            ),
            ("a:p/a:pPr/a:spcBef/a:spcPts{val=600}", None, "a:p/a:pPr"),
        ]
    )
    def before_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ('a:p/a:r/a:t"foo"', "a:p"),
            ('a:p/(a:br,a:r/a:t"foo")', "a:p"),
            ('a:p/(a:fld,a:br,a:r/a:t"foo")', "a:p"),
        ]
    )
    def clear_fixture(self, request):
        p_cxml, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, expected_xml

    @pytest.fixture(params=[("a:p", 0), ("a:p/a:pPr{lvl=2}", 2)])
    def level_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", 1, "a:p/a:pPr{lvl=1}"),
            ("a:p/a:pPr{lvl=1}", 2, "a:p/a:pPr{lvl=2}"),
            ("a:p/a:pPr{lvl=2}", 0, "a:p/a:pPr"),
        ]
    )
    def level_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:p", 0), ("a:p/a:pPr{marL=10000}", 10000)])
    def left_margin_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", 1, "a:p/a:pPr{marL=1}"),
            ("a:p/a:pPr{marL=1}", 2, "a:p/a:pPr{marL=2}"),
            ("a:p/a:pPr{marL=2}", 0, "a:p/a:pPr"),
        ]
    )
    def left_margin_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:p", 0), ("a:p/a:pPr{indent=10000}", 10000)])
    def indent_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", 1, "a:p/a:pPr{indent=1}"),
            ("a:p/a:pPr{indent=1}", 2, "a:p/a:pPr{indent=2}"),
            ("a:p/a:pPr{indent=2}", 0, "a:p/a:pPr"),
        ]
    )
    def indent_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:p", 0), ("a:p/a:pPr{marR=10000}", 10000)])
    def right_margin_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", 1, "a:p/a:pPr{marR=1}"),
            ("a:p/a:pPr{marR=1}", 2, "a:p/a:pPr{marR=2}"),
            ("a:p/a:pPr{marR=2}", 0, "a:p/a:pPr"),
        ]
    )
    def right_margin_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:p", "a:p/a:br"),
            ("a:p/a:r", "a:p/(a:r,a:br)"),
            ("a:p/a:br", "a:p/(a:br,a:br)"),
        ]
    )
    def line_break_fixture(self, request):
        cxml, expected_cxml = request.param
        paragraph = _Paragraph(element(cxml), None)
        expected_xml = xml(expected_cxml)
        return paragraph, expected_xml

    @pytest.fixture
    def runs_fixture(self):
        p_cxml = 'a:p/(a:r/a:t"Foo",a:r/a:t"Bar",a:r/a:t"Baz")'
        paragraph = _Paragraph(element(p_cxml), None)
        expected_text = ("Foo", "Bar", "Baz")
        return paragraph, expected_text

    @pytest.fixture(
        params=[
            ("a:p", None),
            ("a:p/a:pPr", None),
            ("a:p/a:pPr/a:lnSpc/a:spcPts{val=1800}", 228600),
            ("a:p/a:pPr/a:lnSpc/a:spcPct{val=142000}", 1.42),
            ("a:p/a:pPr/a:lnSpc/a:spcPct{val=124.64%}", 1.2464),
        ]
    )
    def spacing_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:p", 1.42, "a:p/a:pPr/a:lnSpc/a:spcPct{val=142000}"),
            ("a:p", Pt(42), "a:p/a:pPr/a:lnSpc/a:spcPts{val=4200}"),
            (
                "a:p/a:pPr/a:lnSpc/a:spcPct{val=110000}",
                0.875,
                "a:p/a:pPr/a:lnSpc/a:spcPct{val=87500}",
            ),
            (
                "a:p/a:pPr/a:lnSpc/a:spcPts{val=600}",
                Pt(42),
                "a:p/a:pPr/a:lnSpc/a:spcPts{val=4200}",
            ),
            (
                "a:p/a:pPr/a:lnSpc/a:spcPts{val=1900}",
                0.925,
                "a:p/a:pPr/a:lnSpc/a:spcPct{val=92500}",
            ),
            (
                "a:p/a:pPr/a:lnSpc/a:spcPct{val=150000}",
                Pt(24),
                "a:p/a:pPr/a:lnSpc/a:spcPts{val=2400}",
            ),
            ("a:p/a:pPr/a:lnSpc/a:spcPts{val=600}", None, "a:p/a:pPr"),
            ("a:p/a:pPr/a:lnSpc/a:spcPct{val=150000}", None, "a:p/a:pPr"),
        ]
    )
    def spacing_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            # ---single-run---
            ('a:p/a:r/a:t"foobar"', "foobar"),
            # ---multiple-runs---
            ('a:p/(a:r/a:t"foo",a:r/a:t"bar")', "foobar"),
            # ---line-break between runs---
            ('a:p/(a:r/a:t"foo",a:br,a:r/a:t"bar")', "foo\vbar"),
            # ---field between runs---
            ('a:p/(a:r/a:t"foo ",a:fld/a:t"42",a:r/a:t" bar")', "foo 42 bar"),
            # ---line-break and field---
            ('a:p/(a:r/a:t" foo",a:br,a:fld/a:t"42")', " foo\v42"),
            # ---other common p child elements included---
            ('a:p/(a:pPr,a:r/a:t"foobar",a:endParaRPr)', "foobar"),
            # ---field by itself---
            ('a:p/a:fld/a:t"42"', "42"),
            # ---line-break by itself---
            ("a:p/a:br", "\v"),
        ]
    )
    def text_get_fixture(self, request):
        p_cxml, expected_value = request.param
        p = element(p_cxml)
        return p, expected_value

    @pytest.fixture(
        params=[
            ('a:p/(a:r/a:t"foo",a:r/a:t"bar")', "foobar", 'a:p/a:r/a:t"foobar"'),
            ("a:p", "", "a:p"),
            ("a:p", "foobar", 'a:p/a:r/a:t"foobar"'),
            ("a:p", "foo\nbar", 'a:p/(a:r/a:t"foo",a:br,a:r/a:t"bar")'),
            ("a:p", "\vfoo\n", 'a:p/(a:br,a:r/a:t"foo",a:br)'),
            ("a:p", "\n\nfoo", 'a:p/(a:br,a:br,a:r/a:t"foo")'),
            ("a:p", "foo\n", 'a:p/(a:r/a:t"foo",a:br)'),
            ("a:p", b"foo\x07\n", 'a:p/(a:r/a:t"foo_x0007_",a:br)'),
            ("a:p", b"7-bit str", 'a:p/a:r/a:t"7-bit str"'),
            ("a:p", b"8-\xc9\x93\xc3\xaf\xc8\xb6 str", 'a:p/a:r/a:t"8-ɓïȶ str"'),
            ("a:p", "ŮŦƑ-8\x1bliteral", 'a:p/a:r/a:t"ŮŦƑ-8_x001B_literal"'),
            (
                "a:p",
                "utf-8 unicode: Hér er texti",
                'a:p/a:r/a:t"utf-8 unicode: Hér er texti"',
            ),
        ]
    )
    def text_set_fixture(self, request):
        p_cxml, value, expected_cxml = request.param
        p = element(p_cxml)
        expected_xml = xml(expected_cxml)
        return p, value, expected_xml

    # fixture components -----------------------------------

    @pytest.fixture
    def Font_(self, request):
        return class_mock(request, "pptx.text.text.Font")

    @pytest.fixture
    def p_bldr(self):
        return a_p().with_nsdecls()

    @pytest.fixture
    def p_with_r_xml(self):
        run_bldr = an_r().with_child(a_t())
        return a_p().with_nsdecls().with_child(run_bldr).xml()

    @pytest.fixture
    def paragraph(self, p_bldr):
        return _Paragraph(p_bldr.element, None)


class Describe_Run(object):
    """Unit-test suite for `pptx.text.text._Run` object."""

    def it_provides_access_to_its_font(self, font_fixture):
        run, rPr, Font_, font_ = font_fixture
        font = run.font
        Font_.assert_called_once_with(rPr)
        assert font == font_

    def it_provides_access_to_a_hyperlink_proxy(self, hyperlink_fixture):
        run, rPr, _Hyperlink_, hlink_ = hyperlink_fixture
        hlink = run.hyperlink
        _Hyperlink_.assert_called_once_with(rPr, run)
        assert hlink is hlink_

    def it_can_get_the_text_of_the_run(self, text_get_fixture):
        run, expected_value = text_get_fixture
        text = run.text
        assert text == expected_value
        assert is_unicode(text)

    @pytest.mark.parametrize(
        "r_cxml, new_value, expected_r_cxml",
        (
            ("a:r/a:t", "barfoo", 'a:r/a:t"barfoo"'),
            ("a:r/a:t", "bar\x1bfoo", 'a:r/a:t"bar_x001B_foo"'),
            ("a:r/a:t", "bar\tfoo", 'a:r/a:t"bar\tfoo"'),
        ),
    )
    def it_can_change_its_text(self, r_cxml, new_value, expected_r_cxml):
        run = _Run(element(r_cxml), None)
        run.text = new_value
        assert run._r.xml == xml(expected_r_cxml)

    # fixtures ---------------------------------------------

    @pytest.fixture
    def font_fixture(self, Font_, font_):
        r = element("a:r/a:rPr")
        rPr = r.rPr
        run = _Run(r, None)
        return run, rPr, Font_, font_

    @pytest.fixture
    def hyperlink_fixture(self, _Hyperlink_, hlink_):
        r = element("a:r/a:rPr")
        rPr = r.rPr
        run = _Run(r, None)
        return run, rPr, _Hyperlink_, hlink_

    @pytest.fixture
    def text_get_fixture(self):
        r = element('a:r/a:t"foobar"')
        run = _Run(r, None)
        return run, "foobar"

    # fixture components -----------------------------------

    @pytest.fixture
    def Font_(self, request, font_):
        return class_mock(request, "pptx.text.text.Font", return_value=font_)

    @pytest.fixture
    def font_(self, request):
        return instance_mock(request, Font)

    @pytest.fixture
    def _Hyperlink_(self, request, hlink_):
        return class_mock(request, "pptx.text.text._Hyperlink", return_value=hlink_)

    @pytest.fixture
    def hlink_(self, request):
        return instance_mock(request, _Hyperlink)


class Describe_TextFont(object):
    def it_knows_its_typeface(self, typeface_get_fixture):
        text_font, expected_value = typeface_get_fixture
        assert text_font.typeface == expected_value

    def it_knows_its_panose(self, panose_get_fixture):
        text_font, expected_value = panose_get_fixture
        assert text_font.panose == expected_value

    def it_knows_its_pitch_family(self, pitch_get_fixture):
        text_font, expected_value = pitch_get_fixture
        assert text_font.pitch_family == expected_value

    def it_knows_its_charset(self, charset_get_fixture):
        text_font, expected_value = charset_get_fixture
        assert text_font.charset == expected_value

    # fixtures ---------------------------------------------


    @pytest.fixture(
        params=[("a:latin", None), ("a:latin{typeface=Foobar}", "Foobar")]
    )
    def typeface_get_fixture(self, request):
        latin_cxml, expected_value = request.param
        text_font = TextFont(element(latin_cxml))
        return text_font, expected_value

    @pytest.fixture(
        params=[("a:latin", None), ("a:latin{panose=020F0502020204030204}", "020F0502020204030204")]
    )
    def panose_get_fixture(self, request):
        latin_cxml, expected_value = request.param
        text_font = TextFont(element(latin_cxml))
        return text_font, expected_value

    @pytest.fixture(
        params=[("a:latin", 0), ("a:latin{pitchFamily=1}", 1), ("a:latin{pitchFamily=0}", 0)]
    )
    def pitch_get_fixture(self, request):
        latin_cxml, expected_value = request.param
        text_font = TextFont(element(latin_cxml))
        return text_font, expected_value

    @pytest.fixture(
        params=[("a:latin", 1), ("a:latin{charset=0}", 0), ("a:latin{charset=1}", 1)]
    )
    def charset_get_fixture(self, request):
        latin_cxml, expected_value = request.param
        text_font = TextFont(element(latin_cxml))
        return text_font, expected_value


class Describe_TextListStyle:
    def it_provides_access_to_its_default_style(self, text_list_style):
        assert isinstance(text_list_style.default, ParagraphProperties)

    def it_provides_access_to_its_level1_style(self, text_list_style):
        assert isinstance(text_list_style.level_1, ParagraphProperties)

    def it_provides_access_to_its_level2_style(self, text_list_style):
        assert isinstance(text_list_style.level_2, ParagraphProperties)

    def it_provides_access_to_its_level3_style(self, text_list_style):
        assert isinstance(text_list_style.level_3, ParagraphProperties)

    def it_provides_access_to_its_level4_style(self, text_list_style):
        assert isinstance(text_list_style.level_4, ParagraphProperties)

    def it_provides_access_to_its_level5_style(self, text_list_style):
        assert isinstance(text_list_style.level_5, ParagraphProperties)

    def it_provides_access_to_its_level6_style(self, text_list_style):
        assert isinstance(text_list_style.level_6, ParagraphProperties)

    def it_provides_access_to_its_level7_style(self, text_list_style):
        assert isinstance(text_list_style.level_7, ParagraphProperties)

    def it_provides_access_to_its_level8_style(self, text_list_style):
        assert isinstance(text_list_style.level_8, ParagraphProperties)

    def it_provides_access_to_its_level9_style(self, text_list_style):
        assert isinstance(text_list_style.level_9, ParagraphProperties)

    # fixture components -----------------------------------

    @pytest.fixture
    def text_list_style(self, request):
        return TextListStyle(element("a:lstStyle"))


class Describe_ParagraphProperties(object):
    """Unit test suite for pptx.text.text.ParagraphProperties object."""

    def it_knows_its_horizontal_alignment(self, alignment_get_fixture):
        paragraph, expected_value = alignment_get_fixture
        assert paragraph.alignment == expected_value

    def it_can_change_its_horizontal_alignment(self, alignment_set_fixture):
        paragraph, new_value, expected_xml = alignment_set_fixture
        paragraph.alignment = new_value
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_the_default_paragraph_font(self, paragraph_properties, Font_):
        font = paragraph_properties.font
        Font_.assert_called_once_with(paragraph_properties._defRPr)
        assert font == Font_.return_value

    def it_knows_its_indentation_level(self, level_get_fixture):
        paragraph, expected_value = level_get_fixture
        assert paragraph.level == expected_value

    def it_can_change_its_indentation_level(self, level_set_fixture):
        paragraph, new_value, expected_xml = level_set_fixture
        paragraph.level = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_left_margin(self, left_margin_get_fixture):
        paragraph, expected_value = left_margin_get_fixture
        assert paragraph.margin_left == expected_value

    def it_can_change_its_left_margin(self, left_margin_set_fixture):
        paragraph, new_value, expected_xml = left_margin_set_fixture
        paragraph.margin_left = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_right_margin(self, right_margin_get_fixture):
        paragraph, expected_value = right_margin_get_fixture
        assert paragraph.margin_right == expected_value

    def it_can_change_its_right_margin(self, right_margin_set_fixture):
        paragraph, new_value, expected_xml = right_margin_set_fixture
        paragraph.margin_right = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_indent(self, indent_get_fixture):
        paragraph, expected_value = indent_get_fixture
        assert paragraph.indent == expected_value

    def it_can_change_its_indent(self, indent_set_fixture):
        paragraph, new_value, expected_xml = indent_set_fixture
        paragraph.indent = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_line_spacing(self, spacing_get_fixture):
        paragraph, expected_value = spacing_get_fixture
        assert paragraph.line_spacing == expected_value

    def it_can_change_its_line_spacing(self, spacing_set_fixture):
        paragraph, new_value, expected_xml = spacing_set_fixture
        paragraph.line_spacing = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_space_after(self, after_get_fixture):
        paragraph, expected_value = after_get_fixture
        assert paragraph.space_after == expected_value

    def it_can_change_its_space_after(self, after_set_fixture):
        paragraph, new_value, expected_xml = after_set_fixture
        paragraph.space_after = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_space_before(self, before_get_fixture):
        paragraph, expected_value = before_get_fixture
        assert paragraph.space_before == expected_value

    def it_can_change_its_space_before(self, before_set_fixture):
        paragraph, new_value, expected_xml = before_set_fixture
        paragraph.space_before = new_value
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_its_bullet_text(self, paragraph_properties):
        assert isinstance(paragraph_properties.bullet_text, TextBullet)

    def it_provides_access_to_its_bullet_color(self, paragraph_properties):
        assert isinstance(paragraph_properties.bullet_color, TextBulletColor)

    def it_provides_access_to_its_bullet_size(self, paragraph_properties):
        assert isinstance(paragraph_properties.bullet_size, TextBulletSize)

    def it_provides_access_to_its_bullet_font(self, paragraph_properties):
        assert isinstance(paragraph_properties.bullet_font, TextBulletTypeface)

    # fixtures ---------------------------------------------

    @pytest.fixture(
        params=[
            ("a:pPr", None),
            ("a:pPr/a:spcAft/a:spcPct{val=150000}", None),
            ("a:pPr/a:spcAft/a:spcPts{val=600}", 76200),
        ]
    )
    def after_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            (
                "a:pPr/a:spcAft/a:spcPts{val=600}",
                Pt(42),
                "a:pPr/a:spcAft/a:spcPts{val=4200}",
            ),
            (
                "a:pPr/a:spcAft/a:spcPct{val=150000}",
                Pt(24),
                "a:pPr/a:spcAft/a:spcPts{val=2400}",
            ),
            ("a:pPr/a:spcAft/a:spcPts{val=600}", None, "a:pPr"),
        ]
    )
    def after_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:pPr{algn=ctr}", PP_ALIGN.CENTER),
            ("a:pPr{algn=r}", PP_ALIGN.RIGHT),
        ]
    )
    def alignment_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", PP_ALIGN.LEFT, "a:pPr{algn=l}"),
            ("a:pPr{algn=l}", PP_ALIGN.JUSTIFY, "a:pPr{algn=just}"),
            ("a:pPr{algn=just}", None, "a:pPr"),
        ]
    )
    def alignment_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:pPr", None),
            ("a:pPr/a:spcBef/a:spcPct{val=150000}", None),
            ("a:pPr/a:spcBef/a:spcPts{val=600}", 76200),
        ]
    )
    def before_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", Pt(8.333), "a:pPr/a:spcBef/a:spcPts{val=833}"),
            (
                "a:pPr/a:spcBef/a:spcPts{val=600}",
                Pt(42),
                "a:pPr/a:spcBef/a:spcPts{val=4200}",
            ),
            (
                "a:pPr/a:spcBef/a:spcPct{val=150000}",
                Pt(24),
                "a:pPr/a:spcBef/a:spcPts{val=2400}",
            ),
            ("a:pPr/a:spcBef/a:spcPts{val=600}", None, "a:pPr"),
        ]
    )
    def before_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:pPr", 0), ("a:pPr{lvl=2}", 2)])
    def level_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", 1, "a:pPr{lvl=1}"),
            ("a:pPr{lvl=1}", 2, "a:pPr{lvl=2}"),
            ("a:pPr{lvl=2}", 0, "a:pPr"),
        ]
    )
    def level_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:pPr", 0), ("a:pPr{marL=10000}", 10000)])
    def left_margin_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", 1, "a:pPr{marL=1}"),
            ("a:pPr{marL=1}", 2, "a:pPr{marL=2}"),
            ("a:pPr{marL=2}", 0, "a:pPr"),
        ]
    )
    def left_margin_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:pPr", 0), ("a:pPr{indent=10000}", 10000)])
    def indent_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", 1, "a:pPr{indent=1}"),
            ("a:pPr{indent=1}", 2, "a:pPr{indent=2}"),
            ("a:pPr{indent=2}", 0, "a:pPr"),
        ]
    )
    def indent_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[("a:pPr", 0), ("a:pPr{marR=10000}", 10000)])
    def right_margin_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", 1, "a:pPr{marR=1}"),
            ("a:pPr{marR=1}", 2, "a:pPr{marR=2}"),
            ("a:pPr{marR=2}", 0, "a:pPr"),
        ]
    )
    def right_margin_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(
        params=[
            ("a:pPr", "a:br"),
            ("a:r", "(a:r,a:br)"),
            ("a:br", "(a:br,a:br)"),
        ]
    )
    def line_break_fixture(self, request):
        cxml, expected_cxml = request.param
        paragraph = ParagraphProperties(element(cxml), None)
        expected_xml = xml(expected_cxml)
        return paragraph, expected_xml

    @pytest.fixture(
        params=[
            ("a:pPr", None),
            ("a:pPr/a:lnSpc/a:spcPts{val=1800}", 228600),
            ("a:pPr/a:lnSpc/a:spcPct{val=142000}", 1.42),
            ("a:pPr/a:lnSpc/a:spcPct{val=124.64%}", 1.2464),
        ]
    )
    def spacing_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        return paragraph, expected_value

    @pytest.fixture(
        params=[
            ("a:pPr", 1.42, "a:pPr/a:lnSpc/a:spcPct{val=142000}"),
            ("a:pPr", Pt(42), "a:pPr/a:lnSpc/a:spcPts{val=4200}"),
            (
                "a:pPr/a:lnSpc/a:spcPct{val=110000}",
                0.875,
                "a:pPr/a:lnSpc/a:spcPct{val=87500}",
            ),
            (
                "a:pPr/a:lnSpc/a:spcPts{val=600}",
                Pt(42),
                "a:pPr/a:lnSpc/a:spcPts{val=4200}",
            ),
            (
                "a:pPr/a:lnSpc/a:spcPts{val=1900}",
                0.925,
                "a:pPr/a:lnSpc/a:spcPct{val=92500}",
            ),
            (
                "a:pPr/a:lnSpc/a:spcPct{val=150000}",
                Pt(24),
                "a:pPr/a:lnSpc/a:spcPts{val=2400}",
            ),
            ("a:pPr/a:lnSpc/a:spcPts{val=600}", None, "a:pPr"),
            ("a:pPr/a:lnSpc/a:spcPct{val=150000}", None, "a:pPr"),
        ]
    )
    def spacing_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = ParagraphProperties(element(p_cxml))
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml


    # fixture components -----------------------------------
    @pytest.fixture
    def pPr_bldr(self):
        return an_pPr().with_nsdecls()

    @pytest.fixture
    def paragraph_properties(self, pPr_bldr):
        return ParagraphProperties(pPr_bldr.element)

    @pytest.fixture
    def Font_(self, request):
        return class_mock(request, "pptx.text.text.Font")
